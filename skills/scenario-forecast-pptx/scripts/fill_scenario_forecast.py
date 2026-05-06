"""
fill_scenario_forecast.py — BDD向けシナリオ別見立てチャートをPPTXネイティブで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

レイアウト:
  - 左チャート: 売上高（Base / Upside / Downside の 3 シナリオ折れ線）
  - 右チャート: 調整後 EBITDA（同上）
  - 期間種別凡例（実績 / 見込 / 計画）をプロットエリア下に配置
  - シリーズ凡例（Base / Upside / Downside）を各チャート右側縦並び
  - 出典（stella=動的 textbox / roleup=Source 3 placeholder）

Usage:
  python fill_scenario_forecast.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/scenario_forecast_data.json \\
    --output {{OUTPUT_DIR}}/ScenarioForecast_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "scenario-forecast-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names (template-structure invariants) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_TABLE = "Table 1"

# Defaults (stella). Reassigned in _apply_theme(theme).
SHAPE_SOURCE = "Source"

LEFT_CHART_X = Inches(0.41)
LEFT_CHART_W = Inches(5.20)
LEFT_LEGEND_X = Inches(5.65)
RIGHT_CHART_X = Inches(6.80)
RIGHT_CHART_W = Inches(5.20)
RIGHT_LEGEND_X = Inches(12.04)

CHART_TITLE_Y = Inches(1.40)
CHART_TITLE_LINE_Y = Inches(1.75)
UNIT_Y = Inches(1.80)
CHART_Y = Inches(1.95)
CHART_H = Inches(4.50)

PERIOD_LEGEND_Y = Inches(6.60)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.00)
SOURCE_H = Inches(0.30)

TITLE_LINE_OFFSET_RATIO = 0.25

PLOT_AREA_X_FRAC = 0.09
PLOT_AREA_W_FRAC = 0.89
PLOT_AREA_Y_FRAC = 0.02
PLOT_AREA_H_FRAC = 0.82

# Stella defaults; reassigned in _apply_theme(theme).
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_SUBTITLE = RGBColor(0x33, 0x33, 0x33)
COLOR_BASE = RGBColor(0x4E, 0x79, 0xA7)
COLOR_UPSIDE = RGBColor(0xED, 0x7D, 0x31)
COLOR_DOWNSIDE = RGBColor(0xA5, 0xA5, 0xA5)

TEXT_HEX = "333333"
COLOR_BASE_HEX = "4E79A7"
COLOR_UPSIDE_HEX = "ED7D31"
COLOR_DOWNSIDE_HEX = "A5A5A5"

FONT_NAME_JP = "Meiryo UI"

CHART_TITLE_PT = 14
CHART_TITLE_BOLD = True
CHART_TITLE_ALIGN = PP_ALIGN.CENTER
CHART_TITLE_UNDERLINE = True
UNIT_FONT_PT = 9
LEGEND_FONT_PT = 9
PERIOD_LEGEND_PT = 10
DATA_LABEL_SZ = "1000"
AXIS_FONT_SZ = "1000"
SOURCE_FONT_PT = 10

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global LEFT_CHART_X, LEFT_CHART_W, LEFT_LEGEND_X
    global RIGHT_CHART_X, RIGHT_CHART_W, RIGHT_LEGEND_X
    global CHART_TITLE_Y, CHART_TITLE_LINE_Y, UNIT_Y, CHART_Y, CHART_H
    global PERIOD_LEGEND_Y, SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, COLOR_SUBTITLE
    global COLOR_BASE, COLOR_UPSIDE, COLOR_DOWNSIDE
    global TEXT_HEX, COLOR_BASE_HEX, COLOR_UPSIDE_HEX, COLOR_DOWNSIDE_HEX
    global FONT_NAME_JP
    global CHART_TITLE_PT, CHART_TITLE_BOLD, CHART_TITLE_ALIGN, CHART_TITLE_UNDERLINE
    global UNIT_FONT_PT, LEGEND_FONT_PT, PERIOD_LEGEND_PT
    global DATA_LABEL_SZ, AXIS_FONT_SZ, SOURCE_FONT_PT

    _THEME = theme
    FONT_NAME_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    LEFT_CHART_X = theme.layout("left_chart_x_in")
    LEFT_CHART_W = theme.layout("left_chart_w_in")
    LEFT_LEGEND_X = theme.layout("left_legend_x_in")
    RIGHT_CHART_X = theme.layout("right_chart_x_in")
    RIGHT_CHART_W = theme.layout("right_chart_w_in")
    RIGHT_LEGEND_X = theme.layout("right_legend_x_in")
    CHART_TITLE_Y = theme.layout("chart_title_y_in")
    CHART_TITLE_LINE_Y = theme.layout("chart_title_line_y_in")
    UNIT_Y = theme.layout("unit_y_in")
    CHART_Y = theme.layout("chart_y_in")
    CHART_H = theme.layout("chart_h_in")
    PERIOD_LEGEND_Y = theme.layout("period_legend_y_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        # stella V1 visual defaults
        COLOR_BASE = RGBColor(0x4E, 0x79, 0xA7)
        COLOR_UPSIDE = RGBColor(0xED, 0x7D, 0x31)
        COLOR_DOWNSIDE = RGBColor(0xA5, 0xA5, 0xA5)
        COLOR_BASE_HEX = "4E79A7"
        COLOR_UPSIDE_HEX = "ED7D31"
        COLOR_DOWNSIDE_HEX = "A5A5A5"
        COLOR_SUBTITLE = COLOR_TEXT
        CHART_TITLE_PT = 14
        CHART_TITLE_BOLD = True
        CHART_TITLE_ALIGN = PP_ALIGN.CENTER
        CHART_TITLE_UNDERLINE = True
        UNIT_FONT_PT = 9
        LEGEND_FONT_PT = 9
        PERIOD_LEGEND_PT = 10
        DATA_LABEL_SZ = "1000"
        AXIS_FONT_SZ = "1000"
        SOURCE_FONT_PT = 10
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        # 3 シナリオは brand chart_palette[0..2]
        palette = list(theme.chart_palette)
        COLOR_BASE_HEX = palette[0].lstrip("#")          # #7C4C2C
        COLOR_UPSIDE_HEX = palette[1].lstrip("#")        # #897141
        COLOR_DOWNSIDE_HEX = theme.hex_no_hash("highlight_other")  # #CDCECE
        COLOR_BASE = _hex2rgb(COLOR_BASE_HEX)
        COLOR_UPSIDE = _hex2rgb(COLOR_UPSIDE_HEX)
        COLOR_DOWNSIDE = _hex2rgb(COLOR_DOWNSIDE_HEX)
        COLOR_SUBTITLE = theme.color("subtitle")  # #897141
        CHART_TITLE_PT = theme.pt_value("font_size_subtitle_pt")  # 12
        CHART_TITLE_BOLD = False
        CHART_TITLE_ALIGN = PP_ALIGN.LEFT
        CHART_TITLE_UNDERLINE = False  # roleup template の object 8 が下線役、ここは追加しない
        UNIT_FONT_PT = theme.pt_value("font_size_body_pt")    # 10 (was 9)
        LEGEND_FONT_PT = theme.pt_value("font_size_body_pt")  # 10 (was 9)
        PERIOD_LEGEND_PT = theme.pt_value("font_size_body_pt")  # 10
        DATA_LABEL_SZ = "1000"
        AXIS_FONT_SZ = "1000"
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6


def _hex2rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _silent_remove_shape(slide, name):
    for shape in list(slide.shapes):
        if shape.name == name:
            slide.shapes._spTree.remove(shape._element)
            return True
    return False


def find_shape(slide, name, warn=True):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    if warn:
        print(f"  Warning: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_textbox(slide, text, left, top, width, height, font_size, bold=False,
                alignment=PP_ALIGN.LEFT, color=None):
    if color is None:
        color = COLOR_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME_JP
    return txBox


def add_horizontal_line(slide, left, top, width, color=None, line_width=Pt(1.0)):
    if color is None:
        color = COLOR_TEXT
    connector = slide.shapes.add_connector(
        1,
        int(left), int(top),
        int(left + width), int(top)
    )
    connector.line.color.rgb = color
    connector.line.width = line_width
    return connector


def build_line_chart(slide, chart_data_dict, left, top, width, height,
                     y_max=None, y_step=None):
    """3シナリオ折れ線チャート"""
    from pptx.chart.data import CategoryChartData

    periods = chart_data_dict["periods"]
    categories = [p["label"] for p in periods]

    base_vals = [p.get("base") for p in periods]
    upside_vals = [p.get("upside") for p in periods]
    downside_vals = [p.get("downside") for p in periods]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Base", base_vals)
    chart_data.add_series("Upside", upside_vals)
    chart_data.add_series("Downside", downside_vals)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        left, top, width, height,
        chart_data
    )
    chart = chart_frame.chart

    plotArea = chart._chartSpace.chart.plotArea
    lineChart = plotArea.findall(qn('c:lineChart'))[0]
    sers = lineChart.findall(qn('c:ser'))

    colors = [COLOR_BASE_HEX, COLOR_UPSIDE_HEX, COLOR_DOWNSIDE_HEX]
    marker_symbols = ['diamond', 'circle', 'circle']

    for i, ser in enumerate(sers):
        color_hex = colors[i]
        marker_sym = marker_symbols[i]

        spPr = ser.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ser, qn('c:spPr'))
        ln = spPr.find(qn('a:ln'))
        if ln is None:
            ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w', '25400')
        sf = ln.find(qn('a:solidFill'))
        if sf is None:
            sf = etree.SubElement(ln, qn('a:solidFill'))
        for child in list(sf):
            sf.remove(child)
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': color_hex})

        marker = ser.find(qn('c:marker'))
        if marker is None:
            marker = etree.SubElement(ser, qn('c:marker'))
        symbol = marker.find(qn('c:symbol'))
        if symbol is None:
            symbol = etree.SubElement(marker, qn('c:symbol'))
        symbol.set('val', marker_sym)
        sz = marker.find(qn('c:size'))
        if sz is None:
            sz = etree.SubElement(marker, qn('c:size'))
        sz.set('val', '7')

        m_spPr = marker.find(qn('c:spPr'))
        if m_spPr is None:
            m_spPr = etree.SubElement(marker, qn('c:spPr'))
        for child in list(m_spPr):
            m_spPr.remove(child)
        m_sf = etree.SubElement(m_spPr, qn('a:solidFill'))
        etree.SubElement(m_sf, qn('a:srgbClr'), attrib={'val': color_hex})
        m_ln = etree.SubElement(m_spPr, qn('a:ln'))
        m_ln_sf = etree.SubElement(m_ln, qn('a:solidFill'))
        etree.SubElement(m_ln_sf, qn('a:srgbClr'), attrib={'val': color_hex})

        add_data_labels_to_ser(ser, num_format='#,##0', font_color=color_hex)

    valAxes = plotArea.findall(qn('c:valAx'))
    if valAxes:
        valAx = valAxes[0]
        if y_max is not None:
            scaling = valAx.find(qn('c:scaling'))
            if scaling is None:
                scaling = etree.SubElement(valAx, qn('c:scaling'))
            max_elem = scaling.find(qn('c:max'))
            if max_elem is None:
                max_elem = etree.SubElement(scaling, qn('c:max'))
            max_elem.set('val', str(y_max))
            min_elem = scaling.find(qn('c:min'))
            if min_elem is None:
                min_elem = etree.SubElement(scaling, qn('c:min'))
            min_elem.set('val', '0')

        if y_step is not None:
            major_unit = valAx.find(qn('c:majorUnit'))
            if major_unit is None:
                major_unit = etree.SubElement(valAx, qn('c:majorUnit'))
            major_unit.set('val', str(y_step))

        _set_axis_font(valAx, AXIS_FONT_SZ)
        numFmt = valAx.find(qn('c:numFmt'))
        if numFmt is None:
            numFmt = etree.SubElement(valAx, qn('c:numFmt'))
        numFmt.set('formatCode', '#,##0')
        numFmt.set('sourceLinked', '0')

    catAxes = plotArea.findall(qn('c:catAx'))
    if catAxes:
        _set_axis_font(catAxes[0], AXIS_FONT_SZ, vertical=True)

    for vax in plotArea.findall(qn('c:valAx')):
        for mg in vax.findall(qn('c:majorGridlines')):
            vax.remove(mg)
        for mg in vax.findall(qn('c:minorGridlines')):
            vax.remove(mg)

    chart.has_legend = False
    chart.has_title = False

    plotArea_spPr = plotArea.find(qn('c:spPr'))
    if plotArea_spPr is None:
        plotArea_spPr = etree.SubElement(plotArea, qn('c:spPr'))
    if plotArea_spPr.find(qn('a:noFill')) is None:
        etree.SubElement(plotArea_spPr, qn('a:noFill'))

    layout = plotArea.find(qn('c:layout'))
    if layout is None:
        layout = etree.SubElement(plotArea, qn('c:layout'))
        plotArea.insert(0, layout)
    for child in list(layout):
        layout.remove(child)
    ml = etree.SubElement(layout, qn('c:manualLayout'))
    etree.SubElement(ml, qn('c:layoutTarget'), attrib={'val': 'inner'})
    etree.SubElement(ml, qn('c:xMode'), attrib={'val': 'edge'})
    etree.SubElement(ml, qn('c:yMode'), attrib={'val': 'edge'})
    etree.SubElement(ml, qn('c:x'), attrib={'val': str(PLOT_AREA_X_FRAC)})
    etree.SubElement(ml, qn('c:y'), attrib={'val': str(PLOT_AREA_Y_FRAC)})
    etree.SubElement(ml, qn('c:w'), attrib={'val': str(PLOT_AREA_W_FRAC)})
    etree.SubElement(ml, qn('c:h'), attrib={'val': str(PLOT_AREA_H_FRAC)})

    plot_abs_left = left + width * PLOT_AREA_X_FRAC
    plot_abs_width = width * PLOT_AREA_W_FRAC

    print(f"  Chart: {len(periods)} periods, 3 series  "
          f"(plot x={plot_abs_left/914400:.2f}in w={plot_abs_width/914400:.2f}in)")
    return chart_frame, plot_abs_left, plot_abs_width


def _set_axis_font(ax_elem, font_size_hundredths='1000', vertical=False):
    txPr = ax_elem.find(qn('c:txPr'))
    if txPr is None:
        txPr = etree.SubElement(ax_elem, qn('c:txPr'))
    bodyPr = txPr.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
        txPr.insert(0, bodyPr)
    if vertical:
        bodyPr.set('rot', '-5400000')
        bodyPr.set('vert', 'horz')
    if txPr.find(qn('a:lstStyle')) is None:
        etree.SubElement(txPr, qn('a:lstStyle'))
    p = txPr.find(qn('a:p'))
    if p is None:
        p = etree.SubElement(txPr, qn('a:p'))
    pPr = p.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p, qn('a:pPr'))
    defRPr = pPr.find(qn('a:defRPr'))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn('a:defRPr'))
    defRPr.set('sz', font_size_hundredths)
    if defRPr.find(qn('a:latin')) is None:
        etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
    if defRPr.find(qn('a:ea')) is None:
        etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})
    sf = defRPr.find(qn('a:solidFill'))
    if sf is None:
        sf = etree.SubElement(defRPr, qn('a:solidFill'))
    for child in list(sf):
        sf.remove(child)
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': TEXT_HEX})


def add_data_labels_to_ser(ser_xml, num_format='#,##0', font_color=None):
    if font_color is None:
        font_color = TEXT_HEX
    dLbls = ser_xml.find(qn('c:dLbls'))
    if dLbls is None:
        dLbls = etree.SubElement(ser_xml, qn('c:dLbls'))
    for child in list(dLbls):
        dLbls.remove(child)
    etree.SubElement(dLbls, qn('c:numFmt'), attrib={
        'formatCode': num_format, 'sourceLinked': '0'
    })
    etree.SubElement(dLbls, qn('c:showLegendKey'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dLbls, qn('c:showCatName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showSerName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showPercent'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showBubbleSize'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:dLblPos'), attrib={'val': 't'})

    txPr = etree.SubElement(dLbls, qn('c:txPr'))
    etree.SubElement(txPr, qn('a:bodyPr'))
    etree.SubElement(txPr, qn('a:lstStyle'))
    p = etree.SubElement(txPr, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    defRPr = etree.SubElement(pPr, qn('a:defRPr'),
                              attrib={'sz': DATA_LABEL_SZ, 'b': '1'})
    etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
    etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})
    sf = etree.SubElement(defRPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': font_color})


def add_series_legend_right(slide, series_labels, legend_x, chart_top, chart_h):
    """シリーズ凡例をチャート右側に縦並びで配置"""
    colors = [COLOR_BASE, COLOR_UPSIDE, COLOR_DOWNSIDE]
    markers = ['diamond', 'circle', 'circle']
    labels = [
        series_labels.get("base", "Base"),
        series_labels.get("upside", "Upside"),
        series_labels.get("downside", "Downside"),
    ]

    item_h = Inches(0.30)
    total_legend_h = item_h * len(labels)
    start_y = chart_top + (chart_h - total_legend_h) / 2

    for i, (label, color, marker_type) in enumerate(zip(labels, colors, markers)):
        y = start_y + item_h * i
        x = legend_x

        marker_size = Inches(0.10)
        marker_y = int(y + (item_h - marker_size) / 2)
        if marker_type == 'diamond':
            sym = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND, int(x), marker_y,
                Inches(0.12), Inches(0.12)
            )
        else:
            sym = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, int(x), marker_y,
                marker_size, marker_size
            )
        sym.fill.solid()
        sym.fill.fore_color.rgb = color
        sym.line.fill.background()

        line_x = int(x + Inches(0.15))
        line_y = int(y + item_h / 2)
        connector = slide.shapes.add_connector(
            1, line_x, line_y, int(line_x + Inches(0.25)), line_y
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(1.5)

        txBox = slide.shapes.add_textbox(
            int(x + Inches(0.44)), int(y + Inches(0.02)),
            Inches(0.65), item_h
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label
        run.font.size = Pt(LEGEND_FONT_PT)
        run.font.color.rgb = COLOR_TEXT
        run.font.name = FONT_NAME_JP

    print(f"  Legend (right): {', '.join(labels)}")


def add_period_type_legend(slide, periods, period_type_labels, left, top, width):
    """期間種別凡例（実績|見込|計画）"""
    type_ranges = []
    current_type = None
    start_idx = 0

    for i, p in enumerate(periods):
        ptype = p["type"]
        if ptype != current_type:
            if current_type is not None:
                type_ranges.append((current_type, start_idx, i - 1))
            current_type = ptype
            start_idx = i
    if current_type is not None:
        type_ranges.append((current_type, start_idx, len(periods) - 1))

    n_periods = len(periods)
    legend_h = Inches(0.30)
    period_w = width / n_periods
    margin = Inches(0.05)

    for ptype, s_idx, e_idx in type_ranges:
        label = period_type_labels.get(ptype, ptype)
        box_x = left + period_w * s_idx
        box_w = period_w * (e_idx - s_idx + 1)

        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(box_x + margin), int(top),
            int(box_w - margin * 2), int(legend_h)
        )
        rect.fill.background()
        rect.line.color.rgb = COLOR_TEXT
        rect.line.width = Pt(0.75)

        tf = rect.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(PERIOD_LEGEND_PT)
        run.font.color.rgb = COLOR_TEXT
        run.font.name = FONT_NAME_JP

    print(f"  Period legend: {[r[0] for r in type_ranges]}")


def add_source_label(slide, text):
    """出典: roleup なら Source 3 placeholder、stella なら textbox。"""
    src_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
    if src_shape is not None and src_shape.has_text_frame:
        tf = src_shape.text_frame
        tf.word_wrap = True
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(SOURCE_FONT_PT)
        run.font.color.rgb = COLOR_SOURCE
        run.font.name = FONT_NAME_JP
        print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {text[:50]}")
        return

    txBox = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SOURCE_FONT_PT)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP
    print(f"  ✓ 出典 (textbox): {text[:50]}")


def build_chart_panel(slide, chart_data, series_labels, period_type_labels,
                      chart_x, chart_w, legend_x):
    """左 or 右のチャートパネル全体を構築"""

    title_w = chart_w + (legend_x - chart_x - chart_w) + Inches(0.80)
    add_textbox(slide, chart_data.get("title", ""),
                chart_x, CHART_TITLE_Y, title_w, Inches(0.22),
                Pt(CHART_TITLE_PT), bold=CHART_TITLE_BOLD,
                alignment=CHART_TITLE_ALIGN, color=COLOR_SUBTITLE)

    if CHART_TITLE_UNDERLINE:
        full_line_w = chart_w + Inches(0.80)
        line_w = full_line_w * 0.70
        line_x = chart_x + (full_line_w - line_w) * TITLE_LINE_OFFSET_RATIO
        add_horizontal_line(slide, line_x, CHART_TITLE_LINE_Y, line_w,
                            color=COLOR_TEXT, line_width=Pt(0.75))

    add_textbox(slide, chart_data.get("unit", ""),
                chart_x, UNIT_Y, Inches(2.00), Inches(0.18),
                Pt(UNIT_FONT_PT), alignment=PP_ALIGN.LEFT,
                color=COLOR_SOURCE)

    _, plot_left, plot_width = build_line_chart(
        slide, chart_data, chart_x, CHART_Y, chart_w, CHART_H,
        y_max=chart_data.get("y_max"),
        y_step=chart_data.get("y_step"))

    add_series_legend_right(slide, series_labels, legend_x, CHART_Y, CHART_H)

    add_period_type_legend(slide, chart_data["periods"], period_type_labels,
                           plot_left, PERIOD_LEGEND_Y, plot_width)


def main():
    parser = argparse.ArgumentParser(description="BDD Scenario Forecast Chart Generator")
    parser.add_argument("--data", required=True)
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand.",
    )
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "scenario-forecast")

    print(f"=== シナリオ別見立てチャート (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    # chart_title のデフォルトを data に埋めて、brand 別の top/subtitle 解決で
    # どちらの placeholder field に書かれても落ちない状態にする。
    if not data.get("chart_title"):
        data["chart_title"] = "今後の見立て ー ケース別 ー"

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top placeholder (brand-aware)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:50]}")

    # Stella V1: Table 1 placeholder removal (silent for roleup template)
    _silent_remove_shape(slide, SHAPE_TABLE)

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    series_labels = data.get("series_labels", {
        "base": "Base", "upside": "Upside", "downside": "Downside"
    })
    period_type_labels = data.get("period_type_labels", {
        "actual": "実績", "forecast": "見込", "plan": "計画"
    })

    left_data = data.get("left_chart", {})
    if left_data and left_data.get("periods"):
        build_chart_panel(slide, left_data, series_labels, period_type_labels,
                          LEFT_CHART_X, LEFT_CHART_W, LEFT_LEGEND_X)

    right_data = data.get("right_chart", {})
    if right_data and right_data.get("periods"):
        build_chart_panel(slide, right_data, series_labels, period_type_labels,
                          RIGHT_CHART_X, RIGHT_CHART_W, RIGHT_LEGEND_X)

    source = data.get("source", "")
    if source:
        body = source if source.startswith("出典") else f"出典：{source}"
        add_source_label(slide, body)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
