"""
fill_company_history.py — 会社沿革データをPPTXネイティブテーブルとして生成するスクリプト

テンプレート構造（company-history-template.pptx）:
  - Title 1   (TEXT_BOX): スライドタイトル（デフォルト: 会社沿革）
  - Table 1   (TABLE):    2列テーブル（年 | 概要）。ヘッダー行+データ行1行をテンプレートとして保持

使い方:
  python fill_company_history.py \
    --data /home/claude/company_history_data.json \
    --output /mnt/user-data/outputs/CompanyHistory_output.pptx \
    [--brand stellar_aiz|roleup] [--template <path>]

`--brand` (default: stellar_aiz) selects the output format. `--template`
is optional; if omitted it is resolved via brand_resolver.template_path()
which falls back to the stella default template when the requested brand
does not yet have a curated template. Cell rPr/tcPr styling is driven by
the template itself (no per-brand color/font overrides in this script).
"""

import argparse
import copy
import json
import os
import sys

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── brand_resolver bootstrap (skills/_common/lib/brand_resolver.py) ──
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "company-history-pptx"

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──────────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_TABLE        = "Table 1"
SHAPE_SOURCE       = "Source 3"  # roleup template placeholder; stella では存在しない可能性あり

SKILL_ID = "company-history-pptx"

# ── Brand-aware module global (slide height) ──
# Default = stella's 7.5in for backward compat; reassigned in main() from
# the actual loaded template's prs.slide_height so layout calculations stay
# correct when a non-stella brand introduces a different slide size.
SLIDE_H_DEFAULT = Inches(7.5)
# ────────────────────────────────────────────────────────


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def _silent_remove_shape(slide, name):
    """find_shape の warning を出さずに削除を試みる(brand 別 shape 名フォールバック用)"""
    for shape in slide.shapes:
        if shape.name == name:
            slide.shapes._spTree.remove(shape._element)
            return True
    return False


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def rebuild_history_table(slide, history_data, slide_height=None):
    """
    テンプレートのテーブルを削除し、沿革データに応じた行数でネイティブテーブルを再構築する。
    ヘッダー行・データ行のセルスタイルはテンプレートから複製する。
    """
    table_shape = find_shape(slide, SHAPE_TABLE)
    if table_shape is None:
        print("  ⚠ WARNING: Table shape not found", file=sys.stderr)
        return

    old_table = table_shape.table

    # セルスタイル（tcPr）をテンプレートからコピー
    header_tcPr_0 = copy.deepcopy(old_table.cell(0, 0)._tc.find(qn("a:tcPr")))
    header_tcPr_1 = copy.deepcopy(old_table.cell(0, 1)._tc.find(qn("a:tcPr")))
    data_tcPr_0   = copy.deepcopy(old_table.cell(1, 0)._tc.find(qn("a:tcPr")))
    data_tcPr_1   = copy.deepcopy(old_table.cell(1, 1)._tc.find(qn("a:tcPr")))

    # runスタイル（rPr）をテンプレートからコピー
    def get_rPr(cell):
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.find(qn("a:rPr"))
                if rPr is not None:
                    return copy.deepcopy(rPr)
        return None

    header_rPr_0 = get_rPr(old_table.cell(0, 0))
    header_rPr_1 = get_rPr(old_table.cell(0, 1))
    data_rPr_0   = get_rPr(old_table.cell(1, 0))
    data_rPr_1   = get_rPr(old_table.cell(1, 1))

    # 位置・サイズ・列幅を保存
    tbl_left   = table_shape.left
    tbl_top    = table_shape.top
    tbl_width  = table_shape.width
    tbl_height = table_shape.height
    col0_width = old_table.columns[0].width
    col1_width = old_table.columns[1].width

    # 既存テーブルを削除
    sp_tree = slide.shapes._spTree
    sp_tree.remove(table_shape._element)

    # テーブル配置の計算
    SLIDE_HEIGHT      = slide_height if slide_height is not None else SLIDE_H_DEFAULT
    BOTTOM_MARGIN     = Inches(0.15)
    HEADER_ROW_HEIGHT = Inches(0.35)
    MAX_TABLE_HEIGHT  = SLIDE_HEIGHT - tbl_top - BOTTOM_MARGIN

    n_rows = len(history_data) + 1  # ヘッダー行 + データ行
    n_cols = 2

    # データ行の高さを動的算出（残りスペースを均等割り、上限0.4inで制限）
    available_for_data = MAX_TABLE_HEIGHT - HEADER_ROW_HEIGHT
    data_row_height = int(available_for_data / len(history_data))
    MAX_DATA_ROW = Inches(0.4)
    if data_row_height > MAX_DATA_ROW:
        data_row_height = MAX_DATA_ROW
    dynamic_height = HEADER_ROW_HEIGHT + data_row_height * len(history_data)

    new_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, dynamic_height)
    new_shape.name = SHAPE_TABLE
    new_table = new_shape.table

    # 列幅を復元
    new_table.columns[0].width = col0_width
    new_table.columns[1].width = col1_width

    # 行の高さを明示的に設定
    tbl_xml = new_shape._element.find('.//' + qn('a:tbl'))
    for i, tr in enumerate(tbl_xml.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(HEADER_ROW_HEIGHT))
        else:
            tr.set('h', str(data_row_height))

    print(f"  ✓ テーブル高さ: {dynamic_height/914400:.2f}in (データ行: {data_row_height/914400:.2f}in × {len(history_data)}行)")

    # tblPrを設定（デフォルトスタイル除去）
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tbl_elem = new_shape._element.find('.//a:tbl', ns)
    old_tblPr = tbl_elem.find('a:tblPr', ns)
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    def apply_cell_single(cell, text, tcPr_tmpl, rPr_tmpl):
        """セルに単一テキストを設定（スタイルをテンプレートから複製）"""
        tc = cell._tc
        txBody = tc.find(qn("a:txBody"))
        if txBody is None:
            txBody = etree.SubElement(tc, qn("a:txBody"))
            etree.SubElement(txBody, qn("a:bodyPr"))
            etree.SubElement(txBody, qn("a:lstStyle"))

        # 既存の段落をすべて削除
        for p in txBody.findall(qn("a:p")):
            txBody.remove(p)

        # 段落を追加
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", "l")  # 左寄せ
        r_elem = etree.SubElement(p_elem, qn("a:r"))
        if rPr_tmpl is not None:
            r_elem.append(copy.deepcopy(rPr_tmpl))
        else:
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP", "sz": "1400"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = str(text)

        # tcPrを適用
        old_tcPr = tc.find(qn("a:tcPr"))
        if old_tcPr is not None:
            tc.remove(old_tcPr)
        if tcPr_tmpl is not None:
            tc.append(copy.deepcopy(tcPr_tmpl))

    def apply_cell_multiline(cell, lines, tcPr_tmpl, rPr_tmpl):
        """セルに複数行テキストを設定（各行を別段落として追加）"""
        tc = cell._tc
        txBody = tc.find(qn("a:txBody"))
        if txBody is None:
            txBody = etree.SubElement(tc, qn("a:txBody"))
            etree.SubElement(txBody, qn("a:bodyPr"))
            etree.SubElement(txBody, qn("a:lstStyle"))

        # 既存の段落をすべて削除
        for p in txBody.findall(qn("a:p")):
            txBody.remove(p)

        # 各行を個別の段落として追加
        for line in lines:
            p_elem = etree.SubElement(txBody, qn("a:p"))
            pPr = etree.SubElement(p_elem, qn("a:pPr"))
            pPr.set("algn", "l")  # 左寄せ
            r_elem = etree.SubElement(p_elem, qn("a:r"))
            if rPr_tmpl is not None:
                r_elem.append(copy.deepcopy(rPr_tmpl))
            else:
                etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP", "sz": "1400"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = str(line)

        # tcPrを適用
        old_tcPr = tc.find(qn("a:tcPr"))
        if old_tcPr is not None:
            tc.remove(old_tcPr)
        if tcPr_tmpl is not None:
            tc.append(copy.deepcopy(tcPr_tmpl))

    # ── ヘッダー行 ──
    apply_cell_single(new_table.cell(0, 0), "年",   header_tcPr_0, header_rPr_0)
    apply_cell_single(new_table.cell(0, 1), "概要", header_tcPr_1, header_rPr_1)
    print("  ✓ ヘッダー行: 年 | 概要")

    # ── データ行 ──
    for r_idx, item in enumerate(history_data):
        year = str(item.get("year", ""))
        events = item.get("events", [])
        if isinstance(events, str):
            events = [events]

        # 年セル（単一テキスト）
        apply_cell_single(new_table.cell(r_idx + 1, 0), year, data_tcPr_0, data_rPr_0)

        # 概要セル（複数イベントは「、」で結合して1行にする）
        events_text = "、".join(events)
        apply_cell_single(new_table.cell(r_idx + 1, 1), events_text, data_tcPr_1, data_rPr_1)

        print(f"  ✓ [{year}] {events_text[:60]}{'...' if len(events_text) > 60 else ''}")

    print(f"  ✓ テーブル生成完了: {n_rows}行 × {n_cols}列")


def main():
    parser = argparse.ArgumentParser(description="会社沿革 PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True, help="JSONデータファイルのパス")
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    parser.add_argument("--output", required=True, help="出力PPTXファイルのパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    template_path = args.template or theme.template_path(SKILL_DIR, "company-history")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    # JSONデータ読み込み
    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    print(f"  データ読み込み完了: {len(data.get('history', []))}件の沿革")

    # roleup は出所必須 (theme.layout_rules.source_required = true)。
    # stella は no-op で従来挙動維持。
    require_source(data, theme, skill_id=SKILL_ID)

    # テンプレート読み込み
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 0. roleup 公式テンプレ由来のチャート位置ガイド矩形(茶色 accent2)を出力から除去。
    #    ガイドはチャート寸法の参考にすぎず、出力ではマスター背景(白)が見えるべき。
    #    stella テンプレにはこれらの shape 名が存在しないため silent no-op で安全。
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # 1. メインメッセージ設定
    # Top placeholder (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text}")

    # 2. チャートタイトル設定
    # Subtitle placeholder (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme) or "会社沿革"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text}")

    # 3. テーブル再構築（slide_height はテンプレートから取得 = brand 固有のスライドサイズに対応）
    history = data.get("history", [])
    rebuild_history_table(slide, history, slide_height=prs.slide_height)

    # 4. 出典書き込み (Source 3 placeholder があればそこへ)。stella では shape 不在なら skip。
    src_text = data.get("source", "")
    if src_text:
        src_shape = None
        for nm in (SHAPE_SOURCE, "Source"):
            for shape in slide.shapes:
                if shape.name == nm:
                    src_shape = shape
                    break
            if src_shape is not None:
                break
        if src_shape is not None:
            display = src_text if src_text.startswith("出典") else f"出典：{src_text}"
            set_textbox_text(src_shape, display)
            print(f"  ✓ 出典 ({src_shape.name}): {display}")

    # 5. 出力
    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
