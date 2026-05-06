"""
fill_pyramid_native.py — N段ピラミッドをPowerPointネイティブオブジェクトで動的生成

■ 方式:
  テンプレートの Title 1 / Text Placeholder 2 にテキストを流し込み、
  ピラミッド図部分は python-pptx で Shape を動的に生成・配置する。
  - 各段: Rectangle（上段ほど幅が狭く、下段ほど幅が広い）を中央揃えで配置
  - 段ラベル: Rectangle自体のTextFrameに直接テキストを配置
  - 右側詳細: TextBox を各段に対応する位置に配置
  - アクセント線: 右側詳細カードの左に縦ラインを配置

■ 対応段数: 3〜7段（pyramids配列の要素数で自動判定）

使い方:
  python fill_pyramid_native.py \\
    --data /home/claude/pyramid_data.json \\
    --template <SKILL_DIR>/assets/pyramid-structure-template.pptx \\
    --output /mnt/user-data/outputs/PyramidStructure_output.pptx
"""

import argparse
import os
import json
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "pyramid-structure-pptx"
_THEME = None

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──────────────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_DIAGRAM_AREA = "Pyramid Diagram Area"
# ────────────────────────────────────────────────────────────

# ── レイアウト定数 (stella 13.33×7.50 ベース; _apply_theme で roleup 用にスケール) ──
PYRAMID_TOP      = Inches(2.38)
PYRAMID_BOTTOM   = Inches(6.27)
PYRAMID_CENTER_X = Inches(2.65)
PYRAMID_MAX_W    = Inches(4.10)
TIER_GAP         = Inches(0.04)

PYRAMID_MIN_W_TABLE = {
    3: Inches(2.20),
    4: Inches(1.80),
    5: Inches(1.50),
    6: Inches(1.20),
    7: Inches(1.05),
}

DETAIL_LEFT      = Inches(5.30)
DETAIL_RIGHT     = Inches(12.80)
DETAIL_VLINE_X   = Inches(5.25)
DETAIL_TOP       = Inches(1.83)
DETAIL_BOTTOM    = Inches(6.82)

# Default colors / fonts (stella). _apply_theme(theme) overrides these for roleup.
TIER_GRAD_TOP    = (47, 84, 150)     # 最上段濃色 (stella: 濃い青)
TIER_GRAD_BOTTOM = (185, 205, 230)   # 最下段薄色 (stella: 薄い青)
COLOR_FONT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)   # 上半分用
COLOR_FONT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # 下半分用
COLOR_DETAIL_BG  = RGBColor(0xF3, 0xF5, 0xF8)
COLOR_DETAIL_FG  = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT     = RGBColor(0x44, 0x72, 0xC4)
FONT_NAME        = "Meiryo UI"
FONT_NAME_TITLE  = "Meiryo UI"

# 段数別フォントサイズ (stella: 16/15/13/11pt 系列、roleup: 14/12/10pt 集合内)
FONT_SIZE_TABLE = {
    3: {"label": 16, "detail_title": 16, "detail_comment": 14},
    4: {"label": 14, "detail_title": 15, "detail_comment": 13},
    5: {"label": 13, "detail_title": 14, "detail_comment": 12},
    6: {"label": 11, "detail_title": 13, "detail_comment": 11},
    7: {"label": 10, "detail_title": 12, "detail_comment": 10},
}
# ────────────────────────────────────────────────────────────


def _apply_theme(theme):
    """roleup の場合、レイアウト・色・フォント・フォントサイズを brand 仕様に上書きする。"""
    global PYRAMID_TOP, PYRAMID_BOTTOM, PYRAMID_CENTER_X, PYRAMID_MAX_W
    global PYRAMID_MIN_W_TABLE, DETAIL_LEFT, DETAIL_RIGHT, DETAIL_VLINE_X
    global DETAIL_TOP, DETAIL_BOTTOM
    global TIER_GRAD_TOP, TIER_GRAD_BOTTOM, COLOR_FONT_LIGHT, COLOR_FONT_DARK
    global COLOR_DETAIL_BG, COLOR_DETAIL_FG, COLOR_ACCENT
    global FONT_NAME, FONT_NAME_TITLE, FONT_SIZE_TABLE
    global _THEME
    _THEME = theme

    if theme.id != "roleup":
        return

    # A4 横 (11.69 × 8.27) 用にレイアウトを再計算
    PYRAMID_TOP      = Inches(2.20)
    PYRAMID_BOTTOM   = Inches(6.95)
    PYRAMID_CENTER_X = Inches(2.40)
    PYRAMID_MAX_W    = Inches(3.70)
    PYRAMID_MIN_W_TABLE = {
        3: Inches(2.00), 4: Inches(1.65), 5: Inches(1.40),
        6: Inches(1.10), 7: Inches(0.95),
    }
    DETAIL_LEFT      = Inches(4.65)
    DETAIL_RIGHT     = Inches(11.30)
    DETAIL_VLINE_X   = Inches(4.60)
    DETAIL_TOP       = Inches(1.65)
    DETAIL_BOTTOM    = Inches(7.20)

    # 色: roleup brand 茶系
    palette = theme._defaults.get("chart_palette", []) if theme._defaults else []
    # 最上段濃色 (label_bar = #7C4C2C)
    top_hex = (palette[0] if palette else "#7C4C2C").lstrip("#")
    TIER_GRAD_TOP = (int(top_hex[0:2], 16), int(top_hex[2:4], 16), int(top_hex[4:6], 16))
    # 最下段薄色 (highlight_other = #CDCECE)
    bottom_hex = (theme._colors.get("highlight_other", "#CDCECE")).lstrip("#")
    TIER_GRAD_BOTTOM = (int(bottom_hex[0:2], 16), int(bottom_hex[2:4], 16), int(bottom_hex[4:6], 16))
    COLOR_FONT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_FONT_DARK = theme.color("text")
    COLOR_DETAIL_BG = theme.color("label_bg") if "label_bg" in theme._colors else theme.color("header_bg")
    COLOR_DETAIL_FG = theme.color("text")
    COLOR_ACCENT = theme.color("subtitle")  # #897141

    FONT_NAME = theme._defaults.get("font_name_ja", "Yu Gothic UI")
    FONT_NAME_TITLE = FONT_NAME

    # roleup の C4 許容集合 [22, 14, 12, 10, 6] pt に整合
    FONT_SIZE_TABLE = {
        3: {"label": 14, "detail_title": 14, "detail_comment": 12},
        4: {"label": 14, "detail_title": 14, "detail_comment": 12},
        5: {"label": 12, "detail_title": 12, "detail_comment": 10},
        6: {"label": 10, "detail_title": 12, "detail_comment": 10},
        7: {"label": 10, "detail_title": 10, "detail_comment": 10},
    }


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_placeholder_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


# ═══════════════════════════════════════════════════════════════
#  色の計算
# ═══════════════════════════════════════════════════════════════

def tier_fill_color(i, total):
    """段のインデックス(0=最上段)に応じて濃→薄のグラデーション"""
    r1, g1, b1 = TIER_GRAD_TOP
    r2, g2, b2 = TIER_GRAD_BOTTOM
    ratio = i / max(total - 1, 1)
    r = int(r1 + (r2 - r1) * ratio)
    g = int(g1 + (g2 - g1) * ratio)
    b = int(b1 + (b2 - b1) * ratio)
    return RGBColor(r, g, b)


def tier_font_color(i, total):
    """上半分は白、下半分は濃色"""
    if i < total * 0.55:
        return COLOR_FONT_LIGHT
    return COLOR_FONT_DARK


# ═══════════════════════════════════════════════════════════════
#  ピラミッド座標計算（四角形ベース）
# ═══════════════════════════════════════════════════════════════

def calc_tier_geometry(n_tiers):
    """
    N段ピラミッドの各段の座標を計算する。
    ピラミッド（左）と詳細（右）で独立した高さを持つ。
    Returns: list of dict with pyramid coords + detail coords
    """
    # ── ピラミッド側（70%高さ、中央寄せ）──
    pyr_total_h = PYRAMID_BOTTOM - PYRAMID_TOP
    pyr_total_gap = TIER_GAP * (n_tiers - 1)
    pyr_net_h = pyr_total_h - pyr_total_gap
    pyr_tier_h = int(pyr_net_h / n_tiers)

    # ── 詳細側（元の高さ）──
    det_total_h = DETAIL_BOTTOM - DETAIL_TOP
    det_total_gap = TIER_GAP * (n_tiers - 1)
    det_net_h = det_total_h - det_total_gap
    det_tier_h = int(det_net_h / n_tiers)

    # ── 段数に応じた最上段の幅 ──
    clamped = min(max(n_tiers, 3), 7)
    min_w = PYRAMID_MIN_W_TABLE[clamped]

    tiers = []
    for i in range(n_tiers):
        # ピラミッド座標
        pyr_y_top = PYRAMID_TOP + i * (pyr_tier_h + TIER_GAP)
        ratio = i / max(n_tiers - 1, 1)
        width = int(min_w + (PYRAMID_MAX_W - min_w) * ratio)
        left = PYRAMID_CENTER_X - int(width / 2)

        # 詳細座標
        det_y_top = DETAIL_TOP + i * (det_tier_h + TIER_GAP)

        tiers.append({
            # ピラミッド側
            "pyr_y_top": pyr_y_top,
            "pyr_height": pyr_tier_h,
            "left": left,
            "width": width,
            # 詳細側
            "det_y_top": det_y_top,
            "det_height": det_tier_h,
        })
    return tiers


# ═══════════════════════════════════════════════════════════════
#  段数に応じたフォントサイズテーブル
# ═══════════════════════════════════════════════════════════════

FONT_SIZE_TABLE = {
    3: {"label": 16, "detail_title": 16, "detail_comment": 14},
    4: {"label": 14, "detail_title": 15, "detail_comment": 13},
    5: {"label": 13, "detail_title": 14, "detail_comment": 12},
    6: {"label": 11, "detail_title": 13, "detail_comment": 11},
    7: {"label": 10, "detail_title": 12, "detail_comment": 10},
}


# ═══════════════════════════════════════════════════════════════
#  Rectangle（ピラミッドの段）生成
# ═══════════════════════════════════════════════════════════════

def add_tier_rectangle(slide, tier_geo, tier_index, n_tiers, title_text):
    """
    ピラミッドの1段を角丸Rectangleとして追加する。
    Rectangle自体のTextFrameにラベルテキストも配置する。
    """
    g = tier_geo
    fs_table = FONT_SIZE_TABLE.get(min(max(n_tiers, 3), 7))
    fill_color = tier_fill_color(tier_index, n_tiers)
    font_color = tier_font_color(tier_index, n_tiers)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        g["left"], g["pyr_y_top"], g["width"], g["pyr_height"]
    )
    shape.name = f"PyramidTier_{tier_index + 1}"

    # 塗りつぶし
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # 枠線なし
    shape.line.fill.background()

    # ── ラベルテキスト ──
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = f"{tier_index + 1}. {title_text}"
    run.font.size = Pt(fs_table["label"])
    run.font.bold = True
    run.font.color.rgb = font_color
    run.font.name = FONT_NAME_TITLE

    return shape


# ═══════════════════════════════════════════════════════════════
#  右側詳細TextBox
# ═══════════════════════════════════════════════════════════════

def add_detail_textbox(slide, tier_geo, tier_index, n_tiers, title_text, comments):
    """右側の詳細テキストボックスを配置する。"""
    g = tier_geo
    fs_table = FONT_SIZE_TABLE.get(min(max(n_tiers, 3), 7))
    detail_width = DETAIL_RIGHT - DETAIL_LEFT

    txBox = slide.shapes.add_textbox(DETAIL_LEFT, g["det_y_top"], detail_width, g["det_height"])
    txBox.name = f"PyramidDetail_{tier_index + 1}"

    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")

    # 背景色
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = COLOR_DETAIL_BG

    # タイトル行
    para = tf.paragraphs[0]
    para.space_after = Pt(2)
    run = para.add_run()
    run.text = f"{tier_index + 1}. {title_text}"
    run.font.size = Pt(fs_table["detail_title"])
    run.font.bold = True
    run.font.color.rgb = COLOR_DETAIL_FG
    run.font.name = FONT_NAME

    # コメント行
    for comment in comments:
        p = tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = f"・{comment}"
        r.font.size = Pt(fs_table["detail_comment"])
        r.font.bold = False
        r.font.color.rgb = COLOR_DETAIL_FG
        r.font.name = FONT_NAME

    return txBox


# ═══════════════════════════════════════════════════════════════
#  アクセント縦線
# ═══════════════════════════════════════════════════════════════

def add_accent_line(slide, tier_geo, accent_color):
    """右側詳細カードの左にアクセント縦線を配置する。"""
    g = tier_geo
    line_top = g["det_y_top"] + Inches(0.04)
    line_height = g["det_height"] - Inches(0.08)

    connector = slide.shapes.add_connector(
        1, DETAIL_VLINE_X, line_top, DETAIL_VLINE_X, line_top + line_height
    )
    connector.name = "AccentLine"
    connector.line.color.rgb = accent_color
    connector.line.width = Pt(3.5)
    return connector


# ═══════════════════════════════════════════════════════════════
#  メイン処理
# ═══════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="N段ピラミッドをPowerPointネイティブオブジェクトで動的生成する"
    )
    parser.add_argument("--data",     required=True, help="pyramid_data.json のパス")
    parser.add_argument("--template", required=False, default=None, help="(任意) テンプレートを明示指定")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "pyramids"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "pyramids", "pyramid_type",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "pyramid-structure")

    pyramids = data.get("pyramids", [])
    n_tiers = len(pyramids)

    print(f"📐 {n_tiers}段ピラミッドを生成 (brand={theme.id})...")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    diag = find_shape(slide, SHAPE_DIAGRAM_AREA)
    if diag:
        slide.shapes._spTree.remove(diag._element)

    # Top text (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        set_placeholder_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
        print(f"  [Top]      {top_text[:60]}{'...' if len(top_text) > 60 else ''}")

    # Subtitle (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme).strip()
    if sub_text:
        set_placeholder_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
        print(f"  [Subtitle] {sub_text[:60]}{'...' if len(sub_text) > 60 else ''}")

    tier_geos = calc_tier_geometry(n_tiers)

    for i, (geo, pyr_data) in enumerate(zip(tier_geos, pyramids)):
        title = pyr_data.get("title", "").strip()
        comments = [c.strip() for c in pyr_data.get("comments", [])]

        add_tier_rectangle(slide, geo, i, n_tiers, title)
        add_detail_textbox(slide, geo, i, n_tiers, title, comments)
        add_accent_line(slide, geo, COLOR_ACCENT)

        print(f"  [Tier {i+1}] {title} ({len(comments)} comments)")

    # Source 出典 (roleup: Source 3 placeholder, stella: 動的 textbox)
    source_text = (data.get("source") or data.get("source_label")
                   or data.get("source_text") or "").strip()
    if source_text:
        if theme.is_source_required():
            src_shape = find_shape(slide, "Source 3")
            if src_shape is not None:
                set_placeholder_text(src_shape, f"出典: {source_text}")
                src_size_pt = int(theme._defaults.get("font_size_source_pt", 6))
                src_color = theme.color("source")
                for para in src_shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(src_size_pt)
                        run.font.color.rgb = src_color
                        run.font.name = FONT_NAME
        else:
            tb = slide.shapes.add_textbox(Inches(0.41), Inches(7.10),
                                          Inches(12.50), Inches(0.30))
            tb.text_frame.text = f"出典: {source_text}"
            for run in tb.text_frame.paragraphs[0].runs:
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)
        print(f"  [Source]   {source_text[:60]}")

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
