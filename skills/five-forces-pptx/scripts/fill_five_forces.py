"""
fill_five_forces.py — Five ForcesデータをPPTXテンプレートに流し込むスクリプト

テンプレート構造（five-forces-template.pptx 確認済み）:
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Rectangle 5        (SHAPE):       業界内の競争    — bullet 2段落
  - Rectangle 15       (SHAPE):       売り手の交渉力  — bullet 3段落
  - Rectangle 17       (SHAPE):       買い手の交渉力  — bullet 3段落
  - Rectangle 18       (SHAPE):       新規参入の脅威  — bullet 2段落
  - Rectangle 20       (SHAPE):       代替品の脅威    — bullet 2段落
  - TextBox 30         (TEXT_BOX):    意味合い        — bullet 3段落

使い方:
  python fill_five_forces.py \
    --data /home/claude/five_forces_data.json \
    --template <path>/assets/five-forces-template.pptx \
    --output /mnt/user-data/outputs/FiveForces_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "five-forces-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ─────────────────────────────────────
SHAPE_MAIN_MESSAGE      = "Title 1"
SHAPE_CHART_TITLE       = "Text Placeholder 2"
SHAPE_COMPETITION       = "Rectangle 5"       # 業界内の競争
SHAPE_SUPPLIER          = "Rectangle 15"      # 売り手の交渉力
SHAPE_BUYER             = "Rectangle 17"      # 買い手の交渉力
SHAPE_NEW_ENTRANT       = "Rectangle 18"      # 新規参入の脅威
SHAPE_SUBSTITUTE        = "Rectangle 20"      # 代替品の脅威
SHAPE_IMPLICATIONS      = "TextBox 30"        # 意味合い
SHAPE_SOURCE            = "Source 3"           # roleup: 下端 Source placeholder
SHAPE_GUIDE_RECT_LEFT   = "正方形/長方形 1"   # roleup: 茶色ガイド矩形 (除去対象)
SHAPE_GUIDE_RECT_RIGHT  = "正方形/長方形 8"   # roleup: 茶色ガイド矩形 (除去対象)
# ──────────────────────────────────────────────────────────


def silent_remove_shape(slide, name):
    """指定 shape をエラー出さず削除 (存在しなければ no-op)。"""
    for sh in list(slide.shapes):
        if sh.name == name:
            sh._element.getparent().remove(sh._element)
            return


def find_shape(slide, name):
    """スライドからShape名で検索"""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_placeholder_text(shape, text):
    """PlaceholderのTextFrameにテキストをセット（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_force_box(slide, shape_name, items):
    """
    Five Forcesボックスのbullet項目を差し替える。
    1段落目（ラベル）はそのまま保持し、2段落目以降のbulletを差し替える。
    項目数がテンプレートと異なる場合は段落を追加・削除して対応する。
    """
    shape = find_shape(slide, shape_name)
    if shape is None:
        return

    tf = shape.text_frame
    paras = list(tf.paragraphs)

    if len(paras) < 2:
        print(f"  ⚠ Shape '{shape_name}' has no bullet paragraphs", file=sys.stderr)
        return

    # 1段落目はラベルなのでスキップ
    label_para = paras[0]
    bullet_paras = paras[1:]

    # テンプレートのbullet段落をベースにする（スタイルコピー元）
    template_bullet_p = bullet_paras[0]._p

    # txBody要素を取得
    txBody = shape.text_frame._txBody

    # 既存のbullet段落をすべて削除
    for bp in bullet_paras:
        txBody.remove(bp._p)

    # 新しいbullet段落を追加
    for i, item_text in enumerate(items):
        new_p = copy.deepcopy(template_bullet_p)
        # テキストを差し替え
        runs = new_p.findall(qn("a:r"))
        if runs:
            # 最初のrunのテキストを設定
            t_elem = runs[0].find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = item_text
            # 余分なrunを削除
            for run in runs[1:]:
                new_p.remove(run)
        else:
            # runがなければ新しく作る
            r_elem = etree.SubElement(new_p, qn("a:r"))
            rpr = new_p.find(qn("a:pPr"))
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = item_text

        # endParaRPrがあれば末尾に移動（最後の段落用）
        end_rpr = new_p.find(qn("a:endParaRPr"))
        if end_rpr is not None and i < len(items) - 1:
            new_p.remove(end_rpr)

        txBody.append(new_p)
        print(f"  [{shape_name}] • {item_text[:50]}{'...' if len(item_text) > 50 else ''}")


def fill_implications(slide, items):
    """
    意味合いボックスのbullet項目を差し替える。
    1段落目（「意味合い」ラベル）はそのまま保持し、2段落目以降を差し替える。
    """
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return

    tf = shape.text_frame
    paras = list(tf.paragraphs)

    if len(paras) < 2:
        print(f"  ⚠ Shape '{SHAPE_IMPLICATIONS}' has no bullet paragraphs", file=sys.stderr)
        return

    # 1段落目はラベルなのでスキップ
    bullet_paras = paras[1:]
    template_bullet_p = bullet_paras[0]._p
    txBody = tf._txBody

    # 既存のbullet段落をすべて削除
    for bp in bullet_paras:
        txBody.remove(bp._p)

    # 新しいbullet段落を追加
    for i, item_text in enumerate(items):
        new_p = copy.deepcopy(template_bullet_p)
        runs = new_p.findall(qn("a:r"))
        if runs:
            t_elem = runs[0].find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = item_text
            for run in runs[1:]:
                new_p.remove(run)
        else:
            r_elem = etree.SubElement(new_p, qn("a:r"))
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = item_text

        end_rpr = new_p.find(qn("a:endParaRPr"))
        if end_rpr is not None and i < len(items) - 1:
            new_p.remove(end_rpr)

        txBody.append(new_p)
        print(f"  [意味合い] • {item_text[:50]}{'...' if len(item_text) > 50 else ''}")


def main():
    parser = argparse.ArgumentParser(description="Five ForcesデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="five_forces_data.json のパス")
    parser.add_argument("--template", required=False, default=None, help="(任意) テンプレートを明示指定")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=[
            "main_message",
            "industry_competition", "new_entrants", "substitutes",
            "supplier_power", "buyer_power",
        ],
        allowed_top=[
            "main_message", "chart_title", "source", "implications",
            "industry_competition", "new_entrants", "substitutes",
            "supplier_power", "buyer_power",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "five-forces")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # roleup: 茶色ガイド矩形を除去 (C1 担保)
    silent_remove_shape(slide, SHAPE_GUIDE_RECT_LEFT)
    silent_remove_shape(slide, SHAPE_GUIDE_RECT_RIGHT)

    # ── Top text (stella: main_message / roleup: chart_title) ──
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(shape, top_text)
        print(f"  [Top]      {top_text[:60]}{'...' if len(top_text) > 60 else ''}")

    # ── Subtitle (stella: chart_title / roleup: main_message) ──
    sub_text = resolve_subtitle_text(data, theme).strip()
    if sub_text:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(shape, sub_text)
        print(f"  [Subtitle] {sub_text[:60]}{'...' if len(sub_text) > 60 else ''}")

    # ── Five Forces boxes ─────────────────────────────────
    force_mapping = [
        ("industry_competition", SHAPE_COMPETITION,  "業界内の競争"),
        ("new_entrants",         SHAPE_NEW_ENTRANT,  "新規参入の脅威"),
        ("substitutes",          SHAPE_SUBSTITUTE,   "代替品の脅威"),
        ("supplier_power",       SHAPE_SUPPLIER,     "売り手の交渉力"),
        ("buyer_power",          SHAPE_BUYER,        "買い手の交渉力"),
    ]

    for key, shape_name, label in force_mapping:
        items = data.get(key, [])
        if not items:
            print(f"  ⚠ {label}（{key}）のデータがありません", file=sys.stderr)
            continue
        fill_force_box(slide, shape_name, items)

    # ── 意味合い ──────────────────────────────────────────
    implications = data.get("implications", [])
    if implications:
        fill_implications(slide, implications)
    else:
        print("  ⚠ implications が未設定です", file=sys.stderr)

    # ── Source (出典) ─────────────────────────────────────
    source_text = (data.get("source") or data.get("source_label")
                   or data.get("source_text") or "").strip()
    if source_text:
        if theme.is_source_required():
            # roleup: Source 3 placeholder に流し込む
            src_shape = find_shape(slide, SHAPE_SOURCE)
            if src_shape is not None:
                set_placeholder_text(src_shape, f"出典: {source_text}")
                # フォントサイズ・色を保証 (C4 6pt + 茶系)
                src_size_pt = int(theme._defaults.get("font_size_source_pt", 6))
                src_color = theme.color("source")
                src_font = theme._defaults.get("font_name_ja", "Yu Gothic UI")
                for para in src_shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(src_size_pt)
                        run.font.color.rgb = src_color
                        run.font.name = src_font
        else:
            # stella: 動的 textbox で下端に追加
            tb = slide.shapes.add_textbox(Inches(0.41), Inches(7.10),
                                          Inches(12.50), Inches(0.30))
            tb.text_frame.text = f"出典: {source_text}"
            for run in tb.text_frame.paragraphs[0].runs:
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)
        print(f"  [Source]   {source_text[:60]}")

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
