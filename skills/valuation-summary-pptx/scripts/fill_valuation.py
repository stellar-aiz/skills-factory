"""
fill_valuation.py — M&A DDバリュエーション・財務サマリーをPPTXに生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

3つのchart_type:
  - football_field:    手法別バリュエーションレンジ比較
  - equity_bridge:     EV→株式価値のウォーターフォール
  - financial_summary: 主要財務指標テーブル

Usage:
  python fill_valuation.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/valuation_data.json \\
    --output {{OUTPUT_DIR}}/Valuation_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "valuation-summary-pptx"

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names (template-structure invariants) ──
SHAPE_MM = "Title 1"
SHAPE_CT = "Text Placeholder 2"

# Defaults (stella). Reassigned in _apply_theme(theme).
SHAPE_SOURCE = "Source"
MARGIN_L = 370800
SLIDE_W = 12192000
SLIDE_H = 6858000
CONTENT_W = SLIDE_W - 2 * MARGIN_L
CHART_TOP = 1400000
CHART_BOTTOM = 6400000
CHART_H = CHART_BOTTOM - CHART_TOP
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.00)
SOURCE_H = Inches(0.30)

# Stella defaults; reassigned in _apply_theme(theme).
TEXT_HEX = "1A1A1A"
SUB_TEXT_HEX = "666666"
GRID_HEX = "D0D0D0"
ZERO_LINE_HEX = "AAAAAA"
HEADER_FONT_HEX = "FFFFFF"
ROW_BANDED_HEX = "F5F5F5"
ROW_DIVIDER_HEX = "DDDDDD"
TOTAL_DIVIDER_HEX = "333333"
POSITIVE_BAR_HEX = "2E7D32"  # green
NEGATIVE_BAR_HEX = "C62828"  # red
ACCENT_HEX = "1A3C6E"        # fallback; overridden in main()
SOURCE_HEX = "666666"
FONT_NAME = "Meiryo"

SECTION_HEADER_PT = 11
DATA_PT = 10
LABEL_PT = 12
TICK_PT = 9
RANGE_LABEL_PT = 10
BAR_LABEL_PT = 10
BRIDGE_VALUE_PT = 10
BRIDGE_NAME_PT = 10
TABLE_HDR_PT = 11
TABLE_BODY_PT = 10
SOURCE_FONT_PT = 10

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global MARGIN_L, SLIDE_W, SLIDE_H, CONTENT_W
    global CHART_TOP, CHART_BOTTOM, CHART_H
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global TEXT_HEX, SUB_TEXT_HEX, GRID_HEX, ZERO_LINE_HEX
    global HEADER_FONT_HEX, ROW_BANDED_HEX, ROW_DIVIDER_HEX, TOTAL_DIVIDER_HEX
    global POSITIVE_BAR_HEX, NEGATIVE_BAR_HEX, SOURCE_HEX, FONT_NAME
    global SECTION_HEADER_PT, DATA_PT, LABEL_PT, TICK_PT, RANGE_LABEL_PT
    global BAR_LABEL_PT, BRIDGE_VALUE_PT, BRIDGE_NAME_PT
    global TABLE_HDR_PT, TABLE_BODY_PT, SOURCE_FONT_PT

    _THEME = theme
    FONT_NAME = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    SOURCE_HEX = theme.hex_no_hash("source")

    # Slide size (EMU)
    SLIDE_W = int(theme.slide_w)
    SLIDE_H = int(theme.slide_h)
    # Content margin (EMU)
    MARGIN_L = theme.layout("margin_l_in")
    CONTENT_W = SLIDE_W - 2 * MARGIN_L
    CHART_TOP = theme.layout("chart_top_in")
    CHART_BOTTOM = theme.layout("chart_bottom_in")
    CHART_H = CHART_BOTTOM - CHART_TOP
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        # stella V1 visual defaults
        SUB_TEXT_HEX = "666666"
        GRID_HEX = "D0D0D0"
        ZERO_LINE_HEX = "AAAAAA"
        HEADER_FONT_HEX = "FFFFFF"
        ROW_BANDED_HEX = "F5F5F5"
        ROW_DIVIDER_HEX = "DDDDDD"
        TOTAL_DIVIDER_HEX = "333333"
        POSITIVE_BAR_HEX = "2E7D32"
        NEGATIVE_BAR_HEX = "C62828"
        # Font sizes (raw OOXML 100×pt values used in the original code)
        SECTION_HEADER_PT = 11
        DATA_PT = 10
        LABEL_PT = 12
        TICK_PT = 9
        RANGE_LABEL_PT = 10
        BAR_LABEL_PT = 9
        BRIDGE_VALUE_PT = 10
        BRIDGE_NAME_PT = 10
        TABLE_HDR_PT = 11
        TABLE_BODY_PT = 10
        SOURCE_FONT_PT = 10
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        SUB_TEXT_HEX = theme.hex_no_hash("source")            # #3E3A39
        GRID_HEX = theme.hex_no_hash("highlight_other")        # #CDCECE
        ZERO_LINE_HEX = theme.hex_no_hash("highlight_other")   # #CDCECE
        HEADER_FONT_HEX = "FFFFFF"
        ROW_BANDED_HEX = theme.hex_no_hash("label_bg")         # #F2E8DD
        ROW_DIVIDER_HEX = theme.hex_no_hash("highlight_other") # #CDCECE
        TOTAL_DIVIDER_HEX = TEXT_HEX
        POSITIVE_BAR_HEX = theme.hex_no_hash("accent_revenue_bar")    # #7C4C2C (positive: brand-primary)
        NEGATIVE_BAR_HEX = theme.hex_no_hash("accent_op_margin_line")  # #604C3F (negative: 茶系トーンを差別化)
        # Font sizes — C4 厳守 (10pt 統一が基本)
        body_pt = theme.pt_value("font_size_body_pt")          # 10
        sub_pt = theme.pt_value("font_size_subtitle_pt")       # 12
        SECTION_HEADER_PT = sub_pt
        DATA_PT = body_pt
        LABEL_PT = sub_pt
        TICK_PT = body_pt          # was 9 → 10
        RANGE_LABEL_PT = body_pt   # was 10 → 10
        BAR_LABEL_PT = body_pt     # was 9 → 10
        BRIDGE_VALUE_PT = body_pt
        BRIDGE_NAME_PT = body_pt
        TABLE_HDR_PT = sub_pt      # 11→12 (sub_pt)
        TABLE_BODY_PT = body_pt
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6


def _silent_remove_shape(slide, name):
    for shape in list(slide.shapes):
        if shape.name == name:
            slide.shapes._spTree.remove(shape._element)
            return True
    return False


def find_shape(sl, nm, warn=False):
    for s in sl.shapes:
        if s.name == nm:
            return s
    if warn:
        print(f"  ⚠ WARNING: Shape '{nm}' not found", file=sys.stderr)
    return None


def set_ph(sh, txt):
    if sh is None:
        return
    p = sh.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = txt
        for r in p.runs[1:]:
            r.text = ""
    else:
        r = etree.SubElement(p._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t = etree.SubElement(r, qn("a:t"))
        t.text = txt


def extract_accent2(path):
    """Read accent2 color hex from the template's clrScheme."""
    try:
        prs = Presentation(path)
        for rel in prs.slide_masters[0].part.rels.values():
            if "theme" in rel.reltype:
                te = etree.fromstring(rel.target_part.blob)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                cs = te.find(".//a:clrScheme", ns)
                if cs is not None:
                    a2 = cs.find("a:accent2", ns)
                    if a2 is not None:
                        sr = a2.find("a:srgbClr", ns)
                        if sr is not None:
                            return sr.get("val")
    except Exception:
        pass
    return "1A3C6E"


def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def lighten(c, f):
    r, g, b = hex_to_rgb(c)
    return f"{int(r+(255-r)*f):02x}{int(g+(255-g)*f):02x}{int(b+(255-b)*f):02x}"


def add_rect(sl, l, t, w, h, fill=None, border=None, border_w=6350):
    sp = etree.SubElement(sl.shapes._spTree, qn("p:sp"))
    nv = etree.SubElement(sp, qn("p:nvSpPr"))
    etree.SubElement(nv, qn("p:cNvPr"), attrib={"id": str(id(sp) % 99999 + 100), "name": "R"})
    etree.SubElement(nv, qn("p:cNvSpPr"))
    etree.SubElement(nv, qn("p:nvPr"))
    spr = etree.SubElement(sp, qn("p:spPr"))
    xf = etree.SubElement(spr, qn("a:xfrm"))
    etree.SubElement(xf, qn("a:off"), attrib={"x": str(int(l)), "y": str(int(t))})
    etree.SubElement(xf, qn("a:ext"), attrib={"cx": str(max(1, int(w))), "cy": str(max(1, int(h)))})
    pg = etree.SubElement(spr, qn("a:prstGeom"), attrib={"prst": "rect"})
    etree.SubElement(pg, qn("a:avLst"))
    if fill:
        sf = etree.SubElement(spr, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": fill})
    else:
        etree.SubElement(spr, qn("a:noFill"))
    ln = etree.SubElement(spr, qn("a:ln"))
    if border:
        ln.set("w", str(border_w))
        sf2 = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf2, qn("a:srgbClr"), attrib={"val": border})
    else:
        etree.SubElement(ln, qn("a:noFill"))
    tb = etree.SubElement(sp, qn("p:txBody"))
    etree.SubElement(tb, qn("a:bodyPr"))
    etree.SubElement(tb, qn("a:lstStyle"))
    etree.SubElement(tb, qn("a:p"))
    return sp


def add_txt(sl, l, t, w, h, text, sz_pt=11, bold=False, color=None, align="l",
            shrink=False, wrap="square"):
    """Add a textbox. sz_pt is in pt (int); converted to 100×pt internally."""
    if color is None:
        color = TEXT_HEX
    sp = etree.SubElement(sl.shapes._spTree, qn("p:sp"))
    nv = etree.SubElement(sp, qn("p:nvSpPr"))
    etree.SubElement(nv, qn("p:cNvPr"), attrib={"id": str(id(sp) % 99999 + 200), "name": "T"})
    etree.SubElement(nv, qn("p:cNvSpPr"), attrib={"txBox": "1"})
    etree.SubElement(nv, qn("p:nvPr"))
    spr = etree.SubElement(sp, qn("p:spPr"))
    xf = etree.SubElement(spr, qn("a:xfrm"))
    etree.SubElement(xf, qn("a:off"), attrib={"x": str(int(l)), "y": str(int(t))})
    etree.SubElement(xf, qn("a:ext"), attrib={"cx": str(max(1, int(w))), "cy": str(max(1, int(h)))})
    pg = etree.SubElement(spr, qn("a:prstGeom"), attrib={"prst": "rect"})
    etree.SubElement(pg, qn("a:avLst"))
    etree.SubElement(spr, qn("a:noFill"))
    ln = etree.SubElement(spr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))
    tb = etree.SubElement(sp, qn("p:txBody"))
    bp = etree.SubElement(tb, qn("a:bodyPr"), attrib={
        "wrap": wrap, "lIns": "36000", "rIns": "36000",
        "tIns": "0", "bIns": "0", "anchor": "ctr"
    })
    if shrink:
        etree.SubElement(bp, qn("a:normAutofit"))
    etree.SubElement(tb, qn("a:lstStyle"))
    p = etree.SubElement(tb, qn("a:p"))
    etree.SubElement(p, qn("a:pPr"), attrib={"algn": align})
    r = etree.SubElement(p, qn("a:r"))
    rp = {"kumimoji": "1", "lang": "en-GB", "sz": str(int(sz_pt) * 100)}
    if bold:
        rp["b"] = "1"
    rpr = etree.SubElement(r, qn("a:rPr"), attrib=rp)
    sf = etree.SubElement(rpr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color})
    etree.SubElement(rpr, qn("a:latin"), attrib={"typeface": FONT_NAME})
    etree.SubElement(rpr, qn("a:ea"), attrib={"typeface": FONT_NAME})
    te = etree.SubElement(r, qn("a:t"))
    te.text = text


def add_line(sl, x1, y1, x2, y2, color=None, w=9525, dash=None):
    if color is None:
        color = SUB_TEXT_HEX
    sp = etree.SubElement(sl.shapes._spTree, qn("p:cxnSp"))
    nv = etree.SubElement(sp, qn("p:nvCxnSpPr"))
    etree.SubElement(nv, qn("p:cNvPr"), attrib={"id": str(id(sp) % 99999 + 300), "name": "L"})
    etree.SubElement(nv, qn("p:cNvCxnSpPr"))
    etree.SubElement(nv, qn("p:nvPr"))
    spr = etree.SubElement(sp, qn("p:spPr"))
    xf = etree.SubElement(spr, qn("a:xfrm"))
    etree.SubElement(xf, qn("a:off"), attrib={"x": str(int(min(x1, x2))), "y": str(int(min(y1, y2)))})
    etree.SubElement(xf, qn("a:ext"), attrib={
        "cx": str(max(1, abs(int(x2 - x1)))), "cy": str(max(1, abs(int(y2 - y1))))
    })
    pg = etree.SubElement(spr, qn("a:prstGeom"), attrib={"prst": "line"})
    etree.SubElement(pg, qn("a:avLst"))
    ln = etree.SubElement(spr, qn("a:ln"), attrib={"w": str(w)})
    sf = etree.SubElement(ln, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color})
    if dash:
        etree.SubElement(ln, qn("a:prstDash"), attrib={"val": dash})


def fmt_num(v, unit=""):
    """数値を見やすくフォーマット"""
    if abs(v) >= 1e9:
        return f"{v/1e9:,.1f}B{unit}"
    if abs(v) >= 1e6:
        return f"{v/1e6:,.0f}M{unit}"
    if abs(v) >= 1e4:
        return f"{v:,.0f}{unit}"
    return f"{v:,.1f}{unit}"


# ══════════════════════════════════════════════════════════
# 1. FOOTBALL FIELD CHART
# ══════════════════════════════════════════════════════════
def build_football_field(sl, data, accent_hex):
    ff = data["football_field"]
    methods = ff["methods"]
    unit = ff.get("unit", "")
    n = len(methods)

    all_vals = []
    for m in methods:
        all_vals.extend([m["low"], m["high"]])
    rec = ff.get("recommended_range", None)
    ax_min = ff.get("axis_min", min(all_vals) * 0.85)
    ax_max = ff.get("axis_max", max(all_vals) * 1.15)
    ax_range = ax_max - ax_min if ax_max != ax_min else 1

    label_w = 3400000
    chart_l = MARGIN_L + label_w
    chart_w = CONTENT_W - label_w
    row_h = min(550000, CHART_H // (n + 2))
    bar_h = int(row_h * 0.55)
    chart_area_top = CHART_TOP + row_h

    def val_to_x(v):
        return chart_l + (v - ax_min) / ax_range * chart_w

    if rec:
        rx1 = val_to_x(rec["low"])
        rx2 = val_to_x(rec["high"])
        add_rect(sl, rx1, chart_area_top - row_h // 2,
                 rx2 - rx1, row_h * n + row_h,
                 fill=lighten(accent_hex, 0.85))
        add_txt(sl, rx1, CHART_TOP, rx2 - rx1, int(row_h * 0.7),
                f"推定レンジ: {fmt_num(rec['low'])}{unit} - {fmt_num(rec['high'])}{unit}",
                sz_pt=RANGE_LABEL_PT, bold=True, color=accent_hex, align="ctr")

    n_ticks = 6
    for i in range(n_ticks + 1):
        v = ax_min + ax_range * i / n_ticks
        x = val_to_x(v)
        add_line(sl, x, chart_area_top - row_h // 3,
                 x, chart_area_top + row_h * n, color=GRID_HEX, w=6350, dash="dash")
        add_txt(sl, x - 500000, chart_area_top + row_h * n + 30000,
                1000000, 280000, fmt_num(v) + unit,
                sz_pt=TICK_PT, color=SUB_TEXT_HEX, align="ctr")

    for i, m in enumerate(methods):
        y = chart_area_top + i * row_h
        bar_y = y + (row_h - bar_h) // 2

        add_txt(sl, MARGIN_L, y, label_w - 100000, row_h,
                m["name"], sz_pt=LABEL_PT, bold=True, color=TEXT_HEX, align="r",
                shrink=True, wrap="none")

        x_low = val_to_x(m["low"])
        x_high = val_to_x(m["high"])
        bar_color = accent_hex if i % 2 == 0 else lighten(accent_hex, 0.3)
        add_rect(sl, x_low, bar_y, x_high - x_low, bar_h, fill=bar_color)

        if "mid" in m:
            x_mid = val_to_x(m["mid"])
            add_line(sl, x_mid, bar_y - 30000, x_mid, bar_y + bar_h + 30000,
                     color="FFFFFF", w=19050)

        val_lbl_w = 800000
        bar_actual_w = x_high - x_low
        if bar_actual_w > val_lbl_w * 2:
            add_txt(sl, x_low + 30000, bar_y, val_lbl_w, bar_h,
                    fmt_num(m["low"]), sz_pt=BAR_LABEL_PT, color="FFFFFF", align="l")
            add_txt(sl, x_high - val_lbl_w - 30000, bar_y, val_lbl_w, bar_h,
                    fmt_num(m["high"]), sz_pt=BAR_LABEL_PT, color="FFFFFF", align="r")
        else:
            add_txt(sl, x_low - val_lbl_w - 30000, bar_y, val_lbl_w, bar_h,
                    fmt_num(m["low"]), sz_pt=BAR_LABEL_PT, color=TEXT_HEX, align="r")
            add_txt(sl, x_high + 30000, bar_y, val_lbl_w, bar_h,
                    fmt_num(m["high"]), sz_pt=BAR_LABEL_PT, color=TEXT_HEX, align="l")

    add_txt(sl, chart_l + chart_w - 1500000,
            chart_area_top + row_h * n + 300000,
            1500000, 250000, f"（単位: {unit}）" if unit else "",
            sz_pt=TICK_PT, color=SUB_TEXT_HEX, align="r")

    print(f"  [Football Field] {n} methods, range {fmt_num(ax_min)}-{fmt_num(ax_max)}")


# ══════════════════════════════════════════════════════════
# 2. EQUITY VALUE BRIDGE (Waterfall)
# ══════════════════════════════════════════════════════════
def build_equity_bridge(sl, data, accent_hex):
    eb = data["equity_bridge"]
    items = eb["items"]
    unit = eb.get("unit", "")
    n = len(items)

    cumulative = []
    cum = 0
    for it in items:
        if it["type"] == "start":
            cum = it["value"]
            cumulative.append({"base": 0, "top": cum, "val": cum, **it})
        elif it["type"] == "total":
            cumulative.append({"base": 0, "top": it["value"], "val": it["value"], **it})
        else:
            old_cum = cum
            cum += it["value"]
            if it["value"] >= 0:
                cumulative.append({"base": old_cum, "top": cum, "val": it["value"], **it})
            else:
                cumulative.append({"base": cum, "top": old_cum, "val": it["value"], **it})

    max_val = max(c["top"] for c in cumulative) * 1.15
    min_val = min(min(c["base"] for c in cumulative), 0)

    val_range = max_val - min_val
    if val_range == 0:
        val_range = 1

    bar_area_top = CHART_TOP + 400000
    bar_area_h = CHART_H - 800000
    bar_gap = 120000
    total_bar_w = CONTENT_W - 200000
    bar_w = (total_bar_w - bar_gap * (n - 1)) // n
    bar_w = min(bar_w, 1600000)
    total_used = bar_w * n + bar_gap * (n - 1)
    start_x = MARGIN_L + (CONTENT_W - total_used) // 2

    def val_to_y(v):
        return bar_area_top + bar_area_h - (v - min_val) / val_range * bar_area_h

    y_zero = val_to_y(0)
    add_line(sl, start_x - 100000, y_zero,
             start_x + total_used + 100000, y_zero, color=ZERO_LINE_HEX, w=9525)

    for i, c in enumerate(cumulative):
        x = start_x + i * (bar_w + bar_gap)
        y_top = val_to_y(c["top"])
        y_base = val_to_y(c["base"])
        h = abs(y_base - y_top)
        y = min(y_top, y_base)

        if c["type"] in ("start", "total"):
            fill = accent_hex
        elif c["val"] >= 0:
            fill = POSITIVE_BAR_HEX
        else:
            fill = NEGATIVE_BAR_HEX

        add_rect(sl, x, y, bar_w, max(h, Emu(Pt(2))), fill=fill)

        val_text = fmt_num(c["val"]) + unit
        if c["type"] in ("start", "total"):
            val_text = fmt_num(c["top"]) + unit

        label_y = y - 280000 if c["val"] >= 0 or c["type"] in ("start", "total") else y + h + 30000
        add_txt(sl, x, label_y, bar_w, 260000,
                val_text, sz_pt=BRIDGE_VALUE_PT, bold=True, color=TEXT_HEX, align="ctr")

        add_txt(sl, x - 100000, bar_area_top + bar_area_h + 80000,
                bar_w + 200000, 350000,
                c["name"], sz_pt=BRIDGE_NAME_PT, bold=False, color=TEXT_HEX, align="ctr")

        if i < n - 1 and c["type"] != "total":
            next_x = start_x + (i + 1) * (bar_w + bar_gap)
            conn_y = val_to_y(cumulative[i]["top"] if c["val"] >= 0 or c["type"] == "start" else cumulative[i]["base"])
            if cumulative[i + 1]["type"] == "total":
                pass
            else:
                add_line(sl, x + bar_w, conn_y, next_x, conn_y,
                         color=GRID_HEX, w=6350, dash="dash")

    print(f"  [Equity Bridge] {n} items")


# ══════════════════════════════════════════════════════════
# 3. FINANCIAL SUMMARY TABLE
# ══════════════════════════════════════════════════════════
def build_financial_summary(sl, data, accent_hex):
    fs = data["financial_summary"]
    periods = fs["periods"]
    metrics = fs["metrics"]
    n_cols = len(periods) + 1
    n_rows = len(metrics) + 1

    has_cagr = any("cagr" in m for m in metrics)
    if has_cagr:
        n_cols += 1

    tbl_l = MARGIN_L
    tbl_t = CHART_TOP
    tbl_w = CONTENT_W
    row_h = min(420000, (CHART_H - 100000) // n_rows)
    name_col_w = int(tbl_w * 0.22)
    data_col_w = (tbl_w - name_col_w) // (n_cols - 1)

    add_rect(sl, tbl_l, tbl_t, tbl_w, row_h, fill=accent_hex)
    add_txt(sl, tbl_l, tbl_t, name_col_w, row_h,
            "項目", sz_pt=TABLE_HDR_PT, bold=True, color=HEADER_FONT_HEX, align="ctr")
    for j, per in enumerate(periods):
        x = tbl_l + name_col_w + j * data_col_w
        add_txt(sl, x, tbl_t, data_col_w, row_h,
                per, sz_pt=TABLE_HDR_PT, bold=True, color=HEADER_FONT_HEX, align="ctr")
    if has_cagr:
        x = tbl_l + name_col_w + len(periods) * data_col_w
        add_txt(sl, x, tbl_t, data_col_w, row_h,
                "CAGR", sz_pt=TABLE_HDR_PT, bold=True, color=HEADER_FONT_HEX, align="ctr")

    add_line(sl, tbl_l, tbl_t + row_h, tbl_l + tbl_w, tbl_t + row_h,
             color=TOTAL_DIVIDER_HEX, w=15875)

    for i, m in enumerate(metrics):
        y = tbl_t + (i + 1) * row_h
        is_section = m.get("is_section", False)
        is_total = m.get("is_total", False)

        if is_section:
            add_rect(sl, tbl_l, y, tbl_w, row_h, fill=lighten(accent_hex, 0.88))
        elif i % 2 == 1:
            add_rect(sl, tbl_l, y, tbl_w, row_h, fill=ROW_BANDED_HEX)

        name_color = accent_hex if is_section else TEXT_HEX
        add_txt(sl, tbl_l, y, name_col_w, row_h,
                m["name"], sz_pt=TABLE_BODY_PT,
                bold=is_section or is_total,
                color=name_color, align="l")

        values = m.get("values", [])
        for j, v in enumerate(values):
            x = tbl_l + name_col_w + j * data_col_w
            if v is None or v == "":
                continue
            txt = str(v) if isinstance(v, str) else fmt_num(v, m.get("unit_suffix", ""))
            add_txt(sl, x, y, data_col_w, row_h,
                    txt, sz_pt=TABLE_BODY_PT, bold=is_total,
                    color=TEXT_HEX, align="r")

        if has_cagr and "cagr" in m:
            x = tbl_l + name_col_w + len(periods) * data_col_w
            add_txt(sl, x, y, data_col_w, row_h,
                    m["cagr"], sz_pt=TABLE_BODY_PT, bold=False, color=TEXT_HEX, align="r")

        add_line(sl, tbl_l, y + row_h, tbl_l + tbl_w, y + row_h,
                 color=ROW_DIVIDER_HEX, w=6350)
        if is_total:
            add_line(sl, tbl_l, y, tbl_l + tbl_w, y, color=TOTAL_DIVIDER_HEX, w=12700)

    print(f"  [Financial Summary] {n_rows-1} metrics x {len(periods)} periods")


def add_source_label(sl, text):
    """出典: roleup なら Source 3 placeholder、stella なら textbox。"""
    src_shape = find_shape(sl, SHAPE_SOURCE, warn=False)
    if src_shape is not None and src_shape.has_text_frame:
        tf = src_shape.text_frame
        tf.word_wrap = True
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(SOURCE_FONT_PT)
        rgb_int = int(SOURCE_HEX, 16)
        run.font.color.rgb = RGBColor(
            (rgb_int >> 16) & 0xFF, (rgb_int >> 8) & 0xFF, rgb_int & 0xFF
        )
        run.font.name = FONT_NAME
        print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {text[:50]}")
        return

    # Fallback: dynamic textbox (stella with no Source placeholder)
    add_txt(sl, SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H, text,
            sz_pt=SOURCE_FONT_PT, color=SOURCE_HEX, align="l")
    print(f"  ✓ 出典 (textbox): {text[:50]}")


# ══════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand.",
    )
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    # chart_type で 3 種類 (football_field / equity_bridge / financial_summary) のいずれかを 1 つ持つ
    validate_fill_input(
        data,
        required_top=["main_message", "chart_type"],
        allowed_top=[
            "main_message", "chart_title", "source", "chart_type",
            "football_field", "equity_bridge", "financial_summary",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "valuation")

    print(f"=== バリュエーション・財務サマリー (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    # chart_title のデフォルトを data に埋めて、brand 別の top/subtitle 解決で
    # どちらの placeholder field に書かれても落ちない状態にする。
    chart_type = data.get("chart_type", "football_field")
    if not data.get("chart_title"):
        data["chart_title"] = {
            "football_field": "バリュエーション・レンジ分析",
            "equity_bridge": "株式価値ブリッジ",
            "financial_summary": "財務サマリー",
        }.get(chart_type, "バリュエーション・財務サマリー")

    prs = Presentation(template_path)
    sl = prs.slides[0]

    # roleup: brand chart_palette[0] を accent として優先、
    # stella: テンプレート theme accent2 を優先。
    if theme.id != "stellar_aiz":
        accent_hex = theme.hex_no_hash("accent_revenue_bar")  # #7C4C2C
    else:
        accent_hex = extract_accent2(template_path)
    print(f"  Accent: {accent_hex}")

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(sl, "正方形/長方形 1")
    _silent_remove_shape(sl, "正方形/長方形 8")

    # Top placeholder (brand-aware)
    top_text = resolve_top_text(data, theme)
    set_ph(find_shape(sl, SHAPE_MM), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    sub_text = resolve_subtitle_text(data, theme)
    set_ph(find_shape(sl, SHAPE_CT), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:50]}")

    if chart_type == "football_field":
        build_football_field(sl, data, accent_hex)
    elif chart_type == "equity_bridge":
        build_equity_bridge(sl, data, accent_hex)
    elif chart_type == "financial_summary":
        build_financial_summary(sl, data, accent_hex)
    else:
        print(f"  ERROR: Unknown chart_type '{chart_type}'", file=sys.stderr)
        sys.exit(1)

    # Source
    source = data.get("source", "")
    if source:
        body = source if source.startswith("出典") else f"出典：{source}"
        add_source_label(sl, body)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
