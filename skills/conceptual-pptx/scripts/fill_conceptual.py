"""
fill_conceptual.py — コンセプトデータをPPTXテンプレートに流し込むスクリプト

テンプレート構造（確認済み）:

■ Conceptual3.pptx（3コンセプト版）— 詳細2行
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Oval 4             (AUTO_SHAPE):  Concept1 ラベル（楕円内）
  - Oval 6             (AUTO_SHAPE):  Concept2 ラベル（楕円内）
  - Oval 7             (AUTO_SHAPE):  Concept3 ラベル（楕円内）
  - TextBox 8          (TEXT_BOX):    Concept1 名前(Para[0]) + 詳細1(Para[1]) + 詳細2(Para[2])
  - TextBox 9          (TEXT_BOX):    Concept2 名前(Para[0]) + 詳細1(Para[1]) + 詳細2(Para[2])
  - TextBox 10         (TEXT_BOX):    Concept3 名前(Para[0]) + 詳細1(Para[1]) + 詳細2(Para[2])

■ Conceptual5.pptx（5コンセプト版）— 詳細1行
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Oval 4             (AUTO_SHAPE):  Concept1 ラベル（楕円内）
  - Oval 6             (AUTO_SHAPE):  Concept2 ラベル（楕円内）
  - Oval 7             (AUTO_SHAPE):  Concept3 ラベル（楕円内）
  - Oval 11            (AUTO_SHAPE):  Concept4 ラベル（楕円内）
  - Oval 12            (AUTO_SHAPE):  Concept5 ラベル（楕円内）
  - TextBox 8          (TEXT_BOX):    Concept1 名前(Para[0]) + 詳細(Para[1])
  - TextBox 9          (TEXT_BOX):    Concept2 名前(Para[0]) + 詳細(Para[1])
  - TextBox 10         (TEXT_BOX):    Concept3 名前(Para[0]) + 詳細(Para[1])
  - TextBox 13         (TEXT_BOX):    Concept4 名前(Para[0]) + 詳細(Para[1])
  - TextBox 14         (TEXT_BOX):    Concept5 名前(Para[0]) + 詳細(Para[1])

使い方:
  python fill_conceptual.py \
    --data /home/claude/conceptual_data.json \
    --template3 <skill_path>/assets/Conceptual3.pptx \
    --template5 <skill_path>/assets/Conceptual5.pptx \
    --output /mnt/user-data/outputs/Conceptual_output.pptx
"""

import argparse
import json
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング（確認済み）──────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

# 3コンセプト版: 楕円ラベル
OVAL_NAMES_3 = ["Oval 4", "Oval 6", "Oval 7"]
# 3コンセプト版: テキストボックス（名前＋詳細2行）
TEXTBOX_NAMES_3 = ["TextBox 8", "TextBox 9", "TextBox 10"]

# 5コンセプト版: 楕円ラベル
OVAL_NAMES_5 = ["Oval 4", "Oval 6", "Oval 7", "Oval 11", "Oval 12"]
# 5コンセプト版: テキストボックス（名前＋詳細1行）
TEXTBOX_NAMES_5 = ["TextBox 8", "TextBox 9", "TextBox 10", "TextBox 13", "TextBox 14"]
# ────────────────────────────────────────────────────────────


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_text_in_shape(shape, text):
    """Shape内のテキストを上書き（既存スタイル保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def set_para_text(para, text, fallback_attrs=None):
    """段落のテキストを上書き（既存スタイル保持）"""
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        attrs = fallback_attrs or {"lang": "ja-JP"}
        etree.SubElement(r_elem, qn("a:rPr"), attrib=attrs)
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_concept_textbox_3(slide, textbox_name, name, detail1, detail2):
    """3コンセプト版: TextBoxのPara[0]に名前、Para[1]に詳細1、Para[2]に詳細2を書き込む"""
    shape = find_shape(slide, textbox_name)
    if shape is None:
        return
    tf = shape.text_frame

    # Para[0]: コンセプト名（Bold）
    if len(tf.paragraphs) > 0:
        set_para_text(tf.paragraphs[0], name, {"lang": "ja-JP", "b": "1"})

    # Para[1]: 詳細1
    if len(tf.paragraphs) > 1:
        set_para_text(tf.paragraphs[1], detail1, {"lang": "ja-JP"})

    # Para[2]: 詳細2
    if len(tf.paragraphs) > 2:
        set_para_text(tf.paragraphs[2], detail2, {"lang": "ja-JP"})

    print(f"  [{textbox_name}] {name}")
    print(f"    Detail1: {detail1[:55]}{'...' if len(detail1) > 55 else ''}")
    print(f"    Detail2: {detail2[:55]}{'...' if len(detail2) > 55 else ''}")


def fill_concept_textbox_5(slide, textbox_name, name, detail):
    """5コンセプト版: TextBoxのPara[0]に名前、Para[1]に詳細を書き込む"""
    shape = find_shape(slide, textbox_name)
    if shape is None:
        return
    tf = shape.text_frame

    # Para[0]: コンセプト名（Bold）
    if len(tf.paragraphs) > 0:
        set_para_text(tf.paragraphs[0], name, {"lang": "ja-JP", "b": "1"})

    # Para[1]: 詳細
    if len(tf.paragraphs) > 1:
        set_para_text(tf.paragraphs[1], detail, {"lang": "ja-JP"})

    print(f"  [{textbox_name}] {name}: {detail[:50]}{'...' if len(detail) > 50 else ''}")


def main():
    parser = argparse.ArgumentParser(description="コンセプトデータをPPTXに流し込む")
    parser.add_argument("--data", required=True, help="conceptual_data.json のパス")
    parser.add_argument("--template3", required=True, help="Conceptual3.pptx のパス")
    parser.add_argument("--template5", required=True, help="Conceptual5.pptx のパス")
    parser.add_argument("--output", required=True, help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    variant = data.get("variant", 3)
    concepts = data.get("concepts", [])

    # テンプレート選択
    if variant == 5:
        template_path = args.template5
        oval_names = OVAL_NAMES_5
        textbox_names = TEXTBOX_NAMES_5
        expected_count = 5
    else:
        template_path = args.template3
        oval_names = OVAL_NAMES_3
        textbox_names = TEXTBOX_NAMES_3
        expected_count = 3

    print(f"テンプレート: {template_path} ({expected_count}コンセプト版)")

    if len(concepts) != expected_count:
        print(f"  ⚠ コンセプト数が {len(concepts)} 件です（{expected_count}件必要）", file=sys.stderr)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # ── Main Message ──────────────────────────────────────────
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_text_in_shape(shape, main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # ── Chart Title ───────────────────────────────────────────
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_text_in_shape(shape, chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # ── Concepts ──────────────────────────────────────────────
    for i, concept in enumerate(concepts[:expected_count]):
        name = concept.get("name", "").strip()

        # 楕円内のラベル
        if i < len(oval_names):
            oval_shape = find_shape(slide, oval_names[i])
            set_text_in_shape(oval_shape, name)
            print(f"  [{oval_names[i]}] {name}")

        # 右側のテキストボックス
        if i < len(textbox_names):
            if variant == 3:
                # 3コンセプト版: 詳細2行
                detail1 = concept.get("detail1", "").strip()
                detail2 = concept.get("detail2", "").strip()
                fill_concept_textbox_3(slide, textbox_names[i], name, detail1, detail2)
            else:
                # 5コンセプト版: 詳細1行
                detail = concept.get("detail", "").strip()
                fill_concept_textbox_5(slide, textbox_names[i], name, detail)

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
