"""
fill_value_chain.py — バリューチェーン分析スライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 中央: 5〜7個のバリューチェーン段階を矢印（シェブロン）状に横並び配置
  - 詳細: 各段階の利益プール率バー + 対象会社のポジション(強/中/弱)
  - 下部: バリューチェーン分析からの示唆

Usage:
  python fill_value_chain.py \
    --data /home/claude/value_chain_data.json \
    [--template <path>/value-chain-template.pptx] \
    --output /mnt/user-data/outputs/ValueChain_output.pptx \
    [--brand stellar_aiz | roleup]
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "value-chain-pptx"
_THEME = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants (stella 16:9, 13.33 x 7.5 in; _apply_theme で roleup A4 用にスケール) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_PH = "Source 3"  # roleup placeholder

CHAIN_LEFT = Inches(0.41)
CHAIN_TOP = Inches(1.80)
CHAIN_WIDTH = Inches(12.51)
CHAIN_HEIGHT = Inches(2.50)

DETAIL_TOP = Inches(4.35)
DETAIL_HEIGHT = Inches(1.60)

FOOTER_TOP = Inches(6.05)
FOOTER_HEIGHT = Inches(0.85)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)

# ── Colors (stella defaults; _apply_theme で roleup 用に上書き) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTEXT = RGBColor(0x55, 0x55, 0x55)

CHEVRON_COLORS = [
    RGBColor(0x1E, 0x3A, 0x5F),
    RGBColor(0x2E, 0x4A, 0x6B),
    RGBColor(0x3F, 0x5F, 0x80),
    RGBColor(0x50, 0x75, 0x96),
    RGBColor(0x62, 0x8B, 0xAC),
    RGBColor(0x75, 0xA2, 0xC2),
    RGBColor(0x88, 0xB8, 0xD8),
]

# 利益プール / ポジション色は universal indicator (信号色) として両 brand で維持
COLOR_PROFIT_HIGH = RGBColor(0x1B, 0x7A, 0x3B)
COLOR_PROFIT_MID = RGBColor(0xDA, 0x7A, 0x2D)
COLOR_PROFIT_LOW = RGBColor(0xB8, 0x3A, 0x3A)

COLOR_STRONG = RGBColor(0x1B, 0x7A, 0x3B)
COLOR_MEDIUM = RGBColor(0xDA, 0x7A, 0x2D)
COLOR_WEAK = RGBColor(0xB8, 0x3A, 0x3A)

# 示唆ブレットマーカー色 (brand-aware で上書き)
COLOR_BULLET = RGBColor(0x2E, 0x4A, 0x6B)

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_STAGE_NAME = Pt(13)
FONT_SIZE_ACTIVITY = Pt(10)
FONT_SIZE_PROFIT = Pt(11)
FONT_SIZE_POSITION = Pt(11)
FONT_SIZE_FOOTER = Pt(11)
FONT_SIZE_SOURCE = Pt(10)
FONT_SIZE_SECTION = Pt(13)


def _apply_theme(theme):
    """roleup の場合、レイアウト・色・フォントサイズを brand 仕様に上書きする。"""
    global CHAIN_LEFT, CHAIN_TOP, CHAIN_WIDTH, CHAIN_HEIGHT
    global DETAIL_TOP, DETAIL_HEIGHT, FOOTER_TOP, FOOTER_HEIGHT
    global SOURCE_X, SOURCE_Y, SOURCE_W
    global COLOR_TEXT, COLOR_SOURCE, COLOR_SUBTEXT
    global CHEVRON_COLORS, COLOR_BULLET
    global FONT_NAME_JP
    global FONT_SIZE_STAGE_NAME, FONT_SIZE_ACTIVITY, FONT_SIZE_PROFIT
    global FONT_SIZE_POSITION, FONT_SIZE_FOOTER, FONT_SIZE_SOURCE, FONT_SIZE_SECTION
    global _THEME
    _THEME = theme

    if theme.id != "roleup":
        return

    # A4 横 (11.69 × 8.27) 用にレイアウト再計算
    CHAIN_LEFT = Inches(0.41)
    CHAIN_TOP = Inches(1.55)
    CHAIN_WIDTH = Inches(10.87)
    CHAIN_HEIGHT = Inches(2.30)

    DETAIL_TOP = Inches(3.95)
    DETAIL_HEIGHT = Inches(1.50)

    FOOTER_TOP = Inches(5.60)
    FOOTER_HEIGHT = Inches(1.75)

    SOURCE_X = Inches(0.41)
    SOURCE_Y = Inches(7.45)
    SOURCE_W = Inches(10.87)

    # roleup 茶系トーン
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    COLOR_SUBTEXT = theme.color("source")

    # シェブロン色は roleup chart_palette から茶系グラデーションを構成
    # palette: ["#7C4C2C","#897141","#604C3F","#C78624","#AF7026","#3E3A39","#9C755F","#CDCECE"]
    # 濃→淡で並べ替え: 604C3F, 7C4C2C, AF7026, 9C755F, C78624, 897141, CDCECE
    palette_hex = list(theme.chart_palette) if theme.chart_palette else []
    if palette_hex:
        order = [2, 0, 4, 6, 3, 1, 7]
        chev = []
        for i in order:
            if i < len(palette_hex):
                hx = palette_hex[i].lstrip("#")
                chev.append(RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)))
        if chev:
            CHEVRON_COLORS = chev

    COLOR_BULLET = theme.color("label_bar")  # #7C4C2C

    # フォント
    FONT_NAME_JP = theme.font_ea or "Yu Gothic UI"

    # roleup C4 許容集合 [22, 14, 12, 10, 6] pt
    FONT_SIZE_STAGE_NAME = Pt(12)
    FONT_SIZE_ACTIVITY = Pt(10)
    FONT_SIZE_PROFIT = Pt(12)
    FONT_SIZE_POSITION = Pt(10)
    FONT_SIZE_FOOTER = Pt(10)
    FONT_SIZE_SOURCE = Pt(int(theme._defaults.get("font_size_source_pt", 6)))
    FONT_SIZE_SECTION = Pt(12)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    """汎用テキストボックス (font_name は late-resolve で _apply_theme 後の値を使う)"""
    if font_name is None:
        font_name = FONT_NAME_JP
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def _get_position_info(position):
    if not position:
        return (COLOR_SUBTEXT, "—")
    p = str(position).lower()
    if p in ("strong", "high", "強", "●●●"):
        return (COLOR_STRONG, "●●●  強み")
    elif p in ("medium", "mid", "中", "●●○"):
        return (COLOR_MEDIUM, "●●○  中位")
    elif p in ("weak", "low", "弱", "●○○"):
        return (COLOR_WEAK, "●○○  弱み")
    return (COLOR_SUBTEXT, position)


def _get_profit_color(pool_percent):
    if pool_percent is None:
        return COLOR_SUBTEXT
    if pool_percent >= 25:
        return COLOR_PROFIT_HIGH
    elif pool_percent >= 10:
        return COLOR_PROFIT_MID
    else:
        return COLOR_PROFIT_LOW


def draw_value_chain(slide, stages, left, top, width, height):
    n = len(stages)
    if n == 0:
        return

    gap = Inches(0.05)
    chevron_w = Emu(int((width - gap * (n - 1)) / n))
    chevron_h = height

    for i, stage in enumerate(stages):
        x = left + (chevron_w + gap) * i
        y = top

        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, chevron_w, chevron_h)
        color = CHEVRON_COLORS[i % len(CHEVRON_COLORS)]
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.10); tf.margin_bottom = Inches(0.10)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for p in list(tf.paragraphs):
            p._p.getparent().remove(p._p)

        # 段階名
        p1 = etree.SubElement(tf._txBody, qn("a:p"))
        pPr1 = etree.SubElement(p1, qn("a:pPr"))
        pPr1.set("algn", "ctr")
        r1 = etree.SubElement(p1, qn("a:r"))
        rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_STAGE_NAME.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
        s1 = etree.SubElement(sf1, qn("a:srgbClr"))
        s1.set("val", "FFFFFF")
        t1 = etree.SubElement(r1, qn("a:t"))
        t1.text = stage.get("name", f"段階{i+1}")

        activity = stage.get("activity", "")
        if activity:
            p2 = etree.SubElement(tf._txBody, qn("a:p"))
            pPr2 = etree.SubElement(p2, qn("a:pPr"))
            pPr2.set("algn", "ctr")
            spcBef = etree.SubElement(pPr2, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "300"})
            r2 = etree.SubElement(p2, qn("a:r"))
            rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={
                "lang": "ja-JP",
                "sz": str(int(FONT_SIZE_ACTIVITY.pt * 100)),
                "b": "0",
            })
            etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
            etree.SubElement(rPr2, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
            sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
            s2 = etree.SubElement(sf2, qn("a:srgbClr"))
            s2.set("val", "FFFFFF")
            t2 = etree.SubElement(r2, qn("a:t"))
            t2.text = activity


def draw_detail_panel(slide, stages, left, top, width, height):
    n = len(stages)
    if n == 0:
        return

    gap = Inches(0.05)
    cell_w = Emu(int((width - gap * (n - 1)) / n))

    max_pct = max([s.get("profit_pool_pct", 0) or 0 for s in stages], default=100) or 100

    for i, stage in enumerate(stages):
        x = left + (cell_w + gap) * i
        profit_pct = stage.get("profit_pool_pct")
        position = stage.get("position")
        y = top

        profit_text = f"{profit_pct}%" if profit_pct is not None else "—"
        profit_color = _get_profit_color(profit_pct)

        add_text_box(
            slide, f"利益プール: {profit_text}",
            x, y, cell_w, Inches(0.25),
            FONT_SIZE_PROFIT, bold=True,
            color=profit_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        bar_y = y + Inches(0.30)
        bar_max_w = cell_w - Inches(0.20)
        bar_h = Inches(0.18)
        bar_x_start = x + Inches(0.10)

        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_x_start, bar_y, bar_max_w, bar_h,
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        bg_bar.line.fill.background()
        bg_bar.shadow.inherit = False
        bg_bar.text_frame.text = ""

        if profit_pct is not None and profit_pct > 0:
            fg_w = Emu(int(bar_max_w * min(profit_pct, max_pct) / max_pct))
            if fg_w > Emu(100):
                fg_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    bar_x_start, bar_y, fg_w, bar_h,
                )
                fg_bar.fill.solid()
                fg_bar.fill.fore_color.rgb = profit_color
                fg_bar.line.fill.background()
                fg_bar.shadow.inherit = False
                fg_bar.text_frame.text = ""

        pos_color, pos_label = _get_position_info(position)
        add_text_box(
            slide, pos_label,
            x, y + Inches(0.55), cell_w, Inches(0.25),
            FONT_SIZE_POSITION, bold=True,
            color=pos_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def draw_footer(slide, implications, left, top, width, height):
    if not implications:
        return

    add_text_box(
        slide, "▍ バリューチェーン分析からの示唆",
        left, top, width, Inches(0.30),
        FONT_SIZE_SECTION, bold=True,
        align=PP_ALIGN.LEFT,
    )

    body_top = top + Inches(0.35)
    body_h = height - Inches(0.35)

    tb = slide.shapes.add_textbox(
        left + Inches(0.15), body_top,
        width - Inches(0.30), body_h,
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    bullet_hex = "{:02X}{:02X}{:02X}".format(COLOR_BULLET[0], COLOR_BULLET[1], COLOR_BULLET[2])
    text_hex = "{:02X}{:02X}{:02X}".format(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2])

    for i, item in enumerate(implications):
        if isinstance(item, dict):
            text = item.get("text", "")
        else:
            text = str(item)

        p_elem = etree.SubElement(tf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"), attrib={
            "marL": "180000",
            "indent": "-180000",
        })
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "300"})

        buChar = etree.SubElement(pPr, qn("a:buChar"), attrib={"char": "•"})
        buFont = etree.SubElement(pPr, qn("a:buFont"), attrib={"typeface": "Arial"})
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
        buClrSolid.set("val", bullet_hex)

        r = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_FOOTER.pt * 100)),
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", text_hex)
        t = etree.SubElement(r, qn("a:t"))
        t.text = text


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=False, default=None)
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "stages"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "stages", "implications",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "value-chain")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top text (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
        print(f"  ✓ Top: {top_text[:60]}{'...' if len(top_text) > 60 else ''}")

    # Subtitle (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme).strip()
    if sub_text:
        set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
        print(f"  ✓ Subtitle: {sub_text[:60]}{'...' if len(sub_text) > 60 else ''}")

    stages = data.get("stages", [])
    if not stages:
        print("  ✗ ERROR: 'stages' is required", file=sys.stderr)
        sys.exit(1)

    if len(stages) > 7:
        print(f"  ⚠ WARNING: {len(stages)} stages > 7. Only first 7 will be shown.", file=sys.stderr)
        stages = stages[:7]

    draw_value_chain(slide, stages, CHAIN_LEFT, CHAIN_TOP, CHAIN_WIDTH, CHAIN_HEIGHT)
    print(f"  ✓ バリューチェーン: {len(stages)}段階")

    draw_detail_panel(slide, stages, CHAIN_LEFT, DETAIL_TOP, CHAIN_WIDTH, DETAIL_HEIGHT)
    print(f"  ✓ 詳細パネル")

    implications = data.get("implications", [])
    draw_footer(slide, implications, CHAIN_LEFT, FOOTER_TOP, CHAIN_WIDTH, FOOTER_HEIGHT)
    if implications:
        print(f"  ✓ 示唆: {len(implications)}項目")

    # 出典 (roleup: Source 3 placeholder, stella: 動的 textbox)
    source = (data.get("source") or data.get("source_label")
              or data.get("source_text") or "").strip()
    if source:
        if theme.is_source_required():
            src_shape = find_shape(slide, SHAPE_SOURCE_PH)
            if src_shape is not None:
                set_textbox_text(src_shape, f"出典: {source}")
                for para in src_shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = FONT_SIZE_SOURCE
                        run.font.color.rgb = COLOR_SOURCE
                        run.font.name = FONT_NAME_JP
        else:
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.25),
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
