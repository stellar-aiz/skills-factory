"""
fill_section_divider.py — 中扉（Section Divider）スライドを生成

レイアウト:
  - 左半分: アクセントカラー背景 + 巨大なセクション番号（180pt 白） + "SECTION" ラベル
  - 右半分: タイトル / アクセントライン / サブタイトル / トピックリスト

Brand-aware (Phase 2, ISSUE-010):
  --brand stellar_aiz : 13.33×7.50 in / Meiryo UI / 既存 7 色ローテーション (hardcode 値維持)
  --brand roleup      : 11.69×8.27 in / Yu Gothic UI / theme.chart_palette 8 色ローテ

中扉は装飾スライド (大フォント主体) のため C4 (font size constraint) と
C2/C5/C6 (title/source) は profile から除外し、C1/C7/C8/C11 のみ適用する。

Usage:
  python fill_section_divider.py \
    --data /home/claude/section_divider_data.json \
    --brand stellar_aiz \
    --output /mnt/user-data/outputs/SectionDivider_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402

SKILL_ID = "section-divider-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Template placeholder shape names (cp roleup template / stella legacy 共通) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_ROLEUP = "Source 3"  # roleup template carries; we delete (中扉に出典は不要)

# ── Default values (stella) — reassigned in main() via _apply_theme(theme) ──
# Layout (mirrors V1 hardcoded values for stella regression-zero).
LEFT_PANEL_W = Inches(5.50)
NUM_TOP = Inches(1.50)
NUM_H = Inches(4.50)
LABEL_TOP = Inches(2.60)
LABEL_H = Inches(0.40)
RIGHT_X_OFFSET = Inches(0.50)
RIGHT_PAD_RIGHT = Inches(0.50)
TITLE_TOP = Inches(2.30)
TITLE_H = Inches(1.20)
ACCENT_LINE_TOP = Inches(3.35)
ACCENT_LINE_W = Inches(1.00)
ACCENT_LINE_H = Emu(int(Inches(0.06)))
SUBTITLE_TOP = Inches(3.55)
SUBTITLE_H = Inches(0.50)
TOPICS_TOP = Inches(4.30)
TOPIC_LABEL_H = Inches(0.30)
TOPIC_OFFSET = Inches(0.10)
TOPIC_H = Inches(2.00)

# Slide size (stella default; reassigned to roleup A4 in _apply_theme).
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# Colors (stella default).
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTEXT = RGBColor(0x66, 0x66, 0x66)

# Section accent color rotations (stella legacy: 7 colors).
SECTION_COLORS_STELLA = [
    RGBColor(0x2E, 0x4A, 0x6B),   # 紺
    RGBColor(0x7B, 0x4F, 0xB0),   # 紫
    RGBColor(0x2E, 0x6F, 0xBF),   # 青
    RGBColor(0x3D, 0x8F, 0x5A),   # 緑
    RGBColor(0xDA, 0x7A, 0x2D),   # オレンジ
    RGBColor(0xC0, 0x3A, 0x3A),   # 赤
    RGBColor(0x59, 0x59, 0x59),   # グレー
]

# Fonts (stella defaults).
FONT_NAME_JP = "Meiryo UI"
FONT_NAME_LATIN = "Arial"
FONT_SIZE_BIG_NUMBER = Pt(180)   # 巨大なセクション番号 (装飾、C4 検査対象外)
FONT_SIZE_SECTION_LABEL = Pt(20) # SECTION ラベル
FONT_SIZE_SECTION_TITLE = Pt(36) # セクションタイトル
FONT_SIZE_SUBTITLE = Pt(18)
FONT_SIZE_TOPIC = Pt(13)

# Theme module-global; populated in main() via _apply_theme(theme).
_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme.

    Called once from main() after `--brand` is parsed.
    """
    global _THEME, SLIDE_W, SLIDE_H
    global LEFT_PANEL_W, NUM_TOP, NUM_H, LABEL_TOP, LABEL_H
    global RIGHT_X_OFFSET, RIGHT_PAD_RIGHT, TITLE_TOP, TITLE_H
    global ACCENT_LINE_TOP, ACCENT_LINE_W, ACCENT_LINE_H
    global SUBTITLE_TOP, SUBTITLE_H, TOPICS_TOP, TOPIC_LABEL_H, TOPIC_OFFSET, TOPIC_H
    global COLOR_TEXT, COLOR_SUBTEXT
    global FONT_NAME_JP, FONT_NAME_LATIN
    global FONT_SIZE_BIG_NUMBER, FONT_SIZE_SECTION_LABEL, FONT_SIZE_SECTION_TITLE
    global FONT_SIZE_SUBTITLE, FONT_SIZE_TOPIC

    _THEME = theme

    SLIDE_W = theme.slide_w
    SLIDE_H = theme.slide_h

    LEFT_PANEL_W = theme.layout("left_panel_w_in")
    NUM_TOP = theme.layout("num_top_in")
    NUM_H = theme.layout("num_h_in")
    LABEL_TOP = theme.layout("label_top_in")
    LABEL_H = theme.layout("label_h_in")
    RIGHT_X_OFFSET = theme.layout("right_x_offset_in")
    RIGHT_PAD_RIGHT = theme.layout("right_pad_right_in")
    TITLE_TOP = theme.layout("title_top_in")
    TITLE_H = theme.layout("title_h_in")
    ACCENT_LINE_TOP = theme.layout("accent_line_top_in")
    ACCENT_LINE_W = theme.layout("accent_line_w_in")
    ACCENT_LINE_H = Emu(int(theme.layout("accent_line_h_in")))
    SUBTITLE_TOP = theme.layout("subtitle_top_in")
    SUBTITLE_H = theme.layout("subtitle_h_in")
    TOPICS_TOP = theme.layout("topics_top_in")
    TOPIC_LABEL_H = theme.layout("topic_label_h_in")
    TOPIC_OFFSET = theme.layout("topic_offset_in")
    TOPIC_H = theme.layout("topic_h_in")

    COLOR_TEXT = theme.color("text")
    # subtitle 色は brand 別の subtitle 色を使う (roleup=#897141, stella=#333333)
    COLOR_SUBTEXT = theme.color("subtitle") if "subtitle" in theme._colors else COLOR_TEXT

    FONT_NAME_JP = theme.font_ea
    FONT_NAME_LATIN = theme.font_latin

    # roleup では本文系を許容集合 {22, 14, 12, 10, 6} に収める。
    # 装飾的大フォント (BIG_NUMBER 180pt / SECTION_TITLE) は C4 検査対象外 (profile で除外)。
    if theme.id == "stellar_aiz":
        FONT_SIZE_BIG_NUMBER = Pt(180)
        FONT_SIZE_SECTION_LABEL = Pt(20)
        FONT_SIZE_SECTION_TITLE = Pt(36)
        FONT_SIZE_SUBTITLE = Pt(18)
        FONT_SIZE_TOPIC = Pt(13)
    else:
        # roleup: 巨大数字とセクションタイトルは装飾扱い (C4 除外プロファイル)、
        # SECTION ラベル / サブタイトル / トピックは許容集合内へ。
        FONT_SIZE_BIG_NUMBER = Pt(180)
        FONT_SIZE_SECTION_LABEL = Pt(12)
        FONT_SIZE_SECTION_TITLE = Pt(22)
        FONT_SIZE_SUBTITLE = Pt(12)
        FONT_SIZE_TOPIC = Pt(10)


def _silent_remove_shape(slide, shape_name: str) -> None:
    """Remove a shape by name without printing a warning. No-op if absent."""
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def _section_color(theme, section_number: int) -> RGBColor:
    """Resolve accent color for a section number.

    stella: 7 色 hardcode ローテ (V1 互換)
    roleup: theme.chart_palette (8 色) ローテ
    """
    if theme.id == "stellar_aiz":
        return SECTION_COLORS_STELLA[(section_number - 1) % len(SECTION_COLORS_STELLA)]
    palette = theme.chart_palette
    hex_str = palette[(section_number - 1) % len(palette)]
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    # font_name / color default to current module globals (post-_apply_theme).
    if font_name is None:
        font_name = FONT_NAME_JP
    if color is None:
        color = COLOR_TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def remove_template_shapes(slide):
    """テンプレートのプレースホルダーを削除して中扉を一から組む。

    cp roleup template には Title 1 / Text Placeholder 2 / Source 3 + 茶色ガイド × 2
    が含まれるが、中扉専用デザインのため全て削除する。
    """
    targets = (SHAPE_MAIN_MESSAGE, SHAPE_CHART_TITLE, SHAPE_SOURCE_ROLEUP,
               "正方形/長方形 1", "正方形/長方形 8")
    for sh in list(slide.shapes):
        if sh.name in targets:
            sp_tree = slide.shapes._spTree
            sp_tree.remove(sh._element)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "section-divider")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 中扉専用デザインのため、テンプレ placeholder を全削除して一から組む。
    remove_template_shapes(slide)

    # セクション番号と色を取得
    section_number = data.get("section_number", 1)
    color_hex = data.get("color")
    if color_hex:
        accent_color = hex_to_rgb(color_hex)
    else:
        accent_color = _section_color(theme, section_number)

    # ── 左半分: 巨大なセクション番号 + アクセントカラー背景 ──
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), LEFT_PANEL_W, SLIDE_H,
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = accent_color
    left_bg.line.fill.background()
    left_bg.shadow.inherit = False
    left_bg.text_frame.text = ""

    # 巨大な番号（中央）
    num_tb = slide.shapes.add_textbox(
        Emu(0), NUM_TOP, LEFT_PANEL_W, NUM_H,
    )
    ntf = num_tb.text_frame
    ntf.margin_left = 0; ntf.margin_right = 0
    ntf.margin_top = 0; ntf.margin_bottom = 0
    ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for p in list(ntf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(ntf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr")
    r = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_BIG_NUMBER.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_LATIN})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    s = etree.SubElement(sf, qn("a:srgbClr"))
    s.set("val", "FFFFFF")
    t = etree.SubElement(r, qn("a:t"))
    t.text = f"{section_number:02d}"

    # 番号上の小ラベル "SECTION"
    label_tb = slide.shapes.add_textbox(
        Emu(0), LABEL_TOP, LEFT_PANEL_W, LABEL_H,
    )
    ltf = label_tb.text_frame
    ltf.margin_left = 0; ltf.margin_right = 0
    ltf.margin_top = 0; ltf.margin_bottom = 0
    p_lbl = ltf.paragraphs[0]
    p_lbl.alignment = PP_ALIGN.CENTER
    r_lbl = p_lbl.add_run()
    r_lbl.text = "SECTION"
    r_lbl.font.size = FONT_SIZE_SECTION_LABEL
    r_lbl.font.bold = True
    r_lbl.font.name = FONT_NAME_LATIN
    r_lbl.font.color.rgb = COLOR_WHITE

    # ── 右半分: タイトル・サブタイトル・トピック ──
    right_x = LEFT_PANEL_W + RIGHT_X_OFFSET
    right_w = SLIDE_W - right_x - RIGHT_PAD_RIGHT

    # タイトル
    title = data.get("title", "セクションタイトル")
    add_text_box(
        slide, title,
        right_x, TITLE_TOP, right_w, TITLE_H,
        FONT_SIZE_SECTION_TITLE, bold=True,
        color=COLOR_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # アクセントライン（タイトルの下）
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        right_x, ACCENT_LINE_TOP,
        ACCENT_LINE_W, ACCENT_LINE_H,
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent_color
    accent_line.line.fill.background()

    # サブタイトル
    subtitle = data.get("subtitle", "")
    if subtitle:
        add_text_box(
            slide, subtitle,
            right_x, SUBTITLE_TOP, right_w, SUBTITLE_H,
            FONT_SIZE_SUBTITLE, bold=False,
            color=COLOR_SUBTEXT, align=PP_ALIGN.LEFT,
        )

    # トピックリスト
    topics = data.get("topics", [])
    if topics:
        topic_label_tb = slide.shapes.add_textbox(
            right_x, TOPICS_TOP, right_w, TOPIC_LABEL_H,
        )
        tltf = topic_label_tb.text_frame
        tltf.margin_left = 0; tltf.margin_right = 0
        tltf.margin_top = 0; tltf.margin_bottom = 0
        p_tl = tltf.paragraphs[0]
        p_tl.alignment = PP_ALIGN.LEFT
        r_tl = p_tl.add_run()
        r_tl.text = "▍ このセクションで扱う内容"
        r_tl.font.size = FONT_SIZE_TOPIC
        r_tl.font.bold = True
        r_tl.font.name = FONT_NAME_JP
        r_tl.font.color.rgb = accent_color

        # トピック項目
        topic_tb = slide.shapes.add_textbox(
            right_x + TOPIC_OFFSET, TOPICS_TOP + Inches(0.40),
            right_w - TOPIC_OFFSET, TOPIC_H,
        )
        topic_tf = topic_tb.text_frame
        topic_tf.word_wrap = True
        topic_tf.margin_left = 0; topic_tf.margin_right = 0
        topic_tf.margin_top = 0; topic_tf.margin_bottom = 0

        for p in list(topic_tf.paragraphs):
            p._p.getparent().remove(p._p)

        for i, topic in enumerate(topics):
            p_topic = etree.SubElement(topic_tf._txBody, qn("a:p"))
            pPr = etree.SubElement(p_topic, qn("a:pPr"), attrib={
                "marL": "200000",
                "indent": "-200000",
            })
            if i > 0:
                spcBef = etree.SubElement(pPr, qn("a:spcBef"))
                etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "400"})

            buChar = etree.SubElement(pPr, qn("a:buChar"), attrib={"char": "▸"})
            # bullet typeface: roleup でも Yu Gothic UI に揃える (C8 適用範囲は run のみだが念のため)
            buFont = etree.SubElement(pPr, qn("a:buFont"), attrib={"typeface": FONT_NAME_JP})
            buClr = etree.SubElement(pPr, qn("a:buClr"))
            buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
            buClrSolid.set("val", "{:02X}{:02X}{:02X}".format(accent_color[0], accent_color[1], accent_color[2]))

            r_top = etree.SubElement(p_topic, qn("a:r"))
            rPr_top = etree.SubElement(r_top, qn("a:rPr"), attrib={
                "lang": "ja-JP",
                "sz": str(int(FONT_SIZE_TOPIC.pt * 100)),
            })
            etree.SubElement(rPr_top, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
            etree.SubElement(rPr_top, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
            sf_top = etree.SubElement(rPr_top, qn("a:solidFill"))
            s_top = etree.SubElement(sf_top, qn("a:srgbClr"))
            s_top.set("val", "{:02X}{:02X}{:02X}".format(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2]))
            t_top = etree.SubElement(r_top, qn("a:t"))
            t_top.text = topic

    print(f"  ✓ Section {section_number:02d}: {title}")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
