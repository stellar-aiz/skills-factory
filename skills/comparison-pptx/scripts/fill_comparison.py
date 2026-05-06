"""
fill_comparison.py — Comparison Chart PPTX Generator

Usage:
    python fill_comparison.py --data comparison_data.json \
                              --template assets/comparison-template.pptx \
                              --output output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402
from pptx import Presentation
from pptx.util import Pt
import re

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# Shape name mapping (verified against template)
SHAPE_MAP = {
    "main_message": "Title 1",
    "chart_title": "Text Placeholder 2",
    "axis_label": "TextBox 6",
    "col_a": "TextBox 22",
    "col_b": "TextBox 3",
    "col_c": "TextBox 4",
    "label_1": "Rectangle 5",
    "label_2": "Rectangle 8",
    "label_3": "Rectangle 9",
    "label_4": "Rectangle 12",
    "label_5": "Rectangle 15",
    "implications": "TextBox 29",
    # Cell boxes: [row][col]
    "cell_1_a": "TextBox 30",
    "cell_1_b": "TextBox 31",
    "cell_1_c": "TextBox 32",
    "cell_2_a": "TextBox 34",
    "cell_2_b": "TextBox 35",
    "cell_2_c": "TextBox 36",
    "cell_3_a": "TextBox 37",
    "cell_3_b": "TextBox 38",
    "cell_3_c": "TextBox 39",
    "cell_4_a": "TextBox 40",
    "cell_4_b": "TextBox 41",
    "cell_4_c": "TextBox 42",
    "cell_5_a": "TextBox 43",
    "cell_5_b": "TextBox 44",
    "cell_5_c": "TextBox 45",
}


def get_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_text_preserving_format(text_frame, new_text):
    """Replace all text while preserving paragraph/run formatting of the first run."""
    if not text_frame.paragraphs:
        return
    para = text_frame.paragraphs[0]
    if para.runs:
        run = para.runs[0]
        run.text = new_text
        # Clear extra runs
        for extra_run in para.runs[1:]:
            extra_run.text = ""
    else:
        para.text = new_text


def set_cell_text(text_frame, mark, comment):
    """
    Set a cell's text frame: 
      - Para 0: mark (bold, large)
      - Para 1: comment (bold, small)
    """
    paras = text_frame.paragraphs
    # Para 0 = mark
    if len(paras) > 0:
        p0 = paras[0]
        if p0.runs:
            p0.runs[0].text = mark
            for r in p0.runs[1:]:
                r.text = ""
        else:
            p0.text = mark

    # Para 1 = comment
    if len(paras) > 1:
        p1 = paras[1]
        if p1.runs:
            p1.runs[0].text = comment
            for r in p1.runs[1:]:
                r.text = ""
        else:
            p1.text = comment
    elif comment:
        # Add a new paragraph if missing
        from pptx.oxml.ns import qn
        import copy
        new_para = copy.deepcopy(paras[0]._p)
        text_frame._txBody.append(new_para)
        new_p = text_frame.paragraphs[-1]
        if new_p.runs:
            new_p.runs[0].text = comment


def set_implications(text_frame, implications):
    """
    Set implications text box.
    Para 0: "意味合い" (header)
    Para 1-3: implication texts
    """
    paras = text_frame.paragraphs
    # Para 0 is the header "意味合い" - keep as is
    for i, impl in enumerate(implications):
        para_idx = i + 1
        if para_idx < len(paras):
            p = paras[para_idx]
            if p.runs:
                p.runs[0].text = impl
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = impl


def fill_comparison(data, template_path, output_path):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Main message
    shape = get_shape_by_name(slide, SHAPE_MAP["main_message"])
    if shape:
        set_text_preserving_format(shape.text_frame, data["main_message"])

    # Chart title
    shape = get_shape_by_name(slide, SHAPE_MAP["chart_title"])
    if shape:
        set_text_preserving_format(shape.text_frame, data["chart_title"])

    # Axis label
    shape = get_shape_by_name(slide, SHAPE_MAP["axis_label"])
    if shape:
        set_text_preserving_format(shape.text_frame, data.get("axis_label", "評価軸"))

    # Column headers (patterns)
    patterns = data["patterns"]
    for key, label in [("col_a", patterns[0]), ("col_b", patterns[1]), ("col_c", patterns[2])]:
        shape = get_shape_by_name(slide, SHAPE_MAP[key])
        if shape:
            set_text_preserving_format(shape.text_frame, label)

    # Row labels (criteria)
    criteria = data["criteria"]
    label_keys = ["label_1", "label_2", "label_3", "label_4", "label_5"]
    for i, key in enumerate(label_keys):
        shape = get_shape_by_name(slide, SHAPE_MAP[key])
        if shape and i < len(criteria):
            set_text_preserving_format(shape.text_frame, criteria[i]["label"])

    # Cell content
    cols = ["a", "b", "c"]
    cells = data["cells"]  # cells[row_idx][col_idx] = {"mark": "◎", "comment": "..."}
    for row_idx in range(min(5, len(cells))):
        for col_idx, col_char in enumerate(cols):
            cell_key = f"cell_{row_idx+1}_{col_char}"
            shape = get_shape_by_name(slide, SHAPE_MAP[cell_key])
            if shape and col_idx < len(cells[row_idx]):
                cell_data = cells[row_idx][col_idx]
                set_cell_text(shape.text_frame, cell_data["mark"], cell_data["comment"])

    # Implications
    shape = get_shape_by_name(slide, SHAPE_MAP["implications"])
    if shape:
        set_implications(shape.text_frame, data.get("implications", []))

    prs.save(output_path)

    _finalize_pptx(output_path)

    print(f"Saved: {output_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "patterns", "criteria", "cells"],
        allowed_top=[
            "main_message", "chart_title", "patterns", "criteria",
            "cells", "axis_label", "implications",
        ],
        skill_name="comparison-pptx",
    )

    fill_comparison(data, args.template, args.output)


if __name__ == "__main__":
    main()
