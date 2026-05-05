"""
fill_positioning_map.py — ポジショニングマップスライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル (top/subtitle 配置は brand により入れ替え)
  - 左側: 2軸ポジショニングマップ
      - X軸・Y軸ラベル、両端にlow/highラベル
      - 4象限ラベル（オプション）
      - 各プレイヤーをバブル（円）で配置
      - 対象会社は目立つ色＋太枠でハイライト
      - バブル下に企業名ラベル
  - 右側: 示唆（Implications）のブレット項目
  - 下部: 出典

描画方針: ネイティブチャートではなく MSO_SHAPE.OVAL と直線で手動描画

Usage:
  python fill_positioning_map.py \
    --data /home/claude/positioning_map_data.json \
    --brand stellar_aiz | roleup \
    --output /mnt/user-data/outputs/PositioningMap_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "positioning-map-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants (defaults; reassigned in main() via _apply_theme) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE = "Source 3"

PANEL_Y = Inches(1.55)
PANEL_H = Inches(5.35)

LEFT_X = Inches(0.41)
LEFT_W = Inches(7.70)
RIGHT_X = Inches(8.30)
RIGHT_W = Inches(4.65)

# Map drawing area inside left panel (axis labels etc consume margins)
MAP_MARGIN_LEFT = Inches(0.75)
MAP_MARGIN_RIGHT = Inches(0.20)
MAP_MARGIN_TOP = Inches(0.70)
MAP_MARGIN_BOTTOM = Inches(0.80)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(6.93)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.30)

# ── Colors (defaults — reassigned by _apply_theme) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTITLE = RGBColor(0x33, 0x33, 0x33)
COLOR_FRAME = RGBColor(0x33, 0x33, 0x33)
COLOR_GRID = RGBColor(0xC0, 0xC0, 0xC0)
COLOR_AXIS_END = RGBColor(0x66, 0x66, 0x66)
COLOR_QUADRANT_LABEL = RGBColor(0x99, 0x99, 0x99)
COLOR_TARGET = RGBColor(0xE1, 0x57, 0x59)
COLOR_TARGET_LINE = RGBColor(0x8B, 0x2C, 0x2E)
COLOR_BUBBLE_LINE = RGBColor(0x33, 0x33, 0x33)
COLOR_BULLET = RGBColor(0x2E, 0x4A, 0x6B)

# ─── 共通配色（正本: skills/_common/styles/chart_palette.md） ───
CHART_PALETTE = [
    "#4E79A7", "#F28E2B", "#59A14F", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F",
]
OTHER_COLOR = "#BAB0AC"
TARGET_COLOR = "#E15759"
LABEL_BAR_COLOR = "#4E79A7"
LABEL_BG_COLOR = "#E8EEF5"


def _palette_color(index: int, total: int) -> str:
    if total <= 1:
        return CHART_PALETTE[0]
    return CHART_PALETTE[index % len(CHART_PALETTE)]


DEFAULT_COLORS = CHART_PALETTE

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_SECTION = Pt(16)
FONT_SIZE_AXIS_LABEL = Pt(14)
FONT_SIZE_AXIS_END = Pt(11)
FONT_SIZE_QUADRANT = Pt(12)
FONT_SIZE_BUBBLE = Pt(12)
FONT_SIZE_BUBBLE_TARGET = Pt(13)
FONT_SIZE_ITEM = Pt(13)
FONT_SIZE_SOURCE = Pt(11)

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global PANEL_Y, PANEL_H, LEFT_X, LEFT_W, RIGHT_X, RIGHT_W
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, COLOR_SUBTITLE, COLOR_FRAME
    global COLOR_AXIS_END, COLOR_TARGET, COLOR_TARGET_LINE, COLOR_BUBBLE_LINE
    global COLOR_BULLET
    global CHART_PALETTE, OTHER_COLOR, TARGET_COLOR
    global FONT_NAME_JP, FONT_SIZE_SECTION, FONT_SIZE_AXIS_LABEL, FONT_SIZE_AXIS_END
    global FONT_SIZE_QUADRANT, FONT_SIZE_BUBBLE, FONT_SIZE_BUBBLE_TARGET
    global FONT_SIZE_ITEM, FONT_SIZE_SOURCE

    _THEME = theme

    PANEL_Y = theme.layout("panel_y_in")
    PANEL_H = theme.layout("panel_h_in")
    LEFT_X = theme.layout("left_x_in")
    LEFT_W = theme.layout("left_w_in")
    RIGHT_X = theme.layout("right_x_in")
    RIGHT_W = theme.layout("right_w_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    COLOR_FRAME = theme.color("text")
    COLOR_AXIS_END = theme.color("source")
    COLOR_BUBBLE_LINE = theme.color("text")

    FONT_NAME_JP = theme.font_ea
    body_pt = theme.font_size_body_pt(skill_id=SKILL_ID)

    if theme.id == "roleup":
        COLOR_SUBTITLE = theme.color("subtitle")
        COLOR_TARGET = theme.color("highlight_target")
        COLOR_TARGET_LINE = theme.color("label_bar")
        COLOR_BULLET = theme.color("label_bar")
        CHART_PALETTE = list(theme.chart_palette[:8])
        OTHER_COLOR = theme.hex("highlight_other")
        TARGET_COLOR = theme.hex("highlight_target")
        # roleup C4 許容: [22, 14, 12, 10, 6]
        FONT_SIZE_SECTION = theme.pt("font_size_subtitle_pt")           # 12pt
        FONT_SIZE_AXIS_LABEL = theme.pt("font_size_key_message_pt")     # 14pt
        FONT_SIZE_AXIS_END = body_pt                                    # 10pt
        FONT_SIZE_QUADRANT = theme.pt("font_size_subtitle_pt")          # 12pt
        FONT_SIZE_BUBBLE = theme.pt("font_size_subtitle_pt")            # 12pt
        FONT_SIZE_BUBBLE_TARGET = theme.pt("font_size_subtitle_pt")     # 12pt (太字＋色で強調)
        FONT_SIZE_ITEM = body_pt                                        # 10pt
    else:
        COLOR_SUBTITLE = COLOR_TEXT
        # stella: keep V1 hardcoded values for regression-zero.

    FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")


def _silent_remove_shape(slide, shape_name: str) -> None:
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_section_title(slide, text, left, top, width):
    """セクションタイトル
    stella: 16pt Bold center + 黒下線
    roleup: 12pt Bold left  + subtitle 色, 下線なし
    """
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    align_str = _THEME.layout_rule("subtitle_align", "center") if _THEME is not None else "center"
    p.alignment = PP_ALIGN.LEFT if align_str == "left" else PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SECTION
    run.font.bold = True
    run.font.color.rgb = COLOR_SUBTITLE
    run.font.name = FONT_NAME_JP

    if _THEME is None or _THEME.id == "stellar_aiz":
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_TEXT
        line.line.fill.background()
    return txBox


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 italic=False, font_name=None):
    if font_name is None:
        font_name = FONT_NAME_JP
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def add_rotated_text_box(slide, text, left, top, width, height, font_size,
                          bold=False, color=None, align=PP_ALIGN.CENTER,
                          rotation=-90, font_name=None):
    """Y軸ラベル用: 回転したテキストボックス"""
    if font_name is None:
        font_name = FONT_NAME_JP
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.rotation = rotation
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


# ──────────────────────────────────────────────
# Map Drawing
# ──────────────────────────────────────────────
def draw_positioning_map(slide, data, left, top, width, height):
    """2軸ポジショニングマップを描画する"""
    section_title = data.get("section_title", "ポジショニングマップ")
    add_section_title(slide, section_title, left, top, width)

    map_x = left + MAP_MARGIN_LEFT
    map_y = top + MAP_MARGIN_TOP
    map_w = width - MAP_MARGIN_LEFT - MAP_MARGIN_RIGHT
    map_h = height - MAP_MARGIN_TOP - MAP_MARGIN_BOTTOM

    # マップの外枠
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, map_x, map_y, map_w, map_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    frame.line.color.rgb = COLOR_FRAME
    frame.line.width = Pt(1.0)
    frame.shadow.inherit = False
    frame.text_frame.text = ""

    # 十字ガイドライン（破線）
    v_conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        map_x + map_w // 2, map_y,
        map_x + map_w // 2, map_y + map_h,
    )
    v_conn.line.color.rgb = COLOR_GRID
    v_conn.line.width = Pt(0.75)
    _set_dash_style(v_conn, "dash")

    h_conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        map_x, map_y + map_h // 2,
        map_x + map_w, map_y + map_h // 2,
    )
    h_conn.line.color.rgb = COLOR_GRID
    h_conn.line.width = Pt(0.75)
    _set_dash_style(h_conn, "dash")

    # ── 軸ラベル ──
    x_axis = data.get("x_axis", {})
    y_axis = data.get("y_axis", {})

    x_label = x_axis.get("label", "")
    if x_label:
        add_text_box(
            slide, x_label,
            map_x, map_y + map_h + Inches(0.32),
            map_w, Inches(0.25),
            FONT_SIZE_AXIS_LABEL, bold=True, align=PP_ALIGN.CENTER,
        )
    x_low = x_axis.get("low", "")
    x_high = x_axis.get("high", "")
    if x_low:
        add_text_box(
            slide, f"← {x_low}",
            map_x, map_y + map_h + Inches(0.05),
            Inches(1.8), Inches(0.22),
            FONT_SIZE_AXIS_END, color=COLOR_AXIS_END, align=PP_ALIGN.LEFT,
        )
    if x_high:
        add_text_box(
            slide, f"{x_high} →",
            map_x + map_w - Inches(1.8), map_y + map_h + Inches(0.05),
            Inches(1.8), Inches(0.22),
            FONT_SIZE_AXIS_END, color=COLOR_AXIS_END, align=PP_ALIGN.RIGHT,
        )

    y_label = y_axis.get("label", "")
    if y_label:
        add_rotated_text_box(
            slide, y_label,
            map_x - Inches(1.4), map_y + map_h // 2 - Inches(1.0),
            Inches(2.0), Inches(0.30),
            FONT_SIZE_AXIS_LABEL, bold=True, align=PP_ALIGN.CENTER,
            rotation=-90,
        )

    quadrants = data.get("quadrants", {})
    has_quadrants = bool(quadrants)

    if not has_quadrants:
        y_low = y_axis.get("low", "")
        y_high = y_axis.get("high", "")
        if y_high:
            add_rotated_text_box(
                slide, f"{y_high} →",
                map_x - Inches(0.45), map_y - Inches(0.05),
                Inches(1.2), Inches(0.22),
                FONT_SIZE_AXIS_END, color=COLOR_AXIS_END, align=PP_ALIGN.LEFT,
                rotation=-90,
            )
        if y_low:
            add_rotated_text_box(
                slide, f"← {y_low}",
                map_x - Inches(0.45), map_y + map_h - Inches(1.25),
                Inches(1.2), Inches(0.22),
                FONT_SIZE_AXIS_END, color=COLOR_AXIS_END, align=PP_ALIGN.RIGHT,
                rotation=-90,
            )

    if has_quadrants:
        qlabel_w = Inches(2.2)
        qlabel_h = Inches(0.25)
        qmargin = Inches(0.12)
        # roleup: マップが狭く top 配置のバブルラベルが多発するため、上 quadrant
        # ラベルをマップ枠の外側 (section_title 下) に逃がす。下 quadrant は枠内に維持。
        is_roleup = (_THEME is not None and _THEME.id == "roleup")
        top_q_y = map_y - qlabel_h - Inches(0.02) if is_roleup else map_y + qmargin
        if quadrants.get("top_right"):
            add_text_box(
                slide, quadrants["top_right"],
                map_x + map_w - qlabel_w - qmargin,
                top_q_y,
                qlabel_w, qlabel_h,
                FONT_SIZE_QUADRANT, bold=False, italic=True,
                color=COLOR_QUADRANT_LABEL, align=PP_ALIGN.RIGHT,
            )
        if quadrants.get("top_left"):
            add_text_box(
                slide, quadrants["top_left"],
                map_x + qmargin,
                top_q_y,
                qlabel_w, qlabel_h,
                FONT_SIZE_QUADRANT, bold=False, italic=True,
                color=COLOR_QUADRANT_LABEL, align=PP_ALIGN.LEFT,
            )
        if quadrants.get("bottom_right"):
            add_text_box(
                slide, quadrants["bottom_right"],
                map_x + map_w - qlabel_w - qmargin,
                map_y + map_h - Inches(0.33),
                qlabel_w, qlabel_h,
                FONT_SIZE_QUADRANT, bold=False, italic=True,
                color=COLOR_QUADRANT_LABEL, align=PP_ALIGN.RIGHT,
            )
        if quadrants.get("bottom_left"):
            add_text_box(
                slide, quadrants["bottom_left"],
                map_x + qmargin,
                map_y + map_h - Inches(0.33),
                qlabel_w, qlabel_h,
                FONT_SIZE_QUADRANT, bold=False, italic=True,
                color=COLOR_QUADRANT_LABEL, align=PP_ALIGN.LEFT,
            )

    # ── プレイヤーのバブル配置 ──
    players = data.get("players", [])
    target_company = data.get("target_company")

    x_min = x_axis.get("min", 0)
    x_max = x_axis.get("max", 10)
    y_min = y_axis.get("min", 0)
    y_max = y_axis.get("max", 10)

    sizes = [p.get("size", 1) for p in players]
    min_size = min(sizes) if sizes else 1
    max_size = max(sizes) if sizes else 1
    min_diam = Inches(0.40)
    max_diam = Inches(0.95)

    def compute_diameter(size_val):
        if max_size == min_size:
            return Inches(0.60)
        ratio = (size_val - min_size) / (max_size - min_size)
        return Emu(int(min_diam + (max_diam - min_diam) * ratio))

    for i, p in enumerate(players):
        name = p["name"]
        x_val = p["x"]
        y_val = p["y"]
        size_val = p.get("size", 1)
        is_target = (name == target_company) if target_company else False

        x_ratio = (x_val - x_min) / (x_max - x_min) if x_max != x_min else 0.5
        y_ratio = (y_val - y_min) / (y_max - y_min) if y_max != y_min else 0.5

        diam = compute_diameter(size_val)

        cx = map_x + int(map_w * x_ratio)
        cy = map_y + int(map_h * (1 - y_ratio))

        bubble_x = cx - diam // 2
        bubble_y = cy - diam // 2

        if is_target:
            color = COLOR_TARGET
            line_color = COLOR_TARGET_LINE
            line_w = Pt(2.5)
        else:
            color = hex_to_rgb(_palette_color(i, len(players)))
            line_color = COLOR_BUBBLE_LINE
            line_w = Pt(1.0)

        bubble = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, bubble_x, bubble_y, diam, diam
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = color
        _set_shape_transparency(bubble, 10 if is_target else 30)
        bubble.line.color.rgb = line_color
        bubble.line.width = line_w
        bubble.shadow.inherit = False
        bubble.text_frame.text = ""

        label_pos = p.get("label_position", "bottom")
        # roleup map is narrower (A4 horizontal); shrink labels to avoid overlap.
        if _THEME is not None and _THEME.id == "roleup":
            label_w = Inches(0.95)
            label_h = Inches(0.20)
        else:
            label_w = Inches(1.8)
            label_h = Inches(0.26)
        label_gap = Inches(0.03)

        if label_pos == "top":
            label_x = cx - label_w // 2
            label_y = cy - diam // 2 - label_h - label_gap
        elif label_pos == "left":
            label_x = bubble_x - label_w - label_gap
            label_y = cy - label_h // 2
        elif label_pos == "right":
            label_x = bubble_x + diam + label_gap
            label_y = cy - label_h // 2
        else:
            label_x = cx - label_w // 2
            label_y = cy + diam // 2 + label_gap

        # roleup はバブルクラスタが密 (A4 横の狭いマップ) なため、上方向の境界に
        # 微少な許容を持たせ、"top" ラベルを fallback で他バブルと衝突させない。
        top_tolerance = Inches(0.30) if (_THEME is not None and _THEME.id == "roleup") else Emu(0)
        if label_y + label_h > map_y + map_h:
            label_x = cx - label_w // 2
            label_y = cy - diam // 2 - label_h - label_gap
        if label_y < map_y - top_tolerance:
            label_x = cx - label_w // 2
            label_y = cy + diam // 2 + label_gap

        label_tb = slide.shapes.add_textbox(label_x, label_y, label_w, label_h)
        ltf = label_tb.text_frame
        ltf.word_wrap = False
        ltf.margin_left = 0; ltf.margin_right = 0; ltf.margin_top = 0; ltf.margin_bottom = 0
        lp = ltf.paragraphs[0]
        if label_pos == "left":
            lp.alignment = PP_ALIGN.RIGHT
        elif label_pos == "right":
            lp.alignment = PP_ALIGN.LEFT
        else:
            lp.alignment = PP_ALIGN.CENTER
        lrun = lp.add_run()
        lrun.text = name
        # roleup は target も BUBBLE と同サイズ（C4 許容外の +1pt を避ける）+太字+色で強調
        lrun.font.size = FONT_SIZE_BUBBLE_TARGET if is_target else FONT_SIZE_BUBBLE
        lrun.font.bold = True if is_target else False
        lrun.font.name = FONT_NAME_JP
        lrun.font.color.rgb = COLOR_TARGET_LINE if is_target else COLOR_TEXT

    print(f"  ✓ ポジショニングマップ: {len(players)}プレイヤー配置")


def _set_dash_style(shape, dash_style="dash"):
    ln = shape.line._get_or_add_ln()
    for elem in ln.findall(qn("a:prstDash")):
        ln.remove(elem)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", dash_style)


def _set_shape_transparency(shape, alpha_percent):
    sp_pr = shape.fill._xPr
    solidFill = sp_pr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for elem in srgb.findall(qn("a:alpha")):
        srgb.remove(elem)
    alpha_val = (100 - alpha_percent) * 1000
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(alpha_val))


# ──────────────────────────────────────────────
# Right Panel: Implications
# ──────────────────────────────────────────────
def build_implications_panel(slide, section_title, implications, left, top, width, height):
    """示唆（implications）のブレット項目リスト"""
    add_section_title(slide, section_title, left, top, width)

    if not implications:
        return

    body_top = top + Inches(0.50)
    body_h = height - Inches(0.50)

    tb = slide.shapes.add_textbox(
        left + Inches(0.05), body_top,
        width - Inches(0.10), body_h,
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    bullet_hex = "{:02X}{:02X}{:02X}".format(COLOR_BULLET[0], COLOR_BULLET[1], COLOR_BULLET[2])
    text_hex = "{:02X}{:02X}{:02X}".format(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2])

    for i, item_text in enumerate(implications):
        p_elem = etree.SubElement(tf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"), attrib={
            "marL": "220000",
            "indent": "-220000",
        })
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "600"})

        buChar = etree.SubElement(pPr, qn("a:buChar"), attrib={"char": "●"})
        buFont = etree.SubElement(pPr, qn("a:buFont"), attrib={"typeface": "Arial"})
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
        buClrSolid.set("val", bullet_hex)

        r = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", text_hex)
        t = etree.SubElement(r, qn("a:t"))
        t.text = item_text

    print(f"  ✓ 示唆パネル: {len(implications)}項目")


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "positioning-map")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    require_source(data, theme, skill_id=SKILL_ID)

    _mm = data.get("main_message", "")
    if len(_mm) > 65:
        raise ValueError(
            f"main_message は 65 字以内（受領: {len(_mm)}）: {_mm[:80]}..."
        )

    PLAYERS_MIN = 2
    PLAYERS_MAX = 10
    _players = data.get("players", [])
    if not isinstance(_players, list):
        raise ValueError("players は配列である必要があります")
    if not (PLAYERS_MIN <= len(_players) <= PLAYERS_MAX):
        raise ValueError(
            f"players の要素数は {PLAYERS_MIN}〜{PLAYERS_MAX} の範囲である必要があります"
            f"（受領: {len(_players)}、target_company を含む）"
        )

    prs = Presentation(template_path)
    slide = prs.slides[0]

    top_text = resolve_top_text(data, theme)
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE),
                     sub_text or data.get("chart_title", "ポジショニングマップ"))

    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    draw_positioning_map(slide, data, LEFT_X, PANEL_Y, LEFT_W, PANEL_H)

    implications_title = data.get("implications_title", "ポジショニングからの示唆")
    implications = data.get("implications", [])
    build_implications_panel(
        slide, implications_title, implications,
        RIGHT_X, PANEL_Y, RIGHT_W, PANEL_H,
    )

    source = data.get("source", "")
    if source:
        if theme.id == "stellar_aiz":
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE)
            if source_shape is not None:
                set_textbox_text(source_shape, source)
            else:
                add_text_box(
                    slide, source,
                    SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                    FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                    align=PP_ALIGN.LEFT,
                )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
