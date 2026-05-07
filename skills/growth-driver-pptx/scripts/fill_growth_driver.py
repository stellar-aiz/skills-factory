"""
fill_growth_driver.py — 売上/利益ブリッジ（ウォーターフォール）スライドを生成

Usage:
  python fill_growth_driver.py \
    --data /home/claude/growth_driver_data.json \
    [--template <path>/growth-driver-template.pptx] \
    --output /mnt/user-data/outputs/GrowthDriver_output.pptx \
    [--brand stellar_aiz | roleup]

POSITIVE (緑) / NEGATIVE (赤) はウォーターフォール定番の universal indicator として
両 brand で維持。TOTAL は roleup 時のみ茶系 (#604C3F) に置き換える (brand 一貫性)。
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "growth-driver-pptx"
_THEME = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_PH = "Source 3"  # roleup placeholder

# ── Layout (stella defaults; _apply_theme で roleup A4 横用にスケール) ──
PANEL_Y = Inches(1.55)
PANEL_H = Inches(5.35)

LEFT_X = Inches(0.41)
LEFT_W = Inches(7.80)

RIGHT_X = Inches(8.35)
RIGHT_W = Inches(4.60)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)

# ── Colors (stella defaults; _apply_theme で roleup 用に上書き) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SUBTEXT = RGBColor(0x55, 0x55, 0x55)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Universal indicator (信号色) — POSITIVE / NEGATIVE は両 brand で維持
COLOR_TOTAL = RGBColor(0x2E, 0x4A, 0x6B)  # roleup: #604C3F に上書き (brand 一貫性)
COLOR_POSITIVE = RGBColor(0x1B, 0x7A, 0x3B)
COLOR_NEGATIVE = RGBColor(0xC0, 0x3A, 0x3A)
COLOR_CONNECT = RGBColor(0x99, 0x99, 0x99)

COLOR_PANEL_BG = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_PANEL_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_SECTION = Pt(13)
FONT_SIZE_CATEGORY = Pt(10)
FONT_SIZE_VALUE = Pt(11)
FONT_SIZE_ITEM = Pt(11)
FONT_SIZE_SOURCE = Pt(10)
FONT_SIZE_UNIT = Pt(10)
FONT_SIZE_DETAIL = Pt(10)


def _apply_theme(theme):
    """roleup の場合、レイアウト・色・フォントを brand 仕様に上書きする。"""
    global PANEL_Y, PANEL_H, LEFT_X, LEFT_W, RIGHT_X, RIGHT_W
    global SOURCE_X, SOURCE_Y, SOURCE_W
    global COLOR_TEXT, COLOR_SUBTEXT, COLOR_SOURCE
    global COLOR_TOTAL, COLOR_PANEL_BG, COLOR_PANEL_BORDER
    global FONT_NAME_JP, FONT_SIZE_SECTION, FONT_SIZE_CATEGORY
    global FONT_SIZE_VALUE, FONT_SIZE_ITEM, FONT_SIZE_SOURCE
    global FONT_SIZE_UNIT, FONT_SIZE_DETAIL
    global _THEME
    _THEME = theme

    if theme.id != "roleup":
        return

    # A4 横 (11.69 × 8.27) 用にスケール
    PANEL_Y = Inches(1.55)
    PANEL_H = Inches(5.40)

    LEFT_X = Inches(0.41)
    LEFT_W = Inches(6.40)

    RIGHT_X = Inches(7.05)
    RIGHT_W = Inches(4.23)

    SOURCE_X = Inches(0.41)
    SOURCE_Y = Inches(7.45)
    SOURCE_W = Inches(10.87)

    # roleup 茶系トーン
    COLOR_TEXT = theme.color("text")
    COLOR_SUBTEXT = theme.color("source")
    COLOR_SOURCE = theme.color("source")

    # TOTAL は roleup chart_palette[2] (#604C3F 茶こげ) に置き換え (brand 一貫性)
    palette = list(theme.chart_palette)
    if len(palette) > 2:
        hx = palette[2].lstrip("#")
        COLOR_TOTAL = RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))

    # パネル背景 / 境界線も brand 一貫性で薄茶へ
    COLOR_PANEL_BG = theme.color("label_bg")     # #F2E8DD
    pal_border = palette[7] if len(palette) > 7 else "#CDCECE"
    hx = pal_border.lstrip("#")
    COLOR_PANEL_BORDER = RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))

    # フォント
    FONT_NAME_JP = theme.font_ea or "Yu Gothic UI"

    # roleup C4 許容集合 [22, 14, 12, 10, 6] pt
    FONT_SIZE_SECTION = Pt(12)
    FONT_SIZE_CATEGORY = Pt(10)
    FONT_SIZE_VALUE = Pt(10)
    FONT_SIZE_ITEM = Pt(10)
    FONT_SIZE_UNIT = Pt(10)
    FONT_SIZE_DETAIL = Pt(10)
    FONT_SIZE_SOURCE = Pt(int(theme._defaults.get("font_size_source_pt", 6)))


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    """汎用テキストボックス (font_name は late-resolve で _apply_theme 後の値を使う)"""
    if font_name is None:
        font_name = FONT_NAME_JP
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def add_section_title(slide, text, left, top, width):
    add_text_box(
        slide, text,
        left, top, width, Inches(0.30),
        FONT_SIZE_SECTION, bold=True, align=PP_ALIGN.LEFT,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEXT
    line.line.fill.background()


def _set_dash(shape, dash_style="dash"):
    ln = shape.line._get_or_add_ln()
    for elem in ln.findall(qn("a:prstDash")):
        ln.remove(elem)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", dash_style)


def draw_waterfall(slide, data, left, top, width, height):
    start = data.get("start", {})
    drivers = data.get("drivers", [])
    end = data.get("end", {})
    unit = data.get("unit", "")

    section_title = data.get("section_title", "売上ブリッジ")
    add_section_title(slide, section_title, left, top, width)

    if unit:
        add_text_box(
            slide, f"（単位：{unit}）",
            left, top + Inches(0.35), Inches(3.0), Inches(0.25),
            FONT_SIZE_UNIT, bold=False,
            align=PP_ALIGN.LEFT,
        )

    chart_top = top + Inches(0.65)
    chart_h = height - Inches(0.65) - Inches(0.50)
    chart_left = left + Inches(0.10)
    chart_w = width - Inches(0.20)

    all_bars = [
        {"label": start.get("label", ""), "value": start.get("value", 0), "type": "total"},
    ]
    for d in drivers:
        all_bars.append({"label": d.get("label", ""), "value": d.get("value", 0), "type": "driver"})
    all_bars.append({"label": end.get("label", ""), "value": end.get("value", 0), "type": "total"})

    n_bars = len(all_bars)

    cumulative = [start.get("value", 0)]
    for d in drivers:
        cumulative.append(cumulative[-1] + d.get("value", 0))

    all_values = cumulative + [start.get("value", 0), end.get("value", 0)]
    y_max = max(all_values) * 1.15
    y_min = min(0, min(all_values) * 0.9)
    y_range = y_max - y_min if y_max != y_min else 1

    gap = Emu(int(chart_w / (n_bars * 8)))
    bar_w = Emu(int((chart_w - gap * (n_bars - 1)) / n_bars))

    def y_for(value):
        ratio = (value - y_min) / y_range
        return chart_top + Emu(int(chart_h * (1 - ratio)))

    for i, bar in enumerate(all_bars):
        bar_x = chart_left + (bar_w + gap) * i
        bar_type = bar["type"]
        value = bar["value"]
        label = bar["label"]

        if bar_type == "total":
            top_y = y_for(max(value, 0))
            bottom_y = y_for(min(value, 0))
            h = bottom_y - top_y
            color = COLOR_TOTAL
            label_value_text = f"{value:,.0f}"
        else:
            prev_cum = cumulative[i - 1]
            new_cum = cumulative[i - 1] + value
            top_y = y_for(max(prev_cum, new_cum))
            bottom_y = y_for(min(prev_cum, new_cum))
            h = bottom_y - top_y
            color = COLOR_POSITIVE if value >= 0 else COLOR_NEGATIVE
            label_value_text = f"+{value:,.0f}" if value > 0 else f"{value:,.0f}"

        if h < Emu(500):
            h = Emu(500)

        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, bar_x, top_y, bar_w, h,
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = color
        bar_shape.line.fill.background()
        bar_shape.shadow.inherit = False
        bar_shape.text_frame.text = ""

        label_h = Inches(0.22)
        if bar_type == "total" or value >= 0:
            label_y = top_y - label_h - Inches(0.02)
        else:
            label_y = bottom_y + Inches(0.02)

        lbl_color = color if bar_type != "total" else COLOR_TEXT
        add_text_box(
            slide, label_value_text,
            bar_x, label_y, bar_w, label_h,
            FONT_SIZE_VALUE, bold=True,
            color=lbl_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        cat_y = chart_top + chart_h + Inches(0.05)
        cat_h = Inches(0.40)
        add_text_box(
            slide, label,
            bar_x, cat_y, bar_w, cat_h,
            FONT_SIZE_CATEGORY, bold=(bar_type == "total"),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.TOP,
        )

        # Drop line
        if i < n_bars - 1:
            if bar_type == "total":
                next_ref_value = value
            else:
                next_ref_value = cumulative[i]
            line_y = y_for(next_ref_value)
            line_start_x = bar_x + bar_w
            line_end_x = bar_x + bar_w + gap
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                line_start_x, line_y, line_end_x, line_y,
            )
            line.line.color.rgb = COLOR_CONNECT
            line.line.width = Pt(0.75)
            _set_dash(line, "dash")


def draw_drivers_detail(slide, drivers, left, top, width, height):
    add_section_title(slide, "主要ドライバーの解説", left, top, width)

    panel_top = top + Inches(0.45)
    panel_h = height - Inches(0.45)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, panel_top, width, panel_h,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PANEL_BG
    bg.line.color.rgb = COLOR_PANEL_BORDER
    bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    bg.text_frame.text = ""

    tb = slide.shapes.add_textbox(
        left + Inches(0.15), panel_top + Inches(0.12),
        width - Inches(0.30), panel_h - Inches(0.24),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    text_hex = "{:02X}{:02X}{:02X}".format(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2])
    subtext_hex = "{:02X}{:02X}{:02X}".format(COLOR_SUBTEXT[0], COLOR_SUBTEXT[1], COLOR_SUBTEXT[2])

    for i, d in enumerate(drivers):
        # チャート用のラベルには改行が含まれることがあるが、右パネルでは1行にする
        label = d.get("label", "").replace("\n", " ").strip()
        value = d.get("value", 0)
        detail = d.get("detail", "")

        # POSITIVE / NEGATIVE は universal indicator (信号色) として両 brand で維持
        color_val = "1B7A3B" if value >= 0 else "C03A3A"
        sign = "+" if value > 0 else ""

        p_elem = etree.SubElement(tf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "600"})

        r0 = etree.SubElement(p_elem, qn("a:r"))
        rPr0 = etree.SubElement(r0, qn("a:rPr"), attrib={
            "lang": "en-US",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr0, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr0, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf0 = etree.SubElement(rPr0, qn("a:solidFill"))
        s0 = etree.SubElement(sf0, qn("a:srgbClr"))
        s0.set("val", color_val)
        t0 = etree.SubElement(r0, qn("a:t"))
        t0.text = "● "

        r1 = etree.SubElement(p_elem, qn("a:r"))
        rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
        s1 = etree.SubElement(sf1, qn("a:srgbClr"))
        s1.set("val", text_hex)
        t1 = etree.SubElement(r1, qn("a:t"))
        t1.text = f"{label}  "

        r2 = etree.SubElement(p_elem, qn("a:r"))
        rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={
            "lang": "en-US",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr2, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
        s2 = etree.SubElement(sf2, qn("a:srgbClr"))
        s2.set("val", color_val)
        t2 = etree.SubElement(r2, qn("a:t"))
        t2.text = f"({sign}{value:,.0f})"

        if detail:
            p2 = etree.SubElement(tf._txBody, qn("a:p"))
            pPr2 = etree.SubElement(p2, qn("a:pPr"), attrib={
                "marL": "180000",
                "indent": "0",
            })
            spcBef2 = etree.SubElement(pPr2, qn("a:spcBef"))
            etree.SubElement(spcBef2, qn("a:spcPts"), attrib={"val": "150"})

            r3 = etree.SubElement(p2, qn("a:r"))
            rPr3 = etree.SubElement(r3, qn("a:rPr"), attrib={
                "lang": "ja-JP",
                "sz": str(int(FONT_SIZE_DETAIL.pt * 100)),
                "b": "0",
            })
            etree.SubElement(rPr3, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
            etree.SubElement(rPr3, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
            sf3 = etree.SubElement(rPr3, qn("a:solidFill"))
            s3 = etree.SubElement(sf3, qn("a:srgbClr"))
            s3.set("val", subtext_hex)
            t3 = etree.SubElement(r3, qn("a:t"))
            t3.text = detail

    print(f"  ✓ ドライバー詳細: {len(drivers)}項目")


def write_source(slide, source):
    """出典を Source 3 placeholder (roleup) または dynamic textbox (stella) に書き込む。"""
    if not source:
        return

    if _THEME is not None and _THEME.is_source_required():
        src_shape = find_shape(slide, SHAPE_SOURCE_PH)
        if src_shape is not None:
            set_textbox_text(src_shape, f"出典: {source}")
            for para in src_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = FONT_SIZE_SOURCE
                    run.font.color.rgb = COLOR_SOURCE
                    run.font.name = FONT_NAME_JP
            return

    add_text_box(
        slide, source,
        SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.25),
        FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
        align=PP_ALIGN.LEFT,
    )


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=False, default=None,
                    help="growth-driver-template.pptx のパス (省略時は --brand から自動解決)")
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "waterfall"],
        allowed_top=[
            "main_message", "chart_title", "waterfall",
            "source", "source_label", "source_text",
        ],
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "growth-driver")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top text (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
        print(f"  ✓ Top: {top_text[:60]}{'...' if len(top_text) > 60 else ''}")

    # Subtitle (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme).strip() or "成長ドライバー分析"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle: {sub_text}")

    waterfall = data.get("waterfall", {})
    if not waterfall:
        print("  ✗ ERROR: 'waterfall' is required", file=sys.stderr)
        sys.exit(1)

    draw_waterfall(slide, waterfall, LEFT_X, PANEL_Y, LEFT_W, PANEL_H)
    print(f"  ✓ ウォーターフォールチャート")

    drivers = waterfall.get("drivers", [])
    if drivers:
        draw_drivers_detail(slide, drivers, RIGHT_X, PANEL_Y, RIGHT_W, PANEL_H)

    # Source
    source = (data.get("source") or data.get("source_label")
              or data.get("source_text") or "").strip()
    write_source(slide, source)
    if source:
        print(f"  ✓ Source: {source[:60]}{'...' if len(source) > 60 else ''}")

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
