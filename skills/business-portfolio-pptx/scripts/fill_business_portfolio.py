"""
fill_business_portfolio.py — 事業ポートフォリオスライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 左側: セグメント別売上高の積み上げ棒グラフ（絶対値、複数年）
  - 右側: セグメント別サマリーテーブル（最新期売上、構成比、CAGR、営業利益率）
  - 下部: 出典

Usage:
  python fill_business_portfolio.py \
    --data /home/claude/business_portfolio_data.json \
    --template <path>/business-portfolio-template.pptx \
    --output /mnt/user-data/outputs/BusinessPortfolio_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants (16:9, 13.33 x 7.5 in) ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

PANEL_Y = Inches(1.50)

# Left panel (chart)
LEFT_X = Inches(0.41)
LEFT_W = Inches(6.30)
LEFT_H = Inches(5.20)

# Right panel (table)
RIGHT_X = Inches(7.00)
RIGHT_W = Inches(5.90)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)

# ── Colors ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_HEADER_BG = RGBColor(0x2E, 0x4A, 0x6B)   # Navy header
COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)

# Default segment colors (used when not specified)
DEFAULT_COLORS = [
    "#4E79A7", "#F28E2B", "#59A14F", "#E15759", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F", "#BAB0AC",
]

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_SECTION = Pt(14)
FONT_SIZE_TABLE = Pt(11)
FONT_SIZE_TABLE_HEADER = Pt(11)
FONT_SIZE_SOURCE = Pt(10)
FONT_SIZE_UNIT = Pt(11)


# ──────────────────────────────────────────────
# Utility Functions
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def hex_to_rgb(hex_str):
    """#RRGGBB → RGBColor"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_section_title(slide, text, left, top, width):
    """セクションタイトル（下線付き、Bold、14pt）"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SECTION
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    # 下線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEXT
    line.line.fill.background()
    return txBox


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """汎用テキストボックス"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = FONT_NAME_JP
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


# ──────────────────────────────────────────────
# Left Panel: Stacked Bar Chart
# ──────────────────────────────────────────────
def build_stacked_bar_chart(slide, section_title, chart_data, segments, left, top, width, height):
    """
    セクションタイトル + 単位表記 + 積み上げ棒チャート
    """
    # セクションタイトル
    add_section_title(slide, section_title, left, top, width)

    # 単位表記（左）
    unit_label = chart_data.get("unit_label", "（単位：億円）")
    add_text_box(
        slide, unit_label,
        left, top + Inches(0.40), Inches(2.0), Inches(0.25),
        FONT_SIZE_UNIT, bold=False, align=PP_ALIGN.LEFT,
    )

    # チャートエリア（タイトル・単位の下、残りの高さを使う）
    chart_top = top + Inches(0.70)
    chart_h = height - Inches(0.70)

    # データ準備
    years = chart_data["years"]
    values = chart_data["values"]  # {"事業A": [100, 110, ...], ...}

    cdata = CategoryChartData()
    cdata.categories = years

    seg_names = [seg["name"] for seg in segments]
    for seg_name in seg_names:
        series_vals = values.get(seg_name, [0] * len(years))
        cdata.add_series(seg_name, series_vals)

    # チャート挿入（積み上げ縦棒）
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        left, chart_top, width, chart_h,
        cdata,
    )
    chart = chart_shape.chart

    # 凡例設定
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT_NAME_JP

    # タイトル非表示
    chart.has_title = False

    # 系列ごとに色を設定 + データラベル（全てCENTER配置でシンプルに）
    for idx, series in enumerate(chart.series):
        seg = segments[idx] if idx < len(segments) else {}
        hex_color = seg.get("color") or DEFAULT_COLORS[idx % len(DEFAULT_COLORS)]
        rgb = hex_to_rgb(hex_color)

        # 塗りつぶし色
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = rgb

        # 枠線なし
        series.format.line.fill.background()

        # データラベル（中央、白文字）
        series.data_labels.show_value = True
        series.data_labels.position = XL_LABEL_POSITION.CENTER
        series.data_labels.font.size = Pt(9)
        series.data_labels.font.name = FONT_NAME_JP
        series.data_labels.font.bold = True
        series.data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        series.data_labels.number_format = '#,##0'
        series.data_labels.number_format_is_linked = False

    # 軸フォーマット
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.name = FONT_NAME_JP

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.name = FONT_NAME_JP
    # 目盛線を消す
    val_axis.visible = False
    val_axis.major_tick_mark = 2  # XL_TICK_MARK.NONE

    # Y軸最大値を自動計算（データラベルが切れないよう年ごとの合計×1.20でパディング）
    year_totals = []
    for y_idx in range(len(years)):
        total = sum(values.get(seg_name, [0] * len(years))[y_idx] for seg_name in seg_names)
        year_totals.append(total)
    max_total = max(year_totals) if year_totals else 0
    if max_total > 0:
        val_axis.maximum_scale = max_total * 1.20
        val_axis.minimum_scale = 0

    # 各年の合計を棒の上に表示するため、透明な "合計" 系列を追加するのは複雑なので省略
    # 代わりに積み上げ棒のまま、セグメント値のデータラベルで表現

    print(f"  ✓ 積み上げ棒チャート: {len(seg_names)}セグメント x {len(years)}年")
    return chart_shape


# ──────────────────────────────────────────────
# Right Panel: Summary Table
# ──────────────────────────────────────────────
def build_summary_table(slide, table_data, segments, left, top, width):
    """
    セクションタイトル + サマリーテーブル
    """
    section_title = table_data.get("section_title", "セグメント別サマリー")
    add_section_title(slide, section_title, left, top, width)

    headers = table_data["headers"]
    rows = table_data["rows"]

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    tbl_top = top + Inches(0.55)
    row_h = Inches(0.38)
    tbl_h = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    table = shape.table

    # 列幅: 1列目（セグメント名）を広く、残りは均等
    first_col_w = int(width * 0.32)
    other_col_w = int((width - first_col_w) / (n_cols - 1))
    table.columns[0].width = Emu(first_col_w)
    for i in range(1, n_cols):
        table.columns[i].width = Emu(other_col_w)

    # tblPr設定: バンド無効化
    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # 行の高さ
    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(row_h))

    # ヘッダー行
    for c_idx, h in enumerate(headers):
        cell = table.cell(0, c_idx)
        _style_cell(cell, h, is_header=True, is_alt=False,
                    font_size=FONT_SIZE_TABLE_HEADER)

    # データ行
    seg_color_map = {seg["name"]: seg.get("color") for seg in segments}
    for r_idx, row in enumerate(rows):
        is_alt = (r_idx % 2 == 1)
        for c_idx, h in enumerate(headers):
            # セグメント名列に色のマーカーを付ける場合は name をそのまま使う
            key = _header_to_key(h, c_idx)
            val = row.get(key, "")
            if c_idx == 0:
                # セグメント名列は Bold + セグメント色の小さな四角を前に付ける
                cell = table.cell(r_idx + 1, c_idx)
                _style_cell(cell, val, is_header=False, is_alt=is_alt,
                            font_size=FONT_SIZE_TABLE, bold=True,
                            color_marker_hex=seg_color_map.get(val))
            else:
                cell = table.cell(r_idx + 1, c_idx)
                _style_cell(cell, val, is_header=False, is_alt=is_alt,
                            font_size=FONT_SIZE_TABLE)

    print(f"  ✓ サマリーテーブル: {n_rows}行 x {n_cols}列")
    return shape


def _header_to_key(header, col_idx):
    """ヘッダー名 → row の key。柔軟なキー名を許容する"""
    mapping = {
        "セグメント": "name",
        "事業セグメント": "name",
        "セグメント名": "name",
        "事業": "name",
        "売上": "revenue",
        "売上高": "revenue",
        "最新期売上": "revenue",
        "構成比": "share",
        "シェア": "share",
        "CAGR": "cagr",
        "成長率": "cagr",
        "利益率": "op_margin",
        "営業利益率": "op_margin",
        "営業利益": "op_profit",
    }
    return mapping.get(header, f"col{col_idx}")


def _style_cell(cell, text, is_header=False, is_alt=False, font_size=Pt(11),
                bold=False, color_marker_hex=None):
    """セルのスタイルとテキストを設定"""
    tc = cell._tc

    # 塗りつぶし
    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
    elif is_alt:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ROW_ALT
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # マージン、縦中央寄せ
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # テキスト設定
    tf = cell.text_frame
    tf.word_wrap = True
    # 既存段落をクリア
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    # 新段落追加
    p_elem = etree.SubElement(tf._txBody, qn("a:p"))

    # pPr (段落プロパティ)
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr" if is_header else ("l" if bold else "r"))  # ヘッダー中央、name列左、数値右寄せ
    if not bold and not is_header:
        pPr.set("algn", "r")

    # カラーマーカー（セグメント色の小さな四角）を追加する場合は文字列の前に ■ を付ける
    display_text = text
    if color_marker_hex:
        display_text = f"■ {text}"

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP", "sz": str(font_size.pt * 100)})
    rPr.set("b", "1" if (is_header or bold) else "0")
    # Font
    rFont = etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    rFontEA = etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    # 色
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    if is_header:
        srgb.set("val", "FFFFFF")
    else:
        srgb.set("val", "333333")

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = display_text

    # カラーマーカー部分だけ色を変更する場合、さらに複雑化するので省略
    # (■ の色は上で設定した #333333 になる。将来改良可能)
    if color_marker_hex:
        # ■ 部分だけ色を変えるため、run を2つに分割
        # 既存 r_elem を削除して作り直す
        p_elem.remove(r_elem)

        # Run 1: ■ (セグメント色)
        marker_rgb = color_marker_hex.lstrip("#")
        r1 = etree.SubElement(p_elem, qn("a:r"))
        rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={"lang": "ja-JP", "sz": str(font_size.pt * 100)})
        rPr1.set("b", "1")
        etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
        s1 = etree.SubElement(sf1, qn("a:srgbClr"))
        s1.set("val", marker_rgb)
        t1 = etree.SubElement(r1, qn("a:t"))
        t1.text = "■ "

        # Run 2: text (通常色)
        r2 = etree.SubElement(p_elem, qn("a:r"))
        rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={"lang": "ja-JP", "sz": str(font_size.pt * 100)})
        rPr2.set("b", "1")
        etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr2, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
        s2 = etree.SubElement(sf2, qn("a:srgbClr"))
        s2.set("val", "333333")
        t2 = etree.SubElement(r2, qn("a:t"))
        t2.text = text


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True, help="Path to JSON data file")
    ap.add_argument("--template", required=True, help="Path to template PPTX")
    ap.add_argument("--output", required=True, help="Path to output PPTX")
    add_brand_arg(ap)  # passive: accepted but ignored until brand migration
    args = ap.parse_args()

    # データ読み込み
    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # テンプレート読み込み
    prs = Presentation(args.template)
    slide = prs.slides[0]

    # メインメッセージ
    main_msg = data.get("main_message", "")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg)
    print(f"  ✓ Main Message: {main_msg[:40]}...")

    # チャートタイトル
    chart_title = data.get("chart_title", "事業ポートフォリオ")
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
    print(f"  ✓ Chart Title: {chart_title}")

    # セグメント情報
    segments = data.get("segments", [])
    if not segments:
        print("  ✗ ERROR: 'segments' is required", file=sys.stderr)
        sys.exit(1)

    # デフォルト色を割り当て
    for i, seg in enumerate(segments):
        if not seg.get("color"):
            seg["color"] = DEFAULT_COLORS[i % len(DEFAULT_COLORS)]

    # 左: 積み上げ棒チャート
    chart_data = data.get("chart", {})
    left_section_title = chart_data.get("section_title", "セグメント別売上高推移")
    build_stacked_bar_chart(
        slide, left_section_title, chart_data, segments,
        LEFT_X, PANEL_Y, LEFT_W, LEFT_H,
    )

    # 右: サマリーテーブル
    table_data = data.get("table", {})
    if table_data:
        build_summary_table(slide, table_data, segments, RIGHT_X, PANEL_Y, RIGHT_W)

    # 出典
    source = data.get("source", "")
    if source:
        add_text_box(
            slide, source,
            SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.30),
            FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
            align=PP_ALIGN.LEFT,
        )
        print(f"  ✓ Source: {source[:40]}...")

    # 保存
    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
