"""
fill_executive_summary.py — エグゼクティブサマリースライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 中央: 3〜5個のKey Findingsを縦積みで表示
      各Finding: 番号バッジ + カテゴリラベル + 見出し + 詳細テキスト
  - 下部: 出典

Usage:
  python fill_executive_summary.py \
    --data /home/claude/executive_summary_data.json \
    --template <path>/executive-summary-template.pptx \
    --output /mnt/user-data/outputs/ExecutiveSummary_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "executive-summary-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE = "Source 3"  # roleup template placeholder; stella adds dynamic textbox

# Default values (stella) — reassigned in main() via _apply_theme(theme).
# Findings grid placement.
GRID_LEFT = Inches(0.41)
GRID_TOP = Inches(1.55)
GRID_WIDTH = Inches(12.51)
GRID_HEIGHT = Inches(5.35)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(6.93)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.25)

# ── Colors ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTEXT = RGBColor(0x55, 0x55, 0x55)

# カテゴリ色マッピング（findingの category フィールドで使用）
CATEGORY_COLORS = {
    "対象会社": RGBColor(0x2E, 0x4A, 0x6B),       # 紺
    "company": RGBColor(0x2E, 0x4A, 0x6B),
    "target": RGBColor(0x2E, 0x4A, 0x6B),
    "マクロ環境": RGBColor(0x7B, 0x4F, 0xB0),      # 紫
    "macro": RGBColor(0x7B, 0x4F, 0xB0),
    "pest": RGBColor(0x7B, 0x4F, 0xB0),
    "市場": RGBColor(0x2E, 0x6F, 0xBF),           # 青
    "market": RGBColor(0x2E, 0x6F, 0xBF),
    "競合": RGBColor(0xDA, 0x7A, 0x2D),           # オレンジ
    "competitor": RGBColor(0xDA, 0x7A, 0x2D),
    "財務": RGBColor(0x3D, 0x8F, 0x5A),           # 緑
    "financial": RGBColor(0x3D, 0x8F, 0x5A),
    "リスク": RGBColor(0xB8, 0x3A, 0x3A),         # 赤
    "risk": RGBColor(0xB8, 0x3A, 0x3A),
    "機会": RGBColor(0x1B, 0x7A, 0x3B),           # 濃緑
    "opportunity": RGBColor(0x1B, 0x7A, 0x3B),
    "示唆": RGBColor(0x55, 0x55, 0x55),           # グレー
    "implication": RGBColor(0x55, 0x55, 0x55),
    "結論": RGBColor(0x33, 0x33, 0x33),           # 濃グレー
    "conclusion": RGBColor(0x33, 0x33, 0x33),
}

DEFAULT_COLOR = RGBColor(0x4E, 0x79, 0xA7)  # デフォルト: 紺系

# ─── 共通配色（正本: skills/_common/styles/chart_palette.md） ───
# 編集時は _common/styles/chart_palette.md と他 4 スキルの fill_*.py も同期更新
# CHART_PALETTE には TARGET_COLOR(赤) と OTHER_COLOR(灰) を含めない（palette 外で固定）
CHART_PALETTE = [
    "#4E79A7", "#F28E2B", "#59A14F", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F",
]
OTHER_COLOR = "#BAB0AC"
TARGET_COLOR = "#E15759"
LABEL_BAR_COLOR = "#4E79A7"
LABEL_BG_COLOR = "#E8EEF5"
LABEL_BAR_RGB = RGBColor(0x4E, 0x79, 0xA7)  # P1 の category バー（▍）の単色青


def _palette_color(index: int, total: int) -> str:
    if total <= 1:
        return CHART_PALETTE[0]
    return CHART_PALETTE[index % len(CHART_PALETTE)]

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_BADGE = Pt(20)   # 互換のため残置（draw_number_badgeは未使用）
FONT_SIZE_CATEGORY = Pt(16)
FONT_SIZE_HEADING = Pt(16)
FONT_SIZE_DETAIL = Pt(13)
FONT_SIZE_SOURCE = Pt(11)

# Theme module-global; populated in main() via _apply_theme(theme) so helper
# functions (require_source / resolve_top_text) can read it without an extra
# parameter passed through 200+ lines of fill code.
_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme.

    Called once from main() after `--brand` is parsed. Only invoked from main(),
    so other functions see the post-_apply_theme values.
    """
    global _THEME
    global GRID_LEFT, GRID_TOP, GRID_WIDTH, GRID_HEIGHT
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, COLOR_SUBTEXT, LABEL_BAR_RGB, LABEL_BAR_COLOR, LABEL_BG_COLOR
    global FONT_NAME_JP, FONT_SIZE_CATEGORY, FONT_SIZE_HEADING, FONT_SIZE_DETAIL, FONT_SIZE_SOURCE

    _THEME = theme

    GRID_LEFT = theme.layout("grid_left_in")
    GRID_TOP = theme.layout("grid_top_in")
    GRID_WIDTH = theme.layout("grid_width_in")
    GRID_HEIGHT = theme.layout("grid_height_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    COLOR_SUBTEXT = theme.color("text")  # subtext follows text color in both brands
    LABEL_BAR_RGB = theme.color("label_bar")
    LABEL_BAR_COLOR = theme.hex("label_bar")
    LABEL_BG_COLOR = theme.hex("label_bg")

    FONT_NAME_JP = theme.font_ea
    # Title / chart_title placeholder font sizes are set by template; here we
    # only pin the body / category / detail / source sizes that fill controls.
    # Roleup uses 12pt for executive-summary body (per executive_summary_skill_ids).
    body_pt = theme.font_size_body_pt(skill_id=SKILL_ID)
    FONT_SIZE_CATEGORY = body_pt
    FONT_SIZE_HEADING = body_pt
    FONT_SIZE_DETAIL = body_pt
    FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")


def _silent_remove_shape(slide, shape_name: str) -> None:
    """Remove a shape by name without printing a warning. No-op if absent.

    Used for roleup brown rect guides that the template carries for visual
    consistency but exec-summary's dynamic finding grid does not need.
    """
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 italic=False, font_name=None):
    # font_name / color default to current module globals (post-_apply_theme).
    # Using None sentinel + late binding avoids capturing stella's pre-theme
    # values when the function is defined at module load time.
    if font_name is None:
        font_name = FONT_NAME_JP
    if color is None:
        color = COLOR_TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_category_color(category):
    """カテゴリ名から色を取得。
    新仕様（2026-04-29）: カラフルな意味付き色分けは廃止し、全カテゴリで単色青（LABEL_BAR_COLOR）に統一。
    旧 CATEGORY_COLORS dict は削除済み。
    """
    return LABEL_BAR_RGB


# ──────────────────────────────────────────────
# Number Badge (circular)
# ──────────────────────────────────────────────
def draw_number_badge(slide, number, color, left, top, diameter):
    """番号入りの円形バッジを描画"""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter,
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.shadow.inherit = False

    # テキスト設定（番号）
    tf = badge.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 既存段落クリア
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr")

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_BADGE.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "Arial"})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    s = etree.SubElement(sf, qn("a:srgbClr"))
    s.set("val", "FFFFFF")
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = f"{number:02d}"


# ──────────────────────────────────────────────
# Finding Row
# ──────────────────────────────────────────────
def draw_finding(slide, idx, finding, left, top, width, height):
    """
    1つのFindingを描画する（番号バッジは廃止）
    レイアウト:
      [▎] [カテゴリラベル] [見出し (Bold)]
          [詳細テキスト                ]
    """
    category = finding.get("category", "")
    color = None
    color_hex = finding.get("color")
    if color_hex:
        color = hex_to_rgb(color_hex)
    else:
        color = get_category_color(category)

    # 左側の縦バー（カテゴリ色）— 番号バッジを廃止し、左端から開始
    bar_left = left
    bar_top = top
    bar_w = Inches(0.08)
    bar_h = height - Inches(0.10)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_w, bar_h,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.text_frame.text = ""

    # テキスト領域開始
    content_left = bar_left + bar_w + Inches(0.15)
    content_w = width - (content_left - left)

    # 上段: カテゴリラベル + 見出し
    top_row_top = top + Inches(0.02)
    top_row_h = Inches(0.42)

    # カテゴリラベル（小さな色付きタグ風）
    if category:
        # タグ幅を動的に決定（フォント12ptに合わせて広めに）
        cat_label_w = Inches(1.80)
        cat_tb = slide.shapes.add_textbox(
            content_left, top_row_top, cat_label_w, top_row_h,
        )
        ctf = cat_tb.text_frame
        ctf.margin_left = 0; ctf.margin_right = 0
        ctf.margin_top = 0; ctf.margin_bottom = 0
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for p in list(ctf.paragraphs):
            p._p.getparent().remove(p._p)

        p_elem = etree.SubElement(ctf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", "l")

        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_CATEGORY.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = f"▍ {category}"

        heading_left = content_left + cat_label_w
        heading_w = content_w - cat_label_w
    else:
        heading_left = content_left
        heading_w = content_w

    # 見出し
    heading = finding.get("heading", "")
    if heading:
        add_text_box(
            slide, heading,
            heading_left, top_row_top, heading_w, top_row_h,
            FONT_SIZE_HEADING, bold=True,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

    # 下段: 詳細テキスト
    detail = finding.get("detail", "")
    if detail:
        detail_top = top_row_top + top_row_h + Inches(0.05)
        detail_h = height - (detail_top - top) - Inches(0.10)
        add_text_box(
            slide, detail,
            content_left, detail_top, content_w, detail_h,
            FONT_SIZE_DETAIL, bold=False,
            color=COLOR_SUBTEXT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
# Validation
# ──────────────────────────────────────────────
MAIN_MESSAGE_MAX = 65
CATEGORY_MAX = 8
DETAIL_MAX = 100


def _validate_input(data):
    """main_message ≤65字、category ≤8字＆ユニーク、detail ≤100字。"""
    main_message = data.get("main_message", "")
    if len(main_message) > MAIN_MESSAGE_MAX:
        raise ValueError(
            f"main_message は {MAIN_MESSAGE_MAX} 字以内（受領: {len(main_message)}）: {main_message[:80]}..."
        )
    findings = data.get("findings", [])
    categories = []
    for i, f in enumerate(findings):
        cat = f.get("category", "")
        if cat and len(cat) > CATEGORY_MAX:
            raise ValueError(
                f"findings[{i}].category は {CATEGORY_MAX} 字以内（受領: {len(cat)}）: {cat}"
            )
        categories.append(cat)
        detail = f.get("detail", "")
        if len(detail) > DETAIL_MAX:
            raise ValueError(
                f"findings[{i}].detail は {DETAIL_MAX} 字以内"
                f"（受領: {len(detail)}）: {detail[:80]}..."
            )
    nonempty = [c for c in categories if c]
    if len(set(nonempty)) != len(nonempty):
        dupes = sorted({c for c in nonempty if nonempty.count(c) > 1})
        raise ValueError(
            f"findings の category は全てユニークである必要があります。重複: {dupes} "
            f"（例：「市場」が複数 finding に出る場合は「市場規模」「市場成長」のように区別すること）"
        )


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "executive-summary")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # Roleup: source field is required (hard-fail). Stella: no-op (no requirement).
    require_source(data, theme, skill_id=SKILL_ID)
    _validate_input(data)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top / subtitle placeholder semantics differ between brands:
    #  - stella: Title 1 = main_message (結論文), Text Placeholder 2 = chart_title
    #  - roleup: Title 1 = chart_title (スライドタイトル), Text Placeholder 2 = main_message
    top_text = resolve_top_text(data, theme)
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:40]}")
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:40]}")

    # Roleup: silently remove brown guide rectangles carried by template
    # (matches cp/me/ch convention; stella template has no such shapes → no-op).
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    findings = data.get("findings", [])
    if not findings:
        print("  ✗ ERROR: 'findings' is required", file=sys.stderr)
        sys.exit(1)

    n = len(findings)
    if n > 6:
        print(f"  ⚠ WARNING: {n} findings > 6. Only first 6 will be shown.", file=sys.stderr)
        findings = findings[:6]
        n = 6

    # 各Findingの高さを動的に計算
    # ギャップは等間隔
    total_h = GRID_HEIGHT
    gap_h = Inches(0.12)
    total_gap = gap_h * (n - 1) if n > 1 else Emu(0)
    finding_h = Emu(int((total_h - total_gap) / n))

    for i, f in enumerate(findings):
        top = GRID_TOP + (finding_h + gap_h) * i
        draw_finding(slide, i + 1, f, GRID_LEFT, top, GRID_WIDTH, finding_h)
        print(f"  ✓ Finding {i+1}: {f.get('heading', '')[:40]}...")

    # 出典: roleup は Source 3 placeholder にセット、stella は dynamic textbox を追加
    source = data.get("source", "")
    if source:
        if theme.id == "stellar_aiz":
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE)
            if source_shape is not None:
                set_textbox_text(source_shape, source)
            else:
                # Brand has no Source 3 placeholder — fall back to dynamic textbox.
                add_text_box(
                    slide, source,
                    SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                    FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                    align=PP_ALIGN.LEFT,
                )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
