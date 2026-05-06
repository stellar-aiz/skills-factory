"""
fill_company_overview.py — 会社概要データをPPTXネイティブオブジェクトで生成するスクリプト

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

テンプレート構造:
  Stella (assets/stellar_aiz/company-overview-template.pptx):
    - Title 1            (PLACEHOLDER): Main Message (上段、太字)
    - Text Placeholder 2 (PLACEHOLDER): Chart Title (下段)
    - Overview Table     (TABLE):       会社概要テーブル
    - Photo Caption 1/2  (TEXT_BOX):    写真キャプション
    - Photo Area 1/2     (AUTO_SHAPE):  写真エリア
    - Source             (TEXT_BOX):    出典
  Roleup (assets/roleup/company-overview-template.pptx,
          tools/setup_company_overview_roleup_template.py で生成):
    - Title 1, Text Placeholder 2, Source 3 (PLACEHOLDER)
    - Overview Table, Photo Caption 1/2, Photo Area 1/2 (cp roleup base に追加)
    - 茶色ガイド `正方形/長方形 1` `正方形/長方形 8` は fill が silent_remove

Usage:
  python fill_company_overview.py --brand stellar_aiz \
    --data {{WORK_DIR}}/company_overview_data.json \
    --output {{OUTPUT_DIR}}/CompanyOverview_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "company-overview-pptx-v2"

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names ─────────────────────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_TABLE        = "Overview Table"
SHAPE_CAPTION_1    = "Photo Caption 1"
SHAPE_PHOTO_1      = "Photo Area 1"
SHAPE_CAPTION_2    = "Photo Caption 2"
SHAPE_PHOTO_2      = "Photo Area 2"
# ────────────────────────────────────────────────────────────

# Defaults reassigned in main() via _apply_theme.
SHAPE_SOURCE = "Source"  # stella; roleup uses 'Source 3'

MAIN_MESSAGE_PT = 26
MAIN_MESSAGE_BOLD = True
CHART_TITLE_PT = 18
CHART_TITLE_BOLD = True
SOURCE_PT = 10

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.0)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
FONT_NAME_JP = "Meiryo UI"

CELL_BG_EVEN = "F0F1F5"  # stella V1
CELL_BG_ODD = "FFFFFF"

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global SHAPE_SOURCE
    global MAIN_MESSAGE_PT, MAIN_MESSAGE_BOLD, CHART_TITLE_PT, CHART_TITLE_BOLD, SOURCE_PT
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, FONT_NAME_JP
    global CELL_BG_EVEN, CELL_BG_ODD

    _THEME = theme

    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    FONT_NAME_JP = theme.font_ea

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        MAIN_MESSAGE_PT = 26
        MAIN_MESSAGE_BOLD = True
        CHART_TITLE_PT = 18
        CHART_TITLE_BOLD = True
        SOURCE_PT = 10
        CELL_BG_EVEN = "F0F1F5"
        CELL_BG_ODD = "FFFFFF"
    else:
        SHAPE_SOURCE = "Source 3"
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        MAIN_MESSAGE_PT = theme.pt_value("font_size_title_pt")  # 22
        MAIN_MESSAGE_BOLD = False
        CHART_TITLE_PT = theme.pt_value("font_size_subtitle_pt")  # 12
        CHART_TITLE_BOLD = False
        SOURCE_PT = theme.pt_value("font_size_source_pt")  # 6
        # Cell background uses brand label_bg / white (banded rows for readability).
        CELL_BG_EVEN = theme.hex_no_hash("label_bg").upper()
        CELL_BG_ODD = "FFFFFF"


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text, font_size=None, bold=None):
    """Set textbox text. Optionally pin font_size(pt)/bold via rPr/@sz/@b."""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        run = para.runs[0]
        run.text = text
        rPr = run._r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(run._r, qn("a:rPr"))
            run._r.insert(0, rPr)
        if font_size is not None:
            rPr.set("sz", str(int(font_size * 100)))
        if bold is not None:
            rPr.set("b", "1" if bold else "0")
        for r in para.runs[1:]:
            r.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        attrs = {"lang": "ja-JP"}
        if font_size is not None:
            attrs["sz"] = str(int(font_size * 100))
        if bold is not None:
            attrs["b"] = "1" if bold else "0"
        etree.SubElement(r_elem, qn("a:rPr"), attrib=attrs)
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_dynamic_source_textbox(slide, text):
    """Fallback: dynamic textbox for source when no Source/Source 3 placeholder is present."""
    tb = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SOURCE_PT)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP


def rebuild_table(slide, items):
    """Drop the template's placeholder Overview Table and rebuild with `items` rows.

    Cell rPr is copied from the template's first row, so brand-specific font
    (Yu Gothic UI 10pt for roleup, Meiryo UI 14pt for stella) flows through
    automatically as long as the template's first row has those rPr values.
    """
    table_shape = find_shape(slide, SHAPE_TABLE)
    if table_shape is None:
        print("  ⚠ Table shape not found, cannot rebuild", file=sys.stderr)
        return

    old_table = table_shape.table

    # Copy first row's tcPr / rPr templates from each column.
    label_tcPr = copy.deepcopy(old_table.cell(0, 0)._tc.find(qn("a:tcPr")))
    value_tcPr = copy.deepcopy(old_table.cell(0, 1)._tc.find(qn("a:tcPr")))
    label_rPr = None
    value_rPr = None
    for para in old_table.cell(0, 0).text_frame.paragraphs:
        for run in para.runs:
            label_rPr = copy.deepcopy(run._r.find(qn("a:rPr")))
            break
        break
    for para in old_table.cell(0, 1).text_frame.paragraphs:
        for run in para.runs:
            value_rPr = copy.deepcopy(run._r.find(qn("a:rPr")))
            break
        break

    tbl_left   = table_shape.left
    tbl_top    = table_shape.top
    tbl_width  = table_shape.width
    tbl_height = table_shape.height

    old_col0_width = old_table.columns[0].width
    old_col1_width = old_table.columns[1].width

    sp_tree = slide.shapes._spTree
    sp_tree.remove(table_shape._element)

    n_rows = len(items)
    n_cols = 2

    new_shape = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
    )
    new_shape.name = SHAPE_TABLE
    new_table = new_shape.table

    new_table.columns[0].width = old_col0_width
    new_table.columns[1].width = old_col1_width

    # tblPr: bandRow on
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tbl_elem = new_shape._element.find('.//a:tbl', ns)
    old_tblPr = tbl_elem.find('a:tblPr', ns)
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={'bandRow': '1'})
    tbl_elem.insert(0, tblPr)

    def apply_cell(cell, text, tcPr_tmpl, rPr_tmpl, is_even_row=True, is_label_col=True):
        tc = cell._tc
        txBody = tc.find(qn("a:txBody"))
        if txBody is None:
            txBody = etree.SubElement(tc, qn("a:txBody"))
            etree.SubElement(txBody, qn("a:bodyPr"))
            etree.SubElement(txBody, qn("a:lstStyle"))

        for p in txBody.findall(qn("a:p")):
            txBody.remove(p)

        lines = str(text).split("\n")
        for line in lines:
            p_elem = etree.SubElement(txBody, qn("a:p"))
            pPr = etree.SubElement(p_elem, qn("a:pPr"))
            pPr.set("algn", "l")
            r_elem = etree.SubElement(p_elem, qn("a:r"))
            if rPr_tmpl is not None:
                r_elem.append(copy.deepcopy(rPr_tmpl))
            else:
                etree.SubElement(r_elem, qn("a:rPr"), attrib={
                    "lang": "ja-JP", "sz": "1400"
                })
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = line

        old_tc = tc.find(qn("a:tcPr"))
        if old_tc is not None:
            tc.remove(old_tc)
        if tcPr_tmpl is not None:
            new_tcPr = copy.deepcopy(tcPr_tmpl)
            bg_color = CELL_BG_EVEN if is_even_row else CELL_BG_ODD
            for old_fill in new_tcPr.findall(qn("a:solidFill")):
                new_tcPr.remove(old_fill)
            sf = etree.SubElement(new_tcPr, qn("a:solidFill"))
            etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": bg_color})
            tc.append(new_tcPr)

        tcPr_final = tc.find(qn("a:tcPr"))
        if tcPr_final is not None:
            tcPr_final.set("anchor", "ctr")

    for r_idx, item in enumerate(items):
        label = item.get("label", "")
        value = item.get("value", "")
        is_even = (r_idx % 2 == 0)

        apply_cell(new_table.cell(r_idx, 0), label, label_tcPr, label_rPr, is_even, True)
        apply_cell(new_table.cell(r_idx, 1), value, value_tcPr, value_rPr, is_even, False)

    print(f"  ✓ Table rebuilt: {n_rows} rows x {n_cols} cols")


def insert_photo(slide, photo_info, area_shape_name, caption_shape_name):
    """Update caption + insert image into the photo area (if path exists)."""
    if not photo_info:
        return

    caption = photo_info.get("caption", "")
    if caption:
        cap_shape = find_shape(slide, caption_shape_name)
        if cap_shape:
            set_textbox_text(cap_shape, caption)

    img_path = photo_info.get("url") or photo_info.get("path", "")
    if not img_path or not os.path.exists(img_path):
        print(f"  ℹ No image for {caption_shape_name} (placeholder retained)")
        return

    area_shape = find_shape(slide, area_shape_name)
    if area_shape is None:
        return

    left = area_shape.left
    top = area_shape.top
    width = area_shape.width
    height = area_shape.height

    slide.shapes.add_picture(img_path, left, top, width, height)
    print(f"  ✓ Photo inserted: {img_path}")


def main():
    parser = argparse.ArgumentParser(description="会社概要スライド生成（ネイティブテーブル方式）")
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "company-overview")
    print(f"=== 会社概要スライド生成（brand={theme.id}）===")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Roleup: silently remove brown guide rectangles from cp-roleup-derived template.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # Title 1 / Text Placeholder 2 — brand-aware top/subtitle assignment.
    main_msg = resolve_top_text(data, theme)
    sub_text = resolve_subtitle_text(data, theme) or "対象会社概要：会社概要"

    # data field semantics for company-overview-v2:
    #   stella: 'main_message' = 結論文 → Title 1, 'title' = 見出し → Text Placeholder 2
    #   roleup: 'title' (chart_title 相当) = スライドタイトル → Title 1,
    #           'main_message' = 結論文 → Text Placeholder 2
    # resolve_top_text reads theme.top_placeholder_field() which is
    # 'main_message' for stella and 'chart_title' for roleup. company-overview-v2
    # data uses 'title' instead of 'chart_title', so fall back when 'chart_title'
    # is missing.
    if theme.id != "stellar_aiz" and not data.get("chart_title"):
        # Roleup: top placeholder = 'chart_title', but data has 'title' instead.
        main_msg = data.get("title", "対象会社概要：会社概要")
        sub_text = data.get("main_message", "")

    # main_message length warning (kept from V1, applies to whichever field is the 結論文).
    msg_for_check = data.get("main_message", "")
    if len(msg_for_check) > 65:
        print(f"  ⚠ WARNING: main_message {len(msg_for_check)} 字 > 65 字。スライドに収まらない可能性があります。", file=sys.stderr)

    set_textbox_text(
        find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg,
        font_size=MAIN_MESSAGE_PT, bold=MAIN_MESSAGE_BOLD,
    )
    print(f"  [Title 1] ({len(main_msg)}文字, {MAIN_MESSAGE_PT}pt) {main_msg[:50]}")

    set_textbox_text(
        find_shape(slide, SHAPE_CHART_TITLE), sub_text,
        font_size=CHART_TITLE_PT, bold=CHART_TITLE_BOLD,
    )
    print(f"  [Text Placeholder 2] ({len(sub_text)}文字, {CHART_TITLE_PT}pt) {sub_text[:50]}")

    # Source: stella prefixes "出典：", roleup uses Source 3 placeholder.
    source_text = data.get("source", "")
    if source_text:
        body = source_text if source_text.startswith("出典") else f"出典：{source_text}"
        source_shape = find_shape(slide, SHAPE_SOURCE)
        if source_shape is not None:
            set_textbox_text(source_shape, body, font_size=SOURCE_PT)
        else:
            add_dynamic_source_textbox(slide, body)
        print(f"  [Source]        {body[:50]}")

    # Overview Table
    items = data.get("items", [])
    if items:
        rebuild_table(slide, items)

    # Photos
    photos = data.get("photos", {})
    insert_photo(slide, photos.get("headquarters"), SHAPE_PHOTO_1, SHAPE_CAPTION_1)
    insert_photo(slide, photos.get("product"), SHAPE_PHOTO_2, SHAPE_CAPTION_2)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"  ✓ PPTX saved: {args.output}")
    print("=== 完了 ===")


if __name__ == "__main__":
    main()
