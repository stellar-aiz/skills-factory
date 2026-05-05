"""
fill_market_share.py — 市場シェア分析スライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 左側: 最新年の市場シェア・ドーナツチャート
  - 右側: 市場シェア推移テーブル（ランク、企業名、各年シェア、YoY変化）
    - 対象会社行はイエロー背景でハイライト
  - 下部: 出典

Usage:
  python fill_market_share.py \
    --data /home/claude/market_share_data.json \
    --template <path>/market-share-template.pptx \
    --output /mnt/user-data/outputs/MarketShare_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants (16:9) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

PANEL_Y = Inches(1.55)

# Left panel (doughnut chart)
LEFT_X = Inches(0.41)
LEFT_W = Inches(5.60)
LEFT_H = Inches(5.35)

# Right panel (table)
RIGHT_X = Inches(6.30)
RIGHT_W = Inches(6.60)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(6.93)
SOURCE_W = Inches(12.50)

# ── Colors ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HEADER_BG = RGBColor(0x2E, 0x4A, 0x6B)
COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)
COLOR_HIGHLIGHT = RGBColor(0xFF, 0xF4, 0xC2)   # 対象会社行のハイライト（イエロー）
COLOR_UP = RGBColor(0x1B, 0x7A, 0x3B)          # 上昇（濃緑）
COLOR_DOWN = RGBColor(0xC0, 0x3A, 0x3A)        # 下降（濃赤）
COLOR_FLAT = RGBColor(0x66, 0x66, 0x66)        # 変化なし（グレー）

# ─── 共通配色（正本: skills/_common/styles/chart_palette.md） ───
# 編集時は _common/styles/chart_palette.md と他 4 スキルの fill_*.py も同期更新
# CHART_PALETTE には TARGET_COLOR(赤) と OTHER_COLOR(灰) を含めない（palette 外で固定）
CHART_PALETTE = [
    "#4E79A7", "#F28E2B", "#59A14F", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F",
]
OTHER_COLOR = "#BAB0AC"
TARGET_COLOR = "#E15759"
LABEL_BAR_COLOR = "#4E79A7"
LABEL_BG_COLOR = "#E8EEF5"


def _palette_color(index: int, total: int) -> str:
    if total <= 1:
        return CHART_PALETTE[0]
    return CHART_PALETTE[index % len(CHART_PALETTE)]


def _is_other_player(name: str) -> bool:
    if not name:
        return False
    return (
        name.startswith("その他")
        or name.startswith("Others")
        or name.lower().startswith("other")
    )


# 後方互換のためのエイリアス（既存コードからの参照用、削除する場合は他箇所も更新）
DEFAULT_COLORS = CHART_PALETTE

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_SECTION = Pt(14)
FONT_SIZE_TABLE = Pt(11)
FONT_SIZE_TABLE_HEADER = Pt(11)
FONT_SIZE_SOURCE = Pt(10)
FONT_SIZE_SUB = Pt(11)
FONT_SIZE_TOTAL = Pt(12)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_section_title(slide, text, left, top, width):
    """セクションタイトル（下線付き、14pt Bold）"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SECTION
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEXT
    line.line.fill.background()
    return txBox


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = FONT_NAME_JP
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


# ──────────────────────────────────────────────
# Left Panel: Doughnut Chart
# ──────────────────────────────────────────────
def build_doughnut_chart(slide, section_title, players, year_idx, year_label,
                          total_size_label, left, top, width, height):
    """
    セクションタイトル + ドーナツチャート + 合計ラベル
    """
    # セクションタイトル
    add_section_title(slide, section_title, left, top, width)

    # 合計ラベル（チャート上部）
    sub_y = top + Inches(0.40)
    if total_size_label:
        add_text_box(
            slide, total_size_label,
            left, sub_y, width, Inches(0.30),
            FONT_SIZE_TOTAL, bold=True, align=PP_ALIGN.CENTER,
        )

    # チャートエリア
    chart_top = top + Inches(0.80)
    chart_h = height - Inches(0.80)

    # データ準備
    cdata = CategoryChartData()
    cdata.categories = [p["name"] for p in players]
    shares = [p["shares"][year_idx] for p in players]
    cdata.add_series(f"{year_label}シェア", shares)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        left, chart_top, width, chart_h,
        cdata,
    )
    chart = chart_shape.chart
    chart.has_title = False

    # 凡例: 右側
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT_NAME_JP

    # python-pptxのバグ対策: <c:legendPos/> にval属性を強制設定
    # （val属性がないとOOXML仕様違反となり、PowerPointでスライドが白紙になる）
    legend_elem = chart.legend._element
    legendPos = legend_elem.find(qn("c:legendPos"))
    if legendPos is not None and legendPos.get("val") is None:
        legendPos.set("val", "r")

    # 系列に色を適用 + データラベル
    series = chart.series[0]

    # データラベル: etreeで全構築（python-pptxのsetterは使わない）
    # python-pptxのsetterを使うと <c:dLblPos val="ctr"/> や <c:showLeaderLines>
    # が自動追加され、PowerPoint for Mac で致命的修復エラーになるため。
    # OOXMLスキーマ順（CT_DLbls）で子要素を構築する。
    dLbls_xml = series.data_labels._element
    for child in list(dLbls_xml):
        dLbls_xml.remove(child)

    # 1. numFmt
    numFmt = etree.SubElement(dLbls_xml, qn("c:numFmt"))
    numFmt.set("formatCode", '0.0"%"')
    numFmt.set("sourceLinked", "0")

    # 2. txPr（font size/bold/name/color）
    txPr = etree.SubElement(dLbls_xml, qn("c:txPr"))
    bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
    bodyPr.set("wrap", "square")
    bodyPr.set("anchor", "ctr")
    bodyPr.set("anchorCtr", "1")
    etree.SubElement(txPr, qn("a:lstStyle"))
    p_el = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p_el, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", "1000")  # 10pt (1/100 pt単位)
    defRPr.set("b", "1")
    solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", "FFFFFF")
    for tag in ("a:latin", "a:ea", "a:cs"):
        font_el = etree.SubElement(defRPr, qn(tag))
        font_el.set("typeface", FONT_NAME_JP)
    endParaRPr = etree.SubElement(p_el, qn("a:endParaRPr"))
    endParaRPr.set("lang", "ja-JP")

    # 3. dLblPos は意図的に省略（PowerPoint for Mac の修復エラーの真犯人）

    # 4-9. showLegendKey 〜 showBubbleSize（OOXMLスキーマ順）
    for tag, val in (
        ("c:showLegendKey", "0"),
        ("c:showVal", "1"),
        ("c:showCatName", "0"),
        ("c:showSerName", "0"),
        ("c:showPercent", "0"),
        ("c:showBubbleSize", "0"),
    ):
        etree.SubElement(dLbls_xml, qn(tag)).set("val", val)

    # 念のため、他経路で追加されうる showLeaderLines / leaderLines を除去
    for extra in ("c:showLeaderLines", "c:leaderLines"):
        for e in dLbls_xml.findall(qn(extra)):
            dLbls_xml.remove(e)

    # 各カテゴリごとに色を設定 (data point level)
    # python-pptxではdata points経由で色を個別設定
    for idx, p in enumerate(players):
        # color は事前に main 側で強制割り当て済み（JSON の指定は無視されている）
        hex_color = p["color"]
        rgb = hex_to_rgb(hex_color)
        point = series.points[idx]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = rgb
        point.format.line.color.rgb = COLOR_WHITE
        point.format.line.width = Pt(1.5)

    # ドーナツのホールサイズ (約50%)
    chart_xml = chart._chartSpace
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    plot_area = chart_xml.find(f'.//{{{ns}}}plotArea')
    doughnut = plot_area.find(f'.//{{{ns}}}doughnutChart')
    if doughnut is not None:
        # holeSize要素を追加/更新
        holeSize = doughnut.find(f'{{{ns}}}holeSize')
        if holeSize is None:
            holeSize = etree.SubElement(doughnut, qn("c:holeSize"))
        holeSize.set("val", "55")

    print(f"  ✓ ドーナツチャート: {len(players)}プレイヤー ({year_label}年)")
    return chart_shape


# ──────────────────────────────────────────────
# Right Panel: Trend Summary Table
# ──────────────────────────────────────────────
def build_trend_table(slide, section_title, players, years, target_company,
                       show_ranking, show_yoy_change, left, top, width):
    """
    市場シェア推移テーブル
    列: [ランク (optional)] [企業名] [年1] [年2] ... [年N] [YoY変化 (optional)]
    """
    add_section_title(slide, section_title, left, top, width)

    # 各プレイヤーをランキング順にソート（最新年基準）
    # ただし「その他」系（前方一致）は最下段に固定
    others = [p for p in players if _is_other_player(p.get("name", ""))]
    regular = [p for p in players if not _is_other_player(p.get("name", ""))]
    regular_sorted = sorted(regular, key=lambda p: p["shares"][-1], reverse=True)
    sorted_players = regular_sorted + others

    # 列を構築
    cols = []
    if show_ranking:
        cols.append(("#", "rank"))
    cols.append(("企業名", "name"))
    for y in years:
        cols.append((y, f"share_{y}"))
    if show_yoy_change and len(years) >= 2:
        cols.append(("YoY変化", "yoy"))

    n_cols = len(cols)
    n_rows = len(sorted_players) + 1  # +1 for header

    tbl_top = top + Inches(0.55)
    row_h = Inches(0.38)
    tbl_h = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    table = shape.table

    # 列幅の設定: ランク列は狭く、企業名列は広く、年・変化列は均等
    total_w = width
    rank_w = Inches(0.45) if show_ranking else Emu(0)
    yoy_w = Inches(1.05) if (show_yoy_change and len(years) >= 2) else Emu(0)
    name_w = Inches(2.20)
    remaining = total_w - rank_w - name_w - yoy_w
    year_w = Emu(int(remaining / len(years))) if len(years) > 0 else Inches(1.0)

    col_widths = []
    for hdr, key in cols:
        if key == "rank":
            col_widths.append(rank_w)
        elif key == "name":
            col_widths.append(name_w)
        elif key == "yoy":
            col_widths.append(yoy_w)
        else:
            col_widths.append(year_w)

    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(int(w))

    # tblPr設定
    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(row_h))

    # ヘッダー行
    for c_idx, (hdr, key) in enumerate(cols):
        cell = table.cell(0, c_idx)
        _style_cell(cell, hdr, is_header=True, is_alt=False,
                    font_size=FONT_SIZE_TABLE_HEADER)

    # データ行
    for r_idx, p in enumerate(sorted_players):
        is_target = (p.get("name") == target_company) if target_company else False
        is_alt = (r_idx % 2 == 1)

        for c_idx, (hdr, key) in enumerate(cols):
            cell = table.cell(r_idx + 1, c_idx)

            if key == "rank":
                # 「その他」系はランク表記なし（前方一致）
                if _is_other_player(p.get("name", "")):
                    text = "—"
                else:
                    text = str(regular_sorted.index(p) + 1)
                _style_cell(cell, text, is_header=False, is_alt=is_alt,
                            font_size=FONT_SIZE_TABLE, bold=True,
                            is_target=is_target, align_override="ctr")

            elif key == "name":
                # color は事前に main 側で強制割り当て済み
                color_hex = p["color"]
                _style_cell(cell, p["name"], is_header=False, is_alt=is_alt,
                            font_size=FONT_SIZE_TABLE, bold=True,
                            color_marker_hex=color_hex,
                            is_target=is_target, align_override="l")

            elif key.startswith("share_"):
                y = key.replace("share_", "")
                y_idx = years.index(y)
                val = p["shares"][y_idx]
                text = f"{val:.1f}%"
                _style_cell(cell, text, is_header=False, is_alt=is_alt,
                            font_size=FONT_SIZE_TABLE, is_target=is_target,
                            align_override="r")

            elif key == "yoy":
                delta = p["shares"][-1] - p["shares"][-2]
                if delta > 0.1:
                    text = f"↑ +{delta:.1f}pt"
                    color = COLOR_UP
                elif delta < -0.1:
                    text = f"↓ {delta:.1f}pt"
                    color = COLOR_DOWN
                else:
                    text = f"→ {delta:+.1f}pt"
                    color = COLOR_FLAT
                _style_cell(cell, text, is_header=False, is_alt=is_alt,
                            font_size=FONT_SIZE_TABLE,
                            bold=True, text_color=color,
                            is_target=is_target, align_override="r")

    print(f"  ✓ 市場シェア推移テーブル: {n_rows}行 x {n_cols}列")
    return shape


def _style_cell(cell, text, is_header=False, is_alt=False, font_size=Pt(11),
                bold=False, color_marker_hex=None, is_target=False,
                text_color=None, align_override=None):
    """セルのスタイルとテキストを設定"""
    tc = cell._tc

    # 背景色
    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
    elif is_target:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HIGHLIGHT
    elif is_alt:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_ROW_ALT
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE

    # マージン、縦中央寄せ
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = cell.text_frame
    tf.word_wrap = True

    # 既存段落クリア
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))

    # 配置: is_headerは中央、bold(name列)は左、数値は右、align_overrideで上書き
    if align_override:
        pPr.set("algn", align_override)
    elif is_header:
        pPr.set("algn", "ctr")
    elif bold and not color_marker_hex:
        pPr.set("algn", "l")
    else:
        pPr.set("algn", "r")

    # テキスト色を決定
    if text_color is not None:
        color_val = "{:02X}{:02X}{:02X}".format(text_color[0], text_color[1], text_color[2])
    elif is_header:
        color_val = "FFFFFF"
    else:
        color_val = "333333"

    if color_marker_hex:
        # ■ 色付き + テキスト通常色の2 Run構成
        marker_val = color_marker_hex.lstrip("#")

        r1 = etree.SubElement(p_elem, qn("a:r"))
        rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(font_size.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
        s1 = etree.SubElement(sf1, qn("a:srgbClr"))
        s1.set("val", marker_val)
        t1 = etree.SubElement(r1, qn("a:t"))
        t1.text = "■ "

        r2 = etree.SubElement(p_elem, qn("a:r"))
        rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(font_size.pt * 100)),
            "b": "1",
        })
        etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr2, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
        s2 = etree.SubElement(sf2, qn("a:srgbClr"))
        s2.set("val", color_val)
        t2 = etree.SubElement(r2, qn("a:t"))
        t2.text = text
    else:
        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(font_size.pt * 100)),
            "b": "1" if (is_header or bold) else "0",
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", color_val)
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)  # passive: accepted but ignored until brand migration
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    _mm = data.get("main_message", "")
    if len(_mm) > 65:
        raise ValueError(
            f"main_message は 65 字以内（受領: {len(_mm)}）: {_mm[:80]}..."
        )

    prs = Presentation(args.template)
    slide = prs.slides[0]

    # Main Message & Chart Title
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), data.get("main_message", ""))
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), data.get("chart_title", "市場シェア分析"))

    players = data.get("players", [])
    years = data.get("years", [])
    if not players or not years:
        print("  ✗ ERROR: 'players' and 'years' are required", file=sys.stderr)
        sys.exit(1)

    # 色の強制割り当て（共通パレット使用、JSON の color は無視）
    # 「その他」系エントリは前方一致で OTHER_COLOR(灰) に上書き、
    # それ以外は players 配列の index をそのまま使って palette を引く
    # （P6/P7 で同じ社が同じ色になるよう、target/その他 をスキップせず配列 index で揃える）
    for i, p in enumerate(players):
        if _is_other_player(p.get("name", "")):
            p["color"] = OTHER_COLOR
        else:
            p["color"] = _palette_color(i, len(players))

    # 左: ドーナツチャート（デフォルトは最新年）
    left_panel = data.get("left_panel", {})
    chart_year_idx = data.get("chart_year_index", len(years) - 1)
    chart_year_label = years[chart_year_idx]
    left_section_title = left_panel.get(
        "section_title", f"市場シェア（{chart_year_label}年）"
    )
    build_doughnut_chart(
        slide, left_section_title, players, chart_year_idx, chart_year_label,
        left_panel.get("total_size_label", ""),
        LEFT_X, PANEL_Y, LEFT_W, LEFT_H,
    )

    # 右: 推移テーブル
    right_panel = data.get("right_panel", {})
    right_section_title = right_panel.get("section_title", "主要プレイヤー別シェア推移")
    show_ranking = right_panel.get("show_ranking", True)
    show_yoy_change = right_panel.get("show_yoy_change", True)
    target_company = data.get("target_company")
    build_trend_table(
        slide, right_section_title, players, years, target_company,
        show_ranking, show_yoy_change,
        RIGHT_X, PANEL_Y, RIGHT_W,
    )

    # 出典
    source = data.get("source", "")
    if source:
        add_text_box(
            slide, source,
            SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.30),
            FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
            align=PP_ALIGN.LEFT,
        )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
