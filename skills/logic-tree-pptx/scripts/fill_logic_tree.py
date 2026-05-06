"""
fill_logic_tree.py — ロジックツリーデータをPPTXテンプレートに流し込むスクリプト

テンプレート構造（LogicTree.pptx 確認済み）:
  - Title 1              (PLACEHOLDER): Main Message
  - Text Placeholder 2   (PLACEHOLDER): Chart Title
  - Rectangle 4          (AUTO_SHAPE):  Main Logic   (Para[0]=タイトル, Para[1-3]=3要素)
  - Rectangle 3          (AUTO_SHAPE):  Sub Logic 1  (Para[0]=タイトル, Para[1-3]=3要素)
  - Rectangle 5          (AUTO_SHAPE):  Sub Logic 2  (Para[0]=タイトル, Para[1-3]=3要素)
  - Rectangle 18         (AUTO_SHAPE):  Sub Logic 3  (Para[0]=タイトル, Para[1-3]=3要素)
  - TextBox 23           (TEXT_BOX):    意味合い (Para[0]=ラベル, Para[1-3]=3項目)

使い方:
  python fill_logic_tree.py \
    --data /home/claude/logic_tree_data.json \
    --template <SKILL_DIR>/assets/logic-tree-template.pptx \
    --output /mnt/user-data/outputs/LogicTree_output.pptx
"""

import argparse
import os
import json
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング（確認済み）──────────────────────────────
SHAPE_MAIN_MESSAGE  = "Title 1"
SHAPE_CHART_TITLE   = "Text Placeholder 2"
SHAPE_MAIN_LOGIC    = "Rectangle 4"
SHAPE_SUB_LOGIC_1   = "Rectangle 3"
SHAPE_SUB_LOGIC_2   = "Rectangle 5"
SHAPE_SUB_LOGIC_3   = "Rectangle 18"
SHAPE_IMPLICATIONS  = "TextBox 23"
# ────────────────────────────────────────────────────────────


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_placeholder_text(shape, text):
    """PlaceholderのTextFrameにテキストをセット（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_logic_box(slide, shape_name, title, points):
    """Logic Box（Rectangle）の4段落にタイトル+3要素を書き込む

    Para[0]: タイトル（既存スタイル=太字を保持）
    Para[1-3]: 3つの要素
    """
    shape = find_shape(slide, shape_name)
    if shape is None:
        return
    tf = shape.text_frame

    # Para[0]: タイトル
    para_title = tf.paragraphs[0]
    if para_title.runs:
        para_title.runs[0].text = title
        for run in para_title.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para_title._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP", "b": "1"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = title
    print(f"  [{shape_name} title] {title[:55]}{'...' if len(title) > 55 else ''}")

    # Para[1-3]: 3つの要素
    for i, point in enumerate(points[:3]):
        para_idx = i + 1
        if para_idx < len(tf.paragraphs):
            para = tf.paragraphs[para_idx]
            if para.runs:
                para.runs[0].text = point
                for run in para.runs[1:]:
                    run.text = ""
            else:
                r_elem = etree.SubElement(para._p, qn("a:r"))
                etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = point
        else:
            # 段落が足りない場合は追加
            new_p = etree.SubElement(tf._txBody, qn("a:p"))
            r_elem = etree.SubElement(new_p, qn("a:r"))
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = point
        print(f"  [{shape_name}_{i+1}] {point[:55]}{'...' if len(point) > 55 else ''}")


def fill_implications(slide, items):
    """意味合いTextBoxのPara[1-3]に3つの示唆を書き込む

    Para[0] は「意味合い」ラベル → 変更しない
    Para[1-3] に各示唆を書き込む
    """
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return
    tf = shape.text_frame

    for i, text in enumerate(items[:3]):
        para_idx = i + 1  # Para[0]はラベルなのでスキップ
        if para_idx < len(tf.paragraphs):
            para = tf.paragraphs[para_idx]
            if para.runs:
                para.runs[0].text = text
                for run in para.runs[1:]:
                    run.text = ""
            else:
                r_elem = etree.SubElement(para._p, qn("a:r"))
                etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = text
        else:
            new_p = etree.SubElement(tf._txBody, qn("a:p"))
            r_elem = etree.SubElement(new_p, qn("a:r"))
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text
        print(f"  [Implications_{i+1}] {text[:55]}{'...' if len(text) > 55 else ''}")


def main():
    parser = argparse.ArgumentParser(description="ロジックツリーデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="logic_tree_data.json のパス")
    parser.add_argument("--template", required=True, help="logic-tree-template.pptx のパス")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "main_logic", "sub_logics"],
        allowed_top=[
            "main_message", "chart_title",
            "main_logic", "sub_logics", "implications",
        ],
        skill_name="logic-tree-pptx",
    )

    prs = Presentation(args.template)
    slide = prs.slides[0]

    # ── Main Message ──────────────────────────────────────────
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(shape, main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # ── Chart Title ───────────────────────────────────────────
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(shape, chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # ── Main Logic ────────────────────────────────────────────
    main_logic = data.get("main_logic", {})
    ml_title = main_logic.get("title", "").strip()
    ml_points = main_logic.get("points", [])
    if ml_title:
        fill_logic_box(slide, SHAPE_MAIN_LOGIC, ml_title, ml_points)
    else:
        print("  ⚠ main_logic.title が未設定です", file=sys.stderr)

    # ── Sub Logic 1〜3 ────────────────────────────────────────
    sub_logics = data.get("sub_logics", [])
    sub_shapes = [SHAPE_SUB_LOGIC_1, SHAPE_SUB_LOGIC_2, SHAPE_SUB_LOGIC_3]
    for i, (sl, shape_name) in enumerate(zip(sub_logics[:3], sub_shapes)):
        sl_title = sl.get("title", "").strip()
        sl_points = sl.get("points", [])
        if sl_title:
            fill_logic_box(slide, shape_name, sl_title, sl_points)
        else:
            print(f"  ⚠ sub_logics[{i}].title が未設定です", file=sys.stderr)

    if len(sub_logics) != 3:
        print(f"  ⚠ sub_logics は {len(sub_logics)} 件です（3件必須）", file=sys.stderr)

    # ── Implications ──────────────────────────────────────────
    implications = data.get("implications", [])
    if isinstance(implications, str):
        # 後方互換: 文字列の場合はリストに変換
        implications = [implications]
    if implications:
        fill_implications(slide, implications)
    else:
        print("  ⚠ implications が未設定です", file=sys.stderr)
    if len(implications) != 3:
        print(f"  ⚠ implications は {len(implications)} 件です（3件推奨）", file=sys.stderr)

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
