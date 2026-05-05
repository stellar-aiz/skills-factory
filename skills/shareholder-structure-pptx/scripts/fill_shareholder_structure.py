"""
fill_shareholder_structure.py — 株主構成・役員構成をPPTXネイティブテーブルとして生成

テンプレート構造（shareholder-structure-template.pptx = company-history-template.pptxベース）:
  - Title 1            (PLACEHOLDER): メインメッセージ
  - Text Placeholder 2 (PLACEHOLDER): チャートタイトル
  - Table 1            (TABLE):       テンプレート用テーブル（スタイル複製元→削除）

生成物:
  - ■株主構成 セクションタイトル（TextBox）
  - 株主テーブル（ネイティブTable: 7列）
  - ■役員構成 セクションタイトル（TextBox）
  - 役員テーブル（ネイティブTable: 6列）
  - 出典テキスト（TextBox）

Usage:
  python fill_shareholder_structure.py \
    --data /home/claude/shareholder_structure_data.json \
    --template <SKILL_DIR>/assets/shareholder-structure-template.pptx \
    --output /mnt/user-data/outputs/ShareholderStructure_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_TABLE = "Table 1"

# ── レイアウト定数 ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
TABLE_LEFT = Inches(0.41)
TABLE_WIDTH = Inches(12.52)

# ── スタイル定数 ──
FONT_NAME_JP = "Meiryo UI"
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
HEADER_BG_HEX = "F5F0D0"   # テンプレート準拠ベージュ
TOTAL_BG_HEX = "F0F0F0"    # 合計行背景


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書きし、フォントをMeiryo UIに強制設定"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        # 全runのフォントをMeiryo UIに統一
        for run in para.runs:
            rPr = run._r.find(qn("a:rPr"))
            if rPr is None:
                rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
                run._r.insert(0, rPr)
            # latin/eaを上書き
            for tag in [qn("a:latin"), qn("a:ea")]:
                old = rPr.find(tag)
                if old is not None:
                    rPr.remove(old)
                etree.SubElement(rPr, tag, attrib={"typeface": FONT_NAME_JP})
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def remove_shape(slide, name):
    """名前でShapeを削除"""
    shape = find_shape(slide, name)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def add_section_title(slide, text, left, top, width):
    """セクションタイトル（■株主構成 等）をTextBoxで追加"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP
    # ea（日本語）フォントも明示設定
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        old_ea = rPr.find(qn("a:ea"))
        if old_ea is not None:
            rPr.remove(old_ea)
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    print(f"  ✓ セクションタイトル: {text}")
    return txBox


def add_source_label(slide, text, top):
    """出典テキストを追加"""
    txBox = slide.shapes.add_textbox(TABLE_LEFT, top, Inches(10.0), Inches(0.25))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"出典：{text}" if text else ""
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP
    # ea（日本語）フォントも明示設定
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        old_ea = rPr.find(qn("a:ea"))
        if old_ea is not None:
            rPr.remove(old_ea)
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    print(f"  ✓ 出典: {text}")


def apply_cell(cell, text, is_header=False, bold=False, align="l",
               font_size=1200, bg_hex=None):
    """セルにテキスト・スタイルを設定（ネイティブPPTXセル）"""
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    # bodyPr（中央揃え・マージンはtcPrで管理するので0に）
    old_bodyPr = txBody.find(qn("a:bodyPr"))
    if old_bodyPr is not None:
        txBody.remove(old_bodyPr)
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": "0", "rIns": "0",
        "tIns": "0", "bIns": "0",
        "anchor": "ctr",
    })
    txBody.insert(0, bodyPr)

    if txBody.find(qn("a:lstStyle")) is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.insert(1, lstStyle)

    # 既存段落を削除
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # 新しい段落
    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", align)

    # 行間を詰める
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(font_size),
        "b": "1" if (bold or is_header) else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": "333333"})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    # tcPr（背景色・罫線）
    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)

    tcPr = etree.SubElement(tc, qn("a:tcPr"), attrib={
        "marL": "36576", "marR": "36576",
        "marT": "18288", "marB": "18288",
        "anchor": "ctr",
    })

    # 罫線（薄いグレー）
    for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = etree.SubElement(tcPr, qn(border_name), attrib={"w": "6350", "cmpd": "sng"})
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": "CCCCCC"})

    # 背景色
    if bg_hex:
        cell_fill = etree.SubElement(tcPr, qn("a:solidFill"))
        etree.SubElement(cell_fill, qn("a:srgbClr"), attrib={"val": bg_hex})


def build_table(slide, columns, rows, col_widths_inch, left, top, width,
                font_size=1200, has_total=False, row_height=Inches(0.30)):
    """ネイティブPPTXテーブルを構築する。

    Args:
        row_height: データ行の高さ（Emu）。呼び出し元で決定済み。
    """
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(columns)

    table_height = HEADER_HEIGHT + row_height * len(rows)

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = shape.table

    # 列幅を設定
    for i, w in enumerate(col_widths_inch):
        table.columns[i].width = Inches(w)

    # tblPrを設定（デフォルトスタイル除去）
    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # 行の高さを明示的に設定
    for i, tr in enumerate(tbl_elem.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(HEADER_HEIGHT))
        else:
            tr.set('h', str(row_height))

    print(f"  ✓ テーブル高さ: {table_height/914400:.2f}in "
          f"(データ行: {row_height/914400:.2f}in × {len(rows)}行)")


    # ── ヘッダー行 ──
    for c_idx, col_name in enumerate(columns):
        align = "r" if col_name in ["持株数", "議決権比率(%)", "役員報酬"] else "ctr" if col_name == "#" else "l"
        apply_cell(table.cell(0, c_idx), col_name, is_header=True,
                   align=align, font_size=font_size, bg_hex=HEADER_BG_HEX)

    # ── データ行 ──
    for r_idx, row_data in enumerate(rows):
        cells = row_data["cells"]
        aligns = row_data.get("aligns", ["l"] * n_cols)
        is_total = has_total and (r_idx == len(rows) - 1)
        bg = TOTAL_BG_HEX if is_total else None

        for c_idx, cell_text in enumerate(cells):
            align = aligns[c_idx] if c_idx < len(aligns) else "l"
            apply_cell(table.cell(r_idx + 1, c_idx), cell_text,
                       bold=is_total, align=align, font_size=font_size,
                       bg_hex=bg)

    print(f"  ✓ テーブル生成: {n_rows}行 × {n_cols}列")
    return shape, table_height


def format_shareholder_rows(data):
    """株主構成の行データを整形する"""
    rows = []
    aligns = ["ctr", "l", "l", "l", "r", "r", "l"]

    for row in data.get("rows", []):
        cells = [
            str(row.get("number", "")),
            row.get("name", ""),
            row.get("position", ""),
            row.get("relation", ""),
            row.get("shares", ""),
            f'{row.get("voting_ratio", 0):.1f}',
            row.get("note", ""),
        ]
        rows.append({"cells": cells, "aligns": aligns})

    # 合計行
    total = data.get("total")
    if total:
        cells = [
            "",
            "合計",
            "",
            "",
            total.get("shares", ""),
            f'{total.get("voting_ratio", 100.0):.1f}',
            "",
        ]
        rows.append({"cells": cells, "aligns": aligns})

    return rows


def format_director_rows(data):
    """役員構成の行データを整形する"""
    rows = []
    aligns = ["ctr", "l", "l", "l", "r", "l"]

    for row in data.get("rows", []):
        cells = [
            str(row.get("number", "")),
            row.get("name", ""),
            row.get("position", ""),
            row.get("relation", ""),
            row.get("compensation", ""),
            row.get("note", ""),
        ]
        rows.append({"cells": cells, "aligns": aligns})

    return rows


def required_row_height(font_size_hundredths):
    """フォントサイズから、レンダリング後の実際の行高さ(Emu)を算出する。
    
    PPTXの行高さ(h属性)は最小値であり、レンダリング時にフォント＋マージンに応じて
    自動拡張される。この関数は拡張後の実際の描画高さを返す。
    
    計算式:
      ベース高さ = font_pt × 1.5（行間） ÷ 72（pt→in）
      セルマージン = 0.04in（tcPr: marT + marB）
      レンダリング補正 = 0.05in（PPT/LibreOfficeの追加パディング）
    """
    font_pt = font_size_hundredths / 100
    height_in = (font_pt * 1.5) / 72 + 0.04 + 0.05
    return Inches(height_in)


# フォントサイズ候補（大きい順に試行）
FONT_CANDIDATES = [1300, 1200, 1100, 1000, 900]
MIN_ROW_HEIGHT = Inches(0.25)
HEADER_HEIGHT = Inches(0.35)   # ヘッダーは太字のため少し高め
MAX_ROW_HEIGHT = Inches(0.40)


def main():
    parser = argparse.ArgumentParser(description="株主構成・役員構成 PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument("--template", required=True, help="PPTXテンプレートパス")
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    print("=== 株主構成・役員構成スライド生成 ===")

    prs = Presentation(args.template)
    slide = prs.slides[0]

    # Shape一覧（デバッグ）
    for s in slide.shapes:
        print(f"  Shape: '{s.name}' type={s.shape_type}")

    # 1. メインメッセージ
    main_msg = data.get("main_message", "")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg)
    print(f"  ✓ Main Message: {main_msg}")

    # 2. チャートタイトル
    chart_title = data.get("chart_title", "対象会社概要：株主構成")
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
    print(f"  ✓ Chart Title: {chart_title}")

    # 3. テンプレートのTable 1を削除
    remove_shape(slide, SHAPE_TABLE)

    # 4. データ整形
    shareholders = data.get("shareholders", {})
    directors = data.get("directors", {})
    sh_rows = format_shareholder_rows(shareholders)
    dir_rows = format_director_rows(directors)

    has_total = shareholders.get("total") is not None
    n_data_rows = len(sh_rows) + len(dir_rows)  # ヘッダー除くデータ行数
    n_headers = 2  # 株主・役員各1ヘッダー

    # 5. レイアウト定数
    SLIDE_HEIGHT = Inches(7.50)
    BOTTOM_MARGIN = Inches(0.30)
    SECTION_TITLE_H = Inches(0.28)
    SECTION_GAP = Inches(0.12)

    content_start_y = Inches(1.30)
    total_available = SLIDE_HEIGHT - content_start_y - BOTTOM_MARGIN

    # テーブル以外が使う固定高さ
    overhead = (SECTION_TITLE_H * 2) + SECTION_GAP
    table_budget = total_available - overhead  # テーブル2つに使える高さ

    # 6. ループ: フォントサイズを大→小で試行し、収まるサイズを決定
    chosen_font = None
    chosen_row_h = None

    for font_sz in FONT_CANDIDATES:
        row_h = required_row_height(font_sz)
        # 上限でクランプ
        if row_h > MAX_ROW_HEIGHT:
            row_h = MAX_ROW_HEIGHT

        total_table_h = (HEADER_HEIGHT * n_headers) + (row_h * n_data_rows)

        if total_table_h <= table_budget:
            chosen_font = font_sz
            chosen_row_h = row_h
            print(f"  ✓ フォント {font_sz/100:.0f}pt: "
                  f"行高さ={row_h/914400:.2f}in, "
                  f"テーブル合計={total_table_h/914400:.2f}in ≤ "
                  f"予算{table_budget/914400:.2f}in → OK")
            break
        else:
            print(f"  × フォント {font_sz/100:.0f}pt: "
                  f"行高さ={row_h/914400:.2f}in, "
                  f"テーブル合計={total_table_h/914400:.2f}in > "
                  f"予算{table_budget/914400:.2f}in → 次へ")

    # 全候補で収まらない場合: 最小フォントで行高さを逆算
    if chosen_font is None:
        chosen_font = FONT_CANDIDATES[-1]  # 最小フォント
        available_per_row = (table_budget - HEADER_HEIGHT * n_headers) / n_data_rows
        available_per_row = int(available_per_row)

        if available_per_row < MIN_ROW_HEIGHT:
            # 0.25in未満 → ユーザ判断を仰ぐ
            print(f"\n  ⚠ エラー: データ行数が多すぎます（{n_data_rows}行）")
            print(f"    最小フォント{FONT_CANDIDATES[-1]/100:.0f}ptでも "
                  f"行高さが{available_per_row/914400:.2f}inとなり、")
            print(f"    下限{MIN_ROW_HEIGHT/914400:.2f}inを下回ります。")
            print(f"    以下の対応を検討してください：")
            print(f"    - 株主数または役員数を削減する")
            print(f"    - 2枚のスライドに分割する")
            sys.exit(1)

        chosen_row_h = available_per_row
        print(f"  △ 最小フォント{chosen_font/100:.0f}ptで逆算: "
              f"行高さ={chosen_row_h/914400:.2f}in")

    font_size = chosen_font
    row_height = chosen_row_h

    # 7. 配置（current_yを積み上げ）
    current_y = content_start_y

    # ■株主構成
    sh_section_title = shareholders.get("section_title", "株主構成")
    add_section_title(slide, f"■{sh_section_title}", TABLE_LEFT, current_y, TABLE_WIDTH)
    current_y += SECTION_TITLE_H

    sh_columns = shareholders.get("columns",
        ["#", "株主", "役職", "関係", "持株数", "議決権比率(%)", "備考"])
    sh_col_widths = [0.50, 1.80, 1.30, 1.00, 1.20, 1.50, 5.22]

    sh_shape, sh_table_h = build_table(
        slide, sh_columns, sh_rows, sh_col_widths,
        TABLE_LEFT, current_y, TABLE_WIDTH,
        font_size=font_size, has_total=has_total,
        row_height=row_height
    )
    current_y += sh_table_h + SECTION_GAP

    # ■役員構成
    dir_section_title = directors.get("section_title", "役員構成")
    add_section_title(slide, f"■{dir_section_title}", TABLE_LEFT, current_y, TABLE_WIDTH)
    current_y += SECTION_TITLE_H

    dir_columns = directors.get("columns",
        ["#", "氏名", "役職", "関係", "役員報酬", "備考"])
    dir_col_widths = [0.50, 1.80, 1.30, 1.00, 1.50, 6.42]

    dir_shape, dir_table_h = build_table(
        slide, dir_columns, dir_rows, dir_col_widths,
        TABLE_LEFT, current_y, TABLE_WIDTH,
        font_size=font_size, has_total=False,
        row_height=row_height
    )
    current_y += dir_table_h + Inches(0.05)

    # 8. 出典
    source_text = data.get("source", "")
    if source_text:
        add_source_label(slide, source_text, current_y)

    # 9. 保存
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
