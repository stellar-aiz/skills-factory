"""
fill_sga_breakdown.py — SGA Breakdown Slide (Native PPTX)

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。
"""
import argparse, copy, json, math, os, sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "sga-breakdown-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names ──
SHAPE_MAIN_MSG = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT = "Content Area"
SHAPE_SOURCE = "Source"

# Defaults (stella). Reassigned in _apply_theme(theme) for roleup.
PANEL_Y = Inches(1.50)
CHART_X = Inches(0.50)
CHART_W = Inches(8.50)
CHART_H = Inches(5.00)
CHART_Y_OFFSET = Inches(0.55)
LEGEND_X = Inches(9.30)
LEGEND_W = Inches(3.80)
LEGEND_BOTTOM = Inches(6.30)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SRC = RGBColor(0x66, 0x66, 0x66)
FONT = "Meiryo UI"
TEXT_HEX = "333333"

DATA_LABEL_SZ = "1000"   # OOXML hundredths; was 800 (8pt) — bumped to 10pt for roleup C4
LINE_DATA_LABEL_SZ = "1000"
AXIS_FONT_SZ = "1000"

DEFAULT_LINE_COLOR = "#666666"
SOURCE_FONT_PT = 10
TOTAL_LABEL_PT = 10
TREND_TEXT_PT = 11
UNIT_LABEL_PT = 10

LINE_FALLBACK_HEX = None  # roleup forces line color to brand accent if data omits

_THEME = None


def _apply_theme(theme):
    global _THEME
    global SHAPE_SOURCE
    global PANEL_Y, CHART_X, CHART_W, CHART_H, CHART_Y_OFFSET
    global LEGEND_X, LEGEND_W, LEGEND_BOTTOM
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SRC, FONT, TEXT_HEX
    global DATA_LABEL_SZ, LINE_DATA_LABEL_SZ, AXIS_FONT_SZ
    global SOURCE_FONT_PT, TOTAL_LABEL_PT, TREND_TEXT_PT, UNIT_LABEL_PT
    global LINE_FALLBACK_HEX

    _THEME = theme
    FONT = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    COLOR_TEXT = theme.color("text")
    COLOR_SRC = theme.color("source")

    PANEL_Y = theme.layout("panel_y_in")
    CHART_X = theme.layout("chart_x_in")
    CHART_W = theme.layout("chart_w_in")
    CHART_H = theme.layout("chart_h_in")
    CHART_Y_OFFSET = theme.layout("chart_y_offset_in")
    LEGEND_X = theme.layout("legend_x_in")
    LEGEND_W = theme.layout("legend_w_in")
    LEGEND_BOTTOM = theme.layout("legend_bottom_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        DATA_LABEL_SZ = "800"  # V1: 8pt for stella (no C4 enforcement)
        LINE_DATA_LABEL_SZ = "1000"
        AXIS_FONT_SZ = "1000"
        SOURCE_FONT_PT = 10
        TOTAL_LABEL_PT = 10
        TREND_TEXT_PT = 11
        UNIT_LABEL_PT = 10
        LINE_FALLBACK_HEX = None
    else:
        SHAPE_SOURCE = "Source 3"
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        DATA_LABEL_SZ = "1000"  # 10pt (was 8pt for stella)
        LINE_DATA_LABEL_SZ = "1000"  # 10pt
        AXIS_FONT_SZ = "1000"  # 10pt
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6pt
        TOTAL_LABEL_PT = theme.pt_value("font_size_body_pt")  # 10pt
        TREND_TEXT_PT = theme.pt_value("font_size_subtitle_pt")  # 12pt
        UNIT_LABEL_PT = theme.pt_value("font_size_body_pt")  # 10pt
        LINE_FALLBACK_HEX = theme.hex_no_hash("accent_op_margin_line")  # #604C3F


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def set_text(shape, text):
    if not shape:
        return
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        r = etree.SubElement(p._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t = etree.SubElement(r, qn("a:t"))
        t.text = text


def write_source_placeholder(shape, text, font_size_pt, font_name):
    tf = shape.text_frame
    para = tf.paragraphs[0]
    for r in list(para.runs):
        r.text = ""
    if para.runs:
        run = para.runs[0]
        run.text = text
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text
        run = para.runs[0]
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        run._r.insert(0, rPr)
    rPr.set("sz", str(font_size_pt * 100))
    for tag in [qn("a:latin"), qn("a:ea")]:
        old = rPr.find(tag)
        if old is not None:
            rPr.remove(old)
        etree.SubElement(rPr, tag, attrib={"typeface": font_name})


def hex2rgb(h):
    h = h.replace("#", "")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def build_chart(slide, data, left, top, w, h):
    from pptx.chart.data import CategoryChartData
    cats = data.get("categories", [])
    periods = data.get("periods", [])
    line_cfg = data.get("line")
    threshold = data.get("label_threshold_pct", 5.0)
    n_cats = len(cats)
    n_per = len(periods)

    cd = CategoryChartData()
    cd.categories = [p["label"] for p in periods]
    for ci, cat in enumerate(cats):
        vals = [p["values"][ci] if ci < len(p.get("values", [])) else 0 for p in periods]
        cd.add_series(cat["name"], vals)
    has_line = line_cfg is not None
    if has_line:
        cd.add_series(line_cfg["series_name"], [p.get("line_value", 0) for p in periods])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, left, top, w, h, cd)
    chart = cf.chart
    pa = chart._chartSpace.chart.plotArea
    bc = pa.findall(qn('c:barChart'))[0]

    if has_line:
        sers = bc.findall(qn('c:ser'))
        ls = copy.deepcopy(sers[-1])
        bc.remove(sers[-1])
        lc = etree.SubElement(pa, qn('c:lineChart'))
        etree.SubElement(lc, qn('c:grouping'), attrib={'val': 'standard'})
        etree.SubElement(lc, qn('c:varyColors'), attrib={'val': '0'})
        lc.append(ls)
        mk = ls.find(qn('c:marker'))
        if mk is None:
            mk = etree.SubElement(ls, qn('c:marker'))
        sym = mk.find(qn('c:symbol'))
        if sym is None:
            sym = etree.SubElement(mk, qn('c:symbol'))
        sym.set('val', 'circle')
        sz = mk.find(qn('c:size'))
        if sz is None:
            sz = etree.SubElement(mk, qn('c:size'))
        sz.set('val', '8')
        sec_v = "2094734553"
        sec_c = "2094734554"
        etree.SubElement(lc, qn('c:axId'), attrib={'val': sec_c})
        etree.SubElement(lc, qn('c:axId'), attrib={'val': sec_v})
        sc = etree.SubElement(pa, qn('c:catAx'))
        etree.SubElement(sc, qn('c:axId'), attrib={'val': sec_c})
        s2 = etree.SubElement(sc, qn('c:scaling'))
        etree.SubElement(s2, qn('c:orientation'), attrib={'val': 'minMax'})
        etree.SubElement(sc, qn('c:delete'), attrib={'val': '1'})
        etree.SubElement(sc, qn('c:axPos'), attrib={'val': 'b'})
        etree.SubElement(sc, qn('c:crossAx'), attrib={'val': sec_v})
        sv = etree.SubElement(pa, qn('c:valAx'))
        etree.SubElement(sv, qn('c:axId'), attrib={'val': sec_v})
        s3 = etree.SubElement(sv, qn('c:scaling'))
        etree.SubElement(s3, qn('c:orientation'), attrib={'val': 'minMax'})
        etree.SubElement(sv, qn('c:delete'), attrib={'val': '1'})
        etree.SubElement(sv, qn('c:axPos'), attrib={'val': 'r'})
        etree.SubElement(sv, qn('c:numFmt'), attrib={'formatCode': '0.0"%"', 'sourceLinked': '0'})
        etree.SubElement(sv, qn('c:crossAx'), attrib={'val': sec_c})
        etree.SubElement(sv, qn('c:crosses'), attrib={'val': 'max'})
        # Line color: data-driven, fallback to brand accent if unspecified
        raw_color = line_cfg.get("color") or (
            f"#{LINE_FALLBACK_HEX}" if LINE_FALLBACK_HEX else DEFAULT_LINE_COLOR
        )
        lclr = raw_color.replace("#", "")
        lsp = ls.find(qn('c:spPr'))
        if lsp is None:
            lsp = etree.SubElement(ls, qn('c:spPr'))
        ln = lsp.find(qn('a:ln'))
        if ln is None:
            ln = etree.SubElement(lsp, qn('a:ln'))
        ln.set('w', '19050')
        lsf = ln.find(qn('a:solidFill'))
        if lsf is None:
            lsf = etree.SubElement(ln, qn('a:solidFill'))
        for c in list(lsf):
            lsf.remove(c)
        etree.SubElement(lsf, qn('a:srgbClr'), attrib={'val': lclr})
        msp = mk.find(qn('c:spPr'))
        if msp is None:
            msp = etree.SubElement(mk, qn('c:spPr'))
        msf = etree.SubElement(msp, qn('a:solidFill'))
        etree.SubElement(msf, qn('a:srgbClr'), attrib={'val': lclr})
        _add_line_dlbls(ls)

    # Bar colors + custom labels
    bar_sers = bc.findall(qn('c:ser'))
    for si, bs in enumerate(bar_sers):
        clr = cats[si]["color"].replace("#", "") if si < n_cats else "999999"
        sp = bs.find(qn('c:spPr'))
        if sp is None:
            sp = etree.SubElement(bs, qn('c:spPr'))
        sf = sp.find(qn('a:solidFill'))
        if sf is None:
            sf = etree.SubElement(sp, qn('a:solidFill'))
        for c in list(sf):
            sf.remove(c)
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': clr})
        _add_bar_dlbls(bs, si, periods, threshold)

    gw = bc.find(qn('c:gapWidth'))
    if gw is None:
        gw = etree.SubElement(bc, qn('c:gapWidth'))
    gw.set('val', '80')
    ov = bc.find(qn('c:overlap'))
    if ov is None:
        ov = etree.SubElement(bc, qn('c:overlap'))
    ov.set('val', '100')

    for vx in pa.findall(qn('c:valAx')):
        for mg in vx.findall(qn('c:majorGridlines')):
            vx.remove(mg)
        for mg in vx.findall(qn('c:minorGridlines')):
            vx.remove(mg)
    for cx in pa.findall(qn('c:catAx')):
        for mg in cx.findall(qn('c:majorGridlines')):
            cx.remove(mg)

    for ax in pa.findall(qn('c:catAx')) + pa.findall(qn('c:valAx')):
        de = ax.find(qn('c:delete'))
        if de is not None and de.get('val') == '1':
            continue
        tp = ax.find(qn('c:txPr'))
        if tp is None:
            tp = etree.SubElement(ax, qn('c:txPr'))
        bp = tp.find(qn('a:bodyPr'))
        if bp is None:
            bp = etree.SubElement(tp, qn('a:bodyPr'))
            tp.insert(0, bp)
        if tp.find(qn('a:lstStyle')) is None:
            etree.SubElement(tp, qn('a:lstStyle'))
        p = tp.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(tp, qn('a:p'))
        pp = p.find(qn('a:pPr'))
        if pp is None:
            pp = etree.SubElement(p, qn('a:pPr'))
        dr = pp.find(qn('a:defRPr'))
        if dr is None:
            dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': AXIS_FONT_SZ})
        else:
            dr.set('sz', AXIS_FONT_SZ)
        if dr.find(qn('a:latin')) is None:
            etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT})
        if dr.find(qn('a:ea')) is None:
            etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT})

    chart.has_legend = False
    chart.has_title = False
    psp = pa.find(qn('c:spPr'))
    if psp is None:
        psp = etree.SubElement(pa, qn('c:spPr'))
    if psp.find(qn('a:noFill')) is None:
        etree.SubElement(psp, qn('a:noFill'))
    print(f"  Chart: {n_per} periods, {n_cats} series" + (" + line" if has_line else ""))
    return cf


def _add_bar_dlbls(ser, si, periods, threshold):
    old = ser.find(qn('c:dLbls'))
    if old is not None:
        ser.remove(old)
    dlbls = etree.SubElement(ser, qn('c:dLbls'))
    for pi, per in enumerate(periods):
        vals = per.get("values", [])
        val = vals[si] if si < len(vals) else 0
        total = per.get("total", 0)
        pct = (val / total * 100) if total > 0 else 0
        dlbl = etree.SubElement(dlbls, qn('c:dLbl'))
        etree.SubElement(dlbl, qn('c:idx'), attrib={'val': str(pi)})
        if val == 0 or pct < threshold:
            etree.SubElement(dlbl, qn('c:delete'), attrib={'val': '1'})
        else:
            txt = f"{val}\n({pct:.1f}%)"
            tx = etree.SubElement(dlbl, qn('c:tx'))
            rich = etree.SubElement(tx, qn('c:rich'))
            etree.SubElement(rich, qn('a:bodyPr'))
            etree.SubElement(rich, qn('a:lstStyle'))
            p = etree.SubElement(rich, qn('a:p'))
            pp = etree.SubElement(p, qn('a:pPr'), attrib={'algn': 'ctr'})
            dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': DATA_LABEL_SZ, 'b': '1'})
            sf = etree.SubElement(dr, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': 'FFFFFF'})
            etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT})
            etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT})
            r = etree.SubElement(p, qn('a:r'))
            rPr = etree.SubElement(r, qn('a:rPr'), attrib={
                'lang': 'ja-JP', 'sz': DATA_LABEL_SZ, 'b': '1'
            })
            sf2 = etree.SubElement(rPr, qn('a:solidFill'))
            etree.SubElement(sf2, qn('a:srgbClr'), attrib={'val': 'FFFFFF'})
            etree.SubElement(rPr, qn('a:latin'), attrib={'typeface': FONT})
            etree.SubElement(rPr, qn('a:ea'), attrib={'typeface': FONT})
            t = etree.SubElement(r, qn('a:t'))
            t.text = txt
            etree.SubElement(dlbl, qn('c:dLblPos'), attrib={'val': 'ctr'})
            for k in ['showLegendKey', 'showVal', 'showCatName', 'showSerName', 'showPercent']:
                etree.SubElement(dlbl, qn(f'c:{k}'), attrib={'val': '0'})
    for k in ['showLegendKey', 'showVal', 'showCatName', 'showSerName', 'showPercent']:
        etree.SubElement(dlbls, qn(f'c:{k}'), attrib={'val': '0'})


def _add_line_dlbls(ser):
    dl = ser.find(qn('c:dLbls'))
    if dl is None:
        dl = etree.SubElement(ser, qn('c:dLbls'))
    for c in list(dl):
        dl.remove(c)
    etree.SubElement(dl, qn('c:numFmt'), attrib={'formatCode': '0.0"%"', 'sourceLinked': '0'})
    for k in ['showLegendKey', 'showCatName', 'showSerName', 'showPercent', 'showBubbleSize']:
        etree.SubElement(dl, qn(f'c:{k}'), attrib={'val': '0'})
    etree.SubElement(dl, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dl, qn('c:dLblPos'), attrib={'val': 't'})
    tp = etree.SubElement(dl, qn('c:txPr'))
    etree.SubElement(tp, qn('a:bodyPr'))
    etree.SubElement(tp, qn('a:lstStyle'))
    p = etree.SubElement(tp, qn('a:p'))
    pp = etree.SubElement(p, qn('a:pPr'))
    dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': LINE_DATA_LABEL_SZ, 'b': '1'})
    etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT})
    etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT})
    sf = etree.SubElement(dr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': TEXT_HEX})


def add_total_labels(slide, data, cl, ct, cw, ch):
    periods = data.get("periods", [])
    nc = len(periods)
    if nc == 0:
        return
    plm = cw * 0.06
    prm = cw * 0.04
    ptm = ch * 0.05
    pbm = ch * 0.14
    pl = cl + plm
    pw = cw - plm - prm
    pb = ct + ch - pbm
    ph = pb - (ct + ptm)
    caw = pw / nc
    max_t = max(p.get("total", 0) for p in periods)
    ax_max = max_t * 1.20
    for pi, per in enumerate(periods):
        t = per.get("total", 0)
        cx = pl + (pi + 0.5) * caw
        ty = pb - (t / ax_max) * ph
        lw = Inches(0.80)
        lh = Inches(0.25)
        tb = slide.shapes.add_textbox(int(cx - lw / 2), int(ty - lh - Emu(20000)), lw, lh)
        tb.text_frame.word_wrap = False
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(t)
        r.font.size = Pt(TOTAL_LABEL_PT)
        r.font.bold = True
        r.font.color.rgb = COLOR_TEXT
        r.font.name = FONT
    print(f"  Total labels: {nc}")


def add_trend_arrow(slide, data, cl, ct, cw, ch):
    if not data.get("trend_arrow", False):
        return
    periods = data.get("periods", [])
    line_cfg = data.get("line")
    if not line_cfg or len(periods) < 2:
        return
    nc = len(periods)
    plm = cw * 0.06
    prm = cw * 0.04
    ptm = ch * 0.05
    pbm = ch * 0.14
    pl = cl + plm
    pw = cw - plm - prm
    pb = ct + ch - pbm
    ph = pb - (ct + ptm)
    caw = pw / nc
    max_t = max(p.get("total", 0) for p in periods)
    ax_max = max_t * 1.20
    gap = max(ph * 0.18, Inches(0.50))
    for i in range(nc - 1):
        ti = periods[i].get("total", 0)
        tj = periods[i + 1].get("total", 0)
        lv = periods[i + 1].get("line_value", 0)
        cx_i = pl + (i + 0.5) * caw
        cx_j = pl + (i + 1.5) * caw
        bti = pb - (ti / ax_max) * ph
        btj = pb - (tj / ax_max) * ph
        ay_i = int(bti - gap)
        ay_j = int(btj - gap)
        mn = int(ct - Inches(0.20))
        ay_i = max(ay_i, mn)
        ay_j = max(ay_j, mn)
        conn = slide.shapes.add_connector(1, int(cx_i), ay_i, int(cx_j), ay_j)
        conn.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        conn.line.width = Pt(2.0)
        sp = conn._element.find(qn('p:spPr'))
        if sp is None:
            sp = conn._element.find(qn('a:spPr'))
        if sp is not None:
            ln = sp.find(qn('a:ln'))
            if ln is None:
                ln = etree.SubElement(sp, qn('a:ln'))
            etree.SubElement(ln, qn('a:tailEnd'), attrib={
                'type': 'triangle', 'w': 'med', 'len': 'med'
            })
        mx = (cx_i + cx_j) / 2
        my = (ay_i + ay_j) / 2
        ow = Inches(0.90)
        oh = Inches(0.28)
        tb = slide.shapes.add_textbox(int(mx - ow / 2), int(my - oh / 2), ow, oh)
        tb.text_frame.word_wrap = False
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{lv:.1f}%"
        r.font.size = Pt(TREND_TEXT_PT)
        r.font.bold = True
        r.font.color.rgb = COLOR_TEXT
        r.font.name = FONT
    print(f"  Trend arrows: {nc - 1}")


def add_legend(slide, data, left, top, w):
    cats = data.get("categories", [])
    line_cfg = data.get("line")
    periods = data.get("periods", [])

    avail_h = LEGEND_BOTTOM - top

    active = []
    for ci, cat in enumerate(cats):
        if any((p.get("values", [])[ci] if ci < len(p.get("values", [])) else 0) > 0 for p in periods):
            active.append(cat)
    n_active = len(active)

    line_header_h = Inches(0.45) if line_cfg else 0
    n_rows = math.ceil(n_active / 2)
    bar_area_h = avail_h - line_header_h
    if n_rows > 0:
        lh = min(Inches(0.50), max(Inches(0.28), int(bar_area_h / n_rows)))
    else:
        lh = Inches(0.40)

    # Roleup C4 enforces font sizes ∈ {22,14,12,10,6}; 10pt fits all row heights.
    if _THEME and _THEME.id != "stellar_aiz":
        font_sz = Pt(10)
        line_font_sz = Pt(10)
        sq = Inches(0.16) if lh >= Inches(0.32) else Inches(0.14)
    else:
        if lh >= Inches(0.40):
            font_sz = Pt(12)
            sq = Inches(0.18)
            line_font_sz = Pt(12)
        elif lh >= Inches(0.32):
            font_sz = Pt(11)
            sq = Inches(0.16)
            line_font_sz = Pt(11)
        else:
            font_sz = Pt(10)
            sq = Inches(0.14)
            line_font_sz = Pt(10)

    content_h = line_header_h + n_rows * lh
    y_offset = (avail_h - content_h) / 2
    y = top + y_offset

    cw = w / 2

    if line_cfg:
        raw_line_color = line_cfg.get("color") or (
            f"#{LINE_FALLBACK_HEX}" if LINE_FALLBACK_HEX else DEFAULT_LINE_COLOR
        )
        lclr = hex2rgb(raw_line_color)
        lw = Inches(0.40)
        ly = int(y + line_header_h / 2)
        cn = slide.shapes.add_connector(1, int(left), ly, int(left + lw), ly)
        cn.line.color.rgb = lclr
        cn.line.width = Pt(2.0)
        mk_sz = Inches(0.12)
        mk = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            int(left + (lw - mk_sz) / 2), int(ly - mk_sz / 2), mk_sz, mk_sz)
        mk.fill.solid()
        mk.fill.fore_color.rgb = lclr
        mk.line.fill.background()
        tx = int(left + lw + Inches(0.08))
        tb = slide.shapes.add_textbox(tx, int(y), int(w - lw - Inches(0.08)), int(line_header_h))
        tb.text_frame.word_wrap = False
        tb.text_frame.paragraphs[0].space_before = Pt(0)
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = line_cfg["series_name"]
        r.font.size = line_font_sz
        r.font.color.rgb = COLOR_TEXT
        r.font.name = FONT
        y += line_header_h

    for i, cat in enumerate(active):
        col = i % 2
        row = i // 2
        ix = left + col * cw
        iy = y + row * lh
        sqy = iy + (lh - sq) / 2
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(ix), int(sqy), sq, sq)
        s.fill.solid()
        s.fill.fore_color.rgb = hex2rgb(cat["color"])
        s.line.fill.background()
        tx = int(ix + sq + Inches(0.06))
        tb = slide.shapes.add_textbox(tx, int(iy), int(cw - sq - Inches(0.06)), int(lh))
        tb.text_frame.word_wrap = False
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = cat["name"]
        r.font.size = font_sz
        r.font.color.rgb = COLOR_TEXT
        r.font.name = FONT
    print(f"  Legend: {n_active} items (2-col, {font_sz.pt}pt, centered)")


def add_unit(slide, text, left, top):
    tb = slide.shapes.add_textbox(left, top, Inches(3.0), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(UNIT_LABEL_PT)
    r.font.color.rgb = COLOR_SRC
    r.font.name = FONT


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "sga-breakdown")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "categories", "periods"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "categories", "periods", "line", "trend_arrow",
            "label_threshold_pct", "unit_label",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    print(f"=== SGA Breakdown Slide (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # Top placeholder (brand-aware)
    top_text = resolve_top_text(data, theme)
    set_text(find_shape(slide, SHAPE_MAIN_MSG), top_text)

    sub_text = resolve_subtitle_text(data, theme) or "対売上販管費比率と販管費構成比の推移"
    set_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)

    # Source: stella=Source textbox / roleup=Source 3 placeholder
    src = data.get("source", "")
    if src:
        body = src if src.startswith("出典") else f"出典：{src}"
        src_shape = find_shape(slide, SHAPE_SOURCE)
        if src_shape is not None:
            if theme.id == "stellar_aiz":
                set_text(src_shape, body)
            else:
                write_source_placeholder(src_shape, body, SOURCE_FONT_PT, FONT)
            print(f"  Source ({SHAPE_SOURCE}): {body[:50]}")

    # Stella V1: 'Content Area' placeholder removal (roleup template doesn't have it)
    _silent_remove_shape(slide, SHAPE_CONTENT)

    chart_y = PANEL_Y + CHART_Y_OFFSET

    ul = data.get("unit_label", "")
    if ul:
        add_unit(slide, ul, CHART_X, PANEL_Y)
    build_chart(slide, data, CHART_X, chart_y, CHART_W, CHART_H)
    add_total_labels(slide, data, CHART_X, chart_y, CHART_W, CHART_H)
    add_trend_arrow(slide, data, CHART_X, chart_y, CHART_W, CHART_H)
    add_legend(slide, data, LEGEND_X, PANEL_Y + Inches(0.10), LEGEND_W)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    print(f"  Done: {args.output}")
    _finalize_pptx(args.output)


if __name__ == "__main__":
    main()
