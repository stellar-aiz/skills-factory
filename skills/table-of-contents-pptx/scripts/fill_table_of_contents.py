"""
fill_table_of_contents.py — 目次（Table of Contents）スライドを生成

レイアウト:
  - 上部: タイトル (Title 1) + サブタイトル (Text Placeholder 2)
  - 中央: 3〜7 セクション (色付き番号バッジ + Bold タイトル + サブ項目 + ページ番号)
  - 下部: 出典 (任意)

Brand-aware (Phase 2, ISSUE-010):
  --brand stellar_aiz : 13.33×7.50 in / Meiryo UI / 28-18-11pt / 7 色ローテ
  --brand roleup      : 11.69×8.27 in / Yu Gothic UI / 22-14-10pt / chart_palette 8 色ローテ

TOC は出典が任意のため C6 (source_required) を profile から除外。

Usage:
  python fill_table_of_contents.py \
    --data /home/claude/toc_data.json \
    --brand stellar_aiz \
    --output /mnt/user-data/outputs/TableOfContents_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text  # noqa: E402

SKILL_ID = "table-of-contents-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Template placeholder shape names ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_ROLEUP = "Source 3"

# ── Default values (stella) — reassigned in main() via _apply_theme(theme) ──
TOC_LEFT = Inches(0.80)
TOC_TOP = Inches(1.55)
TOC_WIDTH = Inches(11.73)
TOC_HEIGHT = Inches(5.40)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.25)

BADGE_W = Inches(0.85)
BADGE_OFFSET = Inches(0.20)
PAGE_W = Inches(0.80)

# Colors (stella default).
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTEXT = RGBColor(0x55, 0x55, 0x55)
COLOR_DIVIDER = RGBColor(0xDD, 0xDD, 0xDD)

# Section accent color rotations (stella legacy: 7 colors).
SECTION_COLORS_STELLA = [
    RGBColor(0x2E, 0x4A, 0x6B),   # 紺
    RGBColor(0x7B, 0x4F, 0xB0),   # 紫
    RGBColor(0x2E, 0x6F, 0xBF),   # 青
    RGBColor(0x3D, 0x8F, 0x5A),   # 緑
    RGBColor(0xDA, 0x7A, 0x2D),   # オレンジ
    RGBColor(0xC0, 0x3A, 0x3A),   # 赤
    RGBColor(0x59, 0x59, 0x59),   # グレー
]

# Fonts (stella defaults).
FONT_NAME_JP = "Meiryo UI"
FONT_NAME_LATIN = "Arial"
FONT_SIZE_NUMBER = Pt(28)
FONT_SIZE_SECTION = Pt(18)
FONT_SIZE_SUBITEM = Pt(11)
FONT_SIZE_PAGE = Pt(11)
FONT_SIZE_SOURCE = Pt(10)

# Theme module-global; populated in main() via _apply_theme(theme).
_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global TOC_LEFT, TOC_TOP, TOC_WIDTH, TOC_HEIGHT
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global BADGE_W, BADGE_OFFSET, PAGE_W
    global COLOR_TEXT, COLOR_SOURCE, COLOR_SUBTEXT
    global FONT_NAME_JP, FONT_NAME_LATIN
    global FONT_SIZE_NUMBER, FONT_SIZE_SECTION, FONT_SIZE_SUBITEM, FONT_SIZE_PAGE, FONT_SIZE_SOURCE

    _THEME = theme

    TOC_LEFT = theme.layout("toc_left_in")
    TOC_TOP = theme.layout("toc_top_in")
    TOC_WIDTH = theme.layout("toc_width_in")
    TOC_HEIGHT = theme.layout("toc_height_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")
    BADGE_W = theme.layout("badge_w_in")
    BADGE_OFFSET = theme.layout("badge_offset_in")
    PAGE_W = theme.layout("page_w_in")

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    COLOR_SUBTEXT = theme.color("text")  # subtext follows text color in both brands

    FONT_NAME_JP = theme.font_ea
    FONT_NAME_LATIN = theme.font_latin

    if theme.id == "stellar_aiz":
        FONT_SIZE_NUMBER = Pt(28)
        FONT_SIZE_SECTION = Pt(18)
        FONT_SIZE_SUBITEM = Pt(11)
        FONT_SIZE_PAGE = Pt(11)
        FONT_SIZE_SOURCE = Pt(10)
    else:
        # roleup 許容集合 {22, 14, 12, 10, 6} に揃える
        FONT_SIZE_NUMBER = Pt(22)
        FONT_SIZE_SECTION = Pt(14)
        FONT_SIZE_SUBITEM = Pt(10)
        FONT_SIZE_PAGE = Pt(10)
        FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")  # 6pt


def _silent_remove_shape(slide, shape_name: str) -> None:
    """Remove a shape by name without printing a warning. No-op if absent."""
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def _section_color(theme, section_index: int) -> RGBColor:
    """Resolve accent color for section_index (1-based).

    stella: 7 色 hardcode ローテ (V1 互換)
    roleup: theme.chart_palette (8 色) ローテ
    """
    if theme.id == "stellar_aiz":
        return SECTION_COLORS_STELLA[(section_index - 1) % len(SECTION_COLORS_STELLA)]
    palette = theme.chart_palette
    hex_str = palette[(section_index - 1) % len(palette)]
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    # font_name / color default to current module globals (post-_apply_theme).
    if font_name is None:
        font_name = FONT_NAME_JP
    if color is None:
        color = COLOR_TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def draw_section_row(slide, idx, section, left, top, width, height):
    """1 セクション行を描画"""
    color_hex = section.get("color")
    if color_hex:
        color = hex_to_rgb(color_hex)
    else:
        color = _section_color(_THEME, idx)

    # 番号バッジ（左、四角）
    badge_w = BADGE_W
    badge_h = height - Inches(0.10)
    badge_y = top + Inches(0.05)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, badge_y, badge_w, badge_h,
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.shadow.inherit = False

    tf = badge.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "ctr")
    r = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_NUMBER.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_LATIN})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    s = etree.SubElement(sf, qn("a:srgbClr"))
    s.set("val", "FFFFFF")
    t = etree.SubElement(r, qn("a:t"))
    t.text = f"{idx:02d}"

    # コンテンツ領域
    content_left = left + badge_w + BADGE_OFFSET
    content_w = width - badge_w - BADGE_OFFSET - PAGE_W
    page_x = left + width - PAGE_W

    title = section.get("title", f"セクション{idx}")
    subitems = section.get("subitems", [])
    has_subitems = len(subitems) > 0

    tb = slide.shapes.add_textbox(content_left, top + Inches(0.05), content_w, height - Inches(0.10))
    tf2 = tb.text_frame
    tf2.word_wrap = True
    tf2.margin_left = 0; tf2.margin_right = 0
    tf2.margin_top = 0; tf2.margin_bottom = 0
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE if not has_subitems else MSO_ANCHOR.TOP

    for p in list(tf2.paragraphs):
        p._p.getparent().remove(p._p)

    # タイトル段落
    p_title = etree.SubElement(tf2._txBody, qn("a:p"))
    pPr_t = etree.SubElement(p_title, qn("a:pPr"))
    pPr_t.set("algn", "l")
    r_title = etree.SubElement(p_title, qn("a:r"))
    rPr_title = etree.SubElement(r_title, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(FONT_SIZE_SECTION.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr_title, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr_title, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf_title = etree.SubElement(rPr_title, qn("a:solidFill"))
    s_title = etree.SubElement(sf_title, qn("a:srgbClr"))
    s_title.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))
    t_title = etree.SubElement(r_title, qn("a:t"))
    t_title.text = title

    # サブ項目段落
    if has_subitems:
        sub_color = "{:02X}{:02X}{:02X}".format(COLOR_SUBTEXT[0], COLOR_SUBTEXT[1], COLOR_SUBTEXT[2])
        for sub in subitems:
            p_sub = etree.SubElement(tf2._txBody, qn("a:p"))
            pPr_sub = etree.SubElement(p_sub, qn("a:pPr"), attrib={
                "marL": "180000",
                "indent": "-180000",
            })
            spcBef = etree.SubElement(pPr_sub, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "200"})

            buChar = etree.SubElement(pPr_sub, qn("a:buChar"), attrib={"char": "▸"})
            # bullet typeface: 本文と同じ ea フォントに揃える (C8)
            buFont = etree.SubElement(pPr_sub, qn("a:buFont"), attrib={"typeface": FONT_NAME_JP})
            buClr = etree.SubElement(pPr_sub, qn("a:buClr"))
            buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
            buClrSolid.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))

            r_sub = etree.SubElement(p_sub, qn("a:r"))
            rPr_sub = etree.SubElement(r_sub, qn("a:rPr"), attrib={
                "lang": "ja-JP",
                "sz": str(int(FONT_SIZE_SUBITEM.pt * 100)),
                "b": "0",
            })
            etree.SubElement(rPr_sub, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
            etree.SubElement(rPr_sub, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
            sf_sub = etree.SubElement(rPr_sub, qn("a:solidFill"))
            s_sub = etree.SubElement(sf_sub, qn("a:srgbClr"))
            s_sub.set("val", sub_color)
            t_sub = etree.SubElement(r_sub, qn("a:t"))
            t_sub.text = sub

    # ページ番号（右端）
    page = section.get("page", "")
    if page:
        add_text_box(
            slide, f"P. {page}",
            page_x, top, PAGE_W, height,
            FONT_SIZE_PAGE, bold=True,
            color=color,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )

    # 区切り線（行の下）
    div_y = top + height - Emu(int(Inches(0.01)))
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + badge_w + BADGE_OFFSET, div_y,
        width - badge_w - BADGE_OFFSET, Emu(int(Inches(0.01))),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_DIVIDER
    div.line.fill.background()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "table-of-contents")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    _mm = data.get("main_message", "目次")
    if len(_mm) > 65:
        raise ValueError(
            f"main_message は 65 字以内（受領: {len(_mm)}）: {_mm[:80]}..."
        )

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top / subtitle placeholder semantics differ between brands:
    #  - stella: Title 1 = main_message (結論文), Text Placeholder 2 = chart_title
    #  - roleup: Title 1 = chart_title (スライドタイトル), Text Placeholder 2 = main_message
    top_text = resolve_top_text(data, theme) or data.get("main_message", "目次")
    sub_text = resolve_subtitle_text(data, theme) or data.get("chart_title", "Table of Contents")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:40]}")
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:40]}")

    # Roleup: silently remove brown guide rectangles carried by template
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    sections = data.get("sections", [])
    if not sections:
        print("  ✗ ERROR: 'sections' is required", file=sys.stderr)
        sys.exit(1)

    n = len(sections)
    if n > 7:
        print(f"  ⚠ WARNING: {n} sections > 7. Only first 7 will be shown.", file=sys.stderr)
        sections = sections[:7]
        n = 7

    # 行の高さを計算（セクション数に応じて）
    gap = Inches(0.10)
    section_h = Emu(int((TOC_HEIGHT - gap * (n - 1)) / n))

    for i, section in enumerate(sections):
        y = TOC_TOP + (section_h + gap) * i
        draw_section_row(slide, i + 1, section, TOC_LEFT, y, TOC_WIDTH, section_h)
        print(f"  ✓ Section {i+1}: {section.get('title', '')[:40]}")

    # 出典: roleup なら Source 3 placeholder にセット、stella は dynamic textbox
    source = data.get("source", "")
    if source:
        if theme.id == "stellar_aiz":
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE_ROLEUP)
            if source_shape is not None:
                set_textbox_text(source_shape, source)
            else:
                add_text_box(
                    slide, source,
                    SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                    FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                    align=PP_ALIGN.LEFT,
                )

    # 出典が無いケース: roleup の Source 3 placeholder は空のまま残す
    # → C5 (font size) は PASS (defaultStyle 継承), C6 (source_required) は profile 除外で OK

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
