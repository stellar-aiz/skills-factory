"""
fill_market_environment.py — 市場環境分析チャートをPPTXネイティブオブジェクトで生成

生成するネイティブオブジェクト（すべてPowerPoint上で人間が編集可能）:
  - 積み上げ棒チャート＋折れ線（PowerPointネイティブ複合チャート）
  - CAGR注釈: テキストボックス
  - 成長率注釈: コネクタ＋楕円＋テキスト
  - 期間区切り線: コネクタ（破線）
  - 凡例: Shape＋テキストボックス

Usage:
  python fill_market_environment.py \
    --data /home/claude/market_environment_data.json \
    --template <SKILL_DIR>/assets/market-environment-template.pptx \
    --output /mnt/user-data/outputs/MarketEnvironment_output.pptx
"""

import argparse, copy, json, math, os, re, sys

# ── brand_resolver bootstrap (skills/_common/lib/brand_resolver.py) ──
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import (  # noqa: E402
    format_fiscal_period,
    resolve_top_text,
    resolve_subtitle_text,
)

SKILL_ID = "market-environment-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape名マッピング ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE       = "Source"

# ── Brand-aware module globals ──
# Default values match stella for safety; reassigned in main() via _apply_theme(theme)
# after argparse resolves --brand. SHAPE_* names above are template-structure invariants.
PANEL_Y = Inches(1.50)
CHART_X = Inches(0.50);  CHART_W = Inches(9.00);  CHART_H = Inches(5.00)
CHART_Y = PANEL_Y + Inches(0.55)
CAGR_X  = Inches(9.80);  CAGR_W  = Inches(3.20)

COLOR_TEXT   = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)

# ─── 配色（brand 経由で resolve、stella 既定値を初期値として保持） ───
# V1 以降は brand theme.json (skills/_common/brands/<id>/theme.json) が単一情報源。
# 旧来 skills/_common/styles/chart_palette.md の手動同期運用は stella のみ継続
# （brand-aware にすると Roleup 等は theme.json から自動 resolve され同期負荷が消える）。
CHART_PALETTE = [
    "#4E79A7", "#F28E2B", "#59A14F", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F",
]
OTHER_COLOR = "#BAB0AC"
TARGET_COLOR = "#E15759"
LABEL_BAR_COLOR = "#4E79A7"
LABEL_BG_COLOR = "#E8EEF5"

FONT_JP = "Meiryo UI"
COLOR_SUBTITLE = RGBColor(0x33, 0x33, 0x33)  # default = text; reassigned in _apply_theme

# Phase 4 (ISSUE-011): module-level theme reference so chart helpers can read
# fiscal_period_format / line_height_pt / layout_rule without an extra parameter.
_THEME = None


def _body_pt(stella_fallback_pt):
    """本文・表・チャート要素フォントの brand 別解決ヘルパー。

    Roleup: theme.font_size_body_pt (= 10pt 統一)。ユーザー指示
      「タイトル/キーメッセージ/サブタイトル/出典 以外は 10pt」に準拠。
    Stella: 旧 hardcode 値を維持(regression-zero)。
    """
    if _THEME is None or _THEME.id == "stellar_aiz":
        return Pt(stella_fallback_pt)
    return _THEME.pt("font_size_body_pt")


def _body_sz(stella_fallback_sz_str):
    """OOXML sz 属性 (100×pt の文字列) 用の brand 別解決ヘルパー。"""
    if _THEME is None or _THEME.id == "stellar_aiz":
        return stella_fallback_sz_str
    return str(_THEME.pt_value("font_size_body_pt") * 100)


def _palette_color(index: int, total: int) -> str:
    if total <= 1:
        return CHART_PALETTE[0]
    return CHART_PALETTE[index % len(CHART_PALETTE)]


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme.

    Called once from main() after `--brand` is parsed. Module-load-time
    defaults above remain correct for direct imports / tests that don't
    go through main() (regression safety net).
    """
    global PANEL_Y, CHART_X, CHART_W, CHART_H, CHART_Y, CAGR_X, CAGR_W
    global COLOR_TEXT, COLOR_SOURCE, COLOR_SUBTITLE
    global CHART_PALETTE, OTHER_COLOR, TARGET_COLOR, LABEL_BAR_COLOR, LABEL_BG_COLOR
    global FONT_JP
    global _THEME
    _THEME = theme

    PANEL_Y = theme.layout("panel_y_in")
    CHART_X = theme.layout("chart_x_in")
    CHART_W = theme.layout("chart_w_in")
    CHART_H = theme.layout("chart_h_in")
    CHART_Y = PANEL_Y + theme.layout("chart_y_offset_in")
    CAGR_X  = theme.layout("cagr_x_in")
    CAGR_W  = theme.layout("cagr_w_in")

    COLOR_TEXT     = theme.color("text")
    COLOR_SOURCE   = theme.color("source")
    COLOR_SUBTITLE = theme.color("subtitle")  # roleup #897141 / stella #333333 (= text)

    CHART_PALETTE   = list(theme.chart_palette)
    OTHER_COLOR     = theme.hex("highlight_other")
    TARGET_COLOR    = theme.hex("highlight_target")
    LABEL_BAR_COLOR = theme.hex("label_bar")
    LABEL_BG_COLOR  = theme.hex("label_bg")

    FONT_JP = theme.font_ea

# ── ユーティリティ ──
def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name: return s
    print(f"  ⚠ Shape '{name}' not found", file=sys.stderr); return None

def set_textbox_text(shape, text):
    if shape is None: return
    tf = shape.text_frame; p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]: r.text = ""
    else:
        r = etree.SubElement(p._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang":"ja-JP"})
        t = etree.SubElement(r, qn("a:t")); t.text = text

def remove_shape(slide, name):
    s = find_shape(slide, name)
    if s: slide.shapes._spTree.remove(s._element); print(f"  ✓ Removed '{name}'")

def _silent_remove_shape(slide, name):
    """find_shape の warning を出さずに削除を試みる(brand 別 shape 名フォールバック用)"""
    for s in slide.shapes:
        if s.name == name:
            slide.shapes._spTree.remove(s._element)
            return True
    return False

def _silent_find_shape(slide, name):
    """find_shape の warning を出さずに検索する(brand 別 shape 名フォールバック用)"""
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def _hex2rgb(h):
    h = h.replace("#","")
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

# ── Y 軸スケール計算 ──
def _round_up_nice(x):
    """1.93 → 2.0、4.7 → 5、47 → 50、187 → 200 のように見栄えの良い max にする"""
    if x <= 0:
        return 1
    mag = 10 ** math.floor(math.log10(x))
    n = x / mag
    if   n <= 1.0: nn = 1.0
    elif n <= 1.2: nn = 1.2
    elif n <= 1.5: nn = 1.5
    elif n <= 2.0: nn = 2.0
    elif n <= 2.5: nn = 2.5
    elif n <= 3.0: nn = 3.0
    elif n <= 4.0: nn = 4.0
    elif n <= 5.0: nn = 5.0
    elif n <= 6.0: nn = 6.0
    elif n <= 8.0: nn = 8.0
    else:          nn = 10.0
    return nn * mag

def _calc_primary_axis_max(data, override=None):
    """棒グラフ（積み上げ合計）の最大値から primary Y 軸 max を決める"""
    if override is not None:
        return float(override)
    totals = [sum(d.get("bars", [])) for d in data]
    if not totals:
        return 1.0
    return _round_up_nice(max(totals) * 1.15)

def _calc_secondary_axis_max(data, override=None):
    """折れ線（％想定）の最大値から secondary Y 軸 max を決める。
    値域が 0-100 内なら 120（％的に余裕を持たせる）、
    それ以上なら nice round up を採用。
    """
    if override is not None:
        return float(override)
    vals = [d.get("line_value", 0) for d in data if d.get("line_value") is not None]
    if not vals:
        return 100.0
    mx = max(vals)
    if mx <= 100:
        return 120.0
    return _round_up_nice(mx * 1.15)

def _set_val_axis_scale(val_ax, axis_min, axis_max):
    """既存の c:valAx に明示的な min / max を持つ c:scaling を設定する。
    OOXML schema 順序を保つため、c:scaling 直下に min / max を append する。
    """
    sc = val_ax.find(qn('c:scaling'))
    if sc is None:
        sc = etree.SubElement(val_ax, qn('c:scaling'))
        # axId の直後に scaling を移動
        ax_id = val_ax.find(qn('c:axId'))
        if ax_id is not None:
            val_ax.remove(sc)
            ax_id.addnext(sc)
    # orientation を確実に
    ori = sc.find(qn('c:orientation'))
    if ori is None:
        ori = etree.SubElement(sc, qn('c:orientation'), attrib={'val': 'minMax'})
    # 既存の min / max を撤去して再設定
    for tag in ('c:max', 'c:min'):
        for el in sc.findall(qn(tag)):
            sc.remove(el)
    etree.SubElement(sc, qn('c:max'), attrib={'val': f'{axis_max:g}'})
    etree.SubElement(sc, qn('c:min'), attrib={'val': f'{axis_min:g}'})

def _reorder_plot_area(pa):
    """OOXML CT_PlotArea schema 順序に正規化:
    layout? → chart elements (barChart/lineChart/...) → axes (catAx/valAx/...) → others
    LibreOffice / PowerPoint が schema 違反順序でレンダ崩壊するのを防ぐ。
    """
    chart_tags = {qn(t) for t in (
        'c:areaChart','c:area3DChart','c:bar3DChart','c:barChart','c:bubbleChart',
        'c:doughnutChart','c:line3DChart','c:lineChart','c:ofPieChart',
        'c:pie3DChart','c:pieChart','c:radarChart','c:scatterChart',
        'c:stockChart','c:surface3DChart','c:surfaceChart',
    )}
    axis_tags = {qn(t) for t in ('c:valAx','c:catAx','c:dateAx','c:serAx')}
    layout_tag = qn('c:layout')

    children = list(pa)
    layouts = [c for c in children if c.tag == layout_tag]
    charts  = [c for c in children if c.tag in chart_tags]
    axes    = [c for c in children if c.tag in axis_tags]
    others  = [c for c in children if c not in layouts and c not in charts and c not in axes]

    for c in children:
        pa.remove(c)
    for c in layouts: pa.append(c)
    for c in charts:  pa.append(c)
    for c in axes:    pa.append(c)
    for c in others:  pa.append(c)

def _year_to_int(year) -> int:
    s = str(year)
    digits = re.sub(r"\D", "", s)
    if not digits:
        raise ValueError(f"year value contains no digits: {year!r}")
    return int(digits)

# ── チャート生成 ──
def build_stacked_combo_chart(slide, cfg, left, top, w, h):
    from pptx.chart.data import CategoryChartData
    data      = cfg.get("data", [])
    bars_cfg  = cfg.get("stacked_bars", [])
    line_cfg  = cfg.get("line", None)
    num_fmt   = cfg.get("num_format", "0.0")
    n_ser     = len(bars_cfg)

    cd = CategoryChartData()
    # categories: roleup なら "YY/MM期" (公式 vF 例)、stella は西暦 4 桁を維持
    if _THEME is not None and _THEME.fiscal_period_format():
        fy_month = cfg.get("fiscal_year_end_month", 12)
        cd.categories = [format_fiscal_period(_year_to_int(d["year"]), fy_month, _THEME) for d in data]
    else:
        cd.categories = [d["year"] for d in data]
    for si, sb in enumerate(bars_cfg):
        cd.add_series(sb["series_name"],
                      [d["bars"][si] if si < len(d.get("bars",[])) else 0 for d in data])
    has_line = line_cfg is not None
    if has_line:
        cd.add_series(line_cfg["series_name"], [d.get("line_value",0) for d in data])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, left, top, w, h, cd)
    chart = cf.chart
    pa = chart._chartSpace.chart.plotArea
    bc = pa.findall(qn('c:barChart'))[0]

    # === Y 軸スケール自動計算 ===
    primary_axis_max   = _calc_primary_axis_max(data, cfg.get("primary_y_axis_max"))
    secondary_axis_max = _calc_secondary_axis_max(data, cfg.get("secondary_y_axis_max")) if has_line else None

    # 既存（プライマリ）valAx に明示スケールを設定（バー値域に合わせる）
    primary_val_axes = pa.findall(qn('c:valAx'))
    if primary_val_axes:
        _set_val_axis_scale(primary_val_axes[0], 0, primary_axis_max)

    # 折れ線を分離
    if has_line:
        sers = bc.findall(qn('c:ser'))
        ls = copy.deepcopy(sers[-1]); bc.remove(sers[-1])
        lc = etree.SubElement(pa, qn('c:lineChart'))
        etree.SubElement(lc, qn('c:grouping'), attrib={'val':'standard'})
        etree.SubElement(lc, qn('c:varyColors'), attrib={'val':'0'})
        lc.append(ls)
        # マーカー
        mk = ls.find(qn('c:marker'))
        if mk is None: mk = etree.SubElement(ls, qn('c:marker'))
        sym = mk.find(qn('c:symbol'))
        if sym is None: sym = etree.SubElement(mk, qn('c:symbol'))
        sym.set('val','circle')
        sz = mk.find(qn('c:size'))
        if sz is None: sz = etree.SubElement(mk, qn('c:size'))
        sz.set('val', str(line_cfg.get("marker_size",8)))
        # 第2軸
        sec_v="2094734553"; sec_c="2094734554"
        etree.SubElement(lc, qn('c:axId'), attrib={'val':sec_c})
        etree.SubElement(lc, qn('c:axId'), attrib={'val':sec_v})
        # 二次カテゴリ軸（非表示）
        sc = etree.SubElement(pa, qn('c:catAx'))
        etree.SubElement(sc, qn('c:axId'), attrib={'val':sec_c})
        s2 = etree.SubElement(sc, qn('c:scaling'))
        etree.SubElement(s2, qn('c:orientation'), attrib={'val':'minMax'})
        etree.SubElement(sc, qn('c:delete'), attrib={'val':'1'})
        etree.SubElement(sc, qn('c:axPos'), attrib={'val':'b'})
        etree.SubElement(sc, qn('c:crossAx'), attrib={'val':sec_v})
        # 二次値軸（右側に可視化、明示的な min/max を設定）
        sv = etree.SubElement(pa, qn('c:valAx'))
        etree.SubElement(sv, qn('c:axId'), attrib={'val':sec_v})
        s3 = etree.SubElement(sv, qn('c:scaling'))
        etree.SubElement(s3, qn('c:orientation'), attrib={'val':'minMax'})
        etree.SubElement(s3, qn('c:max'), attrib={'val': f'{secondary_axis_max:g}'})
        etree.SubElement(s3, qn('c:min'), attrib={'val': '0'})
        etree.SubElement(sv, qn('c:delete'), attrib={'val':'0'})
        etree.SubElement(sv, qn('c:axPos'), attrib={'val':'r'})
        line_axis_fmt = line_cfg.get("num_format", "0")
        etree.SubElement(sv, qn('c:numFmt'), attrib={'formatCode':line_axis_fmt,'sourceLinked':'0'})
        etree.SubElement(sv, qn('c:majorTickMark'), attrib={'val':'out'})
        etree.SubElement(sv, qn('c:minorTickMark'), attrib={'val':'none'})
        etree.SubElement(sv, qn('c:tickLblPos'), attrib={'val':'nextTo'})
        etree.SubElement(sv, qn('c:crossAx'), attrib={'val':sec_c})
        etree.SubElement(sv, qn('c:crosses'), attrib={'val':'max'})
        # 折れ線色
        lclr = line_cfg.get("color","#333333").replace("#","")
        lsp = ls.find(qn('c:spPr'))
        if lsp is None: lsp = etree.SubElement(ls, qn('c:spPr'))
        ln = lsp.find(qn('a:ln'))
        if ln is None: ln = etree.SubElement(lsp, qn('a:ln'))
        ln.set('w','19050')
        lsf = ln.find(qn('a:solidFill'))
        if lsf is None: lsf = etree.SubElement(ln, qn('a:solidFill'))
        for c in list(lsf): lsf.remove(c)
        etree.SubElement(lsf, qn('a:srgbClr'), attrib={'val':lclr})
        msp = mk.find(qn('c:spPr'))
        if msp is None: msp = etree.SubElement(mk, qn('c:spPr'))
        msf = etree.SubElement(msp, qn('a:solidFill'))
        etree.SubElement(msf, qn('a:srgbClr'), attrib={'val':lclr})
        _add_dlbls(ls, 't', num_fmt, '333333', 1000)

    # 棒色＋ラベル（スキル固定パレット使用、JSON の color は無視）
    bar_series = bc.findall(qn('c:ser'))
    n_bar_series = len(bar_series)
    for si, bs in enumerate(bar_series):
        clr = _palette_color(si, n_bar_series).lstrip("#")
        sp = bs.find(qn('c:spPr'))
        if sp is None: sp = etree.SubElement(bs, qn('c:spPr'))
        sf = sp.find(qn('a:solidFill'))
        if sf is None: sf = etree.SubElement(sp, qn('a:solidFill'))
        for c in list(sf): sf.remove(c)
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val':clr})
        _add_dlbls(bs, 'ctr', num_fmt, 'FFFFFF', 1000)

    # 棒間隔
    gw = bc.find(qn('c:gapWidth'))
    if gw is None: gw = etree.SubElement(bc, qn('c:gapWidth'))
    gw.set('val','80')
    ov = bc.find(qn('c:overlap'))
    if ov is None: ov = etree.SubElement(bc, qn('c:overlap'))
    ov.set('val','100')

    # グリッド削除
    for vx in pa.findall(qn('c:valAx')):
        for mg in vx.findall(qn('c:majorGridlines')): vx.remove(mg)
        for mg in vx.findall(qn('c:minorGridlines')): vx.remove(mg)
    for cx in pa.findall(qn('c:catAx')):
        for mg in cx.findall(qn('c:majorGridlines')): cx.remove(mg)

    # 軸フォント
    for ax in pa.findall(qn('c:catAx')) + pa.findall(qn('c:valAx')):
        de = ax.find(qn('c:delete'))
        if de is not None and de.get('val')=='1': continue
        tp = ax.find(qn('c:txPr'))
        if tp is None: tp = etree.SubElement(ax, qn('c:txPr'))
        bp = tp.find(qn('a:bodyPr'))
        if bp is None: bp = etree.SubElement(tp, qn('a:bodyPr')); tp.insert(0, bp)
        if tp.find(qn('a:lstStyle')) is None: etree.SubElement(tp, qn('a:lstStyle'))
        p = tp.find(qn('a:p'))
        if p is None: p = etree.SubElement(tp, qn('a:p'))
        pp = p.find(qn('a:pPr'))
        if pp is None: pp = etree.SubElement(p, qn('a:pPr'))
        dr = pp.find(qn('a:defRPr'))
        ax_sz = _body_sz('1100')  # roleup: 10pt 統一 / stella: 11pt 維持
        if dr is None: dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': ax_sz})
        else: dr.set('sz', ax_sz)
        if dr.find(qn('a:latin')) is None: etree.SubElement(dr, qn('a:latin'), attrib={'typeface':FONT_JP})
        if dr.find(qn('a:ea')) is None: etree.SubElement(dr, qn('a:ea'), attrib={'typeface':FONT_JP})

    chart.has_legend = False; chart.has_title = False
    psp = pa.find(qn('c:spPr'))
    if psp is None: psp = etree.SubElement(pa, qn('c:spPr'))
    etree.SubElement(psp, qn('a:noFill'))

    # OOXML CT_PlotArea 順序に正規化（chart 要素 → axes）
    _reorder_plot_area(pa)

    print(f"  ✓ 複合チャート: {len(data)}年, {n_ser}棒系列" + (f" + 折れ線" if has_line else ""))
    print(f"  ✓ Y 軸 自動レンジ: 左軸 max={primary_axis_max:g}" +
          (f" / 右軸 max={secondary_axis_max:g}" if has_line else ""))
    return cf

def _add_dlbls(ser, pos, fmt, fc, sz):
    dl = ser.find(qn('c:dLbls'))
    if dl is None: dl = etree.SubElement(ser, qn('c:dLbls'))
    for c in list(dl): dl.remove(c)
    etree.SubElement(dl, qn('c:numFmt'), attrib={'formatCode':fmt,'sourceLinked':'0'})
    for k in ['showLegendKey','showCatName','showSerName','showPercent','showBubbleSize']:
        etree.SubElement(dl, qn(f'c:{k}'), attrib={'val':'0'})
    etree.SubElement(dl, qn('c:showVal'), attrib={'val':'1'})
    etree.SubElement(dl, qn('c:dLblPos'), attrib={'val':pos})
    tp = etree.SubElement(dl, qn('c:txPr'))
    etree.SubElement(tp, qn('a:bodyPr')); etree.SubElement(tp, qn('a:lstStyle'))
    p = etree.SubElement(tp, qn('a:p'))
    pp = etree.SubElement(p, qn('a:pPr'))
    dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz':str(sz),'b':'1'})
    etree.SubElement(dr, qn('a:latin'), attrib={'typeface':FONT_JP})
    etree.SubElement(dr, qn('a:ea'), attrib={'typeface':FONT_JP})
    sf = etree.SubElement(dr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val':fc})

# ── 注釈・凡例 ──
def add_section_subtitle(slide, text, left, top, width, theme):
    """サブタイトル(セクション見出し)を動的 textbox で配置。

    customer-profile/add_section_title 相当の機能だが、
    market-environment は section_title フィールドが任意のため、
    呼び出し側で `if text` でガードする想定。

    色・align は theme.subtitle / layout_rule(subtitle_align) 由来。
    roleup: 12pt 左寄せ #897141 / stella: 14pt 中央寄せ #333333。
    """
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    align_str = theme.layout_rule("subtitle_align", "center")
    p.alignment = PP_ALIGN.LEFT if align_str == "left" else PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = theme.pt("font_size_section_pt")
    run.font.bold = True
    run.font.color.rgb = theme.color("subtitle")
    run.font.name = theme.font_ea
    return txBox


def add_unit_label(slide, text, left, top, width=None, height=None):
    """Unit label textbox.

    width default: stella regression のため 3.0 inch (旧挙動)。
    height default: stella regression のため 0.25 inch (旧挙動)。
    roleup ではテキスト幅・高さにフィットする値を呼出側で指定。
    """
    if width is None:
        width = Inches(3.0)
    if height is None:
        height = Inches(0.25)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = _body_pt(11)
    r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP

def add_custom_legend(slide, cfg, left, top, w):
    bars = cfg.get("stacked_bars",[]); lc = cfg.get("line",None)
    lh = Inches(0.25); ix = left
    if lc:
        c = _hex2rgb(lc.get("color","#333333"))
        lw = Inches(0.35); ly = int(top + lh/2)
        cn = slide.shapes.add_connector(1, int(ix), ly, int(ix+lw), ly)
        cn.line.color.rgb = c; cn.line.width = Pt(1.5)
        cs = Inches(0.12)
        ci = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(ix+(lw-cs)/2), int(top+(lh-cs)/2), cs, cs)
        ci.fill.solid(); ci.fill.fore_color.rgb = c; ci.line.fill.background()
        tx = int(ix+lw+Inches(0.06))
        tb = slide.shapes.add_textbox(tx, int(top), Inches(1.0), lh)
        tb.text_frame.word_wrap = False
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = lc["series_name"]; r.font.size = _body_pt(11); r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP
        ix = tx + Inches(1.0) + Inches(0.15)
    n_bars_total = len(bars)
    for sb_idx, sb in enumerate(bars):
        c = _hex2rgb(_palette_color(sb_idx, n_bars_total)); ss = Inches(0.15)
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(ix), int(top+(lh-ss)/2), ss, ss)
        sq.fill.solid(); sq.fill.fore_color.rgb = c; sq.line.fill.background()
        tx = int(ix+ss+Inches(0.06))
        tb = slide.shapes.add_textbox(tx, int(top), Inches(1.5), lh)
        tb.text_frame.word_wrap = False
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = sb["series_name"]; r.font.size = _body_pt(11); r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP
        ix = tx + Inches(1.5) + Inches(0.10)
    print("  ✓ 凡例")

def add_cagr_annotations(slide, cfg, left, top, w):
    """右側セグメント別CAGR注釈テキストボックス"""
    anns = cfg.get("cagr_annotations",[])
    if not anns: return
    y = top
    for ann in anns:
        # 期間ラベル（太字、下線付き）
        tb = slide.shapes.add_textbox(left, int(y), w, Inches(0.25))
        tb.text_frame.word_wrap = True
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = ann.get("label",""); r.font.size = _body_pt(11); r.font.bold = True
        r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, int(y+Inches(0.25)), w, Inches(0.015))
        ln.fill.solid(); ln.fill.fore_color.rgb = COLOR_TEXT; ln.line.fill.background()
        y += Inches(0.32)
        for it in ann.get("items",[]):
            nb = slide.shapes.add_textbox(left, int(y), Inches(1.1), Inches(0.28))
            r = nb.text_frame.paragraphs[0].add_run()
            r.text = f"{it['name']}："; r.font.size = _body_pt(11)
            r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP
            vb = slide.shapes.add_textbox(int(left+Inches(1.1)), int(y), int(w-Inches(1.1)), Inches(0.28))
            vb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            r = vb.text_frame.paragraphs[0].add_run()
            r.text = it['value']; r.font.size = _body_pt(15); r.font.bold = True
            r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP
            y += Inches(0.30)
        y += Inches(0.15)
    print(f"  ✓ 右側CAGR注釈: {len(anns)}セクション")

def add_growth_annotations(slide, cfg, cl, ct, cw, ch):
    """CAGR矢印＋楕円注釈を追加（customer-profile-pptx準拠）"""
    gas = cfg.get("growth_annotations",[]); data = cfg.get("data",[])
    if not gas or not data: return
    nc = len(data)

    # プロットエリア推定
    plm = cw*0.06; prm = cw*0.04; ptm = ch*0.05; pbm = ch*0.14
    pl = cl + plm; pw = cw - plm - prm
    pt_ = ct + ptm; pb = ct + ch - pbm; ph = pb - pt_
    caw = pw / nc

    # Y軸スケール推定
    max_total = max(sum(d.get("bars",[])) for d in data)
    axis_max = max_total * 1.20

    for ga in gas:
        si = ga["start_index"]; ei = ga["end_index"]; val = ga["value"]

        # 開始・終了バーの中心X
        x1 = int(pl + (si+0.5)*caw)
        x2 = int(pl + (ei+0.5)*caw)

        # 開始・終了バーの上端Y
        s_total = sum(data[si].get("bars",[]))
        e_total = sum(data[ei].get("bars",[]))
        y1_top = pb - (s_total / axis_max) * ph
        y2_top = pb - (e_total / axis_max) * ph

        # 矢印は棒の上に配置（比例gap: プロット高さの15%、最低0.40"）
        gap = max(ph * 0.15, Inches(0.40))
        ay1 = int(y1_top - gap)
        ay2 = int(y2_top - gap)

        # チャートエリア上端でクランプ（はみ出し防止）
        min_y = int(ct - Inches(0.30))
        ay1 = max(ay1, min_y)
        ay2 = max(ay2, min_y)

        # 矢印（コネクタ＋矢じり）
        conn = slide.shapes.add_connector(1, x1, ay1, x2, ay2)
        conn.line.color.rgb = COLOR_TEXT
        conn.line.width = Pt(1.5)
        cxnSp = conn._element
        spPr = cxnSp.find(qn('p:spPr'))
        if spPr is None: spPr = cxnSp.find(qn('a:spPr'))
        if spPr is not None:
            ln = spPr.find(qn('a:ln'))
            if ln is None: ln = etree.SubElement(spPr, qn('a:ln'))
            etree.SubElement(ln, qn('a:tailEnd'), attrib={
                'type': 'triangle', 'w': 'med', 'len': 'med'
            })

        # 楕円＋CAGRテキスト（矢印の中間点に配置）
        ow = Inches(1.10); oh = Inches(0.35)
        mid_x = (x1 + x2) / 2
        mid_y = (ay1 + ay2) / 2
        ox = int(mid_x - ow/2); oy = int(mid_y - oh/2)
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, ox, oy, ow, oh)
        ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
        ov.line.color.rgb = COLOR_TEXT; ov.line.width = Pt(1.0)
        ov.text_frame.word_wrap = False
        ov.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = ov.text_frame.paragraphs[0].add_run()
        r.text = val; r.font.size = _body_pt(14); r.font.bold = True
        r.font.color.rgb = COLOR_TEXT; r.font.name = FONT_JP
    print(f"  ✓ CAGR注釈: {len(gas)}個")

def add_period_separator(slide, cfg, cl, ct, cw, ch):
    ps = cfg.get("period_separator"); data = cfg.get("data",[])
    if not ps or not data: return
    nc = len(data); si = ps.get("after_index",-1); lb = ps.get("label","")
    if si < 0 or si >= nc: return
    plm = cw*0.06; prm = cw*0.04; ptm = ch*0.05; pbm = ch*0.14
    pl = cl + plm; pw = cw - plm - prm
    pt_ = ct + ptm; pb = ct + ch - pbm; ph = pb - pt_
    caw = pw / nc
    sx = int(pl + (si+1) * caw)
    # 最大バーの上端に合わせる
    max_total = max(sum(d.get("bars",[])) for d in data)
    axis_max = max_total * 1.20
    bar_top = pb - (max_total / axis_max) * ph
    lt = int(bar_top)
    lb_ = int(pb)
    cn = slide.shapes.add_connector(1, sx, lt, sx, lb_)
    cn.line.color.rgb = COLOR_SOURCE; cn.line.width = Pt(1.0)
    sp = cn._element.find(qn('p:spPr'))
    if sp is None: sp = cn._element.find(qn('a:spPr'))
    if sp is not None:
        ln = sp.find(qn('a:ln'))
        if ln is None: ln = etree.SubElement(sp, qn('a:ln'))
        pd = ln.find(qn('a:prstDash'))
        if pd is None: pd = etree.SubElement(ln, qn('a:prstDash'))
        pd.set('val','dash')
    if lb:
        lw = Inches(0.50); lh_ = Inches(0.15)
        lx = int(sx - lw/2); ly = int(lt - Inches(0.18))
        rt = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, lw, lh_)
        rt.fill.solid(); rt.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
        rt.line.color.rgb = COLOR_SOURCE; rt.line.width = Pt(0.5)
        rt.text_frame.word_wrap = False
        rt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        rt.text_frame.margin_top = Pt(1)
        rt.text_frame.margin_bottom = Pt(1)
        rt.text_frame.margin_left = Pt(2)
        rt.text_frame.margin_right = Pt(2)
        r = rt.text_frame.paragraphs[0].add_run()
        r.text = lb; r.font.size = Pt(6); r.font.color.rgb = COLOR_SOURCE; r.font.name = FONT_JP
    print(f"  ✓ 期間区切り線")

# ── メイン ──
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "market-environment")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f: data = json.load(f)

    _mm = data.get("main_message", "")
    if len(_mm) > 65:
        raise ValueError(
            f"main_message は 65 字以内（受領: {len(_mm)}）: {_mm[:80]}..."
        )

    print("=== 市場環境分析スライド生成（ネイティブPPTX）===")
    prs = Presentation(template_path); slide = prs.slides[0]

    # Top placeholder (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    # Subtitle placeholder (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme) or "市場環境分析"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    src = data.get("source","")
    if src:
        # roleup 公式テンプレでは shape 名が "Source 3"。stella は "Source"。
        # 両方を試して見つかった方を更新する(silent fallback)。
        src_shape = _silent_find_shape(slide, SHAPE_SOURCE) or _silent_find_shape(slide, "Source 3")
        if src_shape is not None:
            set_textbox_text(src_shape, f"出典：{src}")
        else:
            print(f"  ⚠ Source shape not found (tried '{SHAPE_SOURCE}', 'Source 3')", file=sys.stderr)
    # ガイド矩形(チャート位置決め用)を除去:
    #   stella: "Content Area" / roleup: "正方形/長方形 1"(茶色 accent2)
    # 出力ではマスター背景(白)が見えるよう両方 silent に削除を試みる。
    _silent_remove_shape(slide, SHAPE_CONTENT_AREA)
    _silent_remove_shape(slide, "正方形/長方形 1")

    cfg = data.get("chart", {})
    ul = cfg.get("unit_label","")

    if theme.id == "stellar_aiz":
        # stella: 既存挙動維持 (生成順序・座標が旧と完全一致 → regression-zero)
        if ul:
            add_unit_label(slide, ul, CHART_X, PANEL_Y)
        add_custom_legend(slide, cfg, CHART_X, PANEL_Y + Inches(0.25), CHART_W)
    else:
        # roleup (ユーザー視覚調整 2026-05-04):
        # サブタイトル: panel_y+0.05 に section_title (12pt 左寄せ #897141)
        # 凡例: panel_y+0.51, chart_x+0.21 (0.62 in)
        # 単位ラベル: 凡例右側 chart 右上 (6.52, panel_y+0.51) に幅 1.39 in
        section_title = data.get("section_title", "")
        if section_title:
            add_section_subtitle(slide, section_title, CHART_X, PANEL_Y + Inches(0.05),
                                 CHART_W + CAGR_W, theme)
        add_custom_legend(slide, cfg, CHART_X + Inches(0.21), PANEL_Y + Inches(0.51), CHART_W)
        if ul:
            add_unit_label(slide, ul, Inches(6.52), PANEL_Y + Inches(0.51),
                           Inches(1.39), height=Inches(0.27))
    build_stacked_combo_chart(slide, cfg, CHART_X, CHART_Y, CHART_W, CHART_H)
    add_period_separator(slide, cfg, CHART_X, CHART_Y, CHART_W, CHART_H)
    add_growth_annotations(slide, cfg, CHART_X, CHART_Y, CHART_W, CHART_H)
    add_cagr_annotations(slide, cfg, CAGR_X, PANEL_Y + Inches(0.60), CAGR_W)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output); print(f"\n  ✅ 出力完了: {args.output}")
    _finalize_pptx(args.output)
if __name__ == "__main__": main()

