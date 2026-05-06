"""
fill_aspiration.py — Aspiration & Trajectory 詳細スライド3枚（コンサル品質）

Visual: 時間軸ロードマップ（過去-現在-将来のイベント配置）
"""

import argparse
import asyncio
import copy
import json
import os
import sys
import tempfile
from html import escape

# validate_fill_input bootstrap (skills/_common/lib)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from validate_fill_input import validate_fill_input  # noqa: E402

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    import shutil, subprocess, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_DIAGRAM_AREA = "Rectangle 4"
SHAPE_IMPLICATIONS = "TextBox 9"

VIEWPORT_WIDTH = 1400
VIEWPORT_HEIGHT = 760
EVIDENCE_VIEWPORT_WIDTH = 2200
EVIDENCE_VIEWPORT_HEIGHT = 760
DEVICE_SCALE = 2

CONFIDENCE_COLOR = {"high": "#52C41A", "medium": "#FAAD14", "low": "#FF4D4F"}
PHASE_COLOR = {"past": "#999999", "current": "#1565C0", "future": "#FFA940"}
PHASE_LABEL = {"past": "過去", "current": "現在", "future": "将来（推定）"}

DIMENSION_TITLE = "Aspiration & Trajectory：経営意図と時間軸"


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_placeholder_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_implications(slide, items):
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return
    tf = shape.text_frame
    for i, item in enumerate(items[:3]):
        label = item.get("label", "").strip()
        detail = item.get("detail", "").strip()
        para_idx = i + 1
        if para_idx >= len(tf.paragraphs):
            continue
        para = tf.paragraphs[para_idx]
        template_rPr = None
        if para.runs:
            template_rPr = para.runs[0]._r.find(qn("a:rPr"))
        for r in list(para._p.findall(qn("a:r"))):
            para._p.remove(r)

        def _make_run(text_str, bold=False):
            r_elem = etree.SubElement(para._p, qn("a:r"))
            if template_rPr is not None:
                new_rPr = copy.deepcopy(template_rPr)
                if bold:
                    new_rPr.set("b", "1")
                else:
                    new_rPr.attrib.pop("b", None)
                r_elem.insert(0, new_rPr)
            else:
                attrib = {"lang": "ja-JP"}
                if bold:
                    attrib["b"] = "1"
                etree.SubElement(r_elem, qn("a:rPr"), attrib=attrib)
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text_str

        if label:
            _make_run(f"{label}：", bold=True)
            _make_run(detail, bold=False)
        else:
            _make_run(detail, bold=False)


def clear_textbox(shape):
    if shape is None or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for r in list(para._p.findall(qn("a:r"))):
            para._p.remove(r)


def dup_slide(prs, tmpl):
    ns = prs.slides.add_slide(tmpl.slide_layout)
    nsp = ns.shapes._spTree
    for c in list(nsp):
        if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr"):
            nsp.remove(c)
    for c in tmpl.shapes._spTree:
        if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr"):
            nsp.append(copy.deepcopy(c))
    return ns


def _esc(text):
    return escape(str(text))


CSS_BASE = """
* { box-sizing: border-box; margin: 0; padding: 0; font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', 'Yu Gothic UI', sans-serif; }
body { width: 1400px; height: 760px; padding: 0; background: #FFF; }
"""


# ──────────────────────────────────────────────────
# Page 1 (Main): タイムライン・ロードマップ SVG
# ──────────────────────────────────────────────────

def render_main_chart(visual):
    """時間軸タイムライン: 横線上にマイルストーンを配置、phase で色分け"""
    milestones = visual.get("milestones", [])
    n = max(len(milestones), 1)

    svg_w, svg_h = 1300, 720
    margin_x = 80
    line_y = svg_h * 0.55
    inner_w = svg_w - 2 * margin_x
    step_w = inner_w / max(n - 1, 1)

    parts = [f'<svg width="{svg_w}" height="{svg_h}" xmlns="http://www.w3.org/2000/svg">']
    parts.append(f'<text x="{svg_w/2}" y="40" text-anchor="middle" font-size="22" font-weight="700" fill="#0A1B3D">経営意図と時間軸ロードマップ</text>')

    # 凡例（過去/現在/将来）
    legend_y = 70
    legend_x = svg_w / 2 - 280
    for i, phase in enumerate(["past", "current", "future"]):
        lx = legend_x + i * 200
        parts.append(f'<circle cx="{lx}" cy="{legend_y}" r="10" fill="{PHASE_COLOR[phase]}"/>')
        parts.append(f'<text x="{lx + 18}" y="{legend_y + 6}" font-size="16" fill="#444" font-weight="600">{PHASE_LABEL[phase]}</text>')

    # 時間軸ベースライン
    parts.append(f'<line x1="{margin_x}" y1="{line_y}" x2="{svg_w - margin_x}" y2="{line_y}" stroke="#444" stroke-width="3"/>')
    parts.append(f'<polygon points="{svg_w - margin_x},{line_y} {svg_w - margin_x - 14},{line_y - 9} {svg_w - margin_x - 14},{line_y + 9}" fill="#444"/>')
    parts.append(f'<text x="{svg_w - margin_x}" y="{line_y + 35}" text-anchor="end" font-size="14" font-style="italic" fill="#666">時間 →</text>')

    # 各マイルストーン（上下交互配置で重ならないように）
    for i, ms in enumerate(milestones):
        cx = margin_x + step_w * i if n > 1 else margin_x + inner_w / 2
        year = _esc(ms.get("year", ""))
        label = _esc(ms.get("label", ""))
        phase = ms.get("phase", "current")
        note = _esc(ms.get("note", ""))
        color = PHASE_COLOR.get(phase, "#1565C0")

        # ベースライン上のマーカー
        parts.append(f'<circle cx="{cx}" cy="{line_y}" r="14" fill="{color}" stroke="#FFF" stroke-width="3"/>')
        parts.append(f'<circle cx="{cx}" cy="{line_y}" r="7" fill="#FFF"/>')

        # 年（ベースライン下）
        parts.append(f'<text x="{cx}" y="{line_y + 50}" text-anchor="middle" font-size="16" font-weight="700" fill="{color}">{year}</text>')

        # ラベルカード（上下交互）
        is_above = (i % 2 == 0)
        if is_above:
            card_y = line_y - 130
            connector_y2 = line_y - 14
            connector_y1 = card_y + 60
        else:
            card_y = line_y + 80
            connector_y1 = line_y + 14
            connector_y2 = card_y

        # 接続線
        parts.append(f'<line x1="{cx}" y1="{connector_y1}" x2="{cx}" y2="{connector_y2}" stroke="{color}" stroke-width="1.5" stroke-dasharray="3 3" opacity="0.5"/>')

        # カード
        card_w = 200
        card_h = 80
        card_x = cx - card_w / 2
        parts.append(f'<rect x="{card_x}" y="{card_y}" width="{card_w}" height="{card_h}" fill="#FFF" stroke="{color}" stroke-width="2.5" rx="6"/>')
        parts.append(f'<rect x="{card_x}" y="{card_y}" width="6" height="{card_h}" fill="{color}" rx="3"/>')
        parts.append(f'<text x="{card_x + 16}" y="{card_y + 30}" font-size="16" font-weight="700" fill="#0A1B3D">{label}</text>')
        if note:
            parts.append(f'<text x="{card_x + 16}" y="{card_y + 55}" font-size="13" fill="#666" font-style="italic">{note}</text>')

    parts.append("</svg>")
    return f'<!DOCTYPE html><html><head><meta charset="utf-8"><style>{CSS_BASE}.center {{ display:flex; align-items:center; justify-content:center; height:760px; }}</style></head><body><div class="center">{"".join(parts)}</div></body></html>'


# ──────────────────────────────────────────────────
# Page 2 (Detail): フェーズ別マトリクス
# ──────────────────────────────────────────────────

def render_detail_chart(visual):
    phases = visual.get("phases", [])
    cols_html = []
    for ph in phases:
        label = _esc(ph.get("label", ""))
        color = ph.get("color", "#1565C0")
        actions = ph.get("actions", [])
        actions_html = "".join(f'<li>{_esc(a)}</li>' for a in actions)
        cols_html.append(f"""
        <div class="col">
          <div class="col-header" style="background:{color};">{label}</div>
          <ul class="actions">{actions_html}</ul>
        </div>
        """)

    css = CSS_BASE + """
    .matrix { padding: 30px 50px; }
    .matrix-title { font-size: 20px; font-weight: 700; color: #0A1B3D; margin-bottom: 16px; padding-bottom: 8px; border-bottom: 3px solid #1565C0; }
    .cols { display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 16px; }
    .col { background: #FFF; border-radius: 6px; box-shadow: 0 1px 3px rgba(0,0,0,0.06); overflow: hidden; }
    .col-header { color: #FFF; font-weight: 700; font-size: 17px; padding: 14px 16px; }
    .actions { list-style: disc; padding: 16px 16px 16px 36px; }
    .actions li { font-size: 15px; color: #333; line-height: 1.7; margin-bottom: 8px; }
    """
    cols = "\n".join(cols_html)
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head>
<body><div class="matrix">
  <div class="matrix-title">フェーズ別の戦略アクション</div>
  <div class="cols">{cols}</div>
</div></body></html>"""


# ──────────────────────────────────────────────────
# Page 3 (Evidence): findings 表（フルワイド）
# ──────────────────────────────────────────────────

def render_evidence_chart(findings):
    rows_html = []
    for f in findings:
        confidence = f.get("confidence", "medium")
        conf_color = CONFIDENCE_COLOR.get(confidence, "#999")
        rows_html.append(f"""
        <tr>
          <td class="td-id">{_esc(f.get("id", ""))}</td>
          <td class="td-agent">{_esc(f.get("agent", ""))}</td>
          <td class="td-source">{_esc(f.get("source", ""))}</td>
          <td><span class="src-type">{_esc(f.get("source_type", ""))}</span></td>
          <td><span class="conf" style="background:{conf_color};">{_esc(confidence)}</span></td>
          <td class="td-excerpt">{_esc(f.get("excerpt", ""))}</td>
        </tr>
        """)

    css = """
    * { box-sizing: border-box; margin: 0; padding: 0; font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', 'Yu Gothic UI', sans-serif; }
    body { width: 2200px; height: 760px; padding: 0; background: #FFF; }
    .ev-wrap { padding: 24px 36px; }
    table { width: 100%; border-collapse: collapse; font-size: 17px; }
    th { background: #0A1B3D; color: #FFF; padding: 14px 12px; text-align: left; font-weight: 700; font-size: 16px; }
    td { padding: 13px 12px; border-bottom: 1px solid #E0E0E0; vertical-align: top; line-height: 1.6; color: #333; }
    tr:nth-child(even) td { background: #FAFAFA; }
    .td-id { font-weight: 700; color: #1565C0; width: 70px; font-size: 18px; }
    .td-agent { color: #666; font-size: 14px; width: 180px; }
    .td-source { font-size: 14px; color: #555; width: 280px; }
    .src-type { display: inline-block; background: #E8F4FF; color: #1565C0; font-size: 13px; padding: 4px 10px; border-radius: 4px; font-weight: 600; }
    .conf { color: #FFF; font-size: 13px; font-weight: 700; padding: 4px 11px; border-radius: 10px; }
    .td-excerpt { font-size: 15px; color: #444; line-height: 1.55; }
    """
    rows = "\n".join(rows_html)
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head>
<body><div class="ev-wrap">
  <table>
    <thead><tr><th>F#</th><th>Agent</th><th>Source</th><th>Type</th><th>Confidence</th><th>抜粋</th></tr></thead>
    <tbody>{rows}</tbody>
  </table>
</div></body></html>"""


async def take_screenshot(html_content, output_path, suffix="", vw=None, vh=None):
    from playwright.async_api import async_playwright
    vw = vw or VIEWPORT_WIDTH
    vh = vh or VIEWPORT_HEIGHT
    html_path = os.path.join(tempfile.gettempdir(), f"aspiration{suffix}.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={"width": vw, "height": vh},
            device_scale_factor=DEVICE_SCALE,
        )
        await page.goto(f"file://{html_path}")
        await page.wait_for_timeout(500)
        await page.screenshot(path=output_path, full_page=False)
        await browser.close()
    print(f"  📸 {output_path} ({os.path.getsize(output_path)} bytes, vp {vw}x{vh})")
    os.unlink(html_path)


def fill_slide(slide, main_msg, chart_subtitle, implications, screenshot_path, skip_implications=False):
    set_placeholder_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg)
    set_placeholder_text(find_shape(slide, SHAPE_CHART_TITLE), chart_subtitle)
    if skip_implications:
        clear_textbox(find_shape(slide, SHAPE_IMPLICATIONS))
    else:
        fill_implications(slide, implications)
    rect_shape = find_shape(slide, SHAPE_DIAGRAM_AREA)
    if rect_shape:
        if rect_shape.has_text_frame:
            for para in rect_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
        padding = Inches(0.05)
        if skip_implications:
            tb9_shape = find_shape(slide, SHAPE_IMPLICATIONS)
            if tb9_shape is not None:
                left = rect_shape.left
                top = min(rect_shape.top, tb9_shape.top)
                right = tb9_shape.left + tb9_shape.width
                bottom = max(rect_shape.top + rect_shape.height, tb9_shape.top + tb9_shape.height)
                slide.shapes.add_picture(screenshot_path, left + padding, top + padding, (right - left) - 2 * padding, (bottom - top) - 2 * padding)
                return
        slide.shapes.add_picture(screenshot_path, rect_shape.left + padding, rect_shape.top + padding, rect_shape.width - 2 * padding, rect_shape.height - 2 * padding)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main", "detail", "evidence"],
        allowed_top=["main", "detail", "evidence"],
        skill_name="smallcap-aspiration-pptx",
    )

    main_data = data.get("main", {})
    detail_data = data.get("detail", {})
    evidence_data = data.get("evidence", {})

    print("📐 Generating 3 chart HTMLs (timeline / phases / table)...")
    html_main = render_main_chart(main_data.get("visual_data", {}))
    html_detail = render_detail_chart(detail_data.get("visual_data", {}))
    html_evidence = render_evidence_chart(evidence_data.get("findings", []))

    paths = [
        os.path.join(tempfile.gettempdir(), "asp_main.png"),
        os.path.join(tempfile.gettempdir(), "asp_detail.png"),
        os.path.join(tempfile.gettempdir(), "asp_evidence.png"),
    ]

    print("📸 Taking 3 screenshots...")
    asyncio.run(take_screenshot(html_main, paths[0], "_main"))
    asyncio.run(take_screenshot(html_detail, paths[1], "_detail"))
    asyncio.run(take_screenshot(html_evidence, paths[2], "_evidence", vw=EVIDENCE_VIEWPORT_WIDTH, vh=EVIDENCE_VIEWPORT_HEIGHT))

    print("📝 Filling PPTX (3 slides)...")
    prs = Presentation(args.template)
    tmpl = prs.slides[0]
    extras = [dup_slide(prs, tmpl) for _ in range(2)]

    fill_slide(tmpl, main_data.get("main_message", ""), main_data.get("chart_title", f"{DIMENSION_TITLE}（1/3）"), main_data.get("implications", []), paths[0])
    print(f"  [1/3] Main")
    fill_slide(extras[0], detail_data.get("main_message", ""), detail_data.get("chart_title", f"{DIMENSION_TITLE}（2/3）"), detail_data.get("implications", []), paths[1])
    print(f"  [2/3] Detail")
    fill_slide(extras[1], evidence_data.get("main_message", ""), evidence_data.get("chart_title", f"{DIMENSION_TITLE}（3/3）"), evidence_data.get("implications", []), paths[2], skip_implications=True)
    print(f"  [3/3] Evidence (chart only)")

    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output} (3 slides)")

    for p in paths:
        if os.path.exists(p):
            os.unlink(p)


if __name__ == "__main__":
    main()
