"""
fill_value_chain_matrix.py — バリューチェーン・ポジショニング・マトリクスを
PPTXネイティブオブジェクトで生成。

レイアウト:
  - 上部: メインメッセージ (28pt Bold) + チャートタイトル (14pt)
  - 中央: シェブロン行 + 4行のマトリクス
    * 行ラベル列（左、1.90"）
    * N段階のシェブロン（5〜7段階、幅自動調整）
    * 行1,2,4: 記号セル (◎ ○ △ ー 空白)
    * 行3: プレーヤーバー（複数列またぎ）
  - 下部: 出典 (10pt グレー)

全オブジェクトはPPT上で編集可能（画像ではない）。

Usage:
  python fill_value_chain_matrix.py \\
    --data /home/claude/value_chain_matrix_data.json \\
    [--template <SKILL_DIR>/assets/<brand>/value-chain-matrix-template.pptx] \\
    --output /mnt/user-data/outputs/ValueChainMatrix_output.pptx \\
    [--brand stellar_aiz | roleup]
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "value-chain-matrix-pptx"
_THEME = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング（テンプレート上） ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_PH = "Source 3"  # roleup placeholder

# ── スライド全体 (stella defaults; _apply_theme で roleup 用にスケール) ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── レイアウト定数 ──
CONTENT_LEFT = Inches(0.41)
CONTENT_WIDTH = Inches(12.52)

ROW_LABEL_WIDTH = Inches(1.90)
STAGES_LEFT = CONTENT_LEFT + ROW_LABEL_WIDTH          # Inches(2.31)
STAGES_WIDTH = CONTENT_WIDTH - ROW_LABEL_WIDTH        # Inches(10.62)

CHEVRON_TOP = Inches(1.50)
CHEVRON_HEIGHT = Inches(0.75)

GAP = Inches(0.06)
DATA_ROW_HEIGHT = Inches(1.00)

DATA_ROWS_TOP = CHEVRON_TOP + CHEVRON_HEIGHT + GAP   # 1.50 + 0.75 + 0.06 = 2.31

SOURCE_X = CONTENT_LEFT
SOURCE_Y = Inches(6.85)
SOURCE_W = CONTENT_WIDTH
SOURCE_H = Inches(0.28)

# ── 色 (stella defaults; _apply_theme で roleup 用に上書き) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_NAVY = RGBColor(0x1E, 0x3A, 0x5F)           # シェブロン枠線・テキスト・記号
COLOR_CHEVRON_FILL = RGBColor(0xCE, 0xE0, 0xED)   # シェブロン塗り（薄いブルー）
COLOR_BAR_FILL = RGBColor(0xB4, 0xC4, 0xD4)       # プレーヤーバー塗り
COLOR_BAR_BORDER = RGBColor(0x5A, 0x70, 0x90)     # プレーヤーバー枠線
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_SEPARATOR = RGBColor(0xBB, 0xBB, 0xBB)      # 行間セパレーターライン
COLOR_SEPARATOR_MAIN = RGBColor(0x66, 0x66, 0x66) # シェブロン下の太いライン

# ── フォント ──
FONT_NAME_JP = "Meiryo UI"
FONT_NAME_LATIN = "Arial"

# ── フォントサイズ (列数に応じて自動決定; _apply_theme で roleup 用に上書き) ──
FONT_SIZE_SOURCE_PT = 10


def calc_font_sizes(n_stages):
    """列数に応じてシェブロン・記号・バーのフォントサイズを自動決定 (stella default)"""
    if n_stages <= 5:
        return {"chevron": 14, "chevron_sub": 10, "symbol": 36, "bar": 14, "row_label": 12}
    elif n_stages == 6:
        return {"chevron": 13, "chevron_sub": 9, "symbol": 32, "bar": 13, "row_label": 12}
    else:  # 7
        return {"chevron": 12, "chevron_sub": 9, "symbol": 28, "bar": 12, "row_label": 11}


def calc_font_sizes_roleup(n_stages):
    """roleup C4 許容集合 [22, 14, 12, 10, 6] pt にハマるフォントサイズ"""
    if n_stages <= 5:
        return {"chevron": 14, "chevron_sub": 10, "symbol": 22, "bar": 12, "row_label": 12}
    else:  # 6 or 7
        return {"chevron": 12, "chevron_sub": 10, "symbol": 14, "bar": 12, "row_label": 10}


def _apply_theme(theme):
    """roleup の場合、レイアウト・色・フォントサイズを brand 仕様に上書きする。"""
    global SLIDE_W, SLIDE_H
    global CONTENT_LEFT, CONTENT_WIDTH
    global ROW_LABEL_WIDTH, STAGES_LEFT, STAGES_WIDTH
    global CHEVRON_TOP, CHEVRON_HEIGHT, GAP, DATA_ROW_HEIGHT, DATA_ROWS_TOP
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_NAVY, COLOR_CHEVRON_FILL
    global COLOR_BAR_FILL, COLOR_BAR_BORDER
    global COLOR_SOURCE, COLOR_SEPARATOR, COLOR_SEPARATOR_MAIN
    global FONT_NAME_JP, FONT_NAME_LATIN, FONT_SIZE_SOURCE_PT
    global _THEME
    _THEME = theme

    if theme.id != "roleup":
        return

    # A4 横 (11.69 × 8.27)
    SLIDE_W = Inches(11.69)
    SLIDE_H = Inches(8.27)

    CONTENT_LEFT = Inches(0.41)
    CONTENT_WIDTH = Inches(10.87)
    ROW_LABEL_WIDTH = Inches(1.70)
    STAGES_LEFT = CONTENT_LEFT + ROW_LABEL_WIDTH
    STAGES_WIDTH = CONTENT_WIDTH - ROW_LABEL_WIDTH

    CHEVRON_TOP = Inches(1.55)
    CHEVRON_HEIGHT = Inches(0.65)
    GAP = Inches(0.06)
    DATA_ROW_HEIGHT = Inches(0.95)
    DATA_ROWS_TOP = CHEVRON_TOP + CHEVRON_HEIGHT + GAP   # 1.55 + 0.65 + 0.06 = 2.26
    # 4 行 = 3.80 → end at 6.06. Source 3 placeholder is at 7.45, so OK.

    SOURCE_X = CONTENT_LEFT
    SOURCE_Y = Inches(7.45)
    SOURCE_W = CONTENT_WIDTH
    SOURCE_H = Inches(0.30)

    # roleup 茶系トーン
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    # シェブロン・バーの色を roleup palette から構成
    # label_bar = #7C4C2C (濃茶) を border/text 色に、label_bg = #F2E8DD (淡褐) を fill に
    COLOR_NAVY = theme.color("label_bar")            # #7C4C2C
    COLOR_CHEVRON_FILL = theme.color("label_bg")     # #F2E8DD
    COLOR_BAR_FILL = theme.color("header_bg")        # #F5EFE5
    COLOR_BAR_BORDER = theme.color("accent_op_margin_line")  # #604C3F

    # フォント
    FONT_NAME_JP = theme.font_ea or "Yu Gothic UI"
    # latin もまとめて Yu Gothic UI に統一 (C8 fail 回避: 旧 Arial ハードコード)
    FONT_NAME_LATIN = theme.font_ea or "Yu Gothic UI"

    # 出典は theme defaults を尊重
    FONT_SIZE_SOURCE_PT = int(theme._defaults.get("font_size_source_pt", 6))


def get_font_sizes(n_stages):
    """brand に応じてフォントサイズセットを返す (late-resolve)"""
    if _THEME is not None and _THEME.id == "roleup":
        return calc_font_sizes_roleup(n_stages)
    return calc_font_sizes(n_stages)


# ════════════════════════════════════════════════════════════════
# ユーティリティ
# ════════════════════════════════════════════════════════════════

def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_placeholder_text(shape, text, font_size_pt=None, bold=None, color=None):
    """既存placeholderのテキストを置き換え（スタイルを指定可能）"""
    if shape is None:
        return
    tf = shape.text_frame
    p_elem = tf.paragraphs[0]._p
    for r in p_elem.findall(qn("a:r")):
        p_elem.remove(r)
    r = etree.SubElement(p_elem, qn("a:r"))
    rPr_attrs = {"lang": "ja-JP"}
    if font_size_pt is not None:
        rPr_attrs["sz"] = str(int(font_size_pt * 100))
    if bold is not None:
        rPr_attrs["b"] = "1" if bold else "0"
    rPr = etree.SubElement(r, qn("a:rPr"), attrib=rPr_attrs)
    if color is not None:
        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr"),
                         attrib={"val": "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_LATIN})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    t = etree.SubElement(r, qn("a:t"))
    t.text = text


def add_textbox(slide, left, top, width, height, text, *,
                font_size_pt=12, bold=False, color=(0x33, 0x33, 0x33),
                align="center", vertical="middle", word_wrap=True,
                font_name_jp=None, font_name_latin=None):
    """指定位置にテキストボックスを追加（汎用、font_name は late-resolve）"""
    if font_name_jp is None:
        font_name_jp = FONT_NAME_JP
    if font_name_latin is None:
        font_name_latin = FONT_NAME_LATIN
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
        bodyPr.set("anchor", anchor_map.get(vertical, "ctr"))

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.CENTER)
        for r in p._p.findall(qn("a:r")):
            p._p.remove(r)
        pPr = p._p.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.SubElement(p._p, qn("a:pPr"))
            p._p.insert(0, pPr)
        for x in pPr.findall(qn("a:lnSpc")):
            pPr.remove(x)
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
        r = etree.SubElement(p._p, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(font_size_pt * 100)),
            "b": "1" if bold else "0",
        })
        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr"),
                         attrib={"val": "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])})
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": font_name_latin})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": font_name_jp})
        t = etree.SubElement(r, qn("a:t"))
        t.text = line
    return tb


def _add_styled_run(p_elem, text, font_size_pt, bold, text_color):
    """段落要素にスタイル付きrunを追加（共通ヘルパー、font_name は late-resolve）"""
    r = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size_pt * 100)),
        "b": "1" if bold else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"),
                     attrib={"val": "{:02X}{:02X}{:02X}".format(text_color[0], text_color[1], text_color[2])})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_LATIN})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    t = etree.SubElement(r, qn("a:t"))
    t.text = text


def _prepare_paragraph(p):
    """段落から既存runを削除し、行間・中央揃え設定を行う"""
    for r in p._p.findall(qn("a:r")):
        p._p.remove(r)
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(p._p, qn("a:pPr"))
        p._p.insert(0, pPr)
    for x in pPr.findall(qn("a:lnSpc")):
        pPr.remove(x)
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    pPr.set("algn", "ctr")


def _set_chevron_point_depth(shape, depth_val=25000):
    """CHEVRON形状の矢尻の鋭さを調整
    0=長方形に近い、50000=デフォルト、100000=最大。小さいほど文字領域が広い。"""
    spPr = shape._element.find(".//" + qn("p:spPr"))
    if spPr is None:
        return
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn("a:avLst"))
    for gd in avLst.findall(qn("a:gd")):
        avLst.remove(gd)
    etree.SubElement(avLst, qn("a:gd"), attrib={"name": "adj", "fmla": f"val {depth_val}"})


def add_shape_with_text(slide, shape_enum, left, top, width, height,
                        text, *, fill_color, border_color, text_color,
                        font_size_pt, bold=True, vertical="middle",
                        border_width_pt=1.0, sub_text=None, sub_font_size_pt=None,
                        margin_h_emu=9000, margin_v_emu=9000,
                        chevron_depth=None):
    """図形を追加してテキストを入れる（シェブロン、長方形 etc.）"""
    sh = slide.shapes.add_shape(shape_enum, left, top, width, height)
    if shape_enum == MSO_SHAPE.CHEVRON:
        _set_chevron_point_depth(sh, chevron_depth if chevron_depth is not None else 25000)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.color.rgb = border_color
    sh.line.width = Pt(border_width_pt)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(margin_h_emu)
    tf.margin_right = Emu(margin_h_emu)
    tf.margin_top = Emu(margin_v_emu)
    tf.margin_bottom = Emu(margin_v_emu)
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
        bodyPr.set("anchor", anchor_map.get(vertical, "ctr"))

    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        _prepare_paragraph(p)
        _add_styled_run(p._p, line, font_size_pt, bold, text_color)

    if sub_text:
        sub_sz = sub_font_size_pt or max(int(font_size_pt * 0.75), 8)
        for line in str(sub_text).split("\n"):
            p2 = tf.add_paragraph()
            _prepare_paragraph(p2)
            _add_styled_run(p2._p, line, sub_sz, False, text_color)

    return sh


def add_horizontal_line(slide, left, top, width, color, thickness_pt=0.75):
    """水平セパレーターライン"""
    sh = slide.shapes.add_connector(1, left, top, left + width, top)  # 1 = STRAIGHT
    sh.line.color.rgb = color
    sh.line.width = Pt(thickness_pt)
    return sh


# ════════════════════════════════════════════════════════════════
# マトリクス構築
# ════════════════════════════════════════════════════════════════

def build_matrix(slide, stages, rows, font_sizes):
    """シェブロン行 + 4行のマトリクスを描画"""
    n_stages = len(stages)
    col_w = STAGES_WIDTH / n_stages

    overlap = Inches(0.12)
    for i, stage in enumerate(stages):
        cx = STAGES_LEFT + col_w * i
        if i < n_stages - 1:
            cw = col_w + overlap
        else:
            cw = col_w
        add_shape_with_text(
            slide, MSO_SHAPE.CHEVRON,
            cx, CHEVRON_TOP, cw, CHEVRON_HEIGHT,
            stage.get("name", ""),
            fill_color=COLOR_CHEVRON_FILL,
            border_color=COLOR_NAVY,
            text_color=(COLOR_NAVY[0], COLOR_NAVY[1], COLOR_NAVY[2]),
            font_size_pt=font_sizes["chevron"],
            bold=True,
            border_width_pt=1.0,
            sub_text=stage.get("sub"),
            sub_font_size_pt=font_sizes["chevron_sub"],
        )

    sep_y = CHEVRON_TOP + CHEVRON_HEIGHT + Inches(0.02)
    add_horizontal_line(slide, CONTENT_LEFT, sep_y, CONTENT_WIDTH,
                        COLOR_SEPARATOR_MAIN, thickness_pt=1.5)

    if len(rows) != 4:
        print(f"  ⚠ WARNING: expected 4 rows, got {len(rows)}", file=sys.stderr)

    current_y = DATA_ROWS_TOP
    for row_idx, row in enumerate(rows):
        row_label = row.get("label", "")
        row_type = row.get("type", "symbols")

        add_textbox(
            slide, CONTENT_LEFT, current_y, ROW_LABEL_WIDTH, DATA_ROW_HEIGHT,
            row_label,
            font_size_pt=font_sizes["row_label"], bold=False,
            color=(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2]),
            align="center", vertical="middle", word_wrap=True,
        )

        if row_type == "symbols":
            cells = row.get("cells", [])
            while len(cells) < n_stages:
                cells.append("")
            for c_idx in range(n_stages):
                cell_text = str(cells[c_idx] or "").strip()
                if cell_text:
                    cx = STAGES_LEFT + col_w * c_idx
                    add_textbox(
                        slide, cx, current_y, col_w, DATA_ROW_HEIGHT,
                        cell_text,
                        font_size_pt=font_sizes["symbol"], bold=True,
                        color=(COLOR_NAVY[0], COLOR_NAVY[1], COLOR_NAVY[2]),
                        align="center", vertical="middle",
                    )

        elif row_type == "bars":
            bars = row.get("bars", [])
            n_bars = len(bars)
            bar_pad_x = Inches(0.08)
            bar_pad_y_outer = Inches(0.10)

            # 複数バーは縦に sub-stripe で分割 (range が重複しても overlap しないように)
            usable_h = DATA_ROW_HEIGHT - bar_pad_y_outer * 2
            stripe_h = usable_h / max(n_bars, 1)
            inner_pad = Inches(0.02) if n_bars > 1 else Inches(0)

            for b_idx, bar in enumerate(bars):
                s = int(bar.get("span_start", 0))
                e = int(bar.get("span_end", s))
                s = max(0, min(s, n_stages - 1))
                e = max(s, min(e, n_stages - 1))
                bar_label = bar.get("label", "")

                bar_x = STAGES_LEFT + col_w * s + bar_pad_x
                bar_w = col_w * (e - s + 1) - bar_pad_x * 2
                bar_y = current_y + bar_pad_y_outer + stripe_h * b_idx + inner_pad
                bar_h = stripe_h - inner_pad * 2

                add_shape_with_text(
                    slide, MSO_SHAPE.RECTANGLE,
                    bar_x, bar_y, bar_w, bar_h,
                    bar_label,
                    fill_color=COLOR_BAR_FILL,
                    border_color=COLOR_BAR_BORDER,
                    text_color=(COLOR_NAVY[0], COLOR_NAVY[1], COLOR_NAVY[2]),
                    font_size_pt=font_sizes["bar"],
                    bold=True,
                    border_width_pt=0.75,
                )

        if row_idx < len(rows) - 1:
            line_y = current_y + DATA_ROW_HEIGHT
            add_horizontal_line(slide, CONTENT_LEFT, line_y, CONTENT_WIDTH,
                                COLOR_SEPARATOR, thickness_pt=0.5)

        current_y = current_y + DATA_ROW_HEIGHT

    add_horizontal_line(slide, CONTENT_LEFT, current_y, CONTENT_WIDTH,
                        COLOR_SEPARATOR_MAIN, thickness_pt=1.0)


# ════════════════════════════════════════════════════════════════
# メイン処理
# ════════════════════════════════════════════════════════════════

def fill_slide(data, theme, template_path, output_path):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # ── Top text (stella: main_message / roleup: chart_title) ──
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        # roleup では Title 1 のテンプレ default (22pt) を尊重するため明示 size 指定なし
        title_size = 22 if theme.id == "roleup" else 28
        msg_shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(msg_shape, top_text,
                             font_size_pt=title_size, bold=True,
                             color=(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2]))
        print(f"  ✓ Top: {top_text[:60]}{'...' if len(top_text) > 60 else ''}")

    # ── Subtitle (stella: chart_title / roleup: main_message) ──
    sub_text = resolve_subtitle_text(data, theme).strip()
    if sub_text:
        sub_size = 12 if theme.id == "roleup" else 14
        title_shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(title_shape, sub_text,
                             font_size_pt=sub_size, bold=False,
                             color=(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2]))
        print(f"  ✓ Subtitle: {sub_text[:60]}{'...' if len(sub_text) > 60 else ''}")

    # ── マトリクス本体 ──
    stages = data.get("stages", [])
    rows = data.get("rows", [])
    if not stages:
        raise ValueError("stages は必須です")
    if len(stages) < 5 or len(stages) > 7:
        print(f"  ⚠ WARNING: stages数は5〜7推奨（現在 {len(stages)}）", file=sys.stderr)
    if not rows:
        raise ValueError("rows は必須です")

    font_sizes = get_font_sizes(len(stages))
    print(f"  ✓ 段階数: {len(stages)} → フォントサイズ: {font_sizes}")

    build_matrix(slide, stages, rows, font_sizes)
    print(f"  ✓ マトリクス構築完了")

    # ── 出典 (roleup: Source 3 placeholder, stella: 動的 textbox) ──
    source = (data.get("source") or data.get("source_label")
              or data.get("source_text") or "").strip()
    if source:
        if theme.is_source_required():
            src_shape = find_shape(slide, SHAPE_SOURCE_PH)
            if src_shape is not None:
                set_placeholder_text(src_shape, f"出典: {source}",
                                     font_size_pt=FONT_SIZE_SOURCE_PT,
                                     bold=False,
                                     color=(COLOR_SOURCE[0], COLOR_SOURCE[1], COLOR_SOURCE[2]))
        else:
            add_textbox(
                slide, SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                source,
                font_size_pt=FONT_SIZE_SOURCE_PT, bold=False,
                color=(COLOR_SOURCE[0], COLOR_SOURCE[1], COLOR_SOURCE[2]),
                align="left", vertical="top",
            )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    _finalize_pptx(output_path)
    print(f"\n✓ 保存: {output_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True, help="JSON data path")
    ap.add_argument("--template", required=False, default=None, help="PPTX template path (省略時 brand から自動解決)")
    ap.add_argument("--output", required=True, help="Output PPTX path")
    add_brand_arg(ap)
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "stages", "rows"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "stages", "rows",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "value-chain-matrix")

    print("=" * 60)
    print("バリューチェーン・ポジショニング・マトリクス生成")
    print("=" * 60)
    fill_slide(data, theme, template_path, args.output)


if __name__ == "__main__":
    main()
