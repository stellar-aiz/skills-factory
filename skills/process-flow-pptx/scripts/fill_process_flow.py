"""
fill_process_flow.py — プロセスフローデータをPPTXテンプレートに流し込むスクリプト

ProcessFlow9.pptx をベーステンプレートとして使用し、
ステップ数（3〜9）に応じて不要なボックス・コネクタを自動削除する。

テンプレート構造（確認済み）:
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Rectangle 4/3/13/24/25/26/40/41/42: Process1〜9
  - TextBox 6: Features (9 paragraphs)
  - Straight Arrow Connector 9/14/27/30/33/27/46/49: 矢印コネクタ

使い方:
  python fill_process_flow.py \
    --data /home/claude/process_flow_data.json \
    --output /mnt/user-data/outputs/ProcessFlow_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── テンプレートパス ──────────────────────────────────────────
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(SKILL_DIR, "assets", "ProcessFlow9.pptx")

# ── Shape名マッピング（確認済み）──────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
FEATURE_SHAPE      = "TextBox 6"

# プロセスステップのShape名（順序通り: Step1〜Step9）
STEP_SHAPES = [
    "Rectangle 4",   # Step1 (Row1-Col1)
    "Rectangle 3",   # Step2 (Row1-Col2)
    "Rectangle 13",  # Step3 (Row1-Col3)
    "Rectangle 24",  # Step4 (Row2-Col1)
    "Rectangle 25",  # Step5 (Row2-Col2)
    "Rectangle 26",  # Step6 (Row2-Col3)
    "Rectangle 40",  # Step7 (Row3-Col1)
    "Rectangle 41",  # Step8 (Row3-Col2)
    "Rectangle 42",  # Step9 (Row3-Col3)
]

# ── 削除ロジック ─────────────────────────────────────────────
# コネクタの対応関係:
#   Connector 9:  Step1 → Step2 の水平矢印
#   Connector 14: Step2 → Step3 の水平矢印
#   Connector 27 (top=1247755): Row1 → Row2 のU字折り返し
#   Connector 30: Step4 → Step5 の水平矢印
#   Connector 33: Step5 → Step6 の水平矢印
#   Connector 27 (top=2731758): Row2 → Row3 のU字折り返し
#   Connector 46: Step7 → Step8 の水平矢印
#   Connector 49: Step8 → Step9 の水平矢印
#
# 削除ルール: ステップNが不要 → そのステップのRectangleと、
#             そのステップへの入力コネクタを削除
#
# (shape_name, top_position_to_match_or_None)
SHAPES_TO_DELETE_IF_STEP_UNUSED = {
    9: [("Rectangle 42", None), ("Straight Arrow Connector 49", None)],
    8: [("Rectangle 41", None), ("Straight Arrow Connector 46", None)],
    7: [("Rectangle 40", None), ("Straight Arrow Connector 27", 2731758)],  # U-turn Row2→Row3
    6: [("Rectangle 26", None), ("Straight Arrow Connector 33", None)],
    5: [("Rectangle 25", None), ("Straight Arrow Connector 30", None)],
    4: [("Rectangle 24", None), ("Straight Arrow Connector 27", 1247755)],  # U-turn Row1→Row2
}
# ────────────────────────────────────────────────────────────


def find_shape(slide, name):
    """スライドから指定名のShapeを検索"""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def find_shape_by_name_and_top(slide, name, top_pos):
    """同名Shapeが複数ある場合、top位置で特定する"""
    for shape in slide.shapes:
        if shape.name == name and shape.top == top_pos:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' (top={top_pos}) not found", file=sys.stderr)
    return None


def delete_shape(slide, shape):
    """スライドからShapeを削除する"""
    if shape is None:
        return
    sp = shape._element
    sp.getparent().remove(sp)


def set_placeholder_text(shape, text):
    """PlaceholderのTextFrameにテキストをセット（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def set_shape_text(shape, text):
    """AUTO_SHAPEのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_features(slide, shape_name, features):
    """特徴テキストボックスの9段落にデータを書き込む

    段落構成:
      para[0]: Feature1 ラベル (bold)
      para[1]: Feature1 Comment1
      para[2]: Feature1 Comment2
      para[3]: Feature2 ラベル (bold)
      para[4]: Feature2 Comment1
      para[5]: Feature2 Comment2
      para[6]: Feature3 ラベル (bold)
      para[7]: Feature3 Comment1
      para[8]: Feature3 Comment2
    """
    shape = find_shape(slide, shape_name)
    if shape is None:
        return

    tf = shape.text_frame

    for i, feature in enumerate(features[:3]):
        label    = feature.get("label", "").strip()
        comment1 = feature.get("comment1", "").strip()
        comment2 = feature.get("comment2", "").strip()

        base_idx = i * 3  # 0, 3, 6

        # ラベル
        if base_idx < len(tf.paragraphs):
            para = tf.paragraphs[base_idx]
            if para.runs:
                para.runs[0].text = label
                for run in para.runs[1:]:
                    run.text = ""
            else:
                r_elem = etree.SubElement(para._p, qn("a:r"))
                etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP", "b": "1"})
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = label
            print(f"  [Feature{i+1} Label] {label}")

        # Comment1
        if base_idx + 1 < len(tf.paragraphs):
            para = tf.paragraphs[base_idx + 1]
            if para.runs:
                para.runs[0].text = comment1
                for run in para.runs[1:]:
                    run.text = ""
            else:
                r_elem = etree.SubElement(para._p, qn("a:r"))
                etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = comment1
            print(f"  [Feature{i+1} C1]    {comment1[:50]}{'...' if len(comment1) > 50 else ''}")

        # Comment2
        if base_idx + 2 < len(tf.paragraphs):
            para = tf.paragraphs[base_idx + 2]
            if para.runs:
                para.runs[0].text = comment2
                for run in para.runs[1:]:
                    run.text = ""
            else:
                r_elem = etree.SubElement(para._p, qn("a:r"))
                etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = comment2
            print(f"  [Feature{i+1} C2]    {comment2[:50]}{'...' if len(comment2) > 50 else ''}")


def remove_unused_shapes(slide, step_count):
    """ステップ数に応じて不要なRectangleとConnectorを削除する"""
    deleted = []
    # Step 9 → 4 の順に、不要なステップのShapeを削除
    for step_num in range(9, 3, -1):  # 9, 8, 7, 6, 5, 4
        if step_num > step_count:
            shapes_to_del = SHAPES_TO_DELETE_IF_STEP_UNUSED.get(step_num, [])
            for shape_name, top_pos in shapes_to_del:
                if top_pos is not None:
                    shape = find_shape_by_name_and_top(slide, shape_name, top_pos)
                else:
                    shape = find_shape(slide, shape_name)
                if shape:
                    delete_shape(slide, shape)
                    deleted.append(f"{shape_name}" + (f" (top={top_pos})" if top_pos else ""))

    if deleted:
        print(f"  🗑 削除したShape ({len(deleted)}個): {', '.join(deleted)}")
    else:
        print(f"  ✓ 全9ステップ使用（削除なし）")


def main():
    parser = argparse.ArgumentParser(description="プロセスフローデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True,  help="process_flow_data.json のパス")
    parser.add_argument("--template", default=None,   help="テンプレートPPTXのパス（省略時はデフォルト）")
    parser.add_argument("--output",   required=True,  help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "steps"],
        allowed_top=["main_message", "chart_title", "steps", "features"],
        skill_name="process-flow-pptx",
    )

    # テンプレート読み込み
    template_path = args.template if args.template else TEMPLATE
    print(f"📋 テンプレート: {os.path.basename(template_path)}")

    prs   = Presentation(template_path)
    slide = prs.slides[0]

    # ── ステップ数の確認 ──────────────────────────────────────
    steps = data.get("steps", [])
    step_count = len(steps)
    if step_count < 3 or step_count > 9:
        print(f"  ❌ ステップ数は3〜9の範囲で指定してください（現在: {step_count}）", file=sys.stderr)
        sys.exit(1)
    print(f"  ステップ数: {step_count}")

    # ── 不要なShapeを削除 ─────────────────────────────────────
    remove_unused_shapes(slide, step_count)

    # ── Main Message ──────────────────────────────────────────
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(shape, main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # ── Chart Title ───────────────────────────────────────────
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(shape, chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # ── プロセスステップ ──────────────────────────────────────
    for i, step_text in enumerate(steps):
        if i < len(STEP_SHAPES):
            shape = find_shape(slide, STEP_SHAPES[i])
            set_shape_text(shape, step_text.strip())
            print(f"  [Step{i+1}] {STEP_SHAPES[i]} → {step_text.strip()}")

    # ── プロセス特徴 ──────────────────────────────────────────
    features = data.get("features", [])
    if len(features) != 3:
        print(f"  ⚠ features は {len(features)} 件です（3件推奨）", file=sys.stderr)
    fill_features(slide, FEATURE_SHAPE, features)

    # ── 保存 ──────────────────────────────────────────────────
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
