"""
fill_workforce_composition.py — 人員構成スライドをPPTXネイティブオブジェクトで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

レイアウト:
  - 左側: 在籍人員数の推移 (ネイティブ棒グラフ 3 系列: total / new_hires / departures)
  - 右側: 部署別人員構成テーブル (ネイティブテーブル + 合計行)
  - 下部: 出典 (stella=動的 textbox / roleup=Source 3 placeholder)

Usage:
  python fill_workforce_composition.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/workforce_composition_data.json \\
    --output {{OUTPUT_DIR}}/WorkforceComposition_output.pptx
"""

import argparse
import json
import math
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "workforce-composition-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_TABLE = "Table 1"

# Defaults (stella). Reassigned in _apply_theme(theme).
SHAPE_SOURCE = "Source"
PANEL_Y = Inches(1.50)
LEFT_X = Inches(0.41)
LEFT_W = Inches(5.80)
RIGHT_X = Inches(6.50)
RIGHT_W = Inches(6.40)
CHART_H = Inches(5.00)
TABLE_MAX_H = Inches(5.20)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.00)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_BAR_TOTAL = "ED7D31"
COLOR_BAR_HIRES = "4472C4"
COLOR_BAR_DEPARTURES = "A5A5A5"
COLOR_HEADER_BG = "F0F0F0"
COLOR_TOTAL_BG = "E8E8E8"
COLOR_EVEN_ROW = "FAFAFA"
BORDER_HEX = "CCCCCC"

FONT_NAME_JP = "Meiryo UI"
TEXT_HEX = "333333"

SECTION_TITLE_PT = 14
SECTION_TITLE_BOLD = True
SECTION_TITLE_HEX = "333333"
SECTION_TITLE_ALIGN = PP_ALIGN.CENTER
TABLE_HEADER_PT = 12
TABLE_BODY_PT = 12
DATA_LABEL_SZ = "1100"
AXIS_FONT_SZ = "1100"
LEGEND_FONT_PT = 10
UNIT_FONT_PT = 11
SOURCE_FONT_PT = 10

_THEME = None


def _apply_theme(theme):
    global _THEME
    global SHAPE_SOURCE
    global PANEL_Y, LEFT_X, LEFT_W, RIGHT_X, RIGHT_W, CHART_H, TABLE_MAX_H
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE
    global COLOR_BAR_TOTAL, COLOR_BAR_HIRES, COLOR_BAR_DEPARTURES
    global COLOR_HEADER_BG, COLOR_TOTAL_BG, COLOR_EVEN_ROW, BORDER_HEX
    global FONT_NAME_JP, TEXT_HEX
    global SECTION_TITLE_PT, SECTION_TITLE_BOLD, SECTION_TITLE_HEX, SECTION_TITLE_ALIGN
    global TABLE_HEADER_PT, TABLE_BODY_PT
    global DATA_LABEL_SZ, AXIS_FONT_SZ
    global LEGEND_FONT_PT, UNIT_FONT_PT, SOURCE_FONT_PT

    _THEME = theme
    FONT_NAME_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    PANEL_Y = theme.layout("panel_y_in")
    LEFT_X = theme.layout("left_x_in")
    LEFT_W = theme.layout("left_w_in")
    RIGHT_X = theme.layout("right_x_in")
    RIGHT_W = theme.layout("right_w_in")
    CHART_H = theme.layout("chart_h_in")
    TABLE_MAX_H = theme.layout("table_max_h_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        COLOR_BAR_TOTAL = "ED7D31"
        COLOR_BAR_HIRES = "4472C4"
        COLOR_BAR_DEPARTURES = "A5A5A5"
        COLOR_HEADER_BG = "F0F0F0"
        COLOR_TOTAL_BG = "E8E8E8"
        COLOR_EVEN_ROW = "FAFAFA"
        BORDER_HEX = "CCCCCC"
        SECTION_TITLE_PT = 14
        SECTION_TITLE_BOLD = True
        SECTION_TITLE_HEX = TEXT_HEX
        SECTION_TITLE_ALIGN = PP_ALIGN.CENTER
        TABLE_HEADER_PT = 12
        TABLE_BODY_PT = 12
        DATA_LABEL_SZ = "1100"
        AXIS_FONT_SZ = "1100"
        LEGEND_FONT_PT = 10
        UNIT_FONT_PT = 11
        SOURCE_FONT_PT = 10
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        # Bar colors from chart_palette[0..2]
        palette = list(theme.chart_palette)
        COLOR_BAR_TOTAL = palette[0].lstrip("#")        # #7C4C2C
        COLOR_BAR_HIRES = palette[1].lstrip("#")        # #897141
        COLOR_BAR_DEPARTURES = theme.hex_no_hash("highlight_other")  # #CDCECE
        COLOR_HEADER_BG = theme.hex_no_hash("header_bg")  # #F5EFE5
        COLOR_TOTAL_BG = theme.hex_no_hash("label_bg")    # #F2E8DD
        COLOR_EVEN_ROW = theme.hex_no_hash("label_bg")    # #F2E8DD (banded)
        BORDER_HEX = theme.hex_no_hash("highlight_other")  # #CDCECE
        SECTION_TITLE_PT = theme.pt_value("font_size_subtitle_pt")  # 12
        SECTION_TITLE_BOLD = False
        SECTION_TITLE_HEX = theme.hex_no_hash("subtitle")  # #897141
        SECTION_TITLE_ALIGN = PP_ALIGN.LEFT
        TABLE_HEADER_PT = theme.pt_value("font_size_body_pt")  # 10
        TABLE_BODY_PT = theme.pt_value("font_size_body_pt")    # 10
        DATA_LABEL_SZ = "1000"  # 10pt (was 11pt for stella)
        AXIS_FONT_SZ = "1000"   # 10pt (was 11pt for stella)
        LEGEND_FONT_PT = theme.pt_value("font_size_body_pt")  # 10
        UNIT_FONT_PT = theme.pt_value("font_size_body_pt")    # 10 (was 11pt for stella)
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name, warn=True):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    if warn:
        print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def write_source_placeholder(shape, text, font_size_pt, font_name):
    tf = shape.text_frame
    para = tf.paragraphs[0]
    for r in list(para.runs):
        r.text = ""
    if para.runs:
        run = para.runs[0]
        run.text = text
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text
        run = para.runs[0]
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        run._r.insert(0, rPr)
    rPr.set("sz", str(font_size_pt * 100))
    for tag in [qn("a:latin"), qn("a:ea")]:
        old = rPr.find(tag)
        if old is not None:
            rPr.remove(old)
        etree.SubElement(rPr, tag, attrib={"typeface": font_name})


def _hex2rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_section_title(slide, text, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = SECTION_TITLE_ALIGN
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SECTION_TITLE_PT)
    run.font.bold = SECTION_TITLE_BOLD
    run.font.color.rgb = _hex2rgb(SECTION_TITLE_HEX)
    run.font.name = FONT_NAME_JP

    # 下線 (装飾、テキストなしのため C11 対象外)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _hex2rgb(SECTION_TITLE_HEX)
    line.line.fill.background()
    return txBox


def add_unit_label(slide, text, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.22))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(UNIT_FONT_PT)
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP


def add_source_dynamic_textbox(slide, text):
    txBox = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(SOURCE_FONT_PT)
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP
    print(f"  ✓ 出典 (dynamic textbox): {text[:50]}")


def build_headcount_chart(slide, trend_data, left, top, width, height):
    from pptx.chart.data import CategoryChartData

    periods = trend_data["periods"]
    labels = trend_data.get("series_labels", {})
    total_label = labels.get("total", "総従業員数")
    hires_label = labels.get("new_hires", "入社人数")
    depart_label = labels.get("departures", "退職人数")

    chart_data = CategoryChartData()
    chart_data.categories = [p["label"] for p in periods]
    chart_data.add_series(total_label, [p["total"] for p in periods])
    chart_data.add_series(hires_label, [p["new_hires"] for p in periods])
    chart_data.add_series(depart_label, [p["departures"] for p in periods])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = chart_frame.chart

    plotArea = chart._chartSpace.chart.plotArea
    barChart = plotArea.findall(qn('c:barChart'))[0]

    gapWidth = barChart.find(qn('c:gapWidth'))
    if gapWidth is None:
        gapWidth = etree.SubElement(barChart, qn('c:gapWidth'))
    gapWidth.set('val', '80')

    overlap = barChart.find(qn('c:overlap'))
    if overlap is None:
        overlap = etree.SubElement(barChart, qn('c:overlap'))
    overlap.set('val', '-20')

    colors = [COLOR_BAR_TOTAL, COLOR_BAR_HIRES, COLOR_BAR_DEPARTURES]
    sers = barChart.findall(qn('c:ser'))

    for i, ser in enumerate(sers):
        spPr = ser.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ser, qn('c:spPr'))
        sf = spPr.find(qn('a:solidFill'))
        if sf is None:
            sf = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(sf):
            sf.remove(child)
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': colors[i]})

        _add_data_labels(ser, num_format='#,##0;-#,##0;""')

    for vax in plotArea.findall(qn('c:valAx')):
        for mg in vax.findall(qn('c:majorGridlines')):
            vax.remove(mg)
        for mg in vax.findall(qn('c:minorGridlines')):
            vax.remove(mg)
    for cax in plotArea.findall(qn('c:catAx')):
        for mg in cax.findall(qn('c:majorGridlines')):
            cax.remove(mg)

    for vax in plotArea.findall(qn('c:valAx')):
        del_elem = vax.find(qn('c:delete'))
        if del_elem is None:
            del_elem = etree.SubElement(vax, qn('c:delete'))
        del_elem.set('val', '1')

    for cax in plotArea.findall(qn('c:catAx')):
        delete_elem = cax.find(qn('c:delete'))
        if delete_elem is not None and delete_elem.get('val') == '1':
            continue
        txPr = cax.find(qn('c:txPr'))
        if txPr is None:
            txPr = etree.SubElement(cax, qn('c:txPr'))
        bodyPr = txPr.find(qn('a:bodyPr'))
        if bodyPr is None:
            bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
            txPr.insert(0, bodyPr)
        if txPr.find(qn('a:lstStyle')) is None:
            etree.SubElement(txPr, qn('a:lstStyle'))
        p = txPr.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(txPr, qn('a:p'))
        pPr = p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p, qn('a:pPr'))
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': AXIS_FONT_SZ})
        else:
            defRPr.set('sz', AXIS_FONT_SZ)
        if defRPr.find(qn('a:latin')) is None:
            etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
        if defRPr.find(qn('a:ea')) is None:
            etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})

    chart.has_title = False
    chart.has_legend = False

    plotArea_spPr = plotArea.find(qn('c:spPr'))
    if plotArea_spPr is None:
        plotArea_spPr = etree.SubElement(plotArea, qn('c:spPr'))
    noFill = plotArea_spPr.find(qn('a:noFill'))
    if noFill is None:
        etree.SubElement(plotArea_spPr, qn('a:noFill'))

    print(f"  ✓ 人員推移チャート生成: {len(periods)}期分")
    return chart_frame


def _add_data_labels(ser_xml, num_format='#,##0'):
    dLbls = ser_xml.find(qn('c:dLbls'))
    if dLbls is None:
        dLbls = etree.SubElement(ser_xml, qn('c:dLbls'))

    for child in list(dLbls):
        dLbls.remove(child)

    etree.SubElement(dLbls, qn('c:numFmt'), attrib={
        'formatCode': num_format, 'sourceLinked': '0'
    })
    etree.SubElement(dLbls, qn('c:showLegendKey'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dLbls, qn('c:showCatName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showSerName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showPercent'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showBubbleSize'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:dLblPos'), attrib={'val': 'outEnd'})

    txPr = etree.SubElement(dLbls, qn('c:txPr'))
    etree.SubElement(txPr, qn('a:bodyPr'))
    etree.SubElement(txPr, qn('a:lstStyle'))
    p = etree.SubElement(txPr, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': DATA_LABEL_SZ, 'b': '1'})
    etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
    etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})
    sf = etree.SubElement(defRPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': TEXT_HEX})


def add_custom_legend_chart(slide, trend_data, left, top, width):
    labels_cfg = trend_data.get("series_labels", {})
    items = [
        (labels_cfg.get("total", "総従業員数"), COLOR_BAR_TOTAL),
        (labels_cfg.get("new_hires", "入社人数"), COLOR_BAR_HIRES),
        (labels_cfg.get("departures", "退職人数"), COLOR_BAR_DEPARTURES),
    ]

    legend_h = Inches(0.22)
    sq_size = Inches(0.14)
    sq_y = top + (legend_h - sq_size) // 2
    cursor_x = left

    for label_text, color_hex in items:
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(cursor_x), int(sq_y), sq_size, sq_size
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        marker.line.fill.background()

        text_x = cursor_x + sq_size + Inches(0.04)
        text_w = Inches(1.10)
        txBox = slide.shapes.add_textbox(int(text_x), top, text_w, legend_h)
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label_text
        run.font.size = Pt(LEGEND_FONT_PT)
        run.font.color.rgb = COLOR_TEXT
        run.font.name = FONT_NAME_JP

        cursor_x = text_x + text_w + Inches(0.08)

    print(f"  ✓ チャート凡例: {', '.join([i[0] for i in items])}")


def _style_table_cell(cell, text, bold=False, font_size_pt=None, align='l',
                      bg_color=None, font_color=None):
    if font_size_pt is None:
        font_size_pt = TABLE_BODY_PT
    if font_color is None:
        font_color = TEXT_HEX

    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    old_bodyPr = txBody.find(qn("a:bodyPr"))
    if old_bodyPr is not None:
        txBody.remove(old_bodyPr)
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": "36576", "rIns": "36576",
        "tIns": "18288", "bIns": "18288",
        "anchor": "ctr",
    })
    txBody.insert(0, bodyPr)

    if txBody.find(qn("a:lstStyle")) is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.insert(1, lstStyle)

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))

    align_map = {'l': 'l', 'c': 'ctr', 'r': 'r'}
    pPr.set("algn", align_map.get(align, 'l'))

    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size_pt * 100)),
        "b": "1" if bold else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": font_color})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text) if text is not None else "-"

    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)
    tcPr = etree.SubElement(tc, qn("a:tcPr"), attrib={
        "marL": "36576", "marR": "36576",
        "marT": "18288", "marB": "18288",
        "anchor": "ctr",
    })

    for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = etree.SubElement(tcPr, qn(border_name), attrib={"w": "6350", "cmpd": "sng"})
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": BORDER_HEX})

    if bg_color:
        fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
        etree.SubElement(fill_elem, qn("a:srgbClr"), attrib={"val": bg_color})


def build_department_table(slide, dept_data, left, top, width, max_height):
    columns = dept_data.get("columns", [
        "部署名", "人数", "平均年齢", "平均勤続年数", "管理職数", "有資格者数"
    ])
    departments = dept_data.get("departments", [])
    show_total = dept_data.get("show_total", True)

    n_cols = len(columns)
    n_data_rows = len(departments)
    n_rows = 1 + n_data_rows + (1 if show_total else 0)

    row_h = min(Inches(0.35), int(max_height / n_rows))
    table_h = row_h * n_rows

    col0_w = int(width * 0.28)
    remaining_w = width - col0_w
    other_col_w = int(remaining_w / (n_cols - 1))
    col_widths = [col0_w] + [other_col_w] * (n_cols - 1)

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = shape.table

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '0', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(row_h))

    for c_idx, col_name in enumerate(columns):
        align = 'c' if c_idx > 0 else 'l'
        _style_table_cell(
            table.cell(0, c_idx), col_name,
            bold=True, font_size_pt=TABLE_HEADER_PT,
            align=align, bg_color=COLOR_HEADER_BG,
        )

    key_map = {
        "部署名": "name",
        "人数": "headcount",
        "平均年齢": "avg_age",
        "平均勤続年数": "avg_tenure",
        "管理職数": "managers",
        "有資格者数": "certified",
    }

    for r_idx, dept in enumerate(departments):
        row_num = r_idx + 1
        bg = COLOR_EVEN_ROW if r_idx % 2 == 1 else None

        for c_idx, col_name in enumerate(columns):
            key = key_map.get(col_name, col_name)
            value = dept.get(key, "-")

            if value is None:
                display_val = "-"
            elif isinstance(value, float):
                display_val = f"{value:.1f}"
            elif isinstance(value, int):
                display_val = str(value)
            else:
                display_val = str(value)

            align = 'c' if c_idx > 0 else 'l'
            _style_table_cell(
                table.cell(row_num, c_idx), display_val,
                bold=False, font_size_pt=TABLE_BODY_PT,
                align=align, bg_color=bg,
            )

    if show_total and departments:
        total_row = n_rows - 1
        total_headcount = sum(d.get("headcount", 0) or 0 for d in departments)

        total_for_age = sum(
            (d.get("headcount", 0) or 0) * (d.get("avg_age", 0) or 0)
            for d in departments if d.get("avg_age") is not None
        )
        count_for_age = sum(
            d.get("headcount", 0) or 0
            for d in departments if d.get("avg_age") is not None
        )
        avg_age_total = total_for_age / count_for_age if count_for_age > 0 else None

        total_for_tenure = sum(
            (d.get("headcount", 0) or 0) * (d.get("avg_tenure", 0) or 0)
            for d in departments if d.get("avg_tenure") is not None
        )
        count_for_tenure = sum(
            d.get("headcount", 0) or 0
            for d in departments if d.get("avg_tenure") is not None
        )
        avg_tenure_total = total_for_tenure / count_for_tenure if count_for_tenure > 0 else None

        total_managers = sum(d.get("managers", 0) or 0 for d in departments)
        total_certified = sum(d.get("certified", 0) or 0 for d in departments)

        total_values = {
            "name": "合計",
            "headcount": total_headcount,
            "avg_age": avg_age_total,
            "avg_tenure": avg_tenure_total,
            "managers": total_managers,
            "certified": total_certified,
        }

        for c_idx, col_name in enumerate(columns):
            key = key_map.get(col_name, col_name)
            value = total_values.get(key, "-")

            if value is None:
                display_val = "-"
            elif isinstance(value, float):
                display_val = f"{value:.1f}"
            elif isinstance(value, int):
                display_val = str(value)
            else:
                display_val = str(value)

            align = 'c' if c_idx > 0 else 'l'
            _style_table_cell(
                table.cell(total_row, c_idx), display_val,
                bold=True, font_size_pt=TABLE_BODY_PT,
                align=align, bg_color=COLOR_TOTAL_BG,
            )

    print(f"  ✓ 部署テーブル: {n_data_rows}部署 + {'合計行' if show_total else 'なし'}")
    return shape


def main():
    parser = argparse.ArgumentParser(description="人員構成 PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True)
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "workforce-composition")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "headcount_trend", "department_table"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "headcount_trend", "department_table",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    print(f"=== 人員構成スライド生成 (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # Top placeholder (brand-aware)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:50]}")

    sub_text = resolve_subtitle_text(data, theme) or "人員構成"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)

    # Stella V1: Table 1 placeholder removal (silent for roleup)
    _silent_remove_shape(slide, SHAPE_TABLE)

    # Left panel: Headcount Trend
    trend = data.get("headcount_trend", {})
    trend_title = trend.get("title", "在籍人員数の推移")
    unit_text = f"（単位：{trend.get('unit', '人')}）"

    add_section_title(slide, trend_title, LEFT_X, PANEL_Y, LEFT_W)

    legend_y = PANEL_Y + Inches(0.35)
    unit_w = Inches(1.50)
    add_unit_label(slide, unit_text, LEFT_X, legend_y, unit_w)
    # 凡例は単位ラベル右側から開始 (C11 bbox 重なり回避)
    legend_x = LEFT_X + unit_w + Inches(0.10)
    legend_w = LEFT_W - unit_w - Inches(0.10)
    add_custom_legend_chart(slide, trend, legend_x, legend_y, legend_w)

    chart_top = PANEL_Y + Inches(0.60)
    build_headcount_chart(slide, trend, LEFT_X, chart_top, LEFT_W, CHART_H)

    # Right panel: Department Table
    dept = data.get("department_table", {})
    dept_title = dept.get("title", "人員構成")

    add_section_title(slide, dept_title, RIGHT_X, PANEL_Y, RIGHT_W)

    table_top = PANEL_Y + Inches(0.40)
    build_department_table(slide, dept, RIGHT_X, table_top, RIGHT_W, TABLE_MAX_H)

    # Source: stella=Source dynamic textbox / roleup=Source 3 placeholder
    source = data.get("source", "")
    if source:
        body = source if source.startswith("出典") else f"出典：{source}"
        src_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
        if src_shape is not None:
            write_source_placeholder(src_shape, body, SOURCE_FONT_PT, FONT_NAME_JP)
            print(f"  ✓ 出典 ({SHAPE_SOURCE} placeholder): {body[:50]}")
        else:
            add_source_dynamic_textbox(slide, body)

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
