"""
fill_cost_breakdown.py — コスト内訳推移チャート (1〜2チャート並列) をPPTXネイティブオブジェクトで生成

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。

生成するネイティブオブジェクト:
  - 積み上げ棒チャート + 折れ線 (PowerPoint ネイティブ複合チャート)
  - 各セグメントの%ラベル (カスタムデータラベル)
  - トータルラベル / 折れ線ラベル: テキストボックス
  - 凡例: Shape + テキストボックス
  - チャートタイトル / 単位ラベル: テキストボックス

Usage:
  python fill_cost_breakdown.py --brand stellar_aiz \\
    --data {{WORK_DIR}}/cost_breakdown_data.json \\
    --output {{OUTPUT_DIR}}/CostBreakdown_output.pptx
"""

import argparse, copy, json, math, os, sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "cost-breakdown-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape names ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE = "Source"

# Defaults (stella). Reassigned in _apply_theme(theme).
PANEL_Y = Inches(1.50)
SINGLE_CHART_X = Inches(0.50)
SINGLE_CHART_W = Inches(12.30)
LEFT_CHART_X = Inches(0.35)
LEFT_CHART_W = Inches(6.15)
RIGHT_CHART_X = Inches(6.75)
RIGHT_CHART_W = Inches(6.15)
CHART_H = Inches(4.20)
CHART_BOTTOM = Inches(6.30)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.30)

CHART_TITLE_H = Inches(0.30)
UNIT_LABEL_H = Inches(0.20)
LEGEND_H = Inches(0.50)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
FONT_JP = "Meiryo UI"
TEXT_HEX = "333333"

DEFAULT_COLORS = [
    "#4E79A7", "#59A14F", "#E8923F", "#F28E2B", "#76B7B2",
    "#EDC948", "#B07AA1", "#FF9DA7", "#9C755F", "#BAB0AC",
]

# Font sizes (OOXML hundredths or pt)
DATA_LABEL_SZ = 800           # %label inside bar segment
LINE_LABEL_SZ = 900           # line series numeric label
LEGEND_FONT_PT = 8
UNIT_FONT_PT = 8
TOTAL_FONT_PT = 10
CHART_TITLE_FONT_PT = 12
CHART_TITLE_HEX = "333333"
CHART_TITLE_BOLD = True
CHART_TITLE_ALIGN = PP_ALIGN.CENTER
SOURCE_FONT_PT = 10
LINE_FALLBACK_HEX = None       # roleup forces line color when data omits

_THEME = None


def _apply_theme(theme):
    global _THEME
    global SHAPE_SOURCE
    global PANEL_Y
    global SINGLE_CHART_X, SINGLE_CHART_W, LEFT_CHART_X, LEFT_CHART_W
    global RIGHT_CHART_X, RIGHT_CHART_W, CHART_H, CHART_BOTTOM
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_SOURCE, FONT_JP, TEXT_HEX
    global DEFAULT_COLORS
    global DATA_LABEL_SZ, LINE_LABEL_SZ, LEGEND_FONT_PT, UNIT_FONT_PT
    global TOTAL_FONT_PT, CHART_TITLE_FONT_PT, CHART_TITLE_HEX, CHART_TITLE_BOLD
    global CHART_TITLE_ALIGN, SOURCE_FONT_PT, LINE_FALLBACK_HEX

    _THEME = theme
    FONT_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    PANEL_Y = theme.layout("panel_y_in")
    SINGLE_CHART_X = theme.layout("single_chart_x_in")
    SINGLE_CHART_W = theme.layout("single_chart_w_in")
    LEFT_CHART_X = theme.layout("left_chart_x_in")
    LEFT_CHART_W = theme.layout("left_chart_w_in")
    RIGHT_CHART_X = theme.layout("right_chart_x_in")
    RIGHT_CHART_W = theme.layout("right_chart_w_in")
    CHART_H = theme.layout("chart_h_in")
    CHART_BOTTOM = theme.layout("chart_bottom_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    DEFAULT_COLORS = list(theme.chart_palette)

    if theme.id == "stellar_aiz":
        SHAPE_SOURCE = "Source"
        DATA_LABEL_SZ = 800     # 8pt V1
        LINE_LABEL_SZ = 900     # 9pt V1
        LEGEND_FONT_PT = 8
        UNIT_FONT_PT = 8
        TOTAL_FONT_PT = 10
        CHART_TITLE_FONT_PT = 12
        CHART_TITLE_HEX = TEXT_HEX
        CHART_TITLE_BOLD = True
        CHART_TITLE_ALIGN = PP_ALIGN.CENTER
        SOURCE_FONT_PT = 10
        LINE_FALLBACK_HEX = None
    else:
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        SHAPE_SOURCE = "Source 3"
        DATA_LABEL_SZ = 1000    # 10pt
        LINE_LABEL_SZ = 1000    # 10pt
        LEGEND_FONT_PT = theme.pt_value("font_size_body_pt")  # 10
        UNIT_FONT_PT = theme.pt_value("font_size_body_pt")    # 10
        TOTAL_FONT_PT = theme.pt_value("font_size_body_pt")   # 10
        CHART_TITLE_FONT_PT = theme.pt_value("font_size_subtitle_pt")  # 12
        CHART_TITLE_HEX = theme.hex_no_hash("subtitle")  # #897141
        CHART_TITLE_BOLD = False
        CHART_TITLE_ALIGN = PP_ALIGN.LEFT
        SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")  # 6
        LINE_FALLBACK_HEX = theme.hex_no_hash("accent_op_margin_line")  # #604C3F


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name, warn=True):
    for s in slide.shapes:
        if s.name == name:
            return s
    if warn:
        print(f"  ⚠ Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        r = etree.SubElement(p._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t = etree.SubElement(r, qn("a:t"))
        t.text = text


def write_source_placeholder(shape, text, font_size_pt, font_name):
    tf = shape.text_frame
    para = tf.paragraphs[0]
    for r in list(para.runs):
        r.text = ""
    if para.runs:
        run = para.runs[0]
        run.text = text
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text
        run = para.runs[0]
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        run._r.insert(0, rPr)
    rPr.set("sz", str(font_size_pt * 100))
    for tag in [qn("a:latin"), qn("a:ea")]:
        old = rPr.find(tag)
        if old is not None:
            rPr.remove(old)
        etree.SubElement(rPr, tag, attrib={"typeface": font_name})


def _hex2rgb(h):
    h = h.replace("#", "")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(slide, left, top, width, height, text, font_size=None,
                bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name=None):
    """汎用テキストボックス追加"""
    if font_size is None:
        font_size = Pt(11)
    if color is None:
        color = COLOR_TEXT
    if font_name is None:
        font_name = FONT_JP
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    r.font.size = font_size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return tb


def _add_custom_pct_labels(ser_xml, shares, min_share=2.0, font_color='FFFFFF'):
    old = ser_xml.find(qn('c:dLbls'))
    if old is not None:
        ser_xml.remove(old)

    dLbls = etree.SubElement(ser_xml, qn('c:dLbls'))

    for i, sv in enumerate(shares):
        dLbl = etree.SubElement(dLbls, qn('c:dLbl'))
        etree.SubElement(dLbl, qn('c:idx'), attrib={'val': str(i)})

        if sv < min_share:
            etree.SubElement(dLbl, qn('c:delete'), attrib={'val': '1'})
            continue

        tx = etree.SubElement(dLbl, qn('c:tx'))
        rich = etree.SubElement(tx, qn('c:rich'))
        etree.SubElement(rich, qn('a:bodyPr'))
        etree.SubElement(rich, qn('a:lstStyle'))
        p = etree.SubElement(rich, qn('a:p'))
        pp = etree.SubElement(p, qn('a:pPr'))
        etree.SubElement(pp, qn('a:defRPr'))
        r = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'), attrib={
            'lang': 'ja-JP', 'sz': str(DATA_LABEL_SZ), 'b': '1'
        })
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': font_color})
        etree.SubElement(rPr, qn('a:latin'), attrib={'typeface': FONT_JP})
        etree.SubElement(rPr, qn('a:ea'), attrib={'typeface': FONT_JP})
        t = etree.SubElement(r, qn('a:t'))
        t.text = f"{sv:.1f}%"

        etree.SubElement(dLbl, qn('c:dLblPos'), attrib={'val': 'ctr'})
        for k in ['showLegendKey', 'showCatName', 'showSerName', 'showPercent', 'showBubbleSize']:
            etree.SubElement(dLbl, qn(f'c:{k}'), attrib={'val': '0'})
        etree.SubElement(dLbl, qn('c:showVal'), attrib={'val': '0'})

    for k in ['showLegendKey', 'showCatName', 'showSerName', 'showPercent', 'showBubbleSize']:
        etree.SubElement(dLbls, qn(f'c:{k}'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showVal'), attrib={'val': '0'})


def _add_line_labels(ser_xml, fmt='0.0', font_color=None):
    if font_color is None:
        font_color = TEXT_HEX
    dl = ser_xml.find(qn('c:dLbls'))
    if dl is None:
        dl = etree.SubElement(ser_xml, qn('c:dLbls'))
    for c in list(dl):
        dl.remove(c)
    etree.SubElement(dl, qn('c:numFmt'), attrib={'formatCode': fmt + '"%"', 'sourceLinked': '0'})
    for k in ['showLegendKey', 'showCatName', 'showSerName', 'showPercent', 'showBubbleSize']:
        etree.SubElement(dl, qn(f'c:{k}'), attrib={'val': '0'})
    etree.SubElement(dl, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dl, qn('c:dLblPos'), attrib={'val': 't'})
    tp = etree.SubElement(dl, qn('c:txPr'))
    etree.SubElement(tp, qn('a:bodyPr'))
    etree.SubElement(tp, qn('a:lstStyle'))
    p = etree.SubElement(tp, qn('a:p'))
    pp = etree.SubElement(p, qn('a:pPr'))
    dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': str(LINE_LABEL_SZ), 'b': '1'})
    etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT_JP})
    etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT_JP})
    sf = etree.SubElement(dr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': font_color})


def build_chart(slide, chart_cfg, chart_left, chart_top, chart_w, chart_h):
    from pptx.chart.data import CategoryChartData

    data = chart_cfg.get("data", [])
    bars_cfg = chart_cfg.get("stacked_bars", [])
    line_cfg = chart_cfg.get("line", None)
    n_bars = len(bars_cfg)
    n_periods = len(data)
    has_line = line_cfg is not None

    cd = CategoryChartData()
    cd.categories = [d["year"] for d in data]
    for si, sb in enumerate(bars_cfg):
        cd.add_series(sb["series_name"],
                      [d["bars"][si] if si < len(d.get("bars", [])) else 0 for d in data])
    if has_line:
        cd.add_series(line_cfg["series_name"], [d.get("line_value", 0) for d in data])

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        int(chart_left), int(chart_top), int(chart_w), int(chart_h),
        cd
    )
    chart = cf.chart
    pa = chart._chartSpace.chart.plotArea
    bc = pa.findall(qn('c:barChart'))[0]

    sec_v = "2094734553"
    sec_c = "2094734554"

    if has_line:
        sers = bc.findall(qn('c:ser'))
        ls = copy.deepcopy(sers[-1])
        bc.remove(sers[-1])
        lc = etree.SubElement(pa, qn('c:lineChart'))
        etree.SubElement(lc, qn('c:grouping'), attrib={'val': 'standard'})
        etree.SubElement(lc, qn('c:varyColors'), attrib={'val': '0'})
        lc.append(ls)

        mk = ls.find(qn('c:marker'))
        if mk is None:
            mk = etree.SubElement(ls, qn('c:marker'))
        sym = mk.find(qn('c:symbol'))
        if sym is None:
            sym = etree.SubElement(mk, qn('c:symbol'))
        sym.set('val', 'circle')
        sz_el = mk.find(qn('c:size'))
        if sz_el is None:
            sz_el = etree.SubElement(mk, qn('c:size'))
        sz_el.set('val', str(line_cfg.get("marker_size", 7)))

        etree.SubElement(lc, qn('c:axId'), attrib={'val': sec_c})
        etree.SubElement(lc, qn('c:axId'), attrib={'val': sec_v})
        sc = etree.SubElement(pa, qn('c:catAx'))
        etree.SubElement(sc, qn('c:axId'), attrib={'val': sec_c})
        s2 = etree.SubElement(sc, qn('c:scaling'))
        etree.SubElement(s2, qn('c:orientation'), attrib={'val': 'minMax'})
        etree.SubElement(sc, qn('c:delete'), attrib={'val': '1'})
        etree.SubElement(sc, qn('c:axPos'), attrib={'val': 'b'})
        etree.SubElement(sc, qn('c:crossAx'), attrib={'val': sec_v})
        sv_ax = etree.SubElement(pa, qn('c:valAx'))
        etree.SubElement(sv_ax, qn('c:axId'), attrib={'val': sec_v})
        s3 = etree.SubElement(sv_ax, qn('c:scaling'))
        etree.SubElement(s3, qn('c:orientation'), attrib={'val': 'minMax'})
        etree.SubElement(sv_ax, qn('c:delete'), attrib={'val': '1'})
        etree.SubElement(sv_ax, qn('c:axPos'), attrib={'val': 'r'})
        etree.SubElement(sv_ax, qn('c:numFmt'), attrib={'formatCode': '0.0', 'sourceLinked': '0'})
        etree.SubElement(sv_ax, qn('c:crossAx'), attrib={'val': sec_c})
        etree.SubElement(sv_ax, qn('c:crosses'), attrib={'val': 'max'})

        raw_color = line_cfg.get("color") or (
            f"#{LINE_FALLBACK_HEX}" if LINE_FALLBACK_HEX else "#4E79A7"
        )
        lclr = raw_color.replace("#", "")
        lsp = ls.find(qn('c:spPr'))
        if lsp is None:
            lsp = etree.SubElement(ls, qn('c:spPr'))
        ln = lsp.find(qn('a:ln'))
        if ln is None:
            ln = etree.SubElement(lsp, qn('a:ln'))
        ln.set('w', '19050')
        lsf = ln.find(qn('a:solidFill'))
        if lsf is None:
            lsf = etree.SubElement(ln, qn('a:solidFill'))
        for c in list(lsf):
            lsf.remove(c)
        etree.SubElement(lsf, qn('a:srgbClr'), attrib={'val': lclr})
        msp = mk.find(qn('c:spPr'))
        if msp is None:
            msp = etree.SubElement(mk, qn('c:spPr'))
        msf = etree.SubElement(msp, qn('a:solidFill'))
        etree.SubElement(msf, qn('a:srgbClr'), attrib={'val': lclr})
        _add_line_labels(ls)

    bar_sers = bc.findall(qn('c:ser'))
    for si, bs in enumerate(bar_sers):
        clr = bars_cfg[si]["color"].replace("#", "") if si < len(bars_cfg) else "999999"
        sp = bs.find(qn('c:spPr'))
        if sp is None:
            sp = etree.SubElement(bs, qn('c:spPr'))
        sf = sp.find(qn('a:solidFill'))
        if sf is None:
            sf = etree.SubElement(sp, qn('a:solidFill'))
        for c in list(sf):
            sf.remove(c)
        etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': clr})

        shares = []
        for d in data:
            s_arr = d.get("shares", [])
            shares.append(s_arr[si] if si < len(s_arr) else 0)
        _add_custom_pct_labels(bs, shares)

    gw = bc.find(qn('c:gapWidth'))
    if gw is None:
        gw = etree.SubElement(bc, qn('c:gapWidth'))
    gw.set('val', '80')
    ov = bc.find(qn('c:overlap'))
    if ov is None:
        ov = etree.SubElement(bc, qn('c:overlap'))
    ov.set('val', '100')

    for vx in pa.findall(qn('c:valAx')):
        axId = vx.find(qn('c:axId'))
        if axId is not None and axId.get('val') == sec_v:
            continue
        de = vx.find(qn('c:delete'))
        if de is None:
            de = etree.SubElement(vx, qn('c:delete'))
        de.set('val', '1')

    for vx in pa.findall(qn('c:valAx')):
        for mg in vx.findall(qn('c:majorGridlines')):
            vx.remove(mg)
        for mg in vx.findall(qn('c:minorGridlines')):
            vx.remove(mg)
    for cx in pa.findall(qn('c:catAx')):
        for mg in cx.findall(qn('c:majorGridlines')):
            cx.remove(mg)

    for ax in pa.findall(qn('c:catAx')):
        de = ax.find(qn('c:delete'))
        if de is not None and de.get('val') == '1':
            continue
        tp = ax.find(qn('c:txPr'))
        if tp is None:
            tp = etree.SubElement(ax, qn('c:txPr'))
        bp = tp.find(qn('a:bodyPr'))
        if bp is None:
            bp = etree.SubElement(tp, qn('a:bodyPr'))
            tp.insert(0, bp)
        if tp.find(qn('a:lstStyle')) is None:
            etree.SubElement(tp, qn('a:lstStyle'))
        p = tp.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(tp, qn('a:p'))
        pp = p.find(qn('a:pPr'))
        if pp is None:
            pp = etree.SubElement(p, qn('a:pPr'))
        dr = pp.find(qn('a:defRPr'))
        if dr is None:
            dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': '1000'})
        else:
            dr.set('sz', '1000')
        if dr.find(qn('a:latin')) is None:
            etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT_JP})
        if dr.find(qn('a:ea')) is None:
            etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT_JP})
        dsf = dr.find(qn('a:solidFill'))
        if dsf is None:
            dsf = etree.SubElement(dr, qn('a:solidFill'))
        for c in list(dsf):
            dsf.remove(c)
        etree.SubElement(dsf, qn('a:srgbClr'), attrib={'val': TEXT_HEX})

    chart.has_legend = False
    chart.has_title = False

    psp = pa.find(qn('c:spPr'))
    if psp is None:
        psp = etree.SubElement(pa, qn('c:spPr'))
    nf = psp.find(qn('a:noFill'))
    if nf is None:
        etree.SubElement(psp, qn('a:noFill'))

    print(f"  ✓ 複合チャート: {n_periods}期, {n_bars}棒系列" + (" + 折れ線" if has_line else ""))
    return cf


def add_total_labels(slide, chart_cfg, cl, ct, cw, ch):
    data = chart_cfg.get("data", [])
    if not data:
        return
    nc = len(data)

    plm = cw * 0.06
    prm = cw * 0.04
    ptm = ch * 0.05
    pbm = ch * 0.14
    pl = cl + plm
    pw = cw - plm - prm
    pb = ct + ch - pbm
    ph = (ct + ch - pbm) - (ct + ptm)

    totals = [sum(d.get("bars", [])) for d in data]
    max_total = max(totals) if totals else 1
    axis_max = max_total * 1.25

    bar_w = pw / nc
    for i, d in enumerate(data):
        tl = d.get("total_label", "")
        if not tl:
            continue
        total = sum(d.get("bars", []))
        bar_cx = pl + (i + 0.5) * bar_w
        bar_top_y = pb - (total / axis_max) * ph

        lw = Inches(0.80)
        lh = Inches(0.25)
        lx = bar_cx - lw / 2
        ly = bar_top_y - lh - Inches(0.02)

        add_textbox(slide, lx, ly, lw, lh, tl,
                    font_size=Pt(TOTAL_FONT_PT), bold=True, color=COLOR_TEXT,
                    alignment=PP_ALIGN.CENTER)
    print(f"  ✓ トータルラベル: {nc}個")


def add_legend(slide, chart_cfg, left, top, max_width):
    bars = chart_cfg.get("stacked_bars", [])
    lc = chart_cfg.get("line", None)
    lh = Inches(0.20)
    ix = left
    row_y = top
    item_count = 0

    if lc:
        line_color = lc.get("color") or (
            f"#{LINE_FALLBACK_HEX}" if LINE_FALLBACK_HEX else "#4E79A7"
        )
        c = _hex2rgb(line_color)
        line_w = Inches(0.30)
        ly = int(row_y + lh / 2)
        cn = slide.shapes.add_connector(1, int(ix), ly, int(ix + line_w), ly)
        cn.line.color.rgb = c
        cn.line.width = Pt(1.5)
        cs = Inches(0.10)
        ci = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(ix + (line_w - cs) / 2), int(row_y + (lh - cs) / 2), int(cs), int(cs)
        )
        ci.fill.solid()
        ci.fill.fore_color.rgb = c
        ci.line.fill.background()
        tx = int(ix + line_w + Inches(0.04))
        name_w = Inches(min(len(lc["series_name"]) * 0.16, 2.4))
        tb = slide.shapes.add_textbox(tx, int(row_y), int(name_w), int(lh))
        tb.text_frame.word_wrap = False
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = lc["series_name"]
        r.font.size = Pt(LEGEND_FONT_PT)
        r.font.color.rgb = COLOR_TEXT
        r.font.name = FONT_JP
        ix = tx + name_w + Inches(0.12)
        item_count += 1

    for sb in bars:
        c = _hex2rgb(sb["color"])
        ss = Inches(0.12)
        # 10pt fonts need wider columns than 8pt; estimate generously.
        per_char_in = 0.13 if LEGEND_FONT_PT >= 10 else 0.11
        name_w = Inches(min(len(sb["series_name"]) * per_char_in + 0.15, 2.2))
        item_w = ss + Inches(0.04) + name_w + Inches(0.06)

        if ix + item_w > left + max_width and item_count > 0:
            row_y += lh + Inches(0.02)
            ix = left
            item_count = 0

        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(ix), int(row_y + (lh - ss) / 2), int(ss), int(ss)
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = c
        sq.line.fill.background()
        tx = int(ix + ss + Inches(0.04))
        tb = slide.shapes.add_textbox(tx, int(row_y), int(name_w), int(lh))
        tb.text_frame.word_wrap = False
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = sb["series_name"]
        r.font.size = Pt(LEGEND_FONT_PT)
        r.font.color.rgb = COLOR_TEXT
        r.font.name = FONT_JP
        ix = tx + name_w + Inches(0.06)
        item_count += 1

    print(f"  ✓ 凡例: {len(bars)}棒" + (" + 折れ線" if lc else ""))
    return row_y + lh


def render_one_chart(slide, chart_cfg, area_left, area_top, area_width):
    y = area_top

    title = chart_cfg.get("title", "")
    if title:
        add_textbox(slide, area_left, y, area_width, CHART_TITLE_H, title,
                    font_size=Pt(CHART_TITLE_FONT_PT), bold=CHART_TITLE_BOLD,
                    color=_hex2rgb(CHART_TITLE_HEX),
                    alignment=CHART_TITLE_ALIGN)
        # 下線 (装飾、テキスト無しのため C11 対象外)
        ul = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(area_left + area_width * 0.10), int(y + CHART_TITLE_H),
            int(area_width * 0.80), Inches(0.012)
        )
        ul.fill.solid()
        ul.fill.fore_color.rgb = _hex2rgb(CHART_TITLE_HEX)
        ul.line.fill.background()
    y += CHART_TITLE_H + Inches(0.05)

    unit = chart_cfg.get("unit_label", "")
    if unit:
        add_textbox(slide, area_left, y, area_width, UNIT_LABEL_H, unit,
                    font_size=Pt(UNIT_FONT_PT), bold=False, color=COLOR_SOURCE)
    y += UNIT_LABEL_H + Inches(0.02)

    legend_bottom = add_legend(slide, chart_cfg, area_left, y, area_width)
    y = legend_bottom + Inches(0.08)

    available_h = CHART_BOTTOM - y
    chart_h = min(available_h, CHART_H)

    build_chart(slide, chart_cfg, area_left, y, area_width, chart_h)
    add_total_labels(slide, chart_cfg, area_left, y, area_width, chart_h)

    print(f"  ✓ チャート領域完了: '{title}'")


def _apply_default_colors(chart_cfg):
    bars = chart_cfg.get("stacked_bars", [])
    for i, sb in enumerate(bars):
        if "color" not in sb or not sb["color"]:
            sb["color"] = DEFAULT_COLORS[i % len(DEFAULT_COLORS)]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "cost-breakdown")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "charts"],
        allowed_top=[
            "main_message", "chart_title", "source", "charts",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    print(f"=== コスト内訳推移スライド生成 (brand={theme.id}) ===")
    print(f"  Template: {template_path}")

    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Roleup: silently remove brown guide rectangles
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # Top placeholder (brand-aware)
    # cost-breakdown は data に chart_title がない場合も多いので charts[0].title へフォールバック
    charts = data.get("charts", [])
    fallback_top = charts[0].get("title", "コスト内訳推移") if charts else "コスト内訳推移"

    if theme.id == "stellar_aiz":
        top_text = data.get("main_message", "")
        sub_text = data.get("chart_title", "")
    else:
        top_text = data.get("chart_title") or fallback_top
        sub_text = data.get("main_message", "")

    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)

    # Source: stella=Source textbox / roleup=Source 3 placeholder
    src = data.get("source", "")
    if src:
        body = src if src.startswith("出典") else f"出典：{src}"
        src_shape = find_shape(slide, SHAPE_SOURCE, warn=False)
        if src_shape is not None:
            if theme.id == "stellar_aiz":
                set_textbox_text(src_shape, body)
            else:
                write_source_placeholder(src_shape, body, SOURCE_FONT_PT, FONT_JP)
            print(f"  Source ({SHAPE_SOURCE}): {body[:50]}")

    # Stella: Content Area placeholder removal (roleup template doesn't have it)
    _silent_remove_shape(slide, SHAPE_CONTENT_AREA)

    n_charts = len(charts)

    if n_charts == 0:
        print("  ⚠ chartsが空です")
    elif n_charts == 1:
        print("\n── 1チャートモード ──")
        _apply_default_colors(charts[0])
        render_one_chart(slide, charts[0], SINGLE_CHART_X, PANEL_Y, SINGLE_CHART_W)
    else:
        print("\n── 2チャートモード ──")
        _apply_default_colors(charts[0])
        _apply_default_colors(charts[1])
        print("\n[左チャート]")
        render_one_chart(slide, charts[0], LEFT_CHART_X, PANEL_Y, LEFT_CHART_W)
        print("\n[右チャート]")
        render_one_chart(slide, charts[1], RIGHT_CHART_X, PANEL_Y, RIGHT_CHART_W)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
