"""
fill_process_arrow.py — プロセス矢羽データをPPTXテンプレートに流し込むスクリプト

テンプレート:
  - ProcessArrow3.pptx (3矢羽)
  - ProcessArrow5.pptx (5矢羽)

使い方:
  python fill_process_arrow.py \
    --data /home/claude/arrow_data.json \
    --template <path-to-template>.pptx \
    --output /mnt/user-data/outputs/ProcessArrow_output.pptx
"""

import argparse
import json
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ────────────────────────────────────────

# 共通
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_LABEL1       = "TextBox 6"
SHAPE_LABEL2       = "TextBox 7"

# 矢羽Shape名 (Pentagon)
ARROW_SHAPES_3 = ["Pentagon 3", "Pentagon 15", "Pentagon 16"]
ARROW_SHAPES_5 = ["Pentagon 3", "Pentagon 15", "Pentagon 16", "Pentagon 20", "Pentagon 21"]

# Label1テキストボックス (各矢羽の上段詳細)
LABEL1_SHAPES_3 = ["TextBox 8", "TextBox 10", "TextBox 12"]
LABEL1_SHAPES_5 = ["TextBox 8", "TextBox 10", "TextBox 12", "TextBox 5", "TextBox 17"]

# Label2テキストボックス (各矢羽の下段詳細)
LABEL2_SHAPES_3 = ["TextBox 9", "TextBox 11", "TextBox 13"]
LABEL2_SHAPES_5 = ["TextBox 9", "TextBox 11", "TextBox 13", "TextBox 14", "TextBox 19"]

# ────────────────────────────────────────────────────────────


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_text(shape, text):
    """TextFrame の最初の段落にテキストをセット（既存スタイル保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def detect_arrow_count(data):
    """JSONのarrowsフィールドの数から3 or 5を判定"""
    arrows = data.get("arrows", [])
    return len(arrows)


def main():
    parser = argparse.ArgumentParser(description="プロセス矢羽データをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="arrow_data.json のパス")
    parser.add_argument("--template", required=True, help="テンプレートPPTX のパス")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(args.template)
    slide = prs.slides[0]

    # ── 矢羽数を判定 ───────────────────────────────────────
    arrows = data.get("arrows", [])
    n = len(arrows)
    if n == 3:
        arrow_shapes = ARROW_SHAPES_3
        label1_shapes = LABEL1_SHAPES_3
        label2_shapes = LABEL2_SHAPES_3
    elif n == 5:
        arrow_shapes = ARROW_SHAPES_5
        label1_shapes = LABEL1_SHAPES_5
        label2_shapes = LABEL2_SHAPES_5
    else:
        print(f"  ❌ ERROR: arrows の数が {n} です（3 または 5 のみ対応）", file=sys.stderr)
        sys.exit(1)

    # ── Main Message ───────────────────────────────────────
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        set_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # ── Chart Title ────────────────────────────────────────
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        set_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # ── Label1 / Label2 ラベル名 ──────────────────────────
    label1 = data.get("label1", "").strip()
    label2 = data.get("label2", "").strip()
    if label1:
        set_text(find_shape(slide, SHAPE_LABEL1), label1)
        print(f"  [Label1] {label1}")
    if label2:
        set_text(find_shape(slide, SHAPE_LABEL2), label2)
        print(f"  [Label2] {label2}")

    # ── 矢羽テキスト + 詳細 ───────────────────────────────
    for i, arrow in enumerate(arrows):
        name = arrow.get("name", "").strip()
        detail1 = arrow.get("label1_detail", "").strip()
        detail2 = arrow.get("label2_detail", "").strip()

        # 矢羽名
        set_text(find_shape(slide, arrow_shapes[i]), name)
        print(f"  [Arrow{i+1}] {name}")

        # Label1 詳細
        if detail1:
            set_text(find_shape(slide, label1_shapes[i]), detail1)
            print(f"    Label1: {detail1[:50]}{'...' if len(detail1) > 50 else ''}")

        # Label2 詳細
        if detail2:
            set_text(find_shape(slide, label2_shapes[i]), detail2)
            print(f"    Label2: {detail2[:50]}{'...' if len(detail2) > 50 else ''}")

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
