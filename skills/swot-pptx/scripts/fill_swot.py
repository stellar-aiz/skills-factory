"""
fill_swot.py — SWOT分析スライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 2×2 SWOTマトリクス:
      [S:強み  ] [W:弱み]   ← 内部要因（上段）
      [O:機会  ] [T:脅威]   ← 外部要因（下段）
  - 下部: 出典

各象限:
  - 上部にカラーヘッダーバー（象限ラベル + 英語名）
  - 下部にブレット項目リスト

Usage:
  python fill_swot.py \
    --data /home/claude/swot_data.json \
    --template <path>/swot-template.pptx \
    --output /mnt/user-data/outputs/SWOT_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "swot-pptx"
_THEME = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants (stella 16:9, 13.33 x 7.5 in; _apply_theme で roleup A4 用にスケール) ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_PH = "Source 3"  # roleup placeholder

GRID_LEFT = Inches(0.41)
GRID_TOP = Inches(1.55)
GRID_WIDTH = Inches(12.51)
GRID_HEIGHT = Inches(5.30)
GAP = Inches(0.15)
CELL_W = (GRID_WIDTH - GAP) / 2
CELL_H = (GRID_HEIGHT - GAP) / 2
HEADER_H = Inches(0.55)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(6.93)
SOURCE_W = Inches(12.50)

# ── Colors (stella defaults; _apply_theme で roleup 用に上書き) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BODY_BG = RGBColor(0xFA, 0xFA, 0xFA)

COLOR_S = RGBColor(0x2E, 0x4A, 0x6B)
COLOR_W = RGBColor(0xE5, 0x7C, 0x52)
COLOR_O = RGBColor(0x5B, 0x8A, 0x3A)
COLOR_T = RGBColor(0xB8, 0x3A, 0x3A)

COLOR_S_LIGHT = RGBColor(0xE8, 0xEE, 0xF4)
COLOR_W_LIGHT = RGBColor(0xFA, 0xEC, 0xE4)
COLOR_O_LIGHT = RGBColor(0xEB, 0xF0, 0xE4)
COLOR_T_LIGHT = RGBColor(0xF6, 0xE2, 0xE2)

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_HEADER = Pt(16)
FONT_SIZE_HEADER_EN = Pt(11)
FONT_SIZE_ITEM = Pt(12)
FONT_SIZE_SOURCE = Pt(10)
FONT_SIZE_AXIS_LABEL = Pt(10)


def _apply_theme(theme):
    """roleup の場合、レイアウト・色・フォントサイズを brand 仕様に上書きする。"""
    global GRID_LEFT, GRID_TOP, GRID_WIDTH, GRID_HEIGHT, GAP, CELL_W, CELL_H
    global HEADER_H, SOURCE_X, SOURCE_Y, SOURCE_W
    global COLOR_TEXT, COLOR_SOURCE
    global COLOR_S, COLOR_W, COLOR_O, COLOR_T
    global COLOR_S_LIGHT, COLOR_W_LIGHT, COLOR_O_LIGHT, COLOR_T_LIGHT
    global FONT_NAME_JP, FONT_SIZE_HEADER, FONT_SIZE_HEADER_EN
    global FONT_SIZE_ITEM, FONT_SIZE_SOURCE, FONT_SIZE_AXIS_LABEL
    global _THEME
    _THEME = theme

    if theme.id != "roleup":
        return

    # A4 横 (11.69 × 8.27) 用にレイアウト再計算
    GRID_LEFT = Inches(0.41)
    GRID_TOP = Inches(1.55)
    GRID_WIDTH = Inches(10.87)
    GRID_HEIGHT = Inches(5.70)
    GAP = Inches(0.15)
    CELL_W = (GRID_WIDTH - GAP) / 2
    CELL_H = (GRID_HEIGHT - GAP) / 2
    HEADER_H = Inches(0.50)
    SOURCE_X = Inches(0.41)
    SOURCE_Y = Inches(7.45)
    SOURCE_W = Inches(10.87)

    # roleup 茶系トーン (SWOT 4 色は universal indicator なのでヘッダー色は維持、
    # light bg は brand 寄りに置換)
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")

    # フォント
    FONT_NAME_JP = theme._defaults.get("font_name_ja", "Yu Gothic UI")

    # roleup C4 許容集合 [22, 14, 12, 10, 6] pt
    FONT_SIZE_HEADER = Pt(14)
    FONT_SIZE_HEADER_EN = Pt(10)
    FONT_SIZE_ITEM = Pt(10)
    FONT_SIZE_SOURCE = Pt(int(theme._defaults.get("font_size_source_pt", 6)))
    FONT_SIZE_AXIS_LABEL = Pt(10)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    """汎用テキストボックス (font_name は late-resolve で _apply_theme 後の値を使う)"""
    if font_name is None:
        font_name = FONT_NAME_JP
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


# ──────────────────────────────────────────────
# SWOT Quadrant Builder
# ──────────────────────────────────────────────
def build_quadrant(slide, label_jp, label_en, items, header_color, body_color,
                   left, top, width, height):
    """
    1つの象限を描画する。
      - 上部: カラーヘッダーバー（ラベル表示）
      - 下部: 薄い背景 + ブレット項目リスト
    """
    # 全体の外枠（uniform border）
    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = body_color
    outer.line.color.rgb = header_color
    outer.line.width = Pt(0.75)
    outer.shadow.inherit = False
    # 文字を空に
    outer.text_frame.text = ""

    # ヘッダーバー（上部）
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, HEADER_H
    )
    header.fill.solid()
    header.fill.fore_color.rgb = header_color
    header.line.fill.background()
    header.shadow.inherit = False

    tf = header.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 既存段落をクリア
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p = tf._txBody.makeelement(qn("a:p"), {})
    tf._txBody.append(p)
    pPr = etree.SubElement(p, qn("a:pPr"))
    pPr.set("algn", "l")

    # Run 1: 日本語ラベル (大)
    r1 = etree.SubElement(p, qn("a:r"))
    rPr1 = etree.SubElement(r1, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(FONT_SIZE_HEADER.pt * 100)),
        "b": "1",
    })
    etree.SubElement(rPr1, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr1, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf1 = etree.SubElement(rPr1, qn("a:solidFill"))
    s1 = etree.SubElement(sf1, qn("a:srgbClr"))
    s1.set("val", "FFFFFF")
    t1 = etree.SubElement(r1, qn("a:t"))
    t1.text = f"{label_jp}  "

    # Run 2: 英語ラベル (小)
    r2 = etree.SubElement(p, qn("a:r"))
    rPr2 = etree.SubElement(r2, qn("a:rPr"), attrib={
        "lang": "en-US",
        "sz": str(int(FONT_SIZE_HEADER_EN.pt * 100)),
        "b": "0",
    })
    etree.SubElement(rPr2, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr2, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf2 = etree.SubElement(rPr2, qn("a:solidFill"))
    s2 = etree.SubElement(sf2, qn("a:srgbClr"))
    s2.set("val", "FFFFFF")
    t2 = etree.SubElement(r2, qn("a:t"))
    t2.text = f"({label_en})"

    # ボディ部分: ブレット項目リスト
    body_top = top + HEADER_H
    body_h = height - HEADER_H

    body_box = slide.shapes.add_textbox(
        left + Inches(0.15), body_top + Inches(0.12),
        width - Inches(0.30), body_h - Inches(0.20),
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.margin_left = 0
    body_tf.margin_right = 0
    body_tf.margin_top = 0
    body_tf.margin_bottom = 0
    body_tf.vertical_anchor = MSO_ANCHOR.TOP

    # 既存段落クリア
    for p in list(body_tf.paragraphs):
        p._p.getparent().remove(p._p)

    # 各項目をブレット形式の段落として追加
    for i, item_text in enumerate(items):
        p_elem = etree.SubElement(body_tf._txBody, qn("a:p"))

        # 段落プロパティ: インデント、行間
        pPr = etree.SubElement(p_elem, qn("a:pPr"), attrib={
            "marL": "180000",   # 左マージン
            "indent": "-180000",  # ハンギングインデント
        })
        # 段落間スペース
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "300"})

        # Bullet character
        buChar = etree.SubElement(pPr, qn("a:buChar"), attrib={"char": "•"})
        buFont = etree.SubElement(pPr, qn("a:buFont"), attrib={"typeface": "Arial"})
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
        buClrSolid.set("val", "{:02X}{:02X}{:02X}".format(
            header_color[0], header_color[1], header_color[2]
        ))

        # Run
        r = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", "333333")
        t = etree.SubElement(r, qn("a:t"))
        t.text = item_text

    print(f"  ✓ 象限 [{label_jp}]: {len(items)}項目")


# ──────────────────────────────────────────────
# Axis Labels (optional small labels around grid)
# ──────────────────────────────────────────────
def add_axis_labels(slide):
    """グリッドの軸ラベル（内部/外部、プラス/マイナス）を追加"""
    # 上部：左半分「内部要因」、右半分「外部要因」
    # (ヘッダーバーにラベルが出ているので、軸ラベルは省略可能。今回は省略。)
    pass


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=False, default=None)
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "swot"],
        allowed_top=[
            "main_message", "chart_title", "source", "swot",
            "title", "subtitle",
        ],
        nested_required={"swot": ["strengths", "weaknesses", "opportunities", "threats"]},
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "swot")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top text (stella: main_message / roleup: chart_title)
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
        print(f"  ✓ Top: {top_text[:60]}{'...' if len(top_text) > 60 else ''}")

    # Subtitle (stella: chart_title / roleup: main_message)
    sub_text = resolve_subtitle_text(data, theme).strip()
    if sub_text:
        set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
        print(f"  ✓ Subtitle: {sub_text[:60]}{'...' if len(sub_text) > 60 else ''}")

    swot = data.get("swot", {})

    # 4象限のコンフィグ
    # レイアウト:
    #   [S:強み]     [W:弱み]     ← 上段（内部要因）
    #   [O:機会]     [T:脅威]     ← 下段（外部要因）
    quadrants = [
        {
            "data": swot.get("strengths", {}),
            "label_jp": "強み",
            "label_en": "Strengths",
            "header_color": COLOR_S,
            "body_color": COLOR_S_LIGHT,
            "left": GRID_LEFT,
            "top": GRID_TOP,
        },
        {
            "data": swot.get("weaknesses", {}),
            "label_jp": "弱み",
            "label_en": "Weaknesses",
            "header_color": COLOR_W,
            "body_color": COLOR_W_LIGHT,
            "left": GRID_LEFT + CELL_W + GAP,
            "top": GRID_TOP,
        },
        {
            "data": swot.get("opportunities", {}),
            "label_jp": "機会",
            "label_en": "Opportunities",
            "header_color": COLOR_O,
            "body_color": COLOR_O_LIGHT,
            "left": GRID_LEFT,
            "top": GRID_TOP + CELL_H + GAP,
        },
        {
            "data": swot.get("threats", {}),
            "label_jp": "脅威",
            "label_en": "Threats",
            "header_color": COLOR_T,
            "body_color": COLOR_T_LIGHT,
            "left": GRID_LEFT + CELL_W + GAP,
            "top": GRID_TOP + CELL_H + GAP,
        },
    ]

    for q in quadrants:
        # ラベルの上書き（dataにカスタムラベルがあればそれを使う）
        label_jp = q["data"].get("label_jp", q["label_jp"])
        label_en = q["data"].get("label_en", q["label_en"])
        items = q["data"].get("items", [])
        build_quadrant(
            slide, label_jp, label_en, items,
            q["header_color"], q["body_color"],
            q["left"], q["top"], CELL_W, CELL_H,
        )

    # 出典 (roleup: Source 3 placeholder, stella: 動的 textbox)
    source = (data.get("source") or data.get("source_label")
              or data.get("source_text") or "").strip()
    if source:
        if theme.is_source_required():
            src_shape = find_shape(slide, SHAPE_SOURCE_PH)
            if src_shape is not None:
                set_textbox_text(src_shape, f"出典: {source}")
                for para in src_shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = FONT_SIZE_SOURCE
                        run.font.color.rgb = COLOR_SOURCE
                        run.font.name = FONT_NAME_JP
        else:
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.30),
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
