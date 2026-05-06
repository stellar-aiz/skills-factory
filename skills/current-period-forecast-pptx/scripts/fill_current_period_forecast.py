"""
fill_current_period_forecast.py — 当期着地見込みテーブルをHTMLで描画→スクリーンショット→PPTXに挿入するスクリプト

BDD（Business Due Diligence）向け。
マネジメント計画（Base）と弊社計画（Downside / Upside）の財務見込みを比較するスライドを生成する。

テンプレート構造（current-period-forecast-template.pptx = company-overviewベース）:
  - Title 1              (PLACEHOLDER): Main Message
  - Text Placeholder 2   (PLACEHOLDER): Chart Title
  - Content Area         (AUTO_SHAPE):  HTML screenshot挿入先
  - Source / Source 3    (TEXT_BOX/PLACEHOLDER): 出典

使い方:
  python fill_current_period_forecast.py \
    --data /home/claude/forecast_data.json \
    --output /mnt/user-data/outputs/CurrentPeriodForecast_output.pptx \
    --brand stellar_aiz
"""

import argparse
import asyncio
import json
import os
import sys
import tempfile
from html import escape as _esc

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── brand_resolver bootstrap ─────────────────────────────────
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "current-period-forecast-pptx"


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape名マッピング ──────────────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE_CANDIDATES = ("Source 3", "Source", "PlaceHolder 3")

DEVICE_SCALE = 2

# ── Theme-controlled module variables (stella defaults) ──
FONT_NAME_JP = "Noto Sans CJK JP, Meiryo UI, Hiragino Sans"
TEXT_HEX = "1A1A2E"
SOURCE_HEX = "666666"
HEADER_BG_HEX = "D9D9D9"
HEADER_DARK_BG_HEX = "BFBFBF"
ROW_EVEN_BG_HEX = "FFFFFF"
ROW_ODD_BG_HEX = "F5F5F0"
BORDER_HEX = "999999"
PCT_TEXT_HEX = "666666"

SOURCE_FONT_PT = 10
TITLE_FONT_PT = 28
SUBTITLE_FONT_PT = 14


def _apply_theme(theme):
    global FONT_NAME_JP
    global TEXT_HEX, SOURCE_HEX, HEADER_BG_HEX, HEADER_DARK_BG_HEX
    global ROW_ODD_BG_HEX, BORDER_HEX, PCT_TEXT_HEX
    global SOURCE_FONT_PT, TITLE_FONT_PT, SUBTITLE_FONT_PT
    FONT_NAME_JP = theme.font_ea
    TEXT_HEX = theme.hex_no_hash("text")
    SOURCE_HEX = theme.hex_no_hash("source")
    HEADER_BG_HEX = theme.hex_no_hash("label_bg")
    HEADER_DARK_BG_HEX = theme.hex_no_hash("highlight_other")
    ROW_ODD_BG_HEX = theme.hex_no_hash("header_bg")
    BORDER_HEX = theme.hex_no_hash("highlight_other")
    PCT_TEXT_HEX = theme.hex_no_hash("source")
    SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")
    TITLE_FONT_PT = theme.pt_value("font_size_title_pt")
    SUBTITLE_FONT_PT = theme.pt_value("font_size_subtitle_pt")


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def find_source_shape(slide):
    for cand in SHAPE_SOURCE_CANDIDATES:
        s = find_shape(slide, cand)
        if s is not None:
            return s
    return None


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def _make_brand_run(para, text_str, font_pt=None, bold=False, color_hex=None):
    for r in list(para._p.findall(qn("a:r"))):
        para._p.remove(r)

    r_elem = etree.SubElement(para._p, qn("a:r"))
    rPr_attrs = {"lang": "ja-JP"}
    if bold:
        rPr_attrs["b"] = "1"
    if font_pt is not None:
        rPr_attrs["sz"] = str(int(font_pt * 100))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib=rPr_attrs)
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color_hex or TEXT_HEX})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text_str


def set_textbox_text(shape, text, font_pt=None, color_hex=None):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    bold = False
    if para.runs:
        existing_rPr = para.runs[0]._r.find(qn("a:rPr"))
        if existing_rPr is not None and existing_rPr.get("b") == "1":
            bold = True
    _make_brand_run(para, text, font_pt=font_pt, bold=bold, color_hex=color_hex)


def generate_html(data, viewport_width, viewport_height):
    """当期着地見込みデータからHTMLテーブルを生成する"""
    period_label = _esc(data.get("period_label", "進行期 見込み"))
    unit_label = _esc(data.get("unit_label", "単位：百万円/(%)"))
    mgmt_label = _esc(data.get("management_plan_label", "マネジメント計画\n(Base)"))
    our_plan_label = _esc(data.get("our_plan_label", "弊社計画"))
    downside_label = _esc(data.get("downside_label", "Downside case"))
    upside_label = _esc(data.get("upside_label", "Upside case"))
    item_header = _esc(data.get("item_header", "項目（調整後）"))

    rows = data.get("rows", [])
    n_rows = len(rows)

    PT = viewport_width / 900.0  # CSS px per pt

    if n_rows <= 5:
        fs_header = int(11 * PT)
        fs_value = int(12 * PT)
        fs_pct = int(8 * PT)
        fs_assumption = int(9 * PT)
        row_pad = int(3 * PT)
    elif n_rows <= 8:
        fs_header = int(10 * PT)
        fs_value = int(11 * PT)
        fs_pct = int(7.5 * PT)
        fs_assumption = int(8 * PT)
        row_pad = int(2 * PT)
    else:
        fs_header = int(9 * PT)
        fs_value = int(9 * PT)
        fs_pct = int(6.5 * PT)
        fs_assumption = int(7 * PT)
        row_pad = int(1.5 * PT)

    HEADER_BG = f"#{HEADER_BG_HEX}"
    HEADER_DARK_BG = f"#{HEADER_DARK_BG_HEX}"
    ROW_EVEN_BG = f"#{ROW_EVEN_BG_HEX}"
    ROW_ODD_BG = f"#{ROW_ODD_BG_HEX}"
    BORDER_COLOR = f"#{BORDER_HEX}"
    TEXT_COLOR = f"#{TEXT_HEX}"
    PCT_COLOR = f"#{PCT_TEXT_HEX}"

    mgmt_label_html = mgmt_label.replace("\n", "<br>")

    header_html = f"""
    <thead>
        <tr>
            <td rowspan="2" class="cell header-cell item-col"
                style="background:{HEADER_DARK_BG};font-weight:700;font-size:{fs_header}px;
                       color:{TEXT_COLOR};text-align:center;vertical-align:middle;">
                <div style="font-size:{int(7*PT)}px;color:{PCT_COLOR};text-align:left;margin-bottom:{int(1*PT)}px;">
                    {unit_label}
                </div>
                {item_header}
            </td>
            <td rowspan="2" class="cell header-cell mgmt-col"
                style="background:{HEADER_BG};font-weight:700;font-size:{fs_header}px;
                       color:{TEXT_COLOR};text-align:center;vertical-align:middle;line-height:1.4;">
                {mgmt_label_html}
            </td>
            <td colspan="4" class="cell header-cell"
                style="background:{HEADER_BG};font-weight:700;font-size:{fs_header}px;
                       color:{TEXT_COLOR};text-align:center;vertical-align:middle;">
                {our_plan_label}
            </td>
        </tr>
        <tr>
            <td colspan="2" class="cell header-cell"
                style="background:{HEADER_BG};font-weight:700;font-size:{fs_header}px;
                       color:{TEXT_COLOR};text-align:center;vertical-align:middle;">
                {downside_label}
            </td>
            <td colspan="2" class="cell header-cell"
                style="background:{HEADER_BG};font-weight:700;font-size:{fs_header}px;
                       color:{TEXT_COLOR};text-align:center;vertical-align:middle;">
                {upside_label}
            </td>
        </tr>
    </thead>"""

    body_rows = ""
    for i, row in enumerate(rows):
        bg = ROW_ODD_BG if i % 2 == 0 else ROW_EVEN_BG
        item_name = _esc(row.get("item", ""))

        mgmt_val = _esc(str(row.get("management_value", "")))
        mgmt_pct = _esc(str(row.get("management_pct", "")))
        ds_val = _esc(str(row.get("downside_value", "")))
        ds_pct = _esc(str(row.get("downside_pct", "")))
        ds_assumption = _esc(str(row.get("downside_assumption", "")))
        us_val = _esc(str(row.get("upside_value", "")))
        us_pct = _esc(str(row.get("upside_pct", "")))
        us_assumption = _esc(str(row.get("upside_assumption", "")))

        def val_cell(val, pct):
            pct_html = (f'<div style="font-size:{fs_pct}px;color:{PCT_COLOR};">({pct})</div>'
                        if pct else "")
            return (f'<div style="font-size:{fs_value}px;font-weight:700;color:{TEXT_COLOR};'
                    f'text-align:center;">{val}</div>{pct_html}')

        def assumption_cell(text):
            if not text:
                return ""
            return (f'<div style="font-size:{fs_assumption}px;color:{TEXT_COLOR};'
                    f'line-height:1.4;text-align:left;">■ {text}</div>')

        body_rows += f"""
        <tr>
            <td class="cell item-col" style="background:{bg};font-weight:700;
                color:{TEXT_COLOR};font-size:{fs_value}px;text-align:center;vertical-align:middle;
                padding:{row_pad}px {int(4*PT)}px;">
                {item_name}
            </td>
            <td class="cell mgmt-col" style="background:{bg};vertical-align:middle;
                padding:{row_pad}px {int(4*PT)}px;text-align:center;">
                {val_cell(mgmt_val, mgmt_pct)}
            </td>
            <td class="cell ds-val-col" style="background:{bg};vertical-align:middle;
                padding:{row_pad}px {int(4*PT)}px;text-align:center;">
                {val_cell(ds_val, ds_pct)}
            </td>
            <td class="cell ds-assumption-col" style="background:{bg};vertical-align:middle;
                padding:{row_pad}px {int(6*PT)}px;">
                {assumption_cell(ds_assumption)}
            </td>
            <td class="cell us-val-col" style="background:{bg};vertical-align:middle;
                padding:{row_pad}px {int(4*PT)}px;text-align:center;">
                {val_cell(us_val, us_pct)}
            </td>
            <td class="cell us-assumption-col" style="background:{bg};vertical-align:middle;
                padding:{row_pad}px {int(6*PT)}px;">
                {assumption_cell(us_assumption)}
            </td>
        </tr>"""

    body_pad = int(4 * PT)
    border_style = f"1px solid {BORDER_COLOR}"

    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
html, body {{
    width:{viewport_width}px;
    height:{viewport_height}px;
    font-family:'{FONT_NAME_JP}','Noto Sans CJK JP','Meiryo UI','Hiragino Sans',sans-serif;
    background:#FFFFFF;
    padding:{body_pad}px;
    overflow:hidden;
}}
table {{
    width:100%;
    height:100%;
    border-collapse:collapse;
    table-layout:fixed;
}}
.cell {{ border:{border_style}; }}
.item-col {{ width:10%; }}
.mgmt-col {{ width:8%; }}
.ds-val-col {{ width:8%; }}
.ds-assumption-col {{ width:32%; }}
.us-val-col {{ width:8%; }}
.us-assumption-col {{ width:34%; }}
</style>
</head>
<body>
<table>
{header_html}
<tbody>
{body_rows}
</tbody>
</table>
</body></html>'''

    return html


async def take_screenshot(html_content, output_path, vw, vh):
    from playwright.async_api import async_playwright

    with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8") as f:
        f.write(html_content)
        html_path = f.name

    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page(
                viewport={"width": vw, "height": vh},
                device_scale_factor=DEVICE_SCALE,
            )
            await page.goto(f"file://{html_path}", wait_until="networkidle")
            await page.wait_for_timeout(500)
            await page.screenshot(path=output_path, full_page=False)
            await browser.close()
    finally:
        os.unlink(html_path)

    print(f"  Screenshot saved: {output_path}")


def insert_screenshot_into_pptx(template_path, data, screenshot_path, output_path):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    for s in slide.shapes:
        print(f"  Shape: '{s.name}' type={s.shape_type}")

    main_msg = data.get("main_message", "")
    msg_shape = find_shape(slide, SHAPE_MAIN_MESSAGE) or find_shape(slide, "PlaceHolder 1")
    set_textbox_text(msg_shape, main_msg, font_pt=TITLE_FONT_PT, color_hex=TEXT_HEX)

    title_text = data.get("chart_title", "")
    title_shape = find_shape(slide, SHAPE_CHART_TITLE) or find_shape(slide, "PlaceHolder 2")
    set_textbox_text(title_shape, title_text, font_pt=SUBTITLE_FONT_PT, color_hex=TEXT_HEX)

    source_text = data.get("source", "受領資料、Q＆A、マネジメントインタビューより当社作成")
    source_shape = find_source_shape(slide)
    if source_shape is not None and source_text:
        set_textbox_text(source_shape, f"出典：{source_text}",
                         font_pt=SOURCE_FONT_PT, color_hex=SOURCE_HEX)

    content_shape = find_shape(slide, SHAPE_CONTENT_AREA)
    if content_shape and os.path.exists(screenshot_path):
        left = content_shape.left
        top = content_shape.top
        width = content_shape.width
        height = content_shape.height
        slide.shapes.add_picture(screenshot_path, left, top, width, height)
        print(f"  Screenshot inserted at ({left},{top}) size=({width},{height})")
        _silent_remove_shape(slide, SHAPE_CONTENT_AREA)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    _finalize_pptx(output_path)
    print(f"  PPTX saved: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="当期着地見込みスライド生成。--brand stellar_aiz / roleup を選択。",
    )
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument("--template", required=False, default=None, help="PPTXテンプレートパス (任意)")
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "current-period-forecast")
    print(f"  ✓ Brand:    {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "rows"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "rows", "item_header", "period_label", "unit_label",
            "management_plan_label", "our_plan_label",
            "upside_label", "downside_label",
            "title", "subtitle",
        ],
        skill_name=SKILL_ID,
    )

    print("=== 当期着地見込みスライド生成 ===")

    # Content Area の実寸から viewport を逆算
    prs_probe = Presentation(template_path)
    content = find_shape(prs_probe.slides[0], SHAPE_CONTENT_AREA)
    if content is not None:
        vw = int(content.width / 914400.0 * 200)
        vh = int(content.height / 914400.0 * 200)
    else:
        vw, vh = 2504, 1080

    print(f"Step 1: HTML生成... (viewport {vw}×{vh})")
    html_content = generate_html(data, vw, vh)

    html_debug = os.path.join(tempfile.gettempdir(), "forecast_debug.html")
    with open(html_debug, "w", encoding="utf-8") as f:
        f.write(html_content)

    print("Step 2: スクリーンショット取得...")
    screenshot_path = os.path.join(tempfile.gettempdir(), "forecast_screenshot.png")
    asyncio.run(take_screenshot(html_content, screenshot_path, vw, vh))

    print("Step 3: PPTX生成...")
    insert_screenshot_into_pptx(template_path, data, screenshot_path, args.output)

    if os.path.exists(screenshot_path):
        os.unlink(screenshot_path)

    print("=== 完了 ===")


if __name__ == "__main__":
    main()
