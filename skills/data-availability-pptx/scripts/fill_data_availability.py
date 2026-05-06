"""
fill_data_availability.py — データアベイラビリティ（調査の網羅度・制約）スライドを生成

レイアウト:
  - 上部: メインメッセージ (Title 1) + チャートタイトル (Text Placeholder 2)
  - 左側: カバレッジテーブル (カテゴリ × 項目 × ステータス × データソース)
       ステータス: ✓取得済 / △一部取得 / ✗未取得
  - 右側: 制約事項パネル (薄色背景 + ⚠ ブレットリスト)
  - 下部: 出典・調査期間

Brand-aware (Phase 2, ISSUE-010):
  --brand stellar_aiz : 16:9 / Meiryo UI / 14-11-10pt / 紺×薄紺 / イエロー panel
  --brand roleup      : A4 横 / Yu Gothic UI / 12-10-6pt / 茶色×label_bg / ベージュ panel

ステータス色 (✓緑/△オレンジ/✗赤) は universal indicator のため両 brand 共通。

Usage:
  python fill_data_availability.py \
    --data /home/claude/data_availability_data.json \
    --brand stellar_aiz \
    --output /mnt/user-data/outputs/DataAvailability_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "data-availability-pptx"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE_ROLEUP = "Source 3"

PANEL_Y = Inches(1.50)
PANEL_H = Inches(5.25)

LEFT_X = Inches(0.41)
LEFT_W = Inches(8.00)

RIGHT_X = Inches(8.55)
RIGHT_W = Inches(4.40)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)
SOURCE_H = Inches(0.25)

TABLE_COL_W = (Inches(2.80), Inches(0.55), Inches(1.35), Inches(3.30))

# ── Colors (stella default) ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HEADER_BG = RGBColor(0x2E, 0x4A, 0x6B)  # 紺
COLOR_ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_CATEGORY_BG = RGBColor(0xE8, 0xEE, 0xF4)  # 薄紺
COLOR_CATEGORY_TEXT = RGBColor(0x2E, 0x4A, 0x6B)

# ステータス色 (両 brand 共通: universal indicator)
COLOR_COMPLETE = RGBColor(0x1B, 0x7A, 0x3B)   # 濃緑
COLOR_PARTIAL = RGBColor(0xDA, 0x7A, 0x2D)    # オレンジ
COLOR_MISSING = RGBColor(0xC0, 0x3A, 0x3A)    # 赤
COLOR_NA = RGBColor(0x99, 0x99, 0x99)         # グレー

# 制約パネル
COLOR_CONSTRAINT_BG = RGBColor(0xFF, 0xF9, 0xE8)  # 薄イエロー
COLOR_CONSTRAINT_BORDER = RGBColor(0xDA, 0x7A, 0x2D)
COLOR_CONSTRAINT_ICON = RGBColor(0xDA, 0x7A, 0x2D)

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_SECTION = Pt(14)
FONT_SIZE_TABLE_HEADER = Pt(11)
FONT_SIZE_TABLE = Pt(10)
FONT_SIZE_CATEGORY = Pt(11)
FONT_SIZE_STATUS = Pt(11)
FONT_SIZE_ITEM = Pt(11)
FONT_SIZE_SOURCE = Pt(10)

_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global PANEL_Y, PANEL_H, LEFT_X, LEFT_W, RIGHT_X, RIGHT_W
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H, TABLE_COL_W
    global COLOR_TEXT, COLOR_SOURCE, COLOR_HEADER_BG, COLOR_CATEGORY_BG, COLOR_CATEGORY_TEXT
    global COLOR_CONSTRAINT_BG, COLOR_CONSTRAINT_BORDER, COLOR_CONSTRAINT_ICON
    global FONT_NAME_JP
    global FONT_SIZE_SECTION, FONT_SIZE_TABLE_HEADER, FONT_SIZE_TABLE
    global FONT_SIZE_CATEGORY, FONT_SIZE_STATUS, FONT_SIZE_ITEM, FONT_SIZE_SOURCE

    _THEME = theme

    PANEL_Y = theme.layout("panel_y_in")
    PANEL_H = theme.layout("panel_h_in")
    LEFT_X = theme.layout("left_x_in")
    LEFT_W = theme.layout("left_w_in")
    RIGHT_X = theme.layout("right_x_in")
    RIGHT_W = theme.layout("right_w_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")
    TABLE_COL_W = (
        theme.layout("table_col1_w_in"),
        theme.layout("table_col2_w_in"),
        theme.layout("table_col3_w_in"),
        theme.layout("table_col4_w_in"),
    )

    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    FONT_NAME_JP = theme.font_ea

    if theme.id == "stellar_aiz":
        COLOR_HEADER_BG = RGBColor(0x2E, 0x4A, 0x6B)
        COLOR_CATEGORY_BG = RGBColor(0xE8, 0xEE, 0xF4)
        COLOR_CATEGORY_TEXT = RGBColor(0x2E, 0x4A, 0x6B)
        COLOR_CONSTRAINT_BG = RGBColor(0xFF, 0xF9, 0xE8)
        COLOR_CONSTRAINT_BORDER = RGBColor(0xDA, 0x7A, 0x2D)
        COLOR_CONSTRAINT_ICON = RGBColor(0xDA, 0x7A, 0x2D)
        FONT_SIZE_SECTION = Pt(14)
        FONT_SIZE_TABLE_HEADER = Pt(11)
        FONT_SIZE_TABLE = Pt(10)
        FONT_SIZE_CATEGORY = Pt(11)
        FONT_SIZE_STATUS = Pt(11)
        FONT_SIZE_ITEM = Pt(11)
        FONT_SIZE_SOURCE = Pt(10)
    else:
        # roleup: 茶色トーン + 許容集合 {22, 14, 12, 10, 6}
        COLOR_HEADER_BG = theme.color("label_bar")          # #7C4C2C
        COLOR_CATEGORY_BG = theme.color("label_bg")         # #F2E8DD
        COLOR_CATEGORY_TEXT = theme.color("label_bar")      # #7C4C2C
        COLOR_CONSTRAINT_BG = theme.color("label_bg")       # #F2E8DD
        COLOR_CONSTRAINT_BORDER = theme.color("highlight_target")  # #C78624
        COLOR_CONSTRAINT_ICON = theme.color("highlight_target")    # #C78624
        FONT_SIZE_SECTION = Pt(12)
        FONT_SIZE_TABLE_HEADER = Pt(10)
        FONT_SIZE_TABLE = Pt(10)
        FONT_SIZE_CATEGORY = Pt(10)
        FONT_SIZE_STATUS = Pt(10)
        FONT_SIZE_ITEM = Pt(10)
        FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")  # 6pt


def _silent_remove_shape(slide, shape_name: str) -> None:
    """Remove a shape by name without printing a warning. No-op if absent."""
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    if font_name is None:
        font_name = FONT_NAME_JP
    if color is None:
        color = COLOR_TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def add_section_title(slide, text, left, top, width):
    """セクションタイトル（下線付き、Bold）。下線色は brand 別。"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    # roleup は左寄せ subtitle 慣習、stella は中央寄せ legacy
    p.alignment = PP_ALIGN.LEFT if (_THEME and _THEME.id == "roleup") else PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SECTION
    run.font.bold = True
    # subtitle 色 (roleup=#897141, stella=#333333)
    if _THEME and "subtitle" in _THEME._colors:
        run.font.color.rgb = _THEME.color("subtitle")
    else:
        run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEXT
    line.line.fill.background()
    return txBox


def _get_status_info(status):
    """ステータス文字列 → (記号, 色, ラベル) を返す"""
    if not status:
        return ("—", COLOR_NA, "未設定")
    s = str(status).lower()
    if s in ("complete", "done", "full", "取得済", "取得済み", "✓", "yes", "ok"):
        return ("✓", COLOR_COMPLETE, "取得済")
    elif s in ("partial", "part", "一部", "△", "limited"):
        return ("△", COLOR_PARTIAL, "一部取得")
    elif s in ("missing", "none", "no", "✗", "x", "未取得", "n/a"):
        return ("✗", COLOR_MISSING, "未取得")
    return ("—", COLOR_NA, status)


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ──────────────────────────────────────────────
# Left Panel: Coverage Table
# ──────────────────────────────────────────────
def build_coverage_table(slide, section_title, categories, left, top, width, height):
    add_section_title(slide, section_title, left, top, width)

    tbl_top = top + Inches(0.50)
    tbl_h = height - Inches(0.50)

    rows = []
    for cat in categories:
        rows.append(("category", cat))
        for item in cat.get("items", []):
            rows.append(("item", item))

    n_rows = len(rows) + 1
    n_cols = 4

    row_h = Emu(int(tbl_h / n_rows))

    shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    table = shape.table

    col_w = list(TABLE_COL_W)
    remainder = width - sum(col_w, Emu(0))
    if remainder != Emu(0):
        col_w[-1] = Emu(col_w[-1] + remainder)
    for i, w in enumerate(col_w):
        table.columns[i].width = w

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(row_h))

    headers = ["項目", "", "ステータス", "データソース"]
    for c_idx, h in enumerate(headers):
        cell = table.cell(0, c_idx)
        _style_cell(cell, h, cell_type="header", font_size=FONT_SIZE_TABLE_HEADER)

    alt_counter = 0
    for r_idx, (row_type, row_data) in enumerate(rows):
        tr_idx = r_idx + 1

        if row_type == "category":
            cat_name = row_data.get("name", "")
            for c_idx in range(n_cols):
                cell = table.cell(tr_idx, c_idx)
                if c_idx == 0:
                    _style_cell(cell, cat_name, cell_type="category",
                                font_size=FONT_SIZE_CATEGORY)
                else:
                    _style_cell(cell, "", cell_type="category",
                                font_size=FONT_SIZE_CATEGORY)
            try:
                start_cell = table.cell(tr_idx, 0)
                end_cell = table.cell(tr_idx, n_cols - 1)
                start_cell.merge(end_cell)
            except Exception:
                pass
            alt_counter = 0

        else:
            label = row_data.get("label", "")
            status = row_data.get("status", "")
            source = row_data.get("source", "")

            symbol, color, status_label = _get_status_info(status)

            is_alt = (alt_counter % 2 == 1)
            alt_counter += 1

            _style_cell(table.cell(tr_idx, 0), label, cell_type="item",
                        is_alt=is_alt, font_size=FONT_SIZE_TABLE, align="l")
            _style_cell(table.cell(tr_idx, 1), symbol, cell_type="status",
                        is_alt=is_alt, font_size=FONT_SIZE_STATUS, align="ctr",
                        bold=True, text_color=color)
            _style_cell(table.cell(tr_idx, 2), status_label, cell_type="item",
                        is_alt=is_alt, font_size=FONT_SIZE_TABLE, align="l",
                        text_color=color, bold=True)
            _style_cell(table.cell(tr_idx, 3), source, cell_type="item",
                        is_alt=is_alt, font_size=FONT_SIZE_TABLE, align="l")

    print(f"  ✓ カバレッジテーブル: {len(categories)}カテゴリ、合計{n_rows - 1}行")


def _style_cell(cell, text, cell_type="item", is_alt=False, font_size=Pt(10),
                align="l", bold=False, text_color=None):
    """cell_type: 'header' | 'category' | 'item' | 'status'"""
    if cell_type == "header":
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
    elif cell_type == "category":
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CATEGORY_BG
    elif is_alt:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE

    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = cell.text_frame
    tf.word_wrap = True

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    p_elem = etree.SubElement(tf._txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))

    if cell_type == "header":
        pPr.set("algn", "ctr")
    elif cell_type == "category":
        pPr.set("algn", "l")
    else:
        pPr.set("algn", align)

    if text_color is not None:
        color_val = "{:02X}{:02X}{:02X}".format(text_color[0], text_color[1], text_color[2])
    elif cell_type == "header":
        color_val = "FFFFFF"
    elif cell_type == "category":
        color_val = "{:02X}{:02X}{:02X}".format(
            COLOR_CATEGORY_TEXT[0], COLOR_CATEGORY_TEXT[1], COLOR_CATEGORY_TEXT[2]
        )
    else:
        color_val = "{:02X}{:02X}{:02X}".format(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2])

    is_bold = bold or (cell_type in ("header", "category"))

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size.pt * 100)),
        "b": "1" if is_bold else "0",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    s = etree.SubElement(sf, qn("a:srgbClr"))
    s.set("val", color_val)
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text


# ──────────────────────────────────────────────
# Right Panel: Constraints
# ──────────────────────────────────────────────
def build_constraints_panel(slide, section_title, constraints, left, top, width, height):
    add_section_title(slide, section_title, left, top, width)

    if not constraints:
        return

    panel_top = top + Inches(0.50)
    panel_h = height - Inches(0.50)
    panel_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, panel_top, width, panel_h,
    )
    panel_bg.fill.solid()
    panel_bg.fill.fore_color.rgb = COLOR_CONSTRAINT_BG
    panel_bg.line.color.rgb = COLOR_CONSTRAINT_BORDER
    panel_bg.line.width = Pt(0.75)
    panel_bg.shadow.inherit = False
    panel_bg.text_frame.text = ""

    tb = slide.shapes.add_textbox(
        left + Inches(0.15), panel_top + Inches(0.12),
        width - Inches(0.30), panel_h - Inches(0.24),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0

    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    icon_color_hex = "{:02X}{:02X}{:02X}".format(
        COLOR_CONSTRAINT_ICON[0], COLOR_CONSTRAINT_ICON[1], COLOR_CONSTRAINT_ICON[2]
    )
    text_color_hex = "{:02X}{:02X}{:02X}".format(COLOR_TEXT[0], COLOR_TEXT[1], COLOR_TEXT[2])

    for i, item in enumerate(constraints):
        if isinstance(item, dict):
            text = item.get("text", "")
        else:
            text = str(item)

        p_elem = etree.SubElement(tf._txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"), attrib={
            "marL": "200000",
            "indent": "-200000",
        })
        if i > 0:
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "500"})

        buChar = etree.SubElement(pPr, qn("a:buChar"), attrib={"char": "⚠"})
        buFont = etree.SubElement(pPr, qn("a:buFont"), attrib={"typeface": FONT_NAME_JP})
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        buClrSolid = etree.SubElement(buClr, qn("a:srgbClr"))
        buClrSolid.set("val", icon_color_hex)

        r = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(int(FONT_SIZE_ITEM.pt * 100)),
        })
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        s = etree.SubElement(sf, qn("a:srgbClr"))
        s.set("val", text_color_hex)
        t_elem = etree.SubElement(r, qn("a:t"))
        t_elem.text = text

    print(f"  ✓ 制約事項パネル: {len(constraints)}項目")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "data-availability")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    # 必須キー欠落で hard-fail、想定外キーで stderr WARN を出す。
    validate_fill_input(
        data,
        required_top=["main_message", "categories"],
        allowed_top=[
            "main_message", "chart_title", "source",
            "categories", "constraints",
            "left_panel", "right_panel",
            # roleup brand で resolve_top_text/subtitle_text が読む可能性のあるキー
            "title", "subtitle",
        ],
        per_item_required={"categories": ["name", "items"]},
        skill_name=SKILL_ID,
    )

    # Roleup: source field is required (hard-fail). Stella: no-op.
    require_source(data, theme, skill_id=SKILL_ID)

    _mm = data.get("main_message", "")
    if len(_mm) > 65:
        raise ValueError(
            f"main_message は 65 字以内（受領: {len(_mm)}）: {_mm[:80]}..."
        )

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top / subtitle placeholder semantics differ between brands.
    top_text = resolve_top_text(data, theme) or data.get("main_message", "")
    sub_text = resolve_subtitle_text(data, theme) or data.get("chart_title", "調査のデータ取得状況")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:40]}")
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:40]}")

    # Roleup: silently remove brown guide rectangles carried by template.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    # 左: カバレッジテーブル
    left_panel = data.get("left_panel", {})
    left_section_title = left_panel.get("section_title", "データ取得カバレッジ")
    categories = data.get("categories", [])
    if categories:
        build_coverage_table(
            slide, left_section_title, categories,
            LEFT_X, PANEL_Y, LEFT_W, PANEL_H,
        )
    else:
        print("  ⚠ WARNING: no categories specified", file=sys.stderr)

    # 右: 制約事項
    right_panel = data.get("right_panel", {})
    right_section_title = right_panel.get("section_title", "調査上の留意事項・制約")
    constraints = data.get("constraints", [])
    build_constraints_panel(
        slide, right_section_title, constraints,
        RIGHT_X, PANEL_Y, RIGHT_W, PANEL_H,
    )

    # 出典: roleup なら Source 3 placeholder にセット、stella は dynamic textbox
    source = data.get("source", "")
    if source:
        if theme.id == "stellar_aiz":
            add_text_box(
                slide, source,
                SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                align=PP_ALIGN.LEFT,
            )
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE_ROLEUP)
            if source_shape is not None:
                set_textbox_text(source_shape, source)
            else:
                add_text_box(
                    slide, source,
                    SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                    FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
                    align=PP_ALIGN.LEFT,
                )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
