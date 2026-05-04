"""
fill_customer_profile.py — 主要顧客プロファイルスライドをPPTXネイティブオブジェクトで生成

テンプレート: customer-profile-template.pptx をベースに、
  - 既存テーブルを削除
  - 左側: 企業の概要テーブル（2列: ラベル | 値）
  - 右側: 業績チャート（ネイティブ複合チャート: 棒=売上高, 折れ線=営業利益率）
  - CAGR矢印＋テキスト注釈
を配置する。

Usage:
  python fill_customer_profile.py \
    --data /home/claude/customer_profile_data.json \
    --output /mnt/user-data/outputs/CustomerProfile_output.pptx \
    [--brand stellar_aiz|roleup] [--template <path>]

`--brand` (default: stellar_aiz) selects the output format. Theme JSON,
layout coordinates, and the per-brand template under assets/<brand>/ are
resolved by skills/_common/lib/brand_resolver.py. `--template` is optional;
if omitted it is resolved from the brand (with fallback to the stella
default template for brands that do not yet have a curated template).
"""

import argparse
import copy
import json
import math
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── brand_resolver bootstrap (skills/_common/lib/brand_resolver.py) ──
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import (  # noqa: E402
    apply_line_spacing,
    require_source,
    resolve_top_text,
    resolve_subtitle_text,
)

SKILL_ID = "customer-profile-pptx"
SHAPE_SOURCE = "Source 3"   # roleup template only; stella falls back to add_textbox

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# Main Message & Chart Title (from template, brand-independent)
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

# ── Brand-aware module globals ──
# Default values match stella for safety; reassigned in main() via _apply_theme(theme)
# after argparse resolves --brand. All visual values (color/font/coordinate) flow
# through theme JSON; SHAPE_* names above are template-structure invariants.
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
PANEL_Y = Inches(1.50)
LEFT_X = Inches(0.41)
LEFT_W = Inches(5.80)
RIGHT_X = Inches(6.50)
RIGHT_W = Inches(6.40)
CHART_H = Inches(4.80)
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.00)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_HEADER_BG = RGBColor(0xF5, 0xF0, 0xD0)
COLOR_BAR = RGBColor(0x4E, 0x79, 0xA7)       # Revenue bars
COLOR_LINE = RGBColor(0x00, 0x33, 0x66)      # Op margin line / legend connector
COLOR_CAGR_ARROW = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_SUBTITLE = RGBColor(0x33, 0x33, 0x33)  # default = text; reassigned in _apply_theme

# Hex strings (no leading '#') for inline OOXML attribute values
TEXT_HEX = "333333"
ACCENT_REVENUE_BAR_HEX = "4E79A7"
ACCENT_OP_MARGIN_LINE_HEX = "003366"

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_LABEL = Pt(14)
FONT_SIZE_VALUE = Pt(14)
FONT_SIZE_SECTION = Pt(14)
FONT_SIZE_CHART_TITLE = Pt(11)
FONT_SIZE_SOURCE = Pt(10)

# Phase 4 (ISSUE-011): module-level theme reference so helpers (e.g. _style_cell)
# can read line_height_pt without an extra parameter. Set in _apply_theme().
_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme.

    Called once from main() after `--brand` is parsed. Only invoked from main(),
    so module-load-time values above remain correct for direct imports / tests
    that don't go through main() (regression safety net).
    """
    global SLIDE_W, SLIDE_H, PANEL_Y, LEFT_X, LEFT_W, RIGHT_X, RIGHT_W, CHART_H
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_HEADER_BG, COLOR_BAR, COLOR_LINE, COLOR_CAGR_ARROW, COLOR_SOURCE, COLOR_SUBTITLE
    global TEXT_HEX, ACCENT_REVENUE_BAR_HEX, ACCENT_OP_MARGIN_LINE_HEX
    global FONT_NAME_JP, FONT_SIZE_LABEL, FONT_SIZE_VALUE, FONT_SIZE_SECTION, FONT_SIZE_CHART_TITLE
    global FONT_SIZE_SOURCE
    global _THEME
    _THEME = theme

    SLIDE_W = theme.slide_w
    SLIDE_H = theme.slide_h
    PANEL_Y = theme.layout("panel_y_in")
    LEFT_X = theme.layout("left_x_in")
    LEFT_W = theme.layout("left_w_in")
    RIGHT_X = theme.layout("right_x_in")
    RIGHT_W = theme.layout("right_w_in")
    CHART_H = theme.layout("chart_h_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_HEADER_BG = theme.color("header_bg")
    COLOR_BAR = theme.color("accent_revenue_bar")
    COLOR_LINE = theme.color("accent_op_margin_line")
    COLOR_CAGR_ARROW = theme.color("cagr_arrow")
    COLOR_SOURCE = theme.color("source")
    COLOR_SUBTITLE = theme.color("subtitle")  # roleup #897141 / stella #333333 (= text)

    TEXT_HEX = theme.hex_no_hash("text")
    ACCENT_REVENUE_BAR_HEX = theme.hex_no_hash("accent_revenue_bar")
    ACCENT_OP_MARGIN_LINE_HEX = theme.hex_no_hash("accent_op_margin_line")

    FONT_NAME_JP = theme.font_ea
    FONT_SIZE_LABEL = theme.pt("font_size_label_pt")
    FONT_SIZE_VALUE = theme.pt("font_size_value_pt")
    FONT_SIZE_SECTION = theme.pt("font_size_section_pt")
    FONT_SIZE_CHART_TITLE = theme.pt("font_size_chart_title_pt")
    FONT_SIZE_SOURCE = theme.pt("font_size_source_pt")


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def remove_shape(slide, name):
    """名前でShapeを削除"""
    shape = find_shape(slide, name)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def add_section_title(slide, text, left, top, width):
    """セクションタイトル（例: 企業の概要、業績）を追加"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    # alignment: stella=CENTER (既存), roleup=LEFT (公式 vF p.4 Subtitle 準拠)
    align_str = _THEME.layout_rule("subtitle_align", "center") if _THEME is not None else "center"
    p.alignment = PP_ALIGN.LEFT if align_str == "left" else PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SECTION
    run.font.bold = True
    run.font.color.rgb = COLOR_SUBTITLE
    run.font.name = FONT_NAME_JP

    # 下線を追加
    txBox2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    txBox2.fill.solid()
    txBox2.fill.fore_color.rgb = COLOR_TEXT
    txBox2.line.fill.background()
    return txBox


def build_overview_table(slide, items, left, top, width):
    """企業の概要テーブルを構築（ブレットポイント形式、枠線なし）"""
    n_rows = len(items)
    n_cols = 2
    col0_w = Inches(1.30)  # "• ラベル" column
    col1_w = width - col0_w

    # 行の最小高さ — PPTがテキスト折り返しに応じて自動拡張する。
    # brand に line_height_pt があれば余白込みで動的計算 (roleup 12pt → 0.20in 程度)、
    # なければ stella の既存値 0.35in (14pt 行) を維持。
    if _THEME is not None and _THEME.line_height_pt() is not None:
        # 12pt 行 + 4pt 余白 = 16pt = 0.222in
        min_row_h = Inches((_THEME.line_height_pt() + 4) / 72.0)
    else:
        min_row_h = Inches(0.35)
    table_h = min_row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = shape.table

    # 列幅設定
    table.columns[0].width = col0_w
    table.columns[1].width = col1_w

    # tblPrを設定（バンド無効化）
    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '0', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # 行の高さを最小値に設定（PPTが内容に応じて自動拡張する）
    for i, tr in enumerate(tbl_elem.findall(qn('a:tr'))):
        tr.set('h', str(min_row_h))

    for r_idx, item in enumerate(items):
        label = item.get("label", "")
        value = item.get("value", "")

        # ラベル列: "• ラベル" 形式
        bullet_label = f"•  {label}"
        _style_cell(table.cell(r_idx, 0), bullet_label, True, FONT_SIZE_LABEL, r_idx)
        _style_cell(table.cell(r_idx, 1), value, False, FONT_SIZE_VALUE, r_idx)

    print(f"  ✓ 企業概要テーブル: {n_rows}行 (ブレットポイント形式)")
    return shape


def _style_cell(cell, text, bold, font_size, row_idx):
    """セルにテキストとスタイルを設定（枠線なし、上揃え、行間詰め）"""
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    # bodyPrを設定/上書き — 上詰め・内部マージン最小
    old_bodyPr = txBody.find(qn("a:bodyPr"))
    if old_bodyPr is not None:
        txBody.remove(old_bodyPr)
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": "0", "rIns": "0",
        "tIns": "27432", "bIns": "27432",
        "anchor": "t",
    })
    # bodyPrを先頭に移動
    txBody.insert(0, bodyPr)

    # lstStyle
    if txBody.find(qn("a:lstStyle")) is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.insert(1, lstStyle)

    # 既存段落を削除
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # 新しい段落を追加
    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "l")

    # 行間: brand に line_height_pt があれば spcPts (roleup=12pt) で固定、
    # なければ既存挙動の spcPct 100% (stella 既定)。
    if _THEME is not None and _THEME.line_height_pt() is not None:
        apply_line_spacing(pPr, _THEME)
    else:
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size.pt * 100)),
        "b": "1" if bold else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": TEXT_HEX})

    # フォント指定
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    # セルスタイル (tcPr) - 上揃え、最小マージン、枠線なし
    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)
    tcPr = etree.SubElement(tc, qn("a:tcPr"), attrib={
        "marL": "45720", "marR": "18288",
        "marT": "27432", "marB": "27432",
        "anchor": "t",
    })

    # 全罫線を無しに設定
    for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = etree.SubElement(tcPr, qn(border_name), attrib={"w": "0", "cmpd": "sng"})
        etree.SubElement(ln, qn("a:noFill"))


def build_combo_chart(slide, perf_data, left, top, width, height):
    """PowerPointネイティブ複合チャート（棒＋折れ線）を作成"""
    from pptx.chart.data import CategoryChartData

    data = perf_data["data"]
    bar_label = perf_data.get("bar_label", "売上高")
    line_label = perf_data.get("line_label", "営業利益率")

    chart_data = CategoryChartData()
    chart_data.categories = [d["year"] for d in data]
    chart_data.add_series(bar_label, [d["revenue"] for d in data])
    chart_data.add_series(line_label, [d["op_margin"] for d in data])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = chart_frame.chart

    # 2番目のシリーズ(営業利益率)を折れ線に変更
    plot = chart.plots[0]

    # XML操作で2番目のシリーズをlineChartに変更
    plotArea = chart._chartSpace.chart.plotArea
    barChart = plotArea.findall(qn('c:barChart'))[0]

    # 2番目のシリーズを取り出す
    sers = barChart.findall(qn('c:ser'))
    line_ser_xml = copy.deepcopy(sers[1])
    barChart.remove(sers[1])

    # lineChartを追加
    lineChart = etree.SubElement(plotArea, qn('c:lineChart'))
    grouping = etree.SubElement(lineChart, qn('c:grouping'), attrib={'val': 'standard'})
    etree.SubElement(lineChart, qn('c:varyColors'), attrib={'val': '0'})
    lineChart.append(line_ser_xml)

    # マーカー追加
    marker_xml = line_ser_xml.find(qn('c:marker'))
    if marker_xml is None:
        marker_xml = etree.SubElement(line_ser_xml, qn('c:marker'))
    symbol = marker_xml.find(qn('c:symbol'))
    if symbol is None:
        symbol = etree.SubElement(marker_xml, qn('c:symbol'))
    symbol.set('val', 'circle')
    sz = marker_xml.find(qn('c:size'))
    if sz is None:
        sz = etree.SubElement(marker_xml, qn('c:size'))
    sz.set('val', '9')

    # 折れ線に第2軸を設定
    axId_elem = line_ser_xml.find(qn('c:axId'))
    # 第2軸追加
    catAx = plotArea.findall(qn('c:catAx'))
    valAx = plotArea.findall(qn('c:valAx'))

    # 既存の軸IDを取得
    primary_catAx_id = catAx[0].find(qn('c:axId')).get('val')
    primary_valAx_id = valAx[0].find(qn('c:axId')).get('val')

    # 第2数値軸を追加
    sec_valAx_id = "2094734553"
    sec_catAx_id = "2094734554"

    # lineChartにaxIdを設定
    etree.SubElement(lineChart, qn('c:axId'), attrib={'val': sec_catAx_id})
    etree.SubElement(lineChart, qn('c:axId'), attrib={'val': sec_valAx_id})

    # 第2カテゴリ軸（非表示）
    sec_catAx = etree.SubElement(plotArea, qn('c:catAx'))
    etree.SubElement(sec_catAx, qn('c:axId'), attrib={'val': sec_catAx_id})
    scaling = etree.SubElement(sec_catAx, qn('c:scaling'))
    etree.SubElement(scaling, qn('c:orientation'), attrib={'val': 'minMax'})
    etree.SubElement(sec_catAx, qn('c:delete'), attrib={'val': '1'})
    etree.SubElement(sec_catAx, qn('c:axPos'), attrib={'val': 'b'})
    etree.SubElement(sec_catAx, qn('c:crossAx'), attrib={'val': sec_valAx_id})

    # 第2数値軸（右側、%用）
    sec_valAx_elem = etree.SubElement(plotArea, qn('c:valAx'))
    etree.SubElement(sec_valAx_elem, qn('c:axId'), attrib={'val': sec_valAx_id})
    scaling2 = etree.SubElement(sec_valAx_elem, qn('c:scaling'))
    etree.SubElement(scaling2, qn('c:orientation'), attrib={'val': 'minMax'})
    etree.SubElement(sec_valAx_elem, qn('c:delete'), attrib={'val': '1'})
    etree.SubElement(sec_valAx_elem, qn('c:axPos'), attrib={'val': 'r'})
    numFmt = etree.SubElement(sec_valAx_elem, qn('c:numFmt'), attrib={
        'formatCode': '0.0"%"', 'sourceLinked': '0'
    })
    etree.SubElement(sec_valAx_elem, qn('c:crossAx'), attrib={'val': sec_catAx_id})
    etree.SubElement(sec_valAx_elem, qn('c:crosses'), attrib={'val': 'max'})

    # === スタイリング ===

    # 棒グラフの色
    bar_ser = barChart.findall(qn('c:ser'))[0]
    spPr = bar_ser.find(qn('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(bar_ser, qn('c:spPr'))
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        sf = etree.SubElement(spPr, qn('a:solidFill'))
    for child in list(sf):
        sf.remove(child)
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': ACCENT_REVENUE_BAR_HEX})

    # 折れ線の色（紺色）
    line_spPr = line_ser_xml.find(qn('c:spPr'))
    if line_spPr is None:
        line_spPr = etree.SubElement(line_ser_xml, qn('c:spPr'))
    ln = line_spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(line_spPr, qn('a:ln'))
    ln.set('w', '19050')  # 1.5pt
    line_sf = ln.find(qn('a:solidFill'))
    if line_sf is None:
        line_sf = etree.SubElement(ln, qn('a:solidFill'))
    for child in list(line_sf):
        line_sf.remove(child)
    etree.SubElement(line_sf, qn('a:srgbClr'), attrib={'val': ACCENT_OP_MARGIN_LINE_HEX})

    # マーカーの色（紺色）・サイズ大きめ
    marker_spPr = marker_xml.find(qn('c:spPr'))
    if marker_spPr is None:
        marker_spPr = etree.SubElement(marker_xml, qn('c:spPr'))
    m_sf = etree.SubElement(marker_spPr, qn('a:solidFill'))
    etree.SubElement(m_sf, qn('a:srgbClr'), attrib={'val': ACCENT_OP_MARGIN_LINE_HEX})

    # データラベル（棒グラフ - 売上高）
    add_data_labels_to_ser(bar_ser, position='outEnd', num_format='0.0', font_color=TEXT_HEX)

    # データラベル（折れ線 - 営業利益率：白文字）
    add_data_labels_to_ser(line_ser_xml, position='t', num_format='0.0', font_color='FFFFFF')

    # 凡例を無効化（カスタム凡例をスライド上に別途配置するため）
    chart.has_legend = False

    # === グリッドライン（目盛線）を削除 ===
    for vax in plotArea.findall(qn('c:valAx')):
        # majorGridlinesを削除
        for mg in vax.findall(qn('c:majorGridlines')):
            vax.remove(mg)
        # minorGridlinesも削除
        for mg in vax.findall(qn('c:minorGridlines')):
            vax.remove(mg)
    # catAxのグリッドラインも削除
    for cax in plotArea.findall(qn('c:catAx')):
        for mg in cax.findall(qn('c:majorGridlines')):
            cax.remove(mg)

    # === 棒グラフの太さを1.2倍に（gapWidthを調整） ===
    gapWidth = barChart.find(qn('c:gapWidth'))
    if gapWidth is None:
        gapWidth = etree.SubElement(barChart, qn('c:gapWidth'))
    gapWidth.set('val', '108')  # デフォルト150→108で約1.2倍太く

    # === 軸のフォント（年ラベルを縦書きに） ===
    for ax in plotArea.findall(qn('c:catAx')) + plotArea.findall(qn('c:valAx')):
        delete_elem = ax.find(qn('c:delete'))
        if delete_elem is not None and delete_elem.get('val') == '1':
            continue
        txPr = ax.find(qn('c:txPr'))
        if txPr is None:
            txPr = etree.SubElement(ax, qn('c:txPr'))
        # bodyPrを設定/更新（縦書き回転）
        bodyPr = txPr.find(qn('a:bodyPr'))
        if bodyPr is None:
            bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
            txPr.insert(0, bodyPr)
        bodyPr.set('rot', '-5400000')  # -90度 = 縦書き
        bodyPr.set('vert', 'horz')

        if txPr.find(qn('a:lstStyle')) is None:
            etree.SubElement(txPr, qn('a:lstStyle'))
        # フォント設定
        p = txPr.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(txPr, qn('a:p'))
        pPr = p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p, qn('a:pPr'))
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': '1100'})
        else:
            defRPr.set('sz', '1100')
        # フォント指定
        if defRPr.find(qn('a:latin')) is None:
            etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
        if defRPr.find(qn('a:ea')) is None:
            etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})

    # プライマリ数値軸を非表示
    for vax in valAx:
        del_elem = vax.find(qn('c:delete'))
        if del_elem is None:
            del_elem = etree.SubElement(vax, qn('c:delete'))
        del_elem.set('val', '1')

    # チャートタイトルなし
    chart.has_title = False

    # プロットエリア背景なし
    plotArea_spPr = plotArea.find(qn('c:spPr'))
    if plotArea_spPr is None:
        plotArea_spPr = etree.SubElement(plotArea, qn('c:spPr'))
    noFill = etree.SubElement(plotArea_spPr, qn('a:noFill'))

    print(f"  ✓ 複合チャート生成: {len(data)}年分")
    return chart_frame


def add_data_labels_to_ser(ser_xml, position='outEnd', num_format='0.0', font_color=None):
    """font_color defaults to brand TEXT_HEX when not provided."""
    if font_color is None:
        font_color = TEXT_HEX
    """シリーズにデータラベルを追加"""
    dLbls = ser_xml.find(qn('c:dLbls'))
    if dLbls is None:
        dLbls = etree.SubElement(ser_xml, qn('c:dLbls'))

    # 既存の子を削除
    for child in list(dLbls):
        dLbls.remove(child)

    numFmt = etree.SubElement(dLbls, qn('c:numFmt'), attrib={
        'formatCode': num_format, 'sourceLinked': '0'
    })
    etree.SubElement(dLbls, qn('c:showLegendKey'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dLbls, qn('c:showCatName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showSerName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showPercent'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showBubbleSize'), attrib={'val': '0'})

    pos_map = {'outEnd': 'outEnd', 't': 't', 'ctr': 'ctr'}
    dLblPos = etree.SubElement(dLbls, qn('c:dLblPos'), attrib={
        'val': pos_map.get(position, 'outEnd')
    })

    # データラベルフォント
    txPr = etree.SubElement(dLbls, qn('c:txPr'))
    etree.SubElement(txPr, qn('a:bodyPr'))
    etree.SubElement(txPr, qn('a:lstStyle'))
    p = etree.SubElement(txPr, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': '1200'})
    latin = etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
    ea = etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})
    sf = etree.SubElement(defRPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': font_color})


def add_cagr_annotation(slide, perf_data, chart_left, chart_top, chart_width, chart_height):
    """CAGR矢印とテキスト注釈を追加（棒グラフの位置に合わせて配置）"""
    data = perf_data["data"]
    if len(data) < 2:
        return

    first_rev = data[0]["revenue"]
    last_rev = data[-1]["revenue"]
    max_rev = max(d["revenue"] for d in data)
    n_years = len(data) - 1
    n_cats = len(data)

    # CAGR計算
    if first_rev > 0 and last_rev > 0 and n_years > 0:
        cagr = (last_rev / first_rev) ** (1.0 / n_years) - 1
        cagr_text = f"+{cagr*100:.1f}%" if cagr >= 0 else f"{cagr*100:.1f}%"
    else:
        cagr_text = "N/A"

    # チャート内のプロットエリアを推定
    # （凡例・軸ラベル分のマージンを考慮）
    plot_left_margin = chart_width * 0.06
    plot_right_margin = chart_width * 0.04
    plot_top_margin = chart_height * 0.18   # 凡例スペース
    plot_bottom_margin = chart_height * 0.14  # カテゴリラベル

    plot_left = chart_left + plot_left_margin
    plot_right = chart_left + chart_width - plot_right_margin
    plot_top = chart_top + plot_top_margin
    plot_bottom = chart_top + chart_height - plot_bottom_margin
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    # 各棒の中心X座標を計算
    cat_width = plot_w / n_cats
    first_bar_cx = plot_left + 0.5 * cat_width
    last_bar_cx = plot_left + (n_cats - 0.5) * cat_width

    # 各棒の上端Y座標を推定（軸最大値 ≈ max_rev * 1.20）
    axis_max = max_rev * 1.20
    first_bar_top_y = plot_bottom - (first_rev / axis_max) * plot_h
    last_bar_top_y = plot_bottom - (last_rev / axis_max) * plot_h

    # 矢印は棒の少し上（0.15in上）
    gap_above = Inches(1.20)
    arrow_start_x = int(first_bar_cx)
    arrow_start_y = int(first_bar_top_y - gap_above)
    arrow_end_x = int(last_bar_cx)
    arrow_end_y = int(last_bar_top_y - gap_above)

    # 矢印Shape
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        arrow_start_x, arrow_start_y,
        arrow_end_x, arrow_end_y
    )
    connector.line.color.rgb = COLOR_CAGR_ARROW
    connector.line.width = Pt(1.5)

    # 矢印の先端を設定
    cxnSp = connector._element
    spPr = cxnSp.find(qn('p:spPr'))
    if spPr is None:
        spPr = cxnSp.find(qn('a:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'), attrib={
        'type': 'triangle', 'w': 'med', 'len': 'med'
    })

    # CAGRテキスト（楕円＋テキスト）— 矢印の真ん中に配置
    text_w = Inches(1.30)
    text_h = Inches(0.40)
    mid_x = (arrow_start_x + arrow_end_x) / 2
    mid_y = (arrow_start_y + arrow_end_y) / 2
    text_x = int(mid_x - text_w / 2)
    text_y = int(mid_y - text_h / 2)

    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, text_x, text_y, text_w, text_h
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    oval.line.color.rgb = COLOR_CAGR_ARROW
    oval.line.width = Pt(1.0)

    tf = oval.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cagr_text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    print(f"  ✓ CAGR注釈: {cagr_text} ({data[0]['year']}→{data[-1]['year']})")


def add_unit_label(slide, text, left, top, width):
    """単位表記テキストを追加（左寄せ）"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.22))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP


def add_custom_legend(slide, perf_data, left, top, width):
    """カスタム凡例を右寄せで配置（■売上高  ●━営業利益率）"""
    bar_label = perf_data.get("bar_label", "売上高")
    line_label = perf_data.get("line_label", "営業利益率")

    legend_w = Inches(2.80)
    legend_h = Inches(0.22)
    legend_x = left + width - legend_w  # 右寄せ

    # 凡例の小さい四角（売上高）
    sq_size = Inches(0.14)
    sq_y = top + (legend_h - sq_size) / 2

    # ■ 売上高マーカー
    bar_marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        legend_x, int(sq_y), sq_size, sq_size
    )
    bar_marker.fill.solid()
    bar_marker.fill.fore_color.rgb = COLOR_BAR
    bar_marker.line.fill.background()

    # "売上高" テキスト
    bar_text_x = legend_x + sq_size + Inches(0.06)
    bar_text_w = Inches(0.70)
    txBox1 = slide.shapes.add_textbox(int(bar_text_x), top, bar_text_w, legend_h)
    tf1 = txBox1.text_frame
    tf1.word_wrap = False
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.add_run()
    run1.text = bar_label
    run1.font.size = Pt(12)
    run1.font.color.rgb = COLOR_TEXT
    run1.font.name = FONT_NAME_JP

    # ●━ 営業利益率マーカー（丸＋線）
    line_section_x = bar_text_x + bar_text_w + Inches(0.20)

    # 線
    line_w = Inches(0.30)
    line_y = top + legend_h / 2
    connector = slide.shapes.add_connector(
        1, int(line_section_x), int(line_y),
        int(line_section_x + line_w), int(line_y)
    )
    connector.line.color.rgb = COLOR_LINE
    connector.line.width = Pt(1.5)

    # 丸マーカー
    circle_size = Inches(0.12)
    circle_x = line_section_x + (line_w - circle_size) / 2
    circle_y = top + (legend_h - circle_size) / 2
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(circle_x), int(circle_y), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_LINE
    circle.line.fill.background()

    # "営業利益率" テキスト
    line_text_x = line_section_x + line_w + Inches(0.06)
    line_text_w = Inches(1.20)
    txBox2 = slide.shapes.add_textbox(int(line_text_x), top, line_text_w, legend_h)
    tf2 = txBox2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = line_label
    run2.font.size = Pt(12)
    run2.font.color.rgb = COLOR_TEXT
    run2.font.name = FONT_NAME_JP

    print(f"  ✓ カスタム凡例: {bar_label} / {line_label}")


def add_source_label(slide, text):
    """スライド左下に出典テキストを追加。

    roleup テンプレに 'Source 3' placeholder があればそこへ書き込み、
    無ければ動的に textbox を生成 (stella の既存挙動)。
    """
    src_shape = None
    for shape in slide.shapes:
        if shape.name == SHAPE_SOURCE:
            src_shape = shape
            break

    if src_shape is not None:
        tf = src_shape.text_frame
        tf.word_wrap = True
        # 既存サンプルテキスト/段落を1段落だけ残してクリア
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        p = tf.paragraphs[0]
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = FONT_SIZE_SOURCE
        run.font.color.rgb = COLOR_SOURCE
        run.font.name = FONT_NAME_JP
        print(f"  ✓ 出典 (Source 3 placeholder): {text}")
        return

    txBox = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SOURCE
    run.font.color.rgb = COLOR_SOURCE
    run.font.name = FONT_NAME_JP
    print(f"  ✓ 出典 (textbox): {text}")


def main():
    parser = argparse.ArgumentParser(description="主要顧客プロファイル PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True)
    parser.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "customer-profile")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # Phase 4 (ISSUE-011): roleup は出所必須。stella は no-op (既存 warning 維持)。
    require_source(data, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 1. Top placeholder (最上部の最大フォント位置)
    #    stella: main_message (結論文) / roleup: chart_title (スライドタイトル見出し)
    top_text = resolve_top_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text}")

    # 2. Subtitle placeholder (副題位置)
    #    stella: chart_title (見出し) / roleup: main_message (結論文)
    sub_text = resolve_subtitle_text(data, theme)
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text}")

    # 3. 既存テーブルを削除
    remove_shape(slide, "Table 1")

    # 4. 左側: 企業の概要
    overview = data.get("company_overview", {})
    section_title_left = overview.get("section_title", "企業の概要")
    add_section_title(slide, section_title_left, LEFT_X, PANEL_Y, LEFT_W)
    items = overview.get("items", [])
    build_overview_table(slide, items, LEFT_X, PANEL_Y + Inches(0.40), LEFT_W)

    # 5. 右側: 業績チャート
    perf = data.get("performance", {})
    section_title_right = perf.get("section_title", "業績")

    # セクションタイトル
    add_section_title(slide, section_title_right, RIGHT_X, PANEL_Y, RIGHT_W)

    # 単位表記（左側）
    unit_label = perf.get("unit_label", "")
    if unit_label:
        add_unit_label(slide, unit_label, RIGHT_X, PANEL_Y + Inches(0.35), Inches(2.50))

    # カスタム凡例（右側）
    add_custom_legend(slide, perf, RIGHT_X, PANEL_Y + Inches(0.35), RIGHT_W)

    # 複合チャート
    chart_top = PANEL_Y + Inches(0.55)
    chart_frame = build_combo_chart(slide, perf, RIGHT_X, chart_top, RIGHT_W, CHART_H)

    # CAGR注釈
    add_cagr_annotation(slide, perf, RIGHT_X, chart_top, RIGHT_W, CHART_H)

    # 7. 出典
    source = data.get("source", "")
    if source:
        add_source_label(slide, source)

    # 8. 保存
    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
