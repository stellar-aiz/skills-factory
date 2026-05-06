"""
fill_issue_risk.py — 課題・リスク一覧データをPPTXテンプレートに流し込むスクリプト

複数ページ対応:
  - 1スライドあたり最大4行（MAX_ROWS_PER_SLIDE、16pt前提）
  - 5行以上の場合は自動的にスライドを追加してバランス分割
  - 複数ページ時のChart Titleに「（1/3）」等のページ番号を付与
  - フォントサイズはデフォルト16pt（JSON側で font_size_header / font_size_data を指定可能）

使い方:
  python fill_issue_risk.py \
    --data /home/claude/issue_risk_data.json \
    --template <SKILL_DIR>/assets/issue-risk-template.pptx \
    --output /mnt/user-data/outputs/IssueRisk_output.pptx
"""

import argparse
import os
import json
import sys
import copy
import math

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── レイアウト定数 ──────────────────────────────────────────
MARGIN_LEFT    = 370800
CONTENT_WIDTH  = 11616153
HEADER_TOP     = 1425630
HEADER_HEIGHT  = 500000
HEADER_SEP_TOP = 1930000
FIRST_ROW_TOP  = 2050000
ROW_HEIGHT     = 900000
ROW_SPACING    = 1100000

MAX_ROWS_PER_SLIDE = 4

# ── フォントサイズデフォルト（16pt = 1600） ──
DEFAULT_HEADER_FONT_SIZE = 1600
DEFAULT_DATA_FONT_SIZE   = 1600

SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"

REMOVE_NAMES = {
    "TextBox 4", "TextBox 5", "TextBox 23", "TextBox 25", "TextBox 26",
    "Straight Connector 28",
    "Group 8", "Group 15", "Group 57", "Group 64", "Group 71", "Group 78",
    "Straight Connector 7", "Straight Connector 51",
    "Straight Connector 63", "Straight Connector 70", "Straight Connector 77",
}


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_placeholder_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def remove_dynamic_shapes(slide):
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if shape.name in REMOVE_NAMES:
            to_remove.append(shape._element)
    for elem in to_remove:
        sp_tree.remove(elem)
    return len(to_remove)


def add_textbox(slide, left, top, width, height, text, font_size=1100,
                bold=False, wrap="square", lang="en-GB"):
    sp_tree = slide.shapes._spTree
    sp = etree.SubElement(sp_tree, qn("p:sp"))

    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    etree.SubElement(nvSpPr, qn("p:cNvPr"), attrib={
        "id": str(id(sp) % 100000 + 100), "name": "DynTextBox"
    })
    etree.SubElement(nvSpPr, qn("p:cNvSpPr"), attrib={"txBox": "1"})
    etree.SubElement(nvSpPr, qn("p:nvPr"))

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    etree.SubElement(xfrm, qn("a:off"), attrib={"x": str(left), "y": str(top)})
    etree.SubElement(xfrm, qn("a:ext"), attrib={"cx": str(width), "cy": str(height)})
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"), attrib={"prst": "rect"})
    etree.SubElement(prstGeom, qn("a:avLst"))
    etree.SubElement(spPr, qn("a:noFill"))

    txBody = etree.SubElement(sp, qn("p:txBody"))
    etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": wrap, "lIns": "0", "rIns": "0", "rtlCol": "0"
    })
    etree.SubElement(txBody, qn("a:lstStyle"))

    p = etree.SubElement(txBody, qn("a:p"))
    etree.SubElement(p, qn("a:pPr"), attrib={"algn": "l"})
    r = etree.SubElement(p, qn("a:r"))

    rPr_attrib = {"kumimoji": "1", "lang": lang, "sz": str(font_size)}
    if bold:
        rPr_attrib["b"] = "1"
    rPr = etree.SubElement(r, qn("a:rPr"), attrib=rPr_attrib)
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "+mn-ea"})

    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return sp


def add_line(slide, left, top, width, line_width=9525, dash="dash"):
    sp_tree = slide.shapes._spTree
    cxnSp = etree.SubElement(sp_tree, qn("p:cxnSp"))

    nvCxnSpPr = etree.SubElement(cxnSp, qn("p:nvCxnSpPr"))
    etree.SubElement(nvCxnSpPr, qn("p:cNvPr"), attrib={
        "id": str(id(cxnSp) % 100000 + 200), "name": "DynConnector"
    })
    etree.SubElement(nvCxnSpPr, qn("p:cNvCxnSpPr"))
    etree.SubElement(nvCxnSpPr, qn("p:nvPr"))

    spPr = etree.SubElement(cxnSp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    etree.SubElement(xfrm, qn("a:off"), attrib={"x": str(left), "y": str(top)})
    etree.SubElement(xfrm, qn("a:ext"), attrib={"cx": str(width), "cy": "0"})
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"), attrib={"prst": "line"})
    etree.SubElement(prstGeom, qn("a:avLst"))

    ln = etree.SubElement(spPr, qn("a:ln"), attrib={"w": str(line_width)})
    sf = etree.SubElement(ln, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:schemeClr"), attrib={"val": "tx1"})
    if dash != "solid":
        etree.SubElement(ln, qn("a:prstDash"), attrib={"val": dash})
    etree.SubElement(spPr, qn("a:effectLst"))
    return cxnSp


def compute_column_positions(columns, total_width, margin_left):
    total_ratio = sum(c.get("width_ratio", 1) for c in columns)
    positions = []
    current_x = margin_left
    for col in columns:
        ratio = col.get("width_ratio", 1)
        col_width = int(total_width * ratio / total_ratio)
        positions.append({"left": current_x, "width": col_width})
        current_x += col_width
    return positions


def build_headers(slide, columns, col_positions, font_size=DEFAULT_HEADER_FONT_SIZE):
    for i, col in enumerate(columns):
        pos = col_positions[i]
        add_textbox(slide, left=pos["left"], top=HEADER_TOP,
                    width=pos["width"], height=HEADER_HEIGHT,
                    text=col["name"], font_size=font_size, bold=True)


def build_data_rows(slide, columns, col_positions, rows, start_row_num=0,
                    font_size=DEFAULT_DATA_FONT_SIZE):
    for r_idx, row in enumerate(rows):
        row_top = FIRST_ROW_TOP + r_idx * ROW_SPACING

        for c_idx, col in enumerate(columns):
            pos = col_positions[c_idx]
            cell_value = row[c_idx] if c_idx < len(row) else ""
            add_textbox(slide, left=pos["left"], top=row_top,
                        width=pos["width"], height=ROW_HEIGHT,
                        text=str(cell_value), font_size=font_size, bold=False)

        global_idx = start_row_num + r_idx + 1
        print(f"  [Row {global_idx}] {' | '.join(str(v) for v in row)}")

        sep_top = row_top + ROW_HEIGHT + (ROW_SPACING - ROW_HEIGHT) // 2
        add_line(slide, left=MARGIN_LEFT, top=sep_top,
                 width=CONTENT_WIDTH, line_width=9525, dash="dash")


def duplicate_slide_from_template(prs, template_slide):
    """テンプレートスライドの「クリーンな状態」を複製して新スライドを追加する"""
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # レイアウト由来のshapeを全削除
    new_sp_tree = new_slide.shapes._spTree
    for child in list(new_sp_tree):
        if child.tag != qn("p:nvGrpSpPr") and child.tag != qn("p:grpSpPr"):
            new_sp_tree.remove(child)

    # テンプレートのshapeをコピー（保存済みのクリーンなXMLから）
    for child in template_slide.shapes._spTree:
        if child.tag != qn("p:nvGrpSpPr") and child.tag != qn("p:grpSpPr"):
            new_sp_tree.append(copy.deepcopy(child))

    return new_slide


def populate_slide(slide, main_msg, chart_title, columns, col_positions, rows,
                   page_num, total_pages, start_row_num,
                   header_font_size=DEFAULT_HEADER_FONT_SIZE,
                   data_font_size=DEFAULT_DATA_FONT_SIZE):
    """1枚のスライドにコンテンツを配置する"""

    shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
    set_placeholder_text(shape, main_msg)

    display_title = chart_title
    if total_pages > 1:
        display_title = f"{chart_title}（{page_num}/{total_pages}）"
    shape = find_shape(slide, SHAPE_CHART_TITLE)
    set_placeholder_text(shape, display_title)

    removed = remove_dynamic_shapes(slide)

    print(f"\n  === Page {page_num}/{total_pages} ===")
    print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    print(f"  [Chart Title]  {display_title}")
    print(f"  Removed {removed} template shapes")

    build_headers(slide, columns, col_positions, font_size=header_font_size)
    print(f"  [Headers] {' | '.join(c['name'] for c in columns)}")

    add_line(slide, left=MARGIN_LEFT, top=HEADER_SEP_TOP,
             width=CONTENT_WIDTH, line_width=15875, dash="solid")

    if rows:
        build_data_rows(slide, columns, col_positions, rows, start_row_num,
                        font_size=data_font_size)
        print(f"  [Subtotal] {len(rows)} rows on this page")


def split_rows_balanced(all_rows, max_per_page):
    """行を複数ページにバランスよく分割する。

    末尾ページに1行だけ残るような不均衡な分割を避け、
    ceil(total / pages) 行を各ページに均等配分する。
    例: 7行・max=4 → [4, 3]（6/1ではなく）
        10行・max=4 → [4, 3, 3]（4/4/2ではなく）
    """
    total = len(all_rows)
    if total == 0:
        return [[]]
    total_pages = max(1, math.ceil(total / max_per_page))
    per_page    = math.ceil(total / total_pages)
    chunks = []
    idx = 0
    for _ in range(total_pages):
        chunks.append(all_rows[idx:idx + per_page])
        idx += per_page
    return [c for c in chunks if c] or [[]]


def main():
    parser = argparse.ArgumentParser(description="課題・リスク一覧データをPPTXに流し込む")
    parser.add_argument("--data",     required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output",   required=True)
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "rows"],
        allowed_top=["main_message", "chart_title", "rows", "columns"],
        skill_name="issue-risk-list-pptx",
    )

    _main_msg = data.get("main_message", "")
    if len(_main_msg) > 70:
        print(
            f"  ⚠ WARNING: main_message is {len(_main_msg)} chars (max 70). "
            f"Text will overflow the textbox. Preview: '{_main_msg[:40]}...'",
            file=sys.stderr,
        )

    prs = Presentation(args.template)

    main_msg    = data.get("main_message", "").strip()
    chart_title = data.get("chart_title", "").strip()
    columns     = data.get("columns", [])
    all_rows    = data.get("rows", [])

    if not columns:
        print("  ERROR: columns is empty", file=sys.stderr)
        sys.exit(1)
    if not main_msg:
        print("  WARNING: main_message is empty", file=sys.stderr)
    if not chart_title:
        print("  WARNING: chart_title is empty", file=sys.stderr)

    # フォントサイズはJSON側で上書き可能（単位: 百分の一pt、16pt = 1600）
    header_font_size = int(data.get("font_size_header", DEFAULT_HEADER_FONT_SIZE))
    data_font_size   = int(data.get("font_size_data",   DEFAULT_DATA_FONT_SIZE))

    col_positions = compute_column_positions(columns, CONTENT_WIDTH, MARGIN_LEFT)

    # ── 行をページごとにバランス分割 ──
    # 末尾ページに1行だけ残るような不均衡な分割を避ける
    row_chunks  = split_rows_balanced(all_rows, MAX_ROWS_PER_SLIDE)
    total_pages = len(row_chunks)

    print(f"  Total: {len(all_rows)} rows -> {total_pages} page(s)")
    print(f"  Split: {[len(c) for c in row_chunks]} rows per page")
    print(f"  Font:  header={header_font_size/100:.1f}pt / data={data_font_size/100:.1f}pt")

    # ── テンプレートスライドのクリーンなXMLを保存 ──
    template_slide = prs.slides[0]
    saved_sp_tree = copy.deepcopy(template_slide.shapes._spTree)

    # ── ページ2以降のスライドを先に複製（テンプレートがクリーンな状態で）──
    extra_slides = []
    for _ in range(1, total_pages):
        new_slide = duplicate_slide_from_template(prs, template_slide)
        extra_slides.append(new_slide)

    # ── ページ1を処理 ──
    populate_slide(
        template_slide, main_msg, chart_title, columns, col_positions,
        row_chunks[0] if row_chunks else [],
        page_num=1, total_pages=total_pages, start_row_num=0,
        header_font_size=header_font_size, data_font_size=data_font_size,
    )

    # ── ページ2以降を処理 ──
    cumulative = len(row_chunks[0]) if row_chunks else 0
    for page_idx, extra_slide in enumerate(extra_slides):
        populate_slide(
            extra_slide, main_msg, chart_title, columns, col_positions,
            row_chunks[page_idx + 1],
            page_num=page_idx + 2,
            total_pages=total_pages,
            start_row_num=cumulative,
            header_font_size=header_font_size, data_font_size=data_font_size,
        )
        cumulative += len(row_chunks[page_idx + 1])

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n  Saved: {args.output} ({total_pages} slide(s))")


if __name__ == "__main__":
    main()
