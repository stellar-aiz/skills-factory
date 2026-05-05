"""
fill_sales_by_customer_native.py — 主要販売先売上高テーブルをPPTXネイティブテーブルで生成

テンプレート構造（sales-by-customer-template.pptx）:
  - Title 1            (PLACEHOLDER): メインメッセージ（上段、太字）
  - Text Placeholder 2 (PLACEHOLDER): チャートタイトル（下段）
  - Content Area       (AUTO_SHAPE):  削除して独立テーブル×N で再構成
  - Source             (TEXT_BOX):    出典（左下）

方式: 期ごとに独立した add_table() でネイティブテーブルを横並び配置。
      Playwright不要。全オブジェクトがPPT上で編集可能。

Usage:
  python fill_sales_by_customer_native.py \
    --data /home/claude/sales_by_customer_data.json \
    --template <SKILL_DIR>/assets/sales-by-customer-template.pptx \
    --output /mnt/user-data/outputs/SalesByCustomer_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE = "Source"

# ── レイアウト定数 ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

CONTENT_LEFT = Inches(0.41)
CONTENT_TOP = Inches(1.50)
CONTENT_WIDTH = Inches(12.52)
CONTENT_BOTTOM = Inches(6.90)     # Source の上

TABLE_GAP = Inches(0.15)          # テーブル間ギャップ
NOTE_HEIGHT = Inches(0.25)        # 注記テキストの高さ
PERIOD_HEADER_HEIGHT = Inches(0.30)  # 期ヘッダーの高さ
BOTTOM_MARGIN = Inches(0.05)

# ── 列幅比率 ──
COL_RATIOS = [0.08, 0.42, 0.30, 0.20]  # #, 企業名, 売上高, 割合

# ── 色定数 ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_ORANGE = RGBColor(0xE6, 0x7E, 0x00)
COLOR_HEADER_BG = RGBColor(0xF0, 0xF0, 0xF0)
COLOR_EVEN_ROW = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_OTHER_ROW = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_BORDER = RGBColor(0xE0, 0xE0, 0xE0)

# ── フォント ──
FONT_NAME = "Meiryo UI"
FONT_NAME_LATIN = "Arial"


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def remove_shape(slide, name):
    shape = find_shape(slide, name)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def get_font_sizes(n_periods):
    """期数に応じたフォントサイズを返す"""
    if n_periods <= 3:
        return {"header": Pt(14), "body": Pt(13), "period": Pt(16), "note": Pt(11)}
    elif n_periods == 4:
        return {"header": Pt(12), "body": Pt(11), "period": Pt(14), "note": Pt(10)}
    else:  # 5
        return {"header": Pt(11), "body": Pt(10), "period": Pt(13), "note": Pt(9)}


def _num_to_kanji(n):
    kanji_map = {2: "二", 3: "三", 4: "四", 5: "五"}
    return kanji_map.get(n, str(n))


def identify_continuous_customers(periods):
    """全期間に出現する企業名を特定"""
    if len(periods) < 2:
        return set()
    name_sets = []
    for period in periods:
        names = set()
        for row in period.get("customers", []):
            name = row.get("name", "")
            if name and name != "その他":
                names.add(name)
        name_sets.append(names)
    result = name_sets[0]
    for ns in name_sets[1:]:
        result = result & ns
    return result


def add_textbox(slide, left, top, width, height, text, font_size,
                font_bold=False, font_color=COLOR_TEXT, alignment=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.bold = font_bold
    run.font.color.rgb = font_color
    # Latin font
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", FONT_NAME_LATIN)
    return txBox


def set_cell_text(cell, text, font_size, font_bold=False, font_color=COLOR_TEXT,
                  alignment=PP_ALIGN.LEFT, bg_color=None):
    """セルにテキストとスタイルを設定"""
    tc = cell._tc

    # --- テキスト設定 ---
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    # bodyPr の設定（上下マージンを小さく）
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    bodyPr.set("lIns", "45720")   # 左マージン 0.05"
    bodyPr.set("rIns", "45720")   # 右マージン 0.05"
    bodyPr.set("tIns", "18288")   # 上マージン 小さめ
    bodyPr.set("bIns", "18288")   # 下マージン 小さめ
    bodyPr.set("anchor", "ctr")   # 垂直中央

    lstStyle = txBody.find(qn("a:lstStyle"))
    if lstStyle is None:
        etree.SubElement(txBody, qn("a:lstStyle"))

    # 既存の段落を削除
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # 段落を追加
    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    align_map = {PP_ALIGN.LEFT: "l", PP_ALIGN.CENTER: "ctr", PP_ALIGN.RIGHT: "r"}
    pPr.set("algn", align_map.get(alignment, "l"))

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size.pt * 100)),
        "b": "1" if font_bold else "0",
        "dirty": "0",
    })
    # フォント色
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{font_color[0]:02X}{font_color[1]:02X}{font_color[2]:02X}"
                if isinstance(font_color, (tuple, list))
                else str(font_color).replace("#", ""))
    # フォント名
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", FONT_NAME_LATIN)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT_NAME)

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    # --- セルプロパティ（背景色・罫線） ---
    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)

    tcPr = etree.SubElement(tc, qn("a:tcPr"))
    tcPr.set("marL", "45720")
    tcPr.set("marR", "45720")
    tcPr.set("marT", "18288")
    tcPr.set("marB", "18288")

    # 罫線（下線のみ薄く）
    for border_name in ["a:lnL", "a:lnR", "a:lnT"]:
        ln = etree.SubElement(tcPr, qn(border_name))
        ln.set("w", "0")
        noFill = etree.SubElement(ln, qn("a:noFill"))

    lnB = etree.SubElement(tcPr, qn("a:lnB"))
    lnB.set("w", "6350")  # 0.5pt
    solidFillB = etree.SubElement(lnB, qn("a:solidFill"))
    srgbClrB = etree.SubElement(solidFillB, qn("a:srgbClr"))
    srgbClrB.set("val", "E0E0E0")

    # 背景色
    if bg_color is not None:
        fill = etree.SubElement(tcPr, qn("a:solidFill"))
        clr = etree.SubElement(fill, qn("a:srgbClr"))
        clr.set("val", f"{bg_color[0]:02X}{bg_color[1]:02X}{bg_color[2]:02X}"
                if isinstance(bg_color, (tuple, list))
                else str(bg_color).replace("#", ""))
    else:
        noFill = etree.SubElement(tcPr, qn("a:noFill"))


def build_period_table(slide, data_period, continuous_names, font_sizes,
                       tbl_left, tbl_top, tbl_width, tbl_height):
    """1期分のネイティブテーブルを構築"""
    customers = data_period.get("customers", [])
    unit = data_period.get("unit", "千円")
    n_rows = len(customers) + 1  # ヘッダー + データ行
    n_cols = 4

    # テーブル追加
    shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = shape.table

    # tblPr設定（デフォルトスタイル除去）
    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # 列幅を設定
    for ci, ratio in enumerate(COL_RATIOS):
        table.columns[ci].width = int(tbl_width * ratio)

    # 行の高さを計算
    header_h = int(tbl_height * 0.09)  # ヘッダー行はやや高め
    data_h = int((tbl_height - header_h) / len(customers))

    tbl_xml = shape._element.find('.//' + qn('a:tbl'))
    for i, tr in enumerate(tbl_xml.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(header_h))
        else:
            tr.set('h', str(data_h))

    # ── ヘッダー行 ──
    headers = ["#", "企業名", f"売上高({unit})", "割合"]
    aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]

    for ci, (hdr, align) in enumerate(zip(headers, aligns)):
        set_cell_text(table.cell(0, ci), hdr,
                      font_size=font_sizes["header"],
                      font_bold=True,
                      font_color=COLOR_TEXT,
                      alignment=align,
                      bg_color=COLOR_HEADER_BG)

    # ── データ行 ──
    for ri, cust in enumerate(customers):
        name = cust.get("name", "")
        revenue = cust.get("revenue", 0)
        share = cust.get("share", 0)
        rank = cust.get("rank", ri + 1)

        is_continuous = name in continuous_names and name != "その他"
        font_color = COLOR_ORANGE if is_continuous else COLOR_TEXT
        font_bold = is_continuous

        # フォーマット
        rank_str = str(rank) if name != "その他" else ""
        rev_str = f"{int(revenue):,}" if isinstance(revenue, (int, float)) else str(revenue)
        share_str = f"{share:.1f}%" if isinstance(share, (int, float)) else str(share)

        # 行背景色
        if name == "その他":
            bg = COLOR_OTHER_ROW
        elif ri % 2 == 1:
            bg = COLOR_EVEN_ROW
        else:
            bg = COLOR_WHITE

        row_idx = ri + 1  # ヘッダー行の分
        set_cell_text(table.cell(row_idx, 0), rank_str,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.CENTER, bg)
        set_cell_text(table.cell(row_idx, 1), name,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.LEFT, bg)
        set_cell_text(table.cell(row_idx, 2), rev_str,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.RIGHT, bg)
        set_cell_text(table.cell(row_idx, 3), share_str,
                      font_sizes["body"], font_bold, font_color, PP_ALIGN.RIGHT, bg)

    return shape


def main():
    parser = argparse.ArgumentParser(description="主要販売先売上高スライド生成（ネイティブテーブル方式）")
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument("--template", required=True, help="PPTXテンプレートパス")
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    print("=== 主要販売先売上高スライド生成（ネイティブテーブル方式） ===")

    # テンプレート読み込み
    prs = Presentation(args.template)
    slide = prs.slides[0]

    periods = data.get("periods", [])
    n_periods = len(periods)
    font_sizes = get_font_sizes(n_periods)
    continuous_names = identify_continuous_customers(periods)

    print(f"  期数: {n_periods}")
    print(f"  継続顧客: {continuous_names}")

    # 1. メインメッセージ
    main_msg = data.get("main_message", "")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg)
    print(f"  ✓ メインメッセージ: {main_msg[:50]}...")

    # 2. チャートタイトル
    chart_title = data.get("chart_title", "主要販売先からの完成工事売上高と割合")
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
    print(f"  ✓ チャートタイトル: {chart_title}")

    # 3. 出典
    source_text = data.get("source", "")
    source_shape = find_shape(slide, SHAPE_SOURCE)
    if source_text:
        set_textbox_text(source_shape, f"出典：{source_text}")
    print(f"  ✓ 出典: {source_text}")

    # 4. Content Area を削除
    remove_shape(slide, SHAPE_CONTENT_AREA)

    # 5. 注記テキストボックスを追加
    subtitle_note = data.get("subtitle_note", "")
    continuous_label = data.get("continuous_label", "継続的な顧客")
    continuous_threshold_label = data.get("continuous_threshold_label", "")

    if not subtitle_note and continuous_names:
        if continuous_threshold_label:
            subtitle_note = f"オレンジ字：{continuous_label}（{continuous_threshold_label}）"
        else:
            subtitle_note = f"オレンジ字：{continuous_label}（{_num_to_kanji(n_periods)}期）"

    if subtitle_note:
        note_box = add_textbox(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=NOTE_HEIGHT,
            text=subtitle_note,
            font_size=font_sizes["note"],
            font_bold=True,
            font_color=COLOR_ORANGE,
            alignment=PP_ALIGN.RIGHT,
        )
        print(f"  ✓ 注記: {subtitle_note}")

    # 6. レイアウト計算
    note_offset = NOTE_HEIGHT if subtitle_note else Inches(0)
    period_header_top = CONTENT_TOP + note_offset
    table_top = period_header_top + PERIOD_HEADER_HEIGHT

    available_width = CONTENT_WIDTH
    total_gaps = TABLE_GAP * (n_periods - 1)
    table_width = int((available_width - total_gaps) / n_periods)

    table_height = CONTENT_BOTTOM - table_top - BOTTOM_MARGIN

    print(f"  テーブル幅: {table_width/914400:.2f}\" × {n_periods}期")
    print(f"  テーブル高さ: {table_height/914400:.2f}\"")

    # 7. 各期のテーブルを生成
    for pi, period in enumerate(periods):
        tbl_left = CONTENT_LEFT + (table_width + TABLE_GAP) * pi

        # 期ヘッダー（テキストボックス）
        period_label = period.get("label", f"期{pi+1}")
        hdr_box = add_textbox(
            slide,
            left=tbl_left,
            top=period_header_top,
            width=table_width,
            height=PERIOD_HEADER_HEIGHT,
            text=period_label,
            font_size=font_sizes["period"],
            font_bold=True,
            font_color=COLOR_TEXT,
            alignment=PP_ALIGN.CENTER,
        )
        # 下線を追加（テキストボックスの下罫線として）
        # テキストに下線を設定
        tf = hdr_box.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.underline = True

        print(f"  ✓ 期ヘッダー: {period_label}")

        # テーブル本体
        tbl_shape = build_period_table(
            slide, period, continuous_names, font_sizes,
            tbl_left=tbl_left,
            tbl_top=table_top,
            tbl_width=table_width,
            tbl_height=table_height,
        )
        print(f"  ✓ テーブル[{period_label}]: {len(period.get('customers', []))+1}行 × 4列")

    # 保存
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
