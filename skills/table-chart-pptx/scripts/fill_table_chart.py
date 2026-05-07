"""
fill_table_chart.py — テーブル＋意味合いスライドをPPTXテンプレートに流し込むスクリプト

テンプレート構造（TableChart.pptx 確認済み）:
  - Title 1              (PLACEHOLDER): Main Message
  - Text Placeholder 2   (PLACEHOLDER): Chart Title
  - TextBox 6            (TEXT_BOX):    テーブルセクションラベル（動的変更可）
  - Straight Connector 26(LINE):        水平線（編集不要）
  - TextBox 29           (TEXT_BOX):    意味合い（タイトル + 1〜5個のBullet）
  - Table 10             (TABLE):       動的テーブル（行・列数はコンテンツに応じて増減）
  - Source 3             (TEXT_BOX):    出典 (roleup brand のみ。stella は dynamic textbox)

使い方:
  python fill_table_chart.py \
    --data /home/claude/table_chart_data.json \
    [--template /path/to/table-chart-template.pptx] \
    --output /mnt/user-data/outputs/TableChart_output.pptx \
    [--brand stellar_aiz | roleup]
"""

import argparse
import os
import json
import sys
import copy

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402

SKILL_ID = "table-chart-pptx"
_THEME = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング（確認済み）──────────────────────────────
SHAPE_MAIN_MESSAGE   = "Title 1"
SHAPE_CHART_TITLE    = "Text Placeholder 2"
SHAPE_TABLE_LABEL    = "TextBox 6"
SHAPE_IMPLICATIONS   = "TextBox 29"
SHAPE_TABLE          = "Table 10"
SHAPE_SOURCE_PH      = "Source 3"  # roleup placeholder
# ────────────────────────────────────────────────────────────

# ── Layout / Color / Font (stella defaults; _apply_theme で roleup 用に上書き) ──
SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(12.50)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_HEADER_FILL = None  # None: テンプレデフォルト (stella 既存色) を維持

FONT_NAME_JP = None  # None: テンプレ default 継承 (stella 既存挙動)
FONT_SIZE_TABLE_HEADER = None  # None: テンプレ既存サイズ維持
FONT_SIZE_TABLE_DATA = None
FONT_SIZE_IMPL_TITLE = None
FONT_SIZE_IMPL_BULLET = None
FONT_SIZE_TABLE_LABEL = None
FONT_SIZE_SOURCE = Pt(10)


def _apply_theme(theme):
    """roleup の場合、フォント名・サイズ・出典座標・ヘッダー色を brand 仕様に上書きする。"""
    global SOURCE_X, SOURCE_Y, SOURCE_W
    global COLOR_TEXT, COLOR_SOURCE, COLOR_HEADER_FILL
    global FONT_NAME_JP, FONT_SIZE_TABLE_HEADER, FONT_SIZE_TABLE_DATA
    global FONT_SIZE_IMPL_TITLE, FONT_SIZE_IMPL_BULLET, FONT_SIZE_TABLE_LABEL
    global FONT_SIZE_SOURCE
    global _THEME
    _THEME = theme

    if theme.id != "roleup":
        return

    # A4 横 (11.69 × 8.27) 用に出典を下端へ
    SOURCE_X = Inches(0.41)
    SOURCE_Y = Inches(7.45)
    SOURCE_W = Inches(10.87)

    # roleup 茶系トーン
    COLOR_TEXT = theme.color("text")
    COLOR_SOURCE = theme.color("source")
    COLOR_HEADER_FILL = theme.color("label_bar")  # #7C4C2C ヘッダー塗り

    # フォント
    FONT_NAME_JP = theme.font_ea or "Yu Gothic UI"

    # roleup C4 許容集合 [22, 14, 12, 10, 6] pt
    FONT_SIZE_TABLE_HEADER = Pt(10)
    FONT_SIZE_TABLE_DATA = Pt(10)
    FONT_SIZE_IMPL_TITLE = Pt(12)
    FONT_SIZE_IMPL_BULLET = Pt(10)
    FONT_SIZE_TABLE_LABEL = Pt(12)
    FONT_SIZE_SOURCE = Pt(int(theme._defaults.get("font_size_source_pt", 6)))


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_placeholder_text(shape, text):
    """PlaceholderのTextFrameにテキストをセット（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def set_textbox_text(shape, text):
    """TextBoxの最初のrunのテキストを上書き（スタイル保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text


def _retypeset_label(shape):
    """Table label (TextBox 6) の font.name / size を _THEME から late-resolve."""
    if shape is None or _THEME is None or _THEME.id != "roleup":
        return
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = FONT_NAME_JP
            run.font.size = FONT_SIZE_TABLE_LABEL
            run.font.color.rgb = COLOR_TEXT
            run.font.bold = True


def fill_implications(slide, items):
    """
    TextBox 29の意味合いセクションを動的に構築する。
    テンプレートの構造:
      para[0]: タイトル「意味合い」(bold)
      para[1..N]: Bullet付きImplication

    items: list of str (1〜5個)
    """
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return

    txBody = shape._element.find(qn("p:txBody"))
    paragraphs = txBody.findall(qn("a:p"))

    if len(paragraphs) < 2:
        print("  WARNING: Implications TextBox has insufficient paragraphs", file=sys.stderr)
        return

    # 1番目の段落 (タイトル) の font を _THEME に合わせて上書き
    if _THEME is not None and _THEME.id == "roleup":
        for run in paragraphs[0].findall(qn("a:r")):
            rPr = run.find(qn("a:rPr"))
            if rPr is None:
                rPr = etree.SubElement(run, qn("a:rPr"))
                run.insert(0, rPr)
            rPr.set("sz", str(int(FONT_SIZE_IMPL_TITLE.pt * 100)))
            rPr.set("b", "1")
            for tag in ("a:latin", "a:ea"):
                el = rPr.find(qn(tag))
                if el is not None:
                    rPr.remove(el)
            etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
            etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    # 2番目の段落（最初のBullet）をテンプレートとして保存
    bullet_template = copy.deepcopy(paragraphs[1])

    # タイトル段落以外をすべて削除
    for p in paragraphs[1:]:
        txBody.remove(p)

    # items数に応じてBullet段落を追加（最大5個）
    for i, text in enumerate(items[:5]):
        new_p = copy.deepcopy(bullet_template)
        runs = new_p.findall(qn("a:r"))
        if runs:
            t_elem = runs[0].find(qn("a:t"))
            if t_elem is not None:
                t_elem.text = text
            for run in runs[1:]:
                t_elem2 = run.find(qn("a:t"))
                if t_elem2 is not None:
                    t_elem2.text = ""
            # roleup: 最初の run の font を上書き
            if _THEME is not None and _THEME.id == "roleup":
                rPr = runs[0].find(qn("a:rPr"))
                if rPr is None:
                    rPr = etree.SubElement(runs[0], qn("a:rPr"))
                    runs[0].insert(0, rPr)
                rPr.set("sz", str(int(FONT_SIZE_IMPL_BULLET.pt * 100)))
                for tag in ("a:latin", "a:ea"):
                    el = rPr.find(qn(tag))
                    if el is not None:
                        rPr.remove(el)
                etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
                etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        else:
            r_elem = etree.SubElement(new_p, qn("a:r"))
            sz_attr = str(int(FONT_SIZE_IMPL_BULLET.pt * 100)) if FONT_SIZE_IMPL_BULLET else "1600"
            rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
                "lang": "en-JP", "sz": sz_attr,
            })
            if FONT_NAME_JP:
                etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
                etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
            else:
                etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "+mn-ea"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text
        txBody.append(new_p)
        print(f"  [Implication {i+1}] {text[:60]}{'...' if len(text) > 60 else ''}")


def _override_rPr_font(rPr):
    """rPr 要素に _THEME のフォント名・サイズを上書きする。roleup brand 限定。"""
    if rPr is None or _THEME is None or _THEME.id != "roleup":
        return
    if FONT_SIZE_TABLE_DATA is not None:
        rPr.set("sz", str(int(FONT_SIZE_TABLE_DATA.pt * 100)))
    for tag in ("a:latin", "a:ea"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
    if FONT_NAME_JP:
        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
        etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})


def _override_tcPr_fill(tcPr, fill_hex):
    """tcPr の solidFill を指定 hex に置き換える。"""
    if tcPr is None or fill_hex is None:
        return
    # 既存の solidFill / gradFill / pattFill / blipFill / noFill を全削除
    for tag in ("a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill", "a:noFill"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": fill_hex})
    # tcPr 内の child 順序: 罫線群 → fill → marL/marR/marT/marB 等。
    # 安全のため tcPr の先頭付近 (罫線の後) に挿入する。
    tcPr.remove(sf)
    # 罫線要素 (lnL/lnR/lnT/lnB) の後に挿入
    insert_idx = 0
    for i, child in enumerate(tcPr):
        if child.tag in (qn("a:lnL"), qn("a:lnR"), qn("a:lnT"), qn("a:lnB"), qn("a:lnTlBr"), qn("a:lnBlToTr")):
            insert_idx = i + 1
    tcPr.insert(insert_idx, sf)


def rebuild_table(slide, headers, rows, template_prs=None):
    """
    既存テーブルを削除し、コンテンツに応じた行列数でテーブルを再構築する。
    ヘッダー行・データ行のセルスタイルはテンプレートから複製する。
    """
    table_shape = find_shape(slide, SHAPE_TABLE)
    if table_shape is None:
        print("  WARNING: Table shape not found, cannot rebuild", file=sys.stderr)
        return

    old_table = table_shape.table

    # セルスタイル（tcPr）をコピー
    header_tcPr = copy.deepcopy(old_table.cell(0, 0)._tc.find(qn("a:tcPr")))
    data_tcPr   = copy.deepcopy(old_table.cell(1, 0)._tc.find(qn("a:tcPr")))

    # roleup: ヘッダー塗り色を label_bar に上書き
    if COLOR_HEADER_FILL is not None and header_tcPr is not None:
        header_hex = "{:02X}{:02X}{:02X}".format(
            COLOR_HEADER_FILL[0], COLOR_HEADER_FILL[1], COLOR_HEADER_FILL[2]
        )
        _override_tcPr_fill(header_tcPr, header_hex)

    # ── ヘッダー行の罫線を白に変更（四辺すべて） ──
    if header_tcPr is not None:
        for ln_tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
            ln = header_tcPr.find(qn(ln_tag))
            if ln is not None:
                # 既存罫線の色を白に変更
                sf = ln.find(qn("a:solidFill"))
                if sf is not None:
                    for child in list(sf):
                        sf.remove(child)
                    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": "FFFFFF"})
            else:
                # 罫線が未定義の場合、白の罫線を新規追加
                ln = etree.SubElement(header_tcPr, qn(ln_tag), attrib={
                    "w": "9525", "cap": "flat", "cmpd": "sng", "algn": "ctr"
                })
                sf = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": "FFFFFF"})
                etree.SubElement(ln, qn("a:prstDash"), attrib={"val": "solid"})
                etree.SubElement(ln, qn("a:round"))

    # run properties（フォントサイズ等）をコピー
    header_rPr = None
    data_rPr   = None
    for para in old_table.cell(0, 0).text_frame.paragraphs:
        for run in para.runs:
            header_rPr = copy.deepcopy(run._r.find(qn("a:rPr")))
            break
        break
    for para in old_table.cell(1, 0).text_frame.paragraphs:
        for run in para.runs:
            data_rPr = copy.deepcopy(run._r.find(qn("a:rPr")))
            break
        break

    # ── ヘッダー行のフォント色を白に変更 ──
    if header_rPr is not None:
        # 既存のsolidFillがあれば削除
        old_fill = header_rPr.find(qn("a:solidFill"))
        if old_fill is not None:
            header_rPr.remove(old_fill)
        # 白色のsolidFillを追加
        sf = etree.SubElement(header_rPr, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": "FFFFFF"})

    # roleup: rPr の font name / size を上書き
    _override_rPr_font(header_rPr)
    _override_rPr_font(data_rPr)

    # 位置・サイズ
    tbl_left  = table_shape.left
    tbl_top   = table_shape.top
    tbl_width = table_shape.width
    tbl_height = table_shape.height

    # 既存テーブルを削除
    sp_tree = slide.shapes._spTree
    sp_tree.remove(table_shape._element)

    # 新テーブルを作成
    n_cols = len(headers)
    n_rows = len(rows) + 1  # ヘッダー行 + データ行

    new_shape = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
    )
    new_table = new_shape.table

    # 列幅を均等配分
    col_width = tbl_width // n_cols
    for col in new_table.columns:
        col.width = col_width

    # tblPrを設定
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tbl_elem = new_shape._element.find('.//a:tbl', ns)
    old_tblPr = tbl_elem.find('a:tblPr', ns)
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '1'
    })
    tbl_elem.insert(0, tblPr)

    def apply_cell(cell, text, tcPr_tmpl, rPr_tmpl):
        """セルにテキストとスタイルを適用"""
        tc = cell._tc
        txBody = tc.find(qn("a:txBody"))
        if txBody is None:
            txBody = etree.SubElement(tc, qn("a:txBody"))
            etree.SubElement(txBody, qn("a:bodyPr"))
            etree.SubElement(txBody, qn("a:lstStyle"))

        for p in txBody.findall(qn("a:p")):
            txBody.remove(p)

        p_elem = etree.SubElement(txBody, qn("a:p"))
        etree.SubElement(p_elem, qn("a:pPr"))
        r_elem = etree.SubElement(p_elem, qn("a:r"))
        if rPr_tmpl is not None:
            r_elem.append(copy.deepcopy(rPr_tmpl))
        else:
            sz_attr = str(int(FONT_SIZE_TABLE_DATA.pt * 100)) if FONT_SIZE_TABLE_DATA else "1200"
            new_rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
                "lang": "en-GB", "sz": sz_attr,
            })
            if FONT_NAME_JP:
                etree.SubElement(new_rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
                etree.SubElement(new_rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = str(text)

        old_tc = tc.find(qn("a:tcPr"))
        if old_tc is not None:
            tc.remove(old_tc)
        if tcPr_tmpl is not None:
            tc.append(copy.deepcopy(tcPr_tmpl))

    # ヘッダー行
    for c_idx, h in enumerate(headers):
        apply_cell(new_table.cell(0, c_idx), h, header_tcPr, header_rPr)
    print(f"  [Table Header] {' | '.join(headers)}")

    # データ行
    for r_idx, row_data in enumerate(rows):
        for c_idx in range(n_cols):
            val = row_data[c_idx] if c_idx < len(row_data) else ""
            apply_cell(new_table.cell(r_idx + 1, c_idx), val, data_tcPr, data_rPr)
        print(f"  [Table Row {r_idx+1}] {' | '.join(str(v) for v in row_data[:n_cols])}")

    print(f"  [Table] {n_rows} rows x {n_cols} cols")


def write_source(slide, source):
    """出典を Source 3 placeholder (roleup) または dynamic textbox (stella) に書き込む。"""
    if not source:
        return

    if _THEME is not None and _THEME.is_source_required():
        src_shape = find_shape(slide, SHAPE_SOURCE_PH)
        if src_shape is not None:
            set_textbox_text(src_shape, f"出典: {source}")
            for para in src_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = FONT_SIZE_SOURCE
                    run.font.color.rgb = COLOR_SOURCE
                    if FONT_NAME_JP:
                        run.font.name = FONT_NAME_JP
            print(f"  [Source] {source[:60]}{'...' if len(source) > 60 else ''}")
            return

    # stella: 動的 textbox
    tb = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.30))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"出典: {source}"
    run.font.size = FONT_SIZE_SOURCE
    run.font.color.rgb = COLOR_SOURCE
    if FONT_NAME_JP:
        run.font.name = FONT_NAME_JP
    print(f"  [Source] {source[:60]}{'...' if len(source) > 60 else ''}")


def main():
    parser = argparse.ArgumentParser(description="テーブルチャートデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="table_chart_data.json のパス")
    parser.add_argument("--template", required=False, default=None, help="table-chart-template.pptx のパス (省略時は brand から解決)")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "table"],
        allowed_top=[
            "main_message", "chart_title",
            "table", "table_label", "implications",
            "source", "source_label", "source_text",
        ],
        skill_name=SKILL_ID,
    )

    # Phase 2: brand-aware
    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    require_source(data, theme, skill_id=SKILL_ID)
    template_path = args.template or theme.template_path(SKILL_DIR, "table-chart")

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # ── Top text (stella: main_message / roleup: chart_title) ──
    top_text = resolve_top_text(data, theme).strip()
    if top_text:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(shape, top_text)
        print(f"  [Top]   {top_text[:60]}{'...' if len(top_text) > 60 else ''}")
    else:
        print("  WARNING: top text is empty", file=sys.stderr)

    # ── Subtitle (stella: chart_title / roleup: main_message) ──
    sub_text = resolve_subtitle_text(data, theme).strip()
    if sub_text:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(shape, sub_text)
        print(f"  [Sub]   {sub_text}")
    else:
        print("  WARNING: subtitle is empty", file=sys.stderr)

    # ── Table Label ──
    table_label = data.get("table_label", "Table").strip()
    shape = find_shape(slide, SHAPE_TABLE_LABEL)
    set_textbox_text(shape, table_label)
    _retypeset_label(shape)
    print(f"  [Table Label]  {table_label}")

    # ── Table (動的行列) ──
    table_data = data.get("table", {})
    table_headers = table_data.get("headers", [])
    table_rows    = table_data.get("rows", [])
    if table_headers:
        rebuild_table(slide, table_headers, table_rows)
    else:
        print("  WARNING: table.headers is empty", file=sys.stderr)

    # ── Implications (1〜5個) ──
    implications = data.get("implications", [])
    if implications:
        fill_implications(slide, implications)
    else:
        print("  WARNING: implications is empty", file=sys.stderr)

    # ── Source ──
    source = (data.get("source") or data.get("source_label")
              or data.get("source_text") or "").strip()
    write_source(slide, source)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n  Saved: {args.output}")


if __name__ == "__main__":
    main()
