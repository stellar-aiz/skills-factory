"""
fill_team.py — プロジェクト体制図データをPPTXテンプレートに流し込むスクリプト

テンプレート構造（確認済み）:

[3WG版 ProjectTeamStructure3.pptx]
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Rectangle 3        (AUTO_SHAPE):  Project Sponsor
  - Rectangle 5        (AUTO_SHAPE):  Project Owner
  - Rectangle 6        (AUTO_SHAPE):  PMO
  - TextBox 9          (TEXT_BOX):    意味合い（para[0]=見出し, para[1-5]=Implication）
  - Rectangle 22       (AUTO_SHAPE):  WG2（左）  → ユーザー入力WG[0]
  - Rectangle 12       (AUTO_SHAPE):  WG1（中央）→ ユーザー入力WG[1]
  - Rectangle 23       (AUTO_SHAPE):  WG3（右）  → ユーザー入力WG[2]

[5WG版 ProjectTeamStructure5.pptx]
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Rectangle 3        (AUTO_SHAPE):  Project Sponsor
  - Rectangle 5        (AUTO_SHAPE):  Project Owner
  - Rectangle 6        (AUTO_SHAPE):  PMO
  - TextBox 9          (TEXT_BOX):    意味合い（para[0]=見出し, para[1-5]=Implication）
  - Rectangle 4        (AUTO_SHAPE):  WG4（最左）→ ユーザー入力WG[0]
  - Rectangle 22       (AUTO_SHAPE):  WG2（左寄）→ ユーザー入力WG[1]
  - Rectangle 12       (AUTO_SHAPE):  WG1（中央）→ ユーザー入力WG[2]
  - Rectangle 23       (AUTO_SHAPE):  WG3（右寄）→ ユーザー入力WG[3]
  - Rectangle 8        (AUTO_SHAPE):  WG5（最右）→ ユーザー入力WG[4]

使い方:
  python fill_team.py \\
    --data /home/claude/team_data.json \\
    --template /path/to/template.pptx \\
    --output /mnt/user-data/outputs/ProjectTeamStructure_output.pptx
"""

import argparse
import os
import json
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──────────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_SPONSOR      = "Rectangle 3"
SHAPE_OWNER        = "Rectangle 5"
SHAPE_PMO          = "Rectangle 6"
SHAPE_IMPLICATIONS = "TextBox 9"

# WG Shape名：左から右への並び順にマッピング
WG_SHAPES_3 = ["Rectangle 22", "Rectangle 12", "Rectangle 23"]                           # 左→中央→右
WG_SHAPES_5 = ["Rectangle 4", "Rectangle 22", "Rectangle 12", "Rectangle 23", "Rectangle 8"]  # 最左→左→中央→右→最右
# ─────────────────────────────────────────────────────────


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_placeholder_text(shape, text, bold=False):
    """PlaceholderのTextFrameにテキストをセット（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        if bold:
            para.runs[0].font.bold = True
        for run in para.runs[1:]:
            run.text = ""
    else:
        attrib = {"lang": "ja-JP"}
        if bold:
            attrib["b"] = "1"
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib=attrib)
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def set_shape_text(slide, shape_name, text, bold=False):
    """AUTO_SHAPE（Rectangle）のpara[0]にテキストをセット"""
    shape = find_shape(slide, shape_name)
    if shape is None:
        return
    set_placeholder_text(shape, text, bold=bold)
    print(f"  [{shape_name}] {text[:60]}{'...' if len(text) > 60 else ''}{' (bold)' if bold else ''}")


def fill_wg(slide, shape_name, wg_data):
    """ワーキンググループのShape（para[0]=名前, para[1-5]=メンバー）を埋める"""
    shape = find_shape(slide, shape_name)
    if shape is None:
        return

    tf = shape.text_frame
    wg_name = wg_data.get("name", "").strip()
    members = wg_data.get("members", [])

    # Ensure exactly 5 members (pad with empty strings)
    while len(members) < 5:
        members.append("")

    # para[0] = WG name (bold)
    if tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = wg_name
        tf.paragraphs[0].runs[0].font.bold = True
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(tf.paragraphs[0]._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP", "b": "1"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = wg_name

    # para[1]-para[5] = members
    for i in range(5):
        para_idx = i + 1
        member = members[i].strip() if i < len(members) else ""
        para = tf.paragraphs[para_idx]
        if para.runs:
            para.runs[0].text = member
            for run in para.runs[1:]:
                run.text = ""
        else:
            r_elem = etree.SubElement(para._p, qn("a:r"))
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = member

    member_str = ", ".join(m for m in members[:5] if m)
    print(f"  [{shape_name}] {wg_name}: {member_str}")


def fill_implications(slide, items):
    """TextBox 9のpara[1]〜para[5]に意味合いを書き込む（para[0]は見出し「意味合い」で編集不要）"""
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return

    tf = shape.text_frame

    # Ensure exactly 5 implications (pad with empty strings)
    while len(items) < 5:
        items.append("")

    for i in range(5):
        para_idx = i + 1  # para[0] is header
        text = items[i].strip() if i < len(items) else ""
        para = tf.paragraphs[para_idx]
        if para.runs:
            para.runs[0].text = text
            for run in para.runs[1:]:
                run.text = ""
        else:
            r_elem = etree.SubElement(para._p, qn("a:r"))
            etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text
        if text:
            print(f"  [Implication {i+1}] {text[:55]}{'...' if len(text) > 55 else ''}")


def main():
    parser = argparse.ArgumentParser(description="プロジェクト体制図データをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="team_data.json のパス")
    parser.add_argument("--template", required=True, help="テンプレートPPTXのパス")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "working_groups"],
        allowed_top=[
            "main_message", "chart_title",
            "project_owner", "project_sponsor", "pmo",
            "working_groups", "implications",
        ],
        skill_name="project-team-structure-pptx",
    )

    prs = Presentation(args.template)
    slide = prs.slides[0]

    # ── Main Message ──────────────────────────────────
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(shape, main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # ── Chart Title ───────────────────────────────────
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(shape, chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # ── Sponsor / Owner / PMO （太字）────────────────
    set_shape_text(slide, SHAPE_SPONSOR, data.get("project_sponsor", ""), bold=True)
    set_shape_text(slide, SHAPE_OWNER, data.get("project_owner", ""), bold=True)
    set_shape_text(slide, SHAPE_PMO, data.get("pmo", ""), bold=True)

    # ── Working Groups ───────────────────────────────
    wgs = data.get("working_groups", [])
    num_wgs = len(wgs)

    if num_wgs <= 3:
        wg_shapes = WG_SHAPES_3
    else:
        wg_shapes = WG_SHAPES_5

    for i, shape_name in enumerate(wg_shapes):
        if i < len(wgs):
            fill_wg(slide, shape_name, wgs[i])
        else:
            # 未使用WGはスライドから削除する
            shape = find_shape(slide, shape_name)
            if shape is not None:
                sp = shape._element
                sp.getparent().remove(sp)
                print(f"  [{shape_name}] 未使用のため削除")

    # ── Implications ─────────────────────────────────
    implications = data.get("implications", [])
    fill_implications(slide, implications)

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
