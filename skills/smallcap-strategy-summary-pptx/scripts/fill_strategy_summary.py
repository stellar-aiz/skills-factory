"""
fill_strategy_summary.py — 戦略仮説サマリースライド（コンサル品質、1スライド）

レイアウト（business-model-template ベース）:
  - Title 1            (PLACEHOLDER): Main Message ← PPTX native
  - Text Placeholder 2 (PLACEHOLDER): Chart Title  ← PPTX native
  - Rectangle 4        (AUTO_SHAPE):  4次元カードのSVG画像 ← HTML→Playwright キャプチャ
  - TextBox 9          (TEXT_BOX):    意味合い 3点（label+detail） ← PPTX native
"""

import argparse
import asyncio
import copy
import json
import os
import sys
import tempfile
from html import escape

# validate_fill_input bootstrap (skills/_common/lib)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from validate_fill_input import validate_fill_input  # noqa: E402

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    import shutil, subprocess, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_DIAGRAM_AREA = "Rectangle 4"
SHAPE_IMPLICATIONS = "TextBox 9"

# Rectangle 4 のアスペクト比に合わせて viewport を設定（横長）
VIEWPORT_WIDTH = 1400
VIEWPORT_HEIGHT = 760
DEVICE_SCALE = 2

CONFIDENCE_COLOR = {"high": "#52C41A", "medium": "#FAAD14", "low": "#FF4D4F"}
CONFIDENCE_LABEL = {"high": "High", "medium": "Medium", "low": "Low"}


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_placeholder_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_implications(slide, items):
    """TextBox 9 の Para[1..3] に意味合いを書き込む（fill_business_model.py 流用）"""
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return
    tf = shape.text_frame
    for i, item in enumerate(items[:3]):
        label = item.get("label", "").strip()
        detail = item.get("detail", "").strip()
        para_idx = i + 1
        if para_idx >= len(tf.paragraphs):
            continue
        para = tf.paragraphs[para_idx]
        template_rPr = None
        if para.runs:
            template_rPr = para.runs[0]._r.find(qn("a:rPr"))
        for r in list(para._p.findall(qn("a:r"))):
            para._p.remove(r)

        def _make_run(text_str, bold=False):
            r_elem = etree.SubElement(para._p, qn("a:r"))
            if template_rPr is not None:
                new_rPr = copy.deepcopy(template_rPr)
                if bold:
                    new_rPr.set("b", "1")
                else:
                    new_rPr.attrib.pop("b", None)
                r_elem.insert(0, new_rPr)
            else:
                attrib = {"lang": "ja-JP"}
                if bold:
                    attrib["b"] = "1"
                etree.SubElement(r_elem, qn("a:rPr"), attrib=attrib)
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text_str

        if label:
            _make_run(f"{label}：", bold=True)
            _make_run(detail, bold=False)
        else:
            _make_run(detail, bold=False)
        print(f"  [Implication {i+1}] {(label + ': ' + detail)[:60]}")


def _esc(text):
    return escape(str(text))


def generate_html(data):
    """4 次元カード（2x2 グリッド）のチャート画像のみ。ヘッダー・出典は含めない（PPTX側で扱う）"""
    dimensions = data.get("dimensions", [])
    cards_html = []
    for dim in dimensions:
        label = _esc(dim.get("label", ""))
        summary = _esc(dim.get("summary", ""))
        confidence = dim.get("confidence", "medium")
        conf_color = CONFIDENCE_COLOR.get(confidence, "#999")
        conf_label = CONFIDENCE_LABEL.get(confidence, "?")
        detail_page = dim.get("detail_page")
        if detail_page:
            page_html = f'<span class="page">→ p.{detail_page}</span>'
        else:
            page_html = '<span class="page-pending">詳細ページ準備中</span>'

        cards_html.append(f"""
        <div class="card">
          <div class="card-header">
            <span class="dim-label">{label}</span>
            <span class="conf-badge" style="background:{conf_color};">{conf_label}</span>
          </div>
          <div class="card-body">{summary}</div>
          <div class="card-footer">{page_html}</div>
        </div>
        """)
    cards_grid = "\n".join(cards_html)

    css = """
    * { box-sizing: border-box; margin: 0; padding: 0; font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', 'Yu Gothic UI', sans-serif; }
    body { width: 1400px; height: 760px; padding: 0; background: #FFF; }
    .grid { display: grid; grid-template-columns: 1fr 1fr; grid-template-rows: 1fr 1fr; gap: 14px; height: 760px; padding: 14px; }
    .card { background: #FFF; border-radius: 6px; padding: 22px 26px; box-shadow: 0 1px 3px rgba(0,0,0,0.1); border-top: 5px solid #1565C0; display: flex; flex-direction: column; }
    .card-header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 14px; }
    .dim-label { font-size: 22px; font-weight: 700; color: #0A1B3D; }
    .conf-badge { color: #FFF; font-size: 13px; font-weight: 600; padding: 4px 12px; border-radius: 12px; }
    .card-body { font-size: 16px; color: #333; line-height: 1.7; flex-grow: 1; }
    .card-footer { margin-top: 12px; font-size: 13px; color: #666; }
    .page { color: #1565C0; font-weight: 600; }
    .page-pending { color: #999; font-style: italic; }
    """
    return f"""<!DOCTYPE html>
<html lang="ja"><head><meta charset="utf-8"><style>{css}</style></head>
<body>
  <div class="grid">{cards_grid}</div>
</body></html>"""


async def take_screenshot(html_content, output_path):
    from playwright.async_api import async_playwright
    html_path = os.path.join(tempfile.gettempdir(), "strategy_summary.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={"width": VIEWPORT_WIDTH, "height": VIEWPORT_HEIGHT},
            device_scale_factor=DEVICE_SCALE,
        )
        await page.goto(f"file://{html_path}")
        await page.wait_for_timeout(500)
        await page.screenshot(path=output_path, full_page=False)
        await browser.close()
    print(f"  📸 {output_path} ({os.path.getsize(output_path)} bytes)")
    os.unlink(html_path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "dimensions"],
        allowed_top=["main_message", "chart_title", "dimensions", "implications"],
        skill_name="smallcap-strategy-summary-pptx",
    )

    print("📐 Generating 4-dimension grid HTML (chart only)...")
    html_content = generate_html(data)

    screenshot_path = os.path.join(tempfile.gettempdir(), "strategy_summary.png")
    print(f"📸 Taking screenshot ({VIEWPORT_WIDTH}x{VIEWPORT_HEIGHT})...")
    asyncio.run(take_screenshot(html_content, screenshot_path))

    print("📝 Filling PPTX template...")
    prs = Presentation(args.template)
    slide = prs.slides[0]

    set_placeholder_text(find_shape(slide, SHAPE_MAIN_MESSAGE), data.get("main_message", "").strip())
    set_placeholder_text(find_shape(slide, SHAPE_CHART_TITLE), data.get("chart_title", "戦略仮説サマリー").strip())

    # Implications を PPTX native で書き込む
    fill_implications(slide, data.get("implications", []))

    # チャート画像を Rectangle 4 に挿入
    rect_shape = find_shape(slide, SHAPE_DIAGRAM_AREA)
    if rect_shape:
        if rect_shape.has_text_frame:
            for para in rect_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
        padding = Inches(0.05)
        slide.shapes.add_picture(
            screenshot_path,
            rect_shape.left + padding,
            rect_shape.top + padding,
            rect_shape.width - 2 * padding,
            rect_shape.height - 2 * padding,
        )
        print("  [Diagram] Image inserted")

    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")

    if os.path.exists(screenshot_path):
        os.unlink(screenshot_path)


if __name__ == "__main__":
    main()
