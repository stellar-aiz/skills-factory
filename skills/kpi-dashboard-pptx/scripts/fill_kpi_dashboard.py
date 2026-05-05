"""
fill_kpi_dashboard.py — KPI・メトリクスダッシュボードをPPTXテンプレートに流し込むスクリプト

改善点:
  - ヘッダーバーはオブジェクト内に直接テキストを記載（TextBox分離をやめる）
  - ヘッダー色はAccent2（スキーマカラー）を使用
  - メトリクス値のフォントサイズ: 24pt
  - カード内その他テキスト: 14pt
  - 単位・目標テキストの色: 黒

使い方:
  python fill_kpi_dashboard.py \
    --data /home/claude/kpi_data.json \
    --template <SKILL_DIR>/assets/kpi-dashboard-template.pptx \
    --output /mnt/user-data/outputs/KPIDashboard_output.pptx
"""

import argparse
import json
import sys
import copy
import math

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── レイアウト定数 ──────────────────────────────────────────
MARGIN_LEFT     = 370800
MARGIN_RIGHT    = 370800
SLIDE_WIDTH     = 12192000
SLIDE_HEIGHT    = 6858000
CONTENT_WIDTH   = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

GRID_TOP        = 1420000
GRID_BOTTOM     = 6550000
COL_GAP         = 200000
ROW_GAP         = 200000
COLS_PER_ROW    = 3
ROWS_PER_PAGE   = 2
MAX_CARDS_PER_PAGE = COLS_PER_ROW * ROWS_PER_PAGE  # 6

CARD_WIDTH      = (CONTENT_WIDTH - (COLS_PER_ROW - 1) * COL_GAP) // COLS_PER_ROW
CARD_HEIGHT     = (GRID_BOTTOM - GRID_TOP - (ROWS_PER_PAGE - 1) * ROW_GAP) // ROWS_PER_PAGE

# カード内部レイアウト
HEADER_BAR_H    = 340000
VALUE_TOP_OFF   = 430000
VALUE_HEIGHT    = 550000
UNIT_TOP_OFF    = 980000
UNIT_HEIGHT     = 280000
TARGET_TOP_OFF  = 1260000
TARGET_HEIGHT   = 280000
STATUS_TOP_OFF  = 1560000
STATUS_HEIGHT   = 280000
TREND_TOP_OFF   = 1840000
TREND_HEIGHT    = 280000

CARD_PADDING    = 140000

# テンプレートShape
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"

# ステータス色マッピング
STATUS_COLORS = {
    "達成":    "2E7D32",
    "on_track": "2E7D32",
    "順調":    "2E7D32",
    "注意":    "F57F17",
    "at_risk": "F57F17",
    "要注意":  "F57F17",
    "未達":    "C62828",
    "off_track":"C62828",
    "遅延":    "C62828",
    "危険":    "C62828",
}


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_placeholder_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_rect_with_text(slide, left, top, width, height, text,
                       fill_scheme=None, fill_srgb=None,
                       font_size=1400, bold=False, font_color=None,
                       font_scheme_color=None, border=False,
                       text_inset_l=100000, text_inset_t=70000,
                       text_anchor="t", align="l"):
    """塗りつぶし矩形にテキストを直接記載する"""
    sp_tree = slide.shapes._spTree
    sp = etree.SubElement(sp_tree, qn("p:sp"))

    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    etree.SubElement(nvSpPr, qn("p:cNvPr"), attrib={
        "id": str(id(sp) % 100000 + 300), "name": "DynRect"
    })
    etree.SubElement(nvSpPr, qn("p:cNvSpPr"))
    etree.SubElement(nvSpPr, qn("p:nvPr"))

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    etree.SubElement(xfrm, qn("a:off"), attrib={"x": str(left), "y": str(top)})
    etree.SubElement(xfrm, qn("a:ext"), attrib={"cx": str(width), "cy": str(height)})
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"), attrib={"prst": "rect"})
    etree.SubElement(prstGeom, qn("a:avLst"))

    # 塗りつぶし
    if fill_scheme:
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:schemeClr"), attrib={"val": fill_scheme})
    elif fill_srgb:
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": fill_srgb})
    else:
        etree.SubElement(spPr, qn("a:noFill"))

    # 枠線
    ln = etree.SubElement(spPr, qn("a:ln"))
    if border:
        ln.set("w", "6350")
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": "D0D0D0"})
    else:
        etree.SubElement(ln, qn("a:noFill"))

    # txBody
    txBody = etree.SubElement(sp, qn("p:txBody"))
    etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": str(text_inset_l),
        "rIns": str(text_inset_l),
        "tIns": str(text_inset_t),
        "bIns": "0",
        "rtlCol": "0",
        "anchor": text_anchor,
    })
    etree.SubElement(txBody, qn("a:lstStyle"))

    if text:
        p = etree.SubElement(txBody, qn("a:p"))
        etree.SubElement(p, qn("a:pPr"), attrib={"algn": align})
        r = etree.SubElement(p, qn("a:r"))

        rPr_attrib = {"kumimoji": "1", "lang": "en-GB", "sz": str(font_size)}
        if bold:
            rPr_attrib["b"] = "1"
        rPr = etree.SubElement(r, qn("a:rPr"), attrib=rPr_attrib)

        if font_color:
            sf = etree.SubElement(rPr, qn("a:solidFill"))
            etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": font_color})
        elif font_scheme_color:
            sf = etree.SubElement(rPr, qn("a:solidFill"))
            etree.SubElement(sf, qn("a:schemeClr"), attrib={"val": font_scheme_color})

        etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "+mn-ea"})
        t = etree.SubElement(r, qn("a:t"))
        t.text = text
    else:
        etree.SubElement(txBody, qn("a:p"))

    return sp


def add_textbox(slide, left, top, width, height, text, font_size=1400,
                bold=False, color=None, align="l"):
    """テキストボックスを追加（カード内テキスト用）"""
    sp_tree = slide.shapes._spTree
    sp = etree.SubElement(sp_tree, qn("p:sp"))

    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    etree.SubElement(nvSpPr, qn("p:cNvPr"), attrib={
        "id": str(id(sp) % 100000 + 100), "name": "DynTextBox"
    })
    etree.SubElement(nvSpPr, qn("p:cNvSpPr"), attrib={"txBox": "1"})
    etree.SubElement(nvSpPr, qn("p:nvPr"))

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    etree.SubElement(xfrm, qn("a:off"), attrib={"x": str(left), "y": str(top)})
    etree.SubElement(xfrm, qn("a:ext"), attrib={"cx": str(width), "cy": str(height)})
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"), attrib={"prst": "rect"})
    etree.SubElement(prstGeom, qn("a:avLst"))
    etree.SubElement(spPr, qn("a:noFill"))
    ln = etree.SubElement(spPr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))

    txBody = etree.SubElement(sp, qn("p:txBody"))
    etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square", "lIns": "0", "rIns": "0",
        "tIns": "0", "bIns": "0", "rtlCol": "0"
    })
    etree.SubElement(txBody, qn("a:lstStyle"))

    p = etree.SubElement(txBody, qn("a:p"))
    etree.SubElement(p, qn("a:pPr"), attrib={"algn": align})
    r = etree.SubElement(p, qn("a:r"))

    rPr_attrib = {"kumimoji": "1", "lang": "en-GB", "sz": str(font_size)}
    if bold:
        rPr_attrib["b"] = "1"
    rPr = etree.SubElement(r, qn("a:rPr"), attrib=rPr_attrib)

    if color:
        sf = etree.SubElement(rPr, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color})

    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": "+mn-ea"})
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return sp


def build_kpi_card(slide, card_left, card_top, kpi):
    """1つのKPIカードを構築"""
    name    = kpi.get("name", "KPI")
    value   = str(kpi.get("value", "—"))
    unit    = kpi.get("unit", "")
    target  = kpi.get("target", "")
    status  = kpi.get("status", "")
    trend   = kpi.get("trend", "")

    inner_left  = card_left + CARD_PADDING
    inner_width = CARD_WIDTH - 2 * CARD_PADDING

    # カード背景（白、グレーボーダー）
    add_rect_with_text(slide, card_left, card_top, CARD_WIDTH, CARD_HEIGHT,
                       text="", fill_srgb="FFFFFF", border=True)

    # ヘッダーバー（Accent2塗り + 白文字テキスト直接記載）
    add_rect_with_text(slide, card_left, card_top, CARD_WIDTH, HEADER_BAR_H,
                       text=name,
                       fill_scheme="accent2",
                       font_size=1200, bold=True,
                       font_color="FFFFFF",
                       text_inset_l=100000, text_inset_t=70000)

    # メトリクス値（24pt, Bold, 黒）
    add_textbox(slide, inner_left, card_top + VALUE_TOP_OFF,
                inner_width, VALUE_HEIGHT,
                text=value, font_size=2400, bold=True, color="1A1A1A", align="l")

    # 単位・説明（14pt, 黒）
    if unit:
        add_textbox(slide, inner_left, card_top + UNIT_TOP_OFF,
                    inner_width, UNIT_HEIGHT,
                    text=unit, font_size=1400, bold=False, color="1A1A1A", align="l")

    # 目標値（14pt, 黒）
    if target:
        target_text = f"目標: {target}" if not target.startswith("目標") else target
        add_textbox(slide, inner_left, card_top + TARGET_TOP_OFF,
                    inner_width, TARGET_HEIGHT,
                    text=target_text, font_size=1400, bold=False, color="1A1A1A", align="l")

    # ステータス（14pt, 色付き, Bold）
    if status:
        status_color = STATUS_COLORS.get(status, "1A1A1A")
        status_text = f"● {status}"
        add_textbox(slide, inner_left, card_top + STATUS_TOP_OFF,
                    inner_width, STATUS_HEIGHT,
                    text=status_text, font_size=1400, bold=True, color=status_color, align="l")

    # トレンド（14pt, 色付き）
    if trend:
        trend_color = "1A1A1A"
        if trend.startswith("+") or trend.startswith("↑"):
            trend_color = "2E7D32"
        elif trend.startswith("-") or trend.startswith("↓"):
            trend_color = "C62828"
        add_textbox(slide, inner_left, card_top + TREND_TOP_OFF,
                    inner_width, TREND_HEIGHT,
                    text=trend, font_size=1400, bold=False, color=trend_color, align="l")


def build_kpi_grid(slide, kpis, page_num, total_pages, start_idx):
    """KPIカードをグリッド配置"""
    for i, kpi in enumerate(kpis):
        col = i % COLS_PER_ROW
        row = i // COLS_PER_ROW

        card_left = MARGIN_LEFT + col * (CARD_WIDTH + COL_GAP)
        card_top  = GRID_TOP + row * (CARD_HEIGHT + ROW_GAP)

        build_kpi_card(slide, card_left, card_top, kpi)
        global_idx = start_idx + i + 1
        print(f"  [KPI {global_idx}] {kpi.get('name', 'KPI')}: {kpi.get('value', '—')}")

    print(f"  [Subtotal] {len(kpis)} cards on this page")


def duplicate_slide(prs, template_slide):
    """テンプレートスライドを複製"""
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    new_sp_tree = new_slide.shapes._spTree
    for child in list(new_sp_tree):
        if child.tag != qn("p:nvGrpSpPr") and child.tag != qn("p:grpSpPr"):
            new_sp_tree.remove(child)

    for child in template_slide.shapes._spTree:
        if child.tag != qn("p:nvGrpSpPr") and child.tag != qn("p:grpSpPr"):
            new_sp_tree.append(copy.deepcopy(child))

    return new_slide


def populate_slide(slide, main_msg, chart_title, kpis, page_num, total_pages, start_idx):
    """1枚のスライドにKPIダッシュボードを配置"""
    shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
    set_placeholder_text(shape, main_msg)

    display_title = chart_title
    if total_pages > 1:
        display_title = f"{chart_title}（{page_num}/{total_pages}）"
    shape = find_shape(slide, SHAPE_CHART_TITLE)
    set_placeholder_text(shape, display_title)

    print(f"\n  === Page {page_num}/{total_pages} ===")
    print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    print(f"  [Chart Title]  {display_title}")

    build_kpi_grid(slide, kpis, page_num, total_pages, start_idx)


def main():
    parser = argparse.ArgumentParser(description="KPIダッシュボードデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output",   required=True)
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(args.template)

    main_msg    = data.get("main_message", "").strip()
    chart_title = data.get("chart_title", "").strip()
    all_kpis    = data.get("kpis", [])

    if not all_kpis:
        print("  ERROR: kpis is empty", file=sys.stderr)
        sys.exit(1)
    if not main_msg:
        print("  WARNING: main_message is empty", file=sys.stderr)
    if not chart_title:
        print("  WARNING: chart_title is empty", file=sys.stderr)

    total_pages = max(1, math.ceil(len(all_kpis) / MAX_CARDS_PER_PAGE))
    kpi_chunks = []
    for i in range(total_pages):
        start = i * MAX_CARDS_PER_PAGE
        end   = start + MAX_CARDS_PER_PAGE
        kpi_chunks.append(all_kpis[start:end])

    print(f"  Total: {len(all_kpis)} KPIs -> {total_pages} page(s)")

    template_slide = prs.slides[0]
    extra_slides = []
    for _ in range(1, total_pages):
        extra_slides.append(duplicate_slide(prs, template_slide))

    populate_slide(template_slide, main_msg, chart_title,
                   kpi_chunks[0], 1, total_pages, 0)

    for page_idx, extra_slide in enumerate(extra_slides):
        populate_slide(extra_slide, main_msg, chart_title,
                       kpi_chunks[page_idx + 1],
                       page_idx + 2, total_pages,
                       (page_idx + 1) * MAX_CARDS_PER_PAGE)

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n  Saved: {args.output} ({total_pages} slide(s))")


if __name__ == "__main__":
    main()
