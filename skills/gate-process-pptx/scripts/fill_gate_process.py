"""
fill_gate_process.py — ゲートプロセスデータをPPTXテンプレートに流し込むスクリプト

テンプレート構造:
  共通:
    - Title 1            (PLACEHOLDER): Main Message
    - Text Placeholder 2 (PLACEHOLDER): Chart Title
    - TextBox 9          (TEXT_BOX):    フィルター詳細（名前 + Feature×2 を繰り返し）
    - TextBox 29         (TEXT_BOX):    意味合い（para[0]固定ラベル、para[1-3]にImplication）

  3フィルター版:
    - スライド直下: Oval 11(Filter1), Oval 3(Filter2), Oval 4(Filter3)

  5フィルター版:
    - スライド直下: Oval 11(Filter1), Oval 3(Filter2), Oval 40(Filter3), Oval 52(Filter4), Oval 4(Filter5)

使い方:
  python fill_gate_process.py \
    --data /home/claude/gate_process_data.json \
    --template /mnt/skills/organization/gate-process-pptx/assets/gate-process-3.pptx \
    --output /mnt/user-data/outputs/GateProcess_output.pptx
"""

import argparse
import os
import json
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
from validate_fill_input import validate_fill_input  # noqa: E402
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング（確認済み）──────────────────────────────
SHAPE_MAIN_MESSAGE    = "Title 1"
SHAPE_CHART_TITLE     = "Text Placeholder 2"
SHAPE_FILTER_DETAIL   = "TextBox 9"
SHAPE_IMPLICATIONS    = "TextBox 29"

# ファネル内フィルター名ラベル（3フィルター版: スライド直下）
FUNNEL_OVALS_3 = ["Oval 11", "Oval 3", "Oval 4"]

# ファネル内フィルター名ラベル（5フィルター版: スライド直下）
FUNNEL_OVALS_5 = ["Oval 11", "Oval 3", "Oval 40", "Oval 52", "Oval 4"]
# ────────────────────────────────────────────────────────────


def find_shape(slide, name):
    """スライド直下からShape名でShapeを検索"""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_text_preserving_style(shape, text):
    """既存のrunスタイルを保持しつつテキストを書き換える"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def set_para_text(para, text):
    """段落のテキストを書き換える（既存runスタイル保持）"""
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_filter_details(slide, filters):
    """TextBox 9にフィルター名 + Feature×2を書き込む"""
    shape = find_shape(slide, SHAPE_FILTER_DETAIL)
    if shape is None:
        print(f"  ⚠ WARNING: Shape '{SHAPE_FILTER_DETAIL}' not found", file=sys.stderr)
        return

    tf = shape.text_frame
    expected_paras = len(filters) * 3
    actual_paras = len(tf.paragraphs)

    if actual_paras < expected_paras:
        print(f"  ⚠ WARNING: TextBox 9 has {actual_paras} paragraphs, expected {expected_paras}", file=sys.stderr)

    for i, filt in enumerate(filters):
        base = i * 3
        name     = filt.get("name", "").strip()
        feature1 = filt.get("feature1", "").strip()
        feature2 = filt.get("feature2", "").strip()

        if base < actual_paras:
            set_para_text(tf.paragraphs[base], name)
            print(f"  [Filter {i+1} Name]     {name}")
        if base + 1 < actual_paras:
            set_para_text(tf.paragraphs[base + 1], feature1)
            print(f"  [Filter {i+1} Feature1] {feature1[:50]}{'...' if len(feature1) > 50 else ''}")
        if base + 2 < actual_paras:
            set_para_text(tf.paragraphs[base + 2], feature2)
            print(f"  [Filter {i+1} Feature2] {feature2[:50]}{'...' if len(feature2) > 50 else ''}")


def fill_funnel_labels(slide, filters, is_3_filter):
    """ファネル図内の楕円にフィルター名を書き込む（両テンプレートともスライド直下）"""
    oval_names = FUNNEL_OVALS_3 if is_3_filter else FUNNEL_OVALS_5
    for i, oval_name in enumerate(oval_names):
        if i < len(filters):
            shape = find_shape(slide, oval_name)
            if shape:
                set_text_preserving_style(shape, filters[i]["name"].strip())
                print(f"  [Funnel Oval {oval_name}] {filters[i]['name'].strip()}")
            else:
                print(f"  ⚠ WARNING: Oval '{oval_name}' not found on slide", file=sys.stderr)


def fill_implications(slide, implications):
    """TextBox 29のpara[1]〜para[3]に意味合いを書き込む"""
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        print(f"  ⚠ WARNING: Shape '{SHAPE_IMPLICATIONS}' not found", file=sys.stderr)
        return

    tf = shape.text_frame
    for i, impl in enumerate(implications[:3]):
        para_idx = i + 1  # para[0]は「意味合い」固定ラベル
        if para_idx < len(tf.paragraphs):
            set_para_text(tf.paragraphs[para_idx], impl.strip())
            print(f"  [Implication {i+1}] {impl.strip()[:55]}{'...' if len(impl.strip()) > 55 else ''}")


def main():
    parser = argparse.ArgumentParser(description="ゲートプロセスデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="gate_process_data.json のパス")
    parser.add_argument("--template", required=True, help="gate-process-3.pptx or gate-process-5.pptx のパス")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main_message", "filters"],
        allowed_top=["main_message", "chart_title", "filters", "implications"],
        skill_name="gate-process-pptx",
    )

    prs = Presentation(args.template)
    slide = prs.slides[0]

    filters = data.get("filters", [])
    num_filters = len(filters)

    if num_filters not in (3, 5):
        print(f"  ⚠ WARNING: フィルター数が {num_filters} です（3 or 5 推奨）", file=sys.stderr)

    is_3_filter = (num_filters <= 3)

    # ── Main Message ──────────────────────────────────────────
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_text_preserving_style(shape, main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # ── Chart Title ───────────────────────────────────────────
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_text_preserving_style(shape, chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # ── Filter Details (TextBox 9) ────────────────────────────
    fill_filter_details(slide, filters)

    # ── Funnel Labels (Ovals) ─────────────────────────────────
    fill_funnel_labels(slide, filters, is_3_filter)

    # ── Implications (TextBox 29) ─────────────────────────────
    implications = data.get("implications", [])
    if len(implications) != 3:
        print(f"  ⚠ implications は {len(implications)} 件です（3件推奨）", file=sys.stderr)
    fill_implications(slide, implications)

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n✅ 保存しました: {args.output}")


if __name__ == "__main__":
    main()
