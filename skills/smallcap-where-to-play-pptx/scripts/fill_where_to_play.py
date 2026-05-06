"""
fill_where_to_play.py — Where to play 詳細スライド3枚（コンサル品質）

各スライド共通レイアウト（business-model-template ベース）:
  - Title 1            (PLACEHOLDER): page-specific Main Message ← PPTX native
  - Text Placeholder 2 (PLACEHOLDER): Chart Title (1/3, 2/3, 3/3) ← PPTX native
  - Rectangle 4        (AUTO_SHAPE):  チャート画像（HTML キャプチャ）
  - TextBox 9          (TEXT_BOX):    意味合い 3点（label+detail）← PPTX native

JSON スキーマ（コンサル品質、narrative 系廃止）:
  {
    "main": {
      "main_message": "...",
      "chart_title": "...",
      "implications": [{"label": "...", "detail": "..."}, ...3個],
      "visual_data": { /* 事業領域マップ用 */ }
    },
    "detail": { /* main と同形式、visual_data は補足図 */ },
    "evidence": { /* main と同形式、visual_data は findings table 用 */ }
  }
"""

import argparse
import asyncio
import copy
import json
import os
import sys
import tempfile
from html import escape

# validate_fill_input bootstrap (skills/_common/lib)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from validate_fill_input import validate_fill_input  # noqa: E402

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    import shutil, subprocess, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_DIAGRAM_AREA = "Rectangle 4"
SHAPE_IMPLICATIONS = "TextBox 9"

# Rectangle 4 のアスペクト比に合わせたチャート専用ビューポート
VIEWPORT_WIDTH = 1400
VIEWPORT_HEIGHT = 760
DEVICE_SCALE = 2

# Evidence ページ用（フルワイド）: TextBox 9 を覆って横長キャンバスに
EVIDENCE_VIEWPORT_WIDTH = 2200
EVIDENCE_VIEWPORT_HEIGHT = 760

CONFIDENCE_COLOR = {"high": "#52C41A", "medium": "#FAAD14", "low": "#FF4D4F"}

DIMENSION_TITLE = "Where to play：事業領域・顧客・地域の選択"


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_placeholder_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_implications(slide, items):
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return
    tf = shape.text_frame
    for i, item in enumerate(items[:3]):
        label = item.get("label", "").strip()
        detail = item.get("detail", "").strip()
        para_idx = i + 1
        if para_idx >= len(tf.paragraphs):
            continue
        para = tf.paragraphs[para_idx]
        template_rPr = None
        if para.runs:
            template_rPr = para.runs[0]._r.find(qn("a:rPr"))
        for r in list(para._p.findall(qn("a:r"))):
            para._p.remove(r)

        def _make_run(text_str, bold=False):
            r_elem = etree.SubElement(para._p, qn("a:r"))
            if template_rPr is not None:
                new_rPr = copy.deepcopy(template_rPr)
                if bold:
                    new_rPr.set("b", "1")
                else:
                    new_rPr.attrib.pop("b", None)
                r_elem.insert(0, new_rPr)
            else:
                attrib = {"lang": "ja-JP"}
                if bold:
                    attrib["b"] = "1"
                etree.SubElement(r_elem, qn("a:rPr"), attrib=attrib)
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text_str

        if label:
            _make_run(f"{label}：", bold=True)
            _make_run(detail, bold=False)
        else:
            _make_run(detail, bold=False)


def dup_slide(prs, tmpl):
    ns = prs.slides.add_slide(tmpl.slide_layout)
    nsp = ns.shapes._spTree
    for c in list(nsp):
        if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr"):
            nsp.remove(c)
    for c in tmpl.shapes._spTree:
        if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr"):
            nsp.append(copy.deepcopy(c))
    return ns


def _esc(text):
    return escape(str(text))


CSS_BASE = """
* { box-sizing: border-box; margin: 0; padding: 0; font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', 'Yu Gothic UI', sans-serif; }
body { width: 1400px; height: 760px; padding: 0; background: #FFF; }
"""


# ──────────────────────────────────────────────────
# Page 1 (Main): 事業領域マップ SVG（チャートのみ）
# ──────────────────────────────────────────────────

def render_main_chart(visual):
    """事業領域マップ（X=BtoB↔BtoC、Y=国内↔海外）のSVGを画面いっぱいに描画"""
    x_left = _esc(visual.get("x_axis_left", "BtoB"))
    x_right = _esc(visual.get("x_axis_right", "BtoC"))
    y_bottom = _esc(visual.get("y_axis_bottom", "国内"))
    y_top = _esc(visual.get("y_axis_top", "海外"))
    x_label = _esc(visual.get("x_axis_label", "顧客タイプ"))
    y_label = _esc(visual.get("y_axis_label", "地理"))
    segments = visual.get("segments", [])

    svg_w, svg_h = 1300, 720
    pad_left, pad_right, pad_top, pad_bottom = 100, 60, 70, 90
    inner_w = svg_w - pad_left - pad_right
    inner_h = svg_h - pad_top - pad_bottom

    parts = [
        f'<svg width="{svg_w}" height="{svg_h}" xmlns="http://www.w3.org/2000/svg">',
        # 4象限の薄色塗り分け
        f'<rect x="{pad_left}" y="{pad_top}" width="{inner_w/2}" height="{inner_h/2}" fill="#F5F5F5"/>',
        f'<rect x="{pad_left+inner_w/2}" y="{pad_top}" width="{inner_w/2}" height="{inner_h/2}" fill="#E8F4FF"/>',
        f'<rect x="{pad_left}" y="{pad_top+inner_h/2}" width="{inner_w/2}" height="{inner_h/2}" fill="#F9F9F9"/>',
        f'<rect x="{pad_left+inner_w/2}" y="{pad_top+inner_h/2}" width="{inner_w/2}" height="{inner_h/2}" fill="#F5F5F5"/>',
        # 中央十字
        f'<line x1="{pad_left+inner_w/2}" y1="{pad_top}" x2="{pad_left+inner_w/2}" y2="{pad_top+inner_h}" stroke="#999" stroke-width="1.5" stroke-dasharray="6 4"/>',
        f'<line x1="{pad_left}" y1="{pad_top+inner_h/2}" x2="{pad_left+inner_w}" y2="{pad_top+inner_h/2}" stroke="#999" stroke-width="1.5" stroke-dasharray="6 4"/>',
        # 外枠
        f'<rect x="{pad_left}" y="{pad_top}" width="{inner_w}" height="{inner_h}" fill="none" stroke="#444" stroke-width="2.5"/>',
        # 軸ラベル
        f'<text x="{svg_w/2}" y="40" text-anchor="middle" font-size="22" font-weight="700" fill="#0A1B3D">{x_label}</text>',
        f'<text x="{pad_left+inner_w/4}" y="{svg_h-30}" text-anchor="middle" font-size="18" font-weight="600" fill="#555">{x_left}</text>',
        f'<text x="{pad_left+inner_w*3/4}" y="{svg_h-30}" text-anchor="middle" font-size="18" font-weight="600" fill="#555">{x_right}</text>',
        f'<text x="{pad_left-20}" y="{pad_top+inner_h/4+8}" text-anchor="end" font-size="18" font-weight="600" fill="#555">{y_top}</text>',
        f'<text x="{pad_left-20}" y="{pad_top+inner_h*3/4+8}" text-anchor="end" font-size="18" font-weight="600" fill="#555">{y_bottom}</text>',
        f'<text x="35" y="{pad_top+inner_h/2}" text-anchor="middle" font-size="22" font-weight="700" fill="#0A1B3D" transform="rotate(-90 35 {pad_top+inner_h/2})">{y_label}</text>',
    ]

    for seg in segments:
        cx = pad_left + max(0.0, min(1.0, float(seg.get("x", 0.5)))) * inner_w
        cy = pad_top + (1 - max(0.0, min(1.0, float(seg.get("y", 0.5))))) * inner_h
        size = seg.get("size", 5)
        radius = max(35, min(85, size * 6))
        highlight = seg.get("highlight", False)
        fill = "#1565C0" if highlight else "#BBBBBB"
        opacity = "0.85" if highlight else "0.45"
        stroke = "#0A1B3D" if highlight else "#777"
        name = _esc(seg.get("name", ""))
        note = _esc(seg.get("note", ""))
        parts.append(f'<circle cx="{cx}" cy="{cy}" r="{radius}" fill="{fill}" fill-opacity="{opacity}" stroke="{stroke}" stroke-width="2.5"/>')
        text_color = "#FFF" if highlight else "#444"
        text_weight = "700" if highlight else "500"
        parts.append(f'<text x="{cx}" y="{cy}" text-anchor="middle" dominant-baseline="middle" font-size="17" font-weight="{text_weight}" fill="{text_color}">{name}</text>')
        if note:
            parts.append(f'<text x="{cx}" y="{cy + radius + 22}" text-anchor="middle" font-size="14" fill="#666" font-style="italic">{note}</text>')

    parts.append("</svg>")
    return f'<!DOCTYPE html><html><head><meta charset="utf-8"><style>{CSS_BASE}.center {{ display:flex; align-items:center; justify-content:center; height:760px; }}</style></head><body><div class="center">{"".join(parts)}</div></body></html>'


# ──────────────────────────────────────────────────
# Page 2 (Detail): 補足図（事業領域別の特徴を構造化した小マトリクス）
# ──────────────────────────────────────────────────

def render_detail_chart(visual):
    """事業領域別の構造（領域名 / 注力度 / 主な特徴）を構造化マトリクスで表示"""
    segments = visual.get("segments", [])
    rows_html = []
    for seg in segments:
        highlight = seg.get("highlight", False)
        name = _esc(seg.get("name", ""))
        note = _esc(seg.get("note", ""))
        focus = "★★★ 注力" if highlight else "△ 後退/縮退"
        focus_color = "#1565C0" if highlight else "#999"
        bg = "#E8F4FF" if highlight else "#F5F5F5"
        rows_html.append(f"""
        <div class="row" style="background:{bg};">
          <div class="cell cell-name">{name}</div>
          <div class="cell cell-focus" style="color:{focus_color}; font-weight:700;">{focus}</div>
          <div class="cell cell-note">{note or "—"}</div>
        </div>
        """)

    css = CSS_BASE + """
    .matrix { padding: 40px 50px; }
    .matrix-title { font-size: 22px; font-weight: 700; color: #0A1B3D; margin-bottom: 20px; padding-bottom: 8px; border-bottom: 3px solid #1565C0; }
    .header-row { display: grid; grid-template-columns: 2fr 1.5fr 3fr; gap: 8px; padding: 12px 16px; background: #0A1B3D; color: #FFF; font-weight: 700; font-size: 16px; border-radius: 4px 4px 0 0; }
    .row { display: grid; grid-template-columns: 2fr 1.5fr 3fr; gap: 8px; padding: 18px 16px; border-bottom: 1px solid #E0E0E0; align-items: center; }
    .cell { font-size: 17px; line-height: 1.6; }
    .cell-name { font-weight: 700; color: #0A1B3D; }
    .cell-focus { font-size: 16px; }
    .cell-note { color: #555; font-style: italic; }
    """
    rows = "\n".join(rows_html)
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head>
<body><div class="matrix">
  <div class="matrix-title">事業領域別の注力度と特徴</div>
  <div class="header-row"><div>事業領域</div><div>注力度</div><div>特徴・補足</div></div>
  {rows}
</div></body></html>"""


# ──────────────────────────────────────────────────
# Page 3 (Evidence): findings 一覧テーブル
# ──────────────────────────────────────────────────

def render_evidence_chart(findings):
    """根拠 finding を表組みで表示（フルワイド：Rect4 + TextBox9 領域を覆う）"""
    rows_html = []
    for f in findings:
        confidence = f.get("confidence", "medium")
        conf_color = CONFIDENCE_COLOR.get(confidence, "#999")
        rows_html.append(f"""
        <tr>
          <td class="td-id">{_esc(f.get("id", ""))}</td>
          <td class="td-agent">{_esc(f.get("agent", ""))}</td>
          <td class="td-source">{_esc(f.get("source", ""))}</td>
          <td><span class="src-type">{_esc(f.get("source_type", ""))}</span></td>
          <td><span class="conf" style="background:{conf_color};">{_esc(confidence)}</span></td>
          <td class="td-excerpt">{_esc(f.get("excerpt", ""))}</td>
        </tr>
        """)

    # Evidence は横長 2200x760 ベース。フォントとセル幅を大きめに
    css = """
    * { box-sizing: border-box; margin: 0; padding: 0; font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', 'Yu Gothic UI', sans-serif; }
    body { width: 2200px; height: 760px; padding: 0; background: #FFF; }
    .ev-wrap { padding: 24px 36px; }
    table { width: 100%; border-collapse: collapse; font-size: 17px; }
    th { background: #0A1B3D; color: #FFF; padding: 14px 12px; text-align: left; font-weight: 700; font-size: 16px; }
    td { padding: 13px 12px; border-bottom: 1px solid #E0E0E0; vertical-align: top; line-height: 1.6; color: #333; }
    tr:nth-child(even) td { background: #FAFAFA; }
    .td-id { font-weight: 700; color: #1565C0; width: 70px; font-size: 18px; }
    .td-agent { color: #666; font-size: 14px; width: 180px; }
    .td-source { font-size: 14px; color: #555; width: 280px; }
    .src-type { display: inline-block; background: #E8F4FF; color: #1565C0; font-size: 13px; padding: 4px 10px; border-radius: 4px; font-weight: 600; }
    .conf { color: #FFF; font-size: 13px; font-weight: 700; padding: 4px 11px; border-radius: 10px; }
    .td-excerpt { font-size: 15px; color: #444; line-height: 1.55; }
    """
    rows = "\n".join(rows_html)
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head>
<body><div class="ev-wrap">
  <table>
    <thead><tr><th>F#</th><th>Agent</th><th>Source</th><th>Type</th><th>Confidence</th><th>抜粋</th></tr></thead>
    <tbody>{rows}</tbody>
  </table>
</div></body></html>"""


# ──────────────────────────────────────────────────
# Playwright
# ──────────────────────────────────────────────────

async def take_screenshot(html_content, output_path, suffix="", vw=None, vh=None):
    from playwright.async_api import async_playwright
    vw = vw or VIEWPORT_WIDTH
    vh = vh or VIEWPORT_HEIGHT
    html_path = os.path.join(tempfile.gettempdir(), f"where_to_play{suffix}.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={"width": vw, "height": vh},
            device_scale_factor=DEVICE_SCALE,
        )
        await page.goto(f"file://{html_path}")
        await page.wait_for_timeout(500)
        await page.screenshot(path=output_path, full_page=False)
        await browser.close()
    print(f"  📸 {output_path} ({os.path.getsize(output_path)} bytes, vp {vw}x{vh})")
    os.unlink(html_path)


def clear_textbox(shape):
    """TextBox の中身を完全に空にする（Implications 不要時の Evidence ページ用）"""
    if shape is None or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for r in list(para._p.findall(qn("a:r"))):
            para._p.remove(r)


def fill_slide(slide, main_msg, chart_subtitle, implications, screenshot_path, skip_implications=False):
    """skip_implications=True なら TextBox 9 を空にしてチャートを横長フルワイドに広げる（Evidence ページ用）"""
    set_placeholder_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_msg)
    set_placeholder_text(find_shape(slide, SHAPE_CHART_TITLE), chart_subtitle)
    if skip_implications:
        clear_textbox(find_shape(slide, SHAPE_IMPLICATIONS))
    else:
        fill_implications(slide, implications)
    rect_shape = find_shape(slide, SHAPE_DIAGRAM_AREA)
    if rect_shape:
        if rect_shape.has_text_frame:
            for para in rect_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
        padding = Inches(0.05)

        # Evidence ページ: Rectangle 4 + TextBox 9 領域を統合してフルワイド表示
        if skip_implications:
            tb9_shape = find_shape(slide, SHAPE_IMPLICATIONS)
            if tb9_shape is not None:
                left = rect_shape.left
                top = min(rect_shape.top, tb9_shape.top)
                right = tb9_shape.left + tb9_shape.width
                bottom = max(rect_shape.top + rect_shape.height, tb9_shape.top + tb9_shape.height)
                full_width = right - left
                full_height = bottom - top
                slide.shapes.add_picture(
                    screenshot_path,
                    left + padding,
                    top + padding,
                    full_width - 2 * padding,
                    full_height - 2 * padding,
                )
                return
        # 通常のページ
        slide.shapes.add_picture(
            screenshot_path,
            rect_shape.left + padding,
            rect_shape.top + padding,
            rect_shape.width - 2 * padding,
            rect_shape.height - 2 * padding,
        )


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ISSUE-012 (2026-05-06): スキーマ齟齬の silent fail 防止
    validate_fill_input(
        data,
        required_top=["main", "detail", "evidence"],
        allowed_top=["main", "detail", "evidence"],
        skill_name="smallcap-where-to-play-pptx",
    )

    main_data = data.get("main", {})
    detail_data = data.get("detail", {})
    evidence_data = data.get("evidence", {})

    # チャート HTML 生成（チャートのみ）
    print("📐 Generating 3 chart HTMLs...")
    html_main = render_main_chart(main_data.get("visual_data", {}))
    html_detail = render_detail_chart(detail_data.get("visual_data", main_data.get("visual_data", {})))
    html_evidence = render_evidence_chart(evidence_data.get("findings", []))

    paths = [
        os.path.join(tempfile.gettempdir(), "wtp_main.png"),
        os.path.join(tempfile.gettempdir(), "wtp_detail.png"),
        os.path.join(tempfile.gettempdir(), "wtp_evidence.png"),
    ]

    print("📸 Taking 3 screenshots...")
    asyncio.run(take_screenshot(html_main, paths[0], "_main"))
    asyncio.run(take_screenshot(html_detail, paths[1], "_detail"))
    asyncio.run(take_screenshot(html_evidence, paths[2], "_evidence", vw=EVIDENCE_VIEWPORT_WIDTH, vh=EVIDENCE_VIEWPORT_HEIGHT))

    print("📝 Filling PPTX (3 slides)...")
    prs = Presentation(args.template)
    tmpl = prs.slides[0]
    extras = [dup_slide(prs, tmpl) for _ in range(2)]

    fill_slide(
        tmpl,
        main_data.get("main_message", ""),
        main_data.get("chart_title", f"{DIMENSION_TITLE}（1/3）"),
        main_data.get("implications", []),
        paths[0],
    )
    print(f"  [1/3] Main")

    fill_slide(
        extras[0],
        detail_data.get("main_message", ""),
        detail_data.get("chart_title", f"{DIMENSION_TITLE}（2/3）"),
        detail_data.get("implications", []),
        paths[1],
    )
    print(f"  [2/3] Detail")

    fill_slide(
        extras[1],
        evidence_data.get("main_message", ""),
        evidence_data.get("chart_title", f"{DIMENSION_TITLE}（3/3）"),
        evidence_data.get("implications", []),
        paths[2],
        skip_implications=True,  # Evidence: チャート（表）のみ、Implications は使わない
    )
    print(f"  [3/3] Evidence (chart only)")

    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output} (3 slides)")

    for p in paths:
        if os.path.exists(p):
            os.unlink(p)


if __name__ == "__main__":
    main()
