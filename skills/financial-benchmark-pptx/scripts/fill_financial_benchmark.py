"""
fill_financial_benchmark.py — 財務ベンチマーク比較スライドをPPTXネイティブオブジェクトで生成

レイアウト:
  - 上部: メインメッセージ + チャートタイトル
  - 2×3グリッド: 6つの小型バーチャート（1指標/1チャート）
      [売上高]    [成長率CAGR]   [営業利益率]
      [EBITDA率]  [ROE]          [自己資本比率]
  - 下部: 出典

各小型チャート:
  - タイトル（指標名 + 単位）
  - 水平バーチャート（MSO_SHAPE.RECTANGLE で手動描画）
  - 対象会社は赤でハイライト、その他は紺
  - 企業名（左）+ 値ラベル（バーの右端）

Usage:
  python fill_financial_benchmark.py \
    --data /home/claude/financial_benchmark_data.json \
    --template <path>/financial-benchmark-template.pptx \
    --output /mnt/user-data/outputs/FinancialBenchmark_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Layout Constants ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

# 6チャートを2x3グリッドで配置
GRID_LEFT = Inches(0.41)
GRID_TOP = Inches(1.50)
GRID_WIDTH = Inches(12.51)
GRID_HEIGHT = Inches(5.40)

N_COLS = 3
N_ROWS = 2
CELL_GAP_X = Inches(0.15)
CELL_GAP_Y = Inches(0.20)

CELL_W = (GRID_WIDTH - CELL_GAP_X * (N_COLS - 1)) / N_COLS
CELL_H = (GRID_HEIGHT - CELL_GAP_Y * (N_ROWS - 1)) / N_ROWS

# 各セル内部のレイアウト
CHART_TITLE_H = Inches(0.35)
BAR_AREA_TOP_MARGIN = Inches(0.10)
BAR_AREA_BOTTOM_MARGIN = Inches(0.05)
BAR_AREA_LEFT_MARGIN = Inches(0.10)
BAR_AREA_RIGHT_MARGIN = Inches(0.10)

# バーのサイズ
LABEL_COL_W = Inches(1.15)   # 企業名列の幅
VALUE_COL_W = Inches(0.85)   # 値ラベル列の幅

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.00)
SOURCE_W = Inches(12.50)

# ── Colors ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CELL_BG = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_CELL_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
COLOR_TITLE_UNDERLINE = RGBColor(0x33, 0x33, 0x33)

COLOR_BAR_DEFAULT = RGBColor(0x4E, 0x79, 0xA7)  # 紺 - 通常企業
COLOR_BAR_TARGET = RGBColor(0xE1, 0x57, 0x59)   # 赤 - 対象会社
COLOR_BAR_NEGATIVE = RGBColor(0xB8, 0x3A, 0x3A) # 濃赤 - マイナス値のバー
COLOR_BAR_TARGET_NEG = RGBColor(0x8B, 0x2C, 0x2E)  # 濃赤 - 対象会社のマイナス値

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_CHART_TITLE = Pt(11)
FONT_SIZE_COMPANY = Pt(10)
FONT_SIZE_VALUE = Pt(10)
FONT_SIZE_SOURCE = Pt(10)


# ──────────────────────────────────────────────
# Utility
# ──────────────────────────────────────────────
def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                 font_name=FONT_NAME_JP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = COLOR_TEXT
    return tb


def format_value(val, unit="", decimals=1, show_sign=False):
    """数値を整形して文字列に変換"""
    if val is None:
        return "—"

    # decimalsに応じてフォーマット
    if decimals == 0:
        num_str = f"{val:,.0f}"
    else:
        num_str = f"{val:,.{decimals}f}"

    # 符号付き表示
    if show_sign and val > 0 and not num_str.startswith("+"):
        num_str = "+" + num_str

    # 単位付加
    if unit:
        return f"{num_str}{unit}"
    return num_str


# ──────────────────────────────────────────────
# Single Bar Chart Cell
# ──────────────────────────────────────────────
def draw_bar_chart_cell(slide, metric, companies, target_company,
                         left, top, width, height):
    """
    1つの指標に対する水平バーチャートを描画する
    """
    # セル背景（薄い枠線付き矩形）
    cell_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    cell_bg.fill.solid()
    cell_bg.fill.fore_color.rgb = COLOR_CELL_BG
    cell_bg.line.color.rgb = COLOR_CELL_BORDER
    cell_bg.line.width = Pt(0.5)
    cell_bg.shadow.inherit = False
    cell_bg.text_frame.text = ""

    # 指標タイトル（上部）
    metric_name = metric.get("name", "指標")
    unit = metric.get("unit", "")
    title_text = f"{metric_name}（{unit}）" if unit else metric_name

    add_text_box(
        slide, title_text,
        left + Inches(0.10), top + Inches(0.05),
        width - Inches(0.20), CHART_TITLE_H,
        FONT_SIZE_CHART_TITLE, bold=True, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    )

    # タイトル下の区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.10), top + CHART_TITLE_H + Inches(0.03),
        width - Inches(0.20), Emu(int(Inches(0.01))),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TITLE_UNDERLINE
    line.line.fill.background()

    # バーエリアの計算
    bar_area_top = top + CHART_TITLE_H + BAR_AREA_TOP_MARGIN + Inches(0.05)
    bar_area_h = height - CHART_TITLE_H - BAR_AREA_TOP_MARGIN - BAR_AREA_BOTTOM_MARGIN - Inches(0.05)
    bar_area_left = left + BAR_AREA_LEFT_MARGIN
    bar_area_right = left + width - BAR_AREA_RIGHT_MARGIN

    # バーの開始位置（企業名列の右）
    bar_start_x = bar_area_left + LABEL_COL_W
    # バーの最大幅（値ラベル列を除く）
    bar_max_w = bar_area_right - bar_start_x - VALUE_COL_W

    # 値と最大値を取得
    values_dict = metric.get("values", {})
    decimals = metric.get("decimals", 1)
    show_sign = metric.get("show_sign", False)

    # 各社の値を取り出し
    rows = []
    for comp in companies:
        name = comp["name"]
        val = values_dict.get(name)
        rows.append({"name": name, "value": val})

    # オプションでソート
    sort_order = metric.get("sort", "keep")  # "keep", "desc", "asc"
    if sort_order == "desc":
        rows.sort(key=lambda r: (r["value"] if r["value"] is not None else float("-inf")), reverse=True)
    elif sort_order == "asc":
        rows.sort(key=lambda r: (r["value"] if r["value"] is not None else float("inf")))

    # 最大絶対値を計算（バー幅の基準）
    vals_numeric = [r["value"] for r in rows if r["value"] is not None]
    if not vals_numeric:
        return
    max_abs = max(abs(v) for v in vals_numeric) or 1
    has_negative = any(v < 0 for v in vals_numeric)

    # 行の高さを計算
    n_rows = len(rows)
    row_h = bar_area_h // n_rows
    bar_h = Emu(int(row_h * 0.60))  # バーの高さは行の60%

    # 負の値がある場合、ゼロラインを中央に設定
    if has_negative:
        zero_x = bar_start_x + bar_max_w // 2
        half_bar_max_w = bar_max_w // 2
    else:
        zero_x = bar_start_x
        half_bar_max_w = bar_max_w

    for i, row in enumerate(rows):
        name = row["name"]
        val = row["value"]
        is_target = (name == target_company) if target_company else False

        row_y = bar_area_top + row_h * i
        cell_y_center = row_y + row_h // 2
        bar_y = cell_y_center - bar_h // 2

        # 企業名ラベル（左）
        label_x = bar_area_left
        label_tb = slide.shapes.add_textbox(
            label_x, row_y, LABEL_COL_W, row_h,
        )
        ltf = label_tb.text_frame
        ltf.word_wrap = False
        ltf.margin_left = 0; ltf.margin_right = Inches(0.05)
        ltf.margin_top = 0; ltf.margin_bottom = 0
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.RIGHT
        lrun = lp.add_run()
        lrun.text = name
        lrun.font.size = FONT_SIZE_COMPANY
        lrun.font.bold = True if is_target else False
        lrun.font.name = FONT_NAME_JP
        lrun.font.color.rgb = COLOR_BAR_TARGET if is_target else COLOR_TEXT

        # バー描画
        if val is None:
            # N/A表示
            na_tb = slide.shapes.add_textbox(
                bar_start_x, row_y, bar_max_w, row_h,
            )
            ntf = na_tb.text_frame
            ntf.margin_left = Inches(0.05); ntf.margin_right = 0
            ntf.margin_top = 0; ntf.margin_bottom = 0
            ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np = ntf.paragraphs[0]
            np.alignment = PP_ALIGN.LEFT
            nrun = np.add_run()
            nrun.text = "N/A"
            nrun.font.size = FONT_SIZE_VALUE
            nrun.font.name = FONT_NAME_JP
            nrun.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            nrun.font.italic = True
            continue

        # バーの長さ
        bar_length = Emu(int(half_bar_max_w * abs(val) / max_abs))

        if val >= 0:
            # 正の値: ゼロから右に伸びる
            bar_x = zero_x
            bar_color = COLOR_BAR_TARGET if is_target else COLOR_BAR_DEFAULT
        else:
            # 負の値: ゼロから左に伸びる
            bar_x = zero_x - bar_length
            bar_color = COLOR_BAR_TARGET_NEG if is_target else COLOR_BAR_NEGATIVE

        if bar_length > Emu(100):  # 極小の場合は描画スキップ
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, bar_y, bar_length, bar_h,
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = bar_color
            bar.line.fill.background()
            bar.shadow.inherit = False
            bar.text_frame.text = ""

        # 値ラベル（バーの右端の右側）
        val_text = format_value(val, unit="", decimals=decimals, show_sign=show_sign)

        # 値ラベルの配置:
        # 正の値: バーの右端の右側
        # 負の値: バーの左端の左側
        if val >= 0:
            vlabel_x = bar_x + bar_length + Inches(0.04)
            vlabel_w = VALUE_COL_W
            vlabel_align = PP_ALIGN.LEFT
        else:
            vlabel_x = bar_x - VALUE_COL_W - Inches(0.04)
            vlabel_w = VALUE_COL_W
            vlabel_align = PP_ALIGN.RIGHT

        # はみ出しチェック
        if vlabel_x + vlabel_w > bar_area_right:
            vlabel_x = bar_area_right - vlabel_w
        if vlabel_x < bar_area_left + LABEL_COL_W:
            vlabel_x = bar_area_left + LABEL_COL_W

        vlabel_tb = slide.shapes.add_textbox(
            vlabel_x, row_y, vlabel_w, row_h,
        )
        vtf = vlabel_tb.text_frame
        vtf.word_wrap = False
        vtf.margin_left = 0; vtf.margin_right = 0
        vtf.margin_top = 0; vtf.margin_bottom = 0
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = vtf.paragraphs[0]
        vp.alignment = vlabel_align
        vrun = vp.add_run()
        vrun.text = val_text
        vrun.font.size = FONT_SIZE_VALUE
        vrun.font.bold = True if is_target else False
        vrun.font.name = FONT_NAME_JP
        vrun.font.color.rgb = COLOR_BAR_TARGET if is_target else COLOR_TEXT

    print(f"    ✓ {metric_name}: {len(rows)}社")


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)  # passive: accepted but ignored until brand migration
    args = ap.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(args.template)
    slide = prs.slides[0]

    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), data.get("main_message", ""))
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), data.get("chart_title", "財務ベンチマーク"))
    print(f"  ✓ Main Message & Chart Title set")

    companies = data.get("companies", [])
    metrics = data.get("metrics", [])
    target_company = data.get("target_company")

    if not companies or not metrics:
        print("  ✗ ERROR: 'companies' and 'metrics' are required", file=sys.stderr)
        sys.exit(1)

    # 指標数に応じてグリッド配置
    n_metrics = len(metrics)
    if n_metrics > N_ROWS * N_COLS:
        print(f"  ⚠ WARNING: {n_metrics} metrics > grid capacity ({N_ROWS * N_COLS}). Only first {N_ROWS * N_COLS} will be shown.", file=sys.stderr)
        metrics = metrics[: N_ROWS * N_COLS]
        n_metrics = len(metrics)

    print(f"\n  各指標のチャート生成:")
    for i, metric in enumerate(metrics):
        row = i // N_COLS
        col = i % N_COLS
        cell_x = GRID_LEFT + (CELL_W + CELL_GAP_X) * col
        cell_y = GRID_TOP + (CELL_H + CELL_GAP_Y) * row

        draw_bar_chart_cell(
            slide, metric, companies, target_company,
            cell_x, cell_y, CELL_W, CELL_H,
        )

    # 出典
    source = data.get("source", "")
    if source:
        add_text_box(
            slide, source,
            SOURCE_X, SOURCE_Y, SOURCE_W, Inches(0.25),
            FONT_SIZE_SOURCE, bold=False, color=COLOR_SOURCE,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        print(f"  ✓ Source: {source[:40]}...")

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ Saved: {args.output}")


if __name__ == "__main__":
    main()
