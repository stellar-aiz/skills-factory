"""
fill_business_model.py — 事業モデルデータをHTMLで描画→スクリーンショット→PPTXに挿入するスクリプト

テンプレート構造（BusinessModel.pptx 確認済み）:
  - Title 1            (PLACEHOLDER): Main Message
  - Text Placeholder 2 (PLACEHOLDER): Chart Title
  - Rectangle 4        (AUTO_SHAPE):  事業モデル図エリア（ここに画像を挿入）
  - TextBox 9          (TEXT_BOX):    意味合い 4段落（見出し + 3項目）

使い方:
  python fill_business_model.py \
    --data /home/claude/business_model_data.json \
    --template <SKILL_DIR>/assets/business-model-template.pptx \
    --output /mnt/user-data/outputs/BusinessModel_output.pptx
"""

import argparse

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
import asyncio
import json
import os
import sys
import tempfile
from html import escape

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング（確認済み）──────────────────────────────
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"
SHAPE_DIAGRAM_AREA = "Rectangle 4"
SHAPE_IMPLICATIONS = "TextBox 9"
# ────────────────────────────────────────────────────────────

# ── HTMLスクリーンショットの設定 ────────────────────────────
VIEWPORT_WIDTH  = 1766   # 8.83 inches * 200 DPI
VIEWPORT_HEIGHT = 1044   # 5.22 inches * 200 DPI
DEVICE_SCALE    = 2      # 高解像度（Retina相当）
# ────────────────────────────────────────────────────────────


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_placeholder_text(shape, text):
    """PlaceholderのTextFrameにテキストをセット（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def fill_implications(slide, items):
    """TextBox 9 の Para[1]〜Para[3] に意味合いを書き込む。
    ラベル部分を太字、詳細部分を通常テキストで2つのrunに分ける。
    """
    shape = find_shape(slide, SHAPE_IMPLICATIONS)
    if shape is None:
        return
    tf = shape.text_frame
    for i, item in enumerate(items[:3]):
        label  = item.get("label", "").strip()
        detail = item.get("detail", "").strip()
        para_idx = i + 1  # Para[0] は「意味合い」見出し
        if para_idx >= len(tf.paragraphs):
            continue
        para = tf.paragraphs[para_idx]

        # 既存のテンプレートrunからフォント情報を取得
        template_rPr = None
        if para.runs:
            template_rPr = para.runs[0]._r.find(qn("a:rPr"))

        # 既存のrunをすべてクリア
        for r in list(para._p.findall(qn("a:r"))):
            para._p.remove(r)

        def _make_run(text_str, bold=False):
            """runを作成し、テンプレートのフォント情報を継承する"""
            r_elem = etree.SubElement(para._p, qn("a:r"))
            if template_rPr is not None:
                # テンプレートの rPr をコピーして太字属性を上書き
                import copy
                new_rPr = copy.deepcopy(template_rPr)
                if bold:
                    new_rPr.set("b", "1")
                else:
                    new_rPr.attrib.pop("b", None)
                r_elem.insert(0, new_rPr)
            else:
                attrib = {"lang": "ja-JP"}
                if bold:
                    attrib["b"] = "1"
                etree.SubElement(r_elem, qn("a:rPr"), attrib=attrib)
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text_str

        if label:
            _make_run(f"{label}：", bold=True)
            _make_run(detail, bold=False)
            display = f"{label}：{detail}"
        else:
            _make_run(detail, bold=False)
            display = detail

        print(f"  [Implication {i+1}] {display[:55]}{'...' if len(display) > 55 else ''}")


def _esc(text):
    """HTMLエスケープのショートカット"""
    return escape(str(text))


def generate_html(data):
    """事業モデルデータからHTML模式図を生成する。

    単一事業モデル（companyのみ）と多事業モデル（businesses配列あり）の両方に対応。
    多事業モデルでは、中央に会社名ヘッダー＋複数の事業ラインを縦に並べ、
    各サプライヤー・顧客が特定の事業ラインと接続される。
    """
    company = data.get("company", {})
    suppliers = data.get("suppliers", [])
    customers = data.get("customers", [])
    businesses = data.get("businesses", [])

    is_multi = len(businesses) > 0

    if is_multi:
        return _generate_html_multi(data, company, suppliers, customers, businesses)
    else:
        return _generate_html_single(data, company, suppliers, customers)


def _generate_html_single(data, company, suppliers, customers):
    """従来の単一事業モデル（自社=1ボックス）のHTML生成"""
    n_suppliers = len(suppliers)
    n_customers = len(customers)
    max_rows = max(n_suppliers, n_customers, 1)

    size_table = {
        1: {"vw": 900, "vh": 530,
            "entity_name": 26, "entity_desc": 20, "company_name": 32, "company_desc": 21,
            "arrow_label": 20, "arrow_head": 9, "box_pad": "14px 16px", "company_pad": "22px 20px",
            "row_gap": 8, "col_gap": 6},
        2: {"vw": 1120, "vh": 660,
            "entity_name": 24, "entity_desc": 19, "company_name": 30, "company_desc": 20,
            "arrow_label": 19, "arrow_head": 9, "box_pad": "12px 14px", "company_pad": "20px 18px",
            "row_gap": 8, "col_gap": 6},
        3: {"vw": 1100, "vh": 650,
            "entity_name": 22, "entity_desc": 18, "company_name": 28, "company_desc": 19,
            "arrow_label": 18, "arrow_head": 8, "box_pad": "10px 12px", "company_pad": "16px 14px",
            "row_gap": 6, "col_gap": 4},
        4: {"vw": 1200, "vh": 710,
            "entity_name": 20, "entity_desc": 17, "company_name": 26, "company_desc": 18,
            "arrow_label": 17, "arrow_head": 8, "box_pad": "8px 10px", "company_pad": "12px 10px",
            "row_gap": 5, "col_gap": 4},
    }
    sz = size_table.get(max_rows, size_table[4])
    vw, vh = sz["vw"], sz["vh"]

    css = _css_common(sz, max_rows)
    css += _css_company_single(sz, max_rows)

    # セル構築
    cells = ""
    supplier_offset = (max_rows - n_suppliers) // 2
    customer_offset = (max_rows - n_customers) // 2

    for row in range(max_rows):
        s_idx = row - supplier_offset
        c_idx = row - customer_offset

        if 0 <= s_idx < n_suppliers:
            s = suppliers[s_idx]
            cells += _html_entity_box(1, row+1, s)
            cells += _html_arrow_cell(2, row+1, s.get("flow_to_company",""), s.get("flow_from_company",""))
        else:
            cells += f'<div class="empty-cell" style="grid-column:1;grid-row:{row+1};"></div>'
            cells += f'<div class="empty-cell" style="grid-column:2;grid-row:{row+1};"></div>'

        if 0 <= c_idx < n_customers:
            c = customers[c_idx]
            cells += _html_arrow_cell(4, row+1, c.get("flow_from_company",""), c.get("flow_to_company",""))
            cells += _html_entity_box(5, row+1, c)
        else:
            cells += f'<div class="empty-cell" style="grid-column:4;grid-row:{row+1};"></div>'
            cells += f'<div class="empty-cell" style="grid-column:5;grid-row:{row+1};"></div>'

    html = _html_wrap(css, f"""
    <div class="company-cell">
        <div class="company-box">
            <div class="company-name">{_esc(company.get('name','自社'))}</div>
            <div class="company-desc">{_esc(company.get('description',''))}</div>
        </div>
    </div>
    {cells}
    """)

    return html, vw, vh


def _generate_html_multi(data, company, suppliers, customers, businesses):
    """多事業モデル（中央に複数の事業ライン）のHTML生成"""
    n_biz = len(businesses)

    size_table = {
        1: {"vw": 900, "vh": 530,
            "entity_name": 26, "entity_desc": 20, "biz_name": 26, "biz_desc": 19,
            "company_name": 30, "arrow_label": 20, "arrow_head": 9,
            "box_pad": "14px 16px", "biz_pad": "10px 14px",
            "row_gap": 8, "col_gap": 6},
        2: {"vw": 1120, "vh": 660,
            "entity_name": 24, "entity_desc": 19, "biz_name": 24, "biz_desc": 18,
            "company_name": 28, "arrow_label": 19, "arrow_head": 9,
            "box_pad": "12px 14px", "biz_pad": "8px 12px",
            "row_gap": 8, "col_gap": 6},
        3: {"vw": 1100, "vh": 650,
            "entity_name": 22, "entity_desc": 18, "biz_name": 22, "biz_desc": 17,
            "company_name": 26, "arrow_label": 18, "arrow_head": 8,
            "box_pad": "10px 12px", "biz_pad": "6px 10px",
            "row_gap": 6, "col_gap": 4},
        4: {"vw": 1200, "vh": 710,
            "entity_name": 20, "entity_desc": 17, "biz_name": 20, "biz_desc": 16,
            "company_name": 24, "arrow_label": 17, "arrow_head": 8,
            "box_pad": "8px 10px", "biz_pad": "5px 8px",
            "row_gap": 5, "col_gap": 4},
    }
    sz = size_table.get(n_biz, size_table[4])
    vw, vh = sz["vw"], sz["vh"]

    # ── サプライヤー/顧客がまたがる行の計算 ──
    def _calc_span(entity_list, key="to_business"):
        """各エンティティのconnections から grid-row の start/end を計算"""
        spans = []
        for ent in entity_list:
            conns = ent.get("connections", [])
            if conns:
                rows = [c.get(key, 0) for c in conns]
                spans.append((min(rows), max(rows)))
            else:
                spans.append((0, 0))
        return spans

    supplier_spans = _calc_span(suppliers, "to_business")
    customer_spans = _calc_span(customers, "from_business")

    # ── CSS ──
    css = f"""
    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
    body {{
        font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', sans-serif;
        background: white;
        width: {vw}px; height: {vh}px;
        display: flex; align-items: stretch; justify-content: stretch;
        padding: 6px 8px; overflow: hidden;
    }}
    .grid {{
        display: grid;
        grid-template-columns: 17fr 16fr 34fr 16fr 17fr;
        grid-template-rows: auto repeat({n_biz}, 1fr);
        align-items: stretch; justify-items: stretch;
        width: 100%; height: 100%;
        gap: {sz['row_gap']}px {sz['col_gap']}px;
    }}

    /* Row 1: ヘッダー行（会社名） */
    .header-label {{
        font-weight: bold;
        font-size: {sz['company_name']}px;
        color: #1A365D;
        text-align: center;
        display: flex; align-items: center; justify-content: center;
        padding: 4px 0;
    }}
    .header-company {{
        font-weight: bold;
        font-size: {sz['company_name']}px;
        color: #1A365D;
        text-align: center;
        display: flex; align-items: center; justify-content: center;
        padding: 4px 0;
    }}

    /* 中央: 事業ラインボックス（自社カラー — ダークネイビー） */
    .biz-box {{
        border: 2px solid #1A365D;
        border-radius: 5px;
        padding: {sz['biz_pad']};
        background: #1A365D;
        display: flex; flex-direction: column;
        justify-content: center; align-items: flex-start;
        overflow: hidden;
    }}
    .biz-name {{
        font-weight: bold;
        font-size: {sz['biz_name']}px;
        color: #FFFFFF;
        line-height: 1.3;
    }}
    .biz-desc {{
        font-size: {sz['biz_desc']}px;
        color: #BEE3F8;
        line-height: 1.3;
    }}
    .biz-metrics {{
        font-size: {sz['biz_desc']}px;
        color: #90CDF4;
        line-height: 1.3;
        margin-top: 2px;
    }}

    /* サプライヤー・顧客ボックス */
    .entity-box {{
        border: 2px solid #4A5568;
        border-radius: 6px;
        padding: {sz['box_pad']};
        text-align: center; width: 100%; height: 100%;
        background: #EBF4FF;
        display: flex; flex-direction: column;
        justify-content: center; align-items: center;
    }}
    .entity-name {{
        font-weight: bold;
        font-size: {sz['entity_name']}px;
        color: #1A365D;
        margin-bottom: 3px; line-height: 1.3;
        word-break: keep-all; overflow-wrap: break-word;
    }}
    .entity-desc {{
        font-size: {sz['entity_desc']}px;
        color: #4A5568; line-height: 1.3;
    }}

    /* 矢印セル */
    .arrow-cell {{
        display: flex; flex-direction: column;
        align-items: stretch; justify-content: center;
        width: 100%; height: 100%;
        gap: 2px; padding: 1px 0;
    }}
    .arrow-row {{
        display: flex; align-items: center; gap: 0;
        height: {sz['arrow_label'] + 8}px;
    }}
    .arrow-label {{
        font-size: {sz['arrow_label']}px;
        line-height: 1.15;
        white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
        flex: 0 1 auto; text-align: center; min-width: 0;
    }}
    .arrow-label-provide {{ color: #C53030; }}
    .arrow-label-payment {{ color: #2B6CB0; }}
    .arrow-line {{
        flex: 1 0 20px; height: 2px; position: relative;
    }}
    .arrow-line-right {{ background: #C53030; }}
    .arrow-line-right::after {{
        content: ''; position: absolute; right: 0; top: 50%;
        transform: translateX(100%) translateY(-50%);
        width: 0; height: 0;
        border-top: {sz['arrow_head']}px solid transparent;
        border-bottom: {sz['arrow_head']}px solid transparent;
        border-left: {sz['arrow_head']+2}px solid #C53030;
    }}
    .arrow-line-left {{ background: #2B6CB0; }}
    .arrow-line-left::before {{
        content: ''; position: absolute; left: 0; top: 50%;
        transform: translateX(-100%) translateY(-50%);
        width: 0; height: 0;
        border-top: {sz['arrow_head']}px solid transparent;
        border-bottom: {sz['arrow_head']}px solid transparent;
        border-right: {sz['arrow_head']+2}px solid #2B6CB0;
    }}
    .empty-cell {{ }}
    """

    # ── セル構築 ──
    cells = ""

    # Row 1 (grid-row:1) = ヘッダー行
    cells += f'<div class="header-label" style="grid-column:1;grid-row:1;">取引先</div>'
    cells += f'<div style="grid-column:2;grid-row:1;"></div>'
    cells += f'<div class="header-company" style="grid-column:3;grid-row:1;">{_esc(company.get("name","自社"))}</div>'
    cells += f'<div style="grid-column:4;grid-row:1;"></div>'
    cells += f'<div class="header-label" style="grid-column:5;grid-row:1;">顧客</div>'

    # ── 中央列 (col 3): 事業ラインボックス ──
    for bi, biz in enumerate(businesses):
        gr = bi + 2  # grid-row (1はヘッダー)
        # 売上・利益率のメトリクス行を構築
        metrics_parts = []
        revenue = biz.get("revenue", "")
        margin = biz.get("margin", "")
        if revenue:
            metrics_parts.append(f"売上{_esc(revenue)}")
        if margin:
            metrics_parts.append(f"利益率{_esc(margin)}")
        metrics_html = ""
        if metrics_parts:
            metrics_html = f'<div class="biz-metrics">（{"／".join(metrics_parts)}）</div>'

        cells += f"""<div class="biz-box" style="grid-column:3;grid-row:{gr};">
            <div class="biz-name">{_esc(biz['name'])}</div>
            <div class="biz-desc">{_esc(biz.get('description',''))}</div>
            {metrics_html}
        </div>"""

    # ── サプライヤー列 (col 1) + 矢印 (col 2) ──
    # 各サプライヤーのボックスをconnection範囲の行にまたがって配置
    supplier_rows_used = set()
    for si, sup in enumerate(suppliers):
        r_start, r_end = supplier_spans[si]
        gr_start = r_start + 2  # grid-row (1-indexed, +1 for header)
        gr_end = r_end + 3      # grid-row end (exclusive)
        cells += f"""<div class="entity-box" style="grid-column:1;grid-row:{gr_start}/{gr_end};">
            <div class="entity-name">{_esc(sup['name'])}</div>
            <div class="entity-desc">{_esc(sup.get('description',''))}</div>
        </div>"""
        for row_i in range(r_start, r_end+1):
            supplier_rows_used.add(row_i)

        # 矢印: 各connectionのto_business行に配置
        for conn in sup.get("connections", []):
            biz_idx = conn.get("to_business", 0)
            gr = biz_idx + 2
            label_to = conn.get("label_to", "")
            label_from = conn.get("label_from", "")
            arrow_rows = ""
            if label_to:
                arrow_rows += f'<div class="arrow-row"><div class="arrow-label arrow-label-provide">{_esc(label_to)}</div><div class="arrow-line arrow-line-right"></div></div>'
            if label_from:
                arrow_rows += f'<div class="arrow-row"><div class="arrow-line arrow-line-left"></div><div class="arrow-label arrow-label-payment">{_esc(label_from)}</div></div>'
            cells += f'<div class="arrow-cell" style="grid-column:2;grid-row:{gr};">{arrow_rows}</div>'

    # 矢印がない行は空セル
    for bi in range(n_biz):
        has_arrow = any(
            any(c.get("to_business") == bi for c in sup.get("connections",[]))
            for sup in suppliers
        )
        if not has_arrow:
            cells += f'<div class="empty-cell" style="grid-column:2;grid-row:{bi+2};"></div>'

    # サプライヤーボックスがない行は空セル
    for bi in range(n_biz):
        if bi not in supplier_rows_used:
            # Check if this row is already covered by a spanning supplier
            covered = False
            for si, (r_start, r_end) in enumerate(supplier_spans):
                if r_start <= bi <= r_end:
                    covered = True
                    break
            if not covered:
                cells += f'<div class="empty-cell" style="grid-column:1;grid-row:{bi+2};"></div>'

    # ── 顧客列 (col 5) + 矢印 (col 4) ──
    customer_rows_used = set()
    for ci, cust in enumerate(customers):
        r_start, r_end = customer_spans[ci]
        gr_start = r_start + 2
        gr_end = r_end + 3
        cells += f"""<div class="entity-box" style="grid-column:5;grid-row:{gr_start}/{gr_end};">
            <div class="entity-name">{_esc(cust['name'])}</div>
            <div class="entity-desc">{_esc(cust.get('description',''))}</div>
        </div>"""
        for row_i in range(r_start, r_end+1):
            customer_rows_used.add(row_i)

        for conn in cust.get("connections", []):
            biz_idx = conn.get("from_business", 0)
            gr = biz_idx + 2
            label_to = conn.get("label_to", "")
            label_from = conn.get("label_from", "")
            arrow_rows = ""
            if label_to:
                arrow_rows += f'<div class="arrow-row"><div class="arrow-label arrow-label-provide">{_esc(label_to)}</div><div class="arrow-line arrow-line-right"></div></div>'
            if label_from:
                arrow_rows += f'<div class="arrow-row"><div class="arrow-line arrow-line-left"></div><div class="arrow-label arrow-label-payment">{_esc(label_from)}</div></div>'
            cells += f'<div class="arrow-cell" style="grid-column:4;grid-row:{gr};">{arrow_rows}</div>'

    for bi in range(n_biz):
        has_arrow = any(
            any(c.get("from_business") == bi for c in cust.get("connections",[]))
            for cust in customers
        )
        if not has_arrow:
            cells += f'<div class="empty-cell" style="grid-column:4;grid-row:{bi+2};"></div>'

    for bi in range(n_biz):
        if bi not in customer_rows_used:
            covered = False
            for ci, (r_start, r_end) in enumerate(customer_spans):
                if r_start <= bi <= r_end:
                    covered = True
                    break
            if not covered:
                cells += f'<div class="empty-cell" style="grid-column:5;grid-row:{bi+2};"></div>'

    html = _html_wrap_raw(css, cells, vw, vh)
    return html, vw, vh


# ── ヘルパー関数 ──

def _css_common(sz, max_rows):
    """単一事業モデル用の共通CSS"""
    return f"""
    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
    body {{
        font-family: 'Noto Sans CJK JP', 'Noto Sans JP', 'Meiryo UI', sans-serif;
        background: white;
        width: {sz['vw']}px; height: {sz['vh']}px;
        display: flex; align-items: stretch; justify-content: stretch;
        padding: 6px 8px; overflow: hidden;
    }}
    .grid {{
        display: grid;
        grid-template-columns: 19fr 18fr 26fr 18fr 19fr;
        grid-template-rows: repeat({max_rows}, 1fr);
        align-items: stretch; justify-items: stretch;
        width: 100%; height: 100%;
        gap: {sz['row_gap']}px {sz['col_gap']}px;
    }}
    .entity-box {{
        border: 2px solid #4A5568; border-radius: 6px;
        padding: {sz['box_pad']}; text-align: center;
        width: 100%; height: 100%; background: #EBF4FF;
        display: flex; flex-direction: column;
        justify-content: center; align-items: center;
    }}
    .entity-name {{
        font-weight: bold; font-size: {sz['entity_name']}px;
        color: #1A365D; margin-bottom: 4px; line-height: 1.3;
        word-break: keep-all; overflow-wrap: break-word;
    }}
    .entity-desc {{
        font-size: {sz['entity_desc']}px; color: #4A5568; line-height: 1.3;
    }}
    .arrow-cell {{
        display: flex; flex-direction: column;
        align-items: stretch; justify-content: center;
        width: 100%; height: 100%; gap: 3px; padding: 2px 0;
    }}
    .arrow-row {{
        display: flex; align-items: center; gap: 0;
        height: {sz['arrow_label']+10}px;
    }}
    .arrow-label {{
        font-size: {sz['arrow_label']}px; line-height: 1.15;
        white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
        flex: 0 1 auto; text-align: center; min-width: 0;
    }}
    .arrow-label-provide {{ color: #C53030; }}
    .arrow-label-payment {{ color: #2B6CB0; }}
    .arrow-line {{
        flex: 1 0 30px; height: 2px; position: relative;
    }}
    .arrow-line-right {{ background: #C53030; }}
    .arrow-line-right::after {{
        content:''; position:absolute; right:0; top:50%;
        transform: translateX(100%) translateY(-50%);
        width:0; height:0;
        border-top:{sz['arrow_head']}px solid transparent;
        border-bottom:{sz['arrow_head']}px solid transparent;
        border-left:{sz['arrow_head']+2}px solid #C53030;
    }}
    .arrow-line-left {{ background: #2B6CB0; }}
    .arrow-line-left::before {{
        content:''; position:absolute; left:0; top:50%;
        transform: translateX(-100%) translateY(-50%);
        width:0; height:0;
        border-top:{sz['arrow_head']}px solid transparent;
        border-bottom:{sz['arrow_head']}px solid transparent;
        border-right:{sz['arrow_head']+2}px solid #2B6CB0;
    }}
    .empty-cell {{ }}
    """

def _css_company_single(sz, max_rows):
    """単一事業モデル用の自社ボックスCSS"""
    return f"""
    .company-cell {{
        grid-column: 3; grid-row: 1 / {max_rows+1};
        display: flex; align-items: stretch; justify-content: stretch;
    }}
    .company-box {{
        border: 3px solid #1A365D; border-radius: 8px;
        padding: {sz['company_pad']}; text-align: center;
        width: 100%; height: 100%; background: #1A365D; color: white;
        display: flex; flex-direction: column;
        justify-content: center; align-items: center;
        box-shadow: 0 3px 10px rgba(0,0,0,0.18);
    }}
    .company-name {{
        font-weight: bold; font-size: {sz['company_name']}px;
        margin-bottom: 8px; line-height: 1.3;
    }}
    .company-desc {{
        font-size: {sz['company_desc']}px; color: #BEE3F8; line-height: 1.4;
    }}
    """

def _html_entity_box(col, row, entity):
    return f"""<div class="entity-box" style="grid-column:{col};grid-row:{row};">
        <div class="entity-name">{_esc(entity['name'])}</div>
        <div class="entity-desc">{_esc(entity.get('description',''))}</div>
    </div>"""

def _html_arrow_cell(col, row, label_provide, label_payment):
    rows = ""
    if label_provide:
        rows += f'<div class="arrow-row"><div class="arrow-label arrow-label-provide">{_esc(label_provide)}</div><div class="arrow-line arrow-line-right"></div></div>'
    if label_payment:
        rows += f'<div class="arrow-row"><div class="arrow-line arrow-line-left"></div><div class="arrow-label arrow-label-payment">{_esc(label_payment)}</div></div>'
    return f'<div class="arrow-cell" style="grid-column:{col};grid-row:{row};">{rows}</div>'

def _html_wrap(css, inner_html):
    return f"""<!DOCTYPE html>
<html lang="ja">
<head><meta charset="utf-8"><style>{css}</style></head>
<body>
<div class="grid">
{inner_html}
</div>
</body>
</html>"""

def _html_wrap_raw(css, cells_html, vw, vh):
    return f"""<!DOCTYPE html>
<html lang="ja">
<head><meta charset="utf-8"><style>{css}</style></head>
<body>
<div class="grid">
{cells_html}
</div>
</body>
</html>"""


async def take_screenshot(html_content, output_path, vp_width=None, vp_height=None):
    """PlaywrightでHTMLをスクリーンショットとして保存する"""
    from playwright.async_api import async_playwright

    vp_width = vp_width or VIEWPORT_WIDTH
    vp_height = vp_height or VIEWPORT_HEIGHT

    # HTMLファイルを一時保存
    html_path = os.path.join(tempfile.gettempdir(), "business_model_diagram.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={"width": vp_width, "height": vp_height},
            device_scale_factor=DEVICE_SCALE,
        )
        await page.goto(f"file://{html_path}")
        # 少し待ってレンダリングを安定させる
        await page.wait_for_timeout(500)
        await page.screenshot(path=output_path, full_page=False)
        await browser.close()

    print(f"  📸 Screenshot saved: {output_path} ({os.path.getsize(output_path)} bytes)")
    # 一時ファイル削除
    os.unlink(html_path)


def main():
    parser = argparse.ArgumentParser(description="事業モデルデータをPPTXに流し込む")
    parser.add_argument("--data",     required=True, help="business_model_data.json のパス")
    parser.add_argument("--template", required=True, help="business-model-template.pptx のパス")
    parser.add_argument("--output",   required=True, help="出力PPTXのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    # ── HTML生成 → スクリーンショット ──────────────────────
    print("📐 Generating business model diagram HTML...")
    html_content, vp_w, vp_h = generate_html(data)

    screenshot_path = os.path.join(tempfile.gettempdir(), "bm_diagram.png")
    print(f"📸 Taking screenshot with Playwright (viewport: {vp_w}x{vp_h})...")
    asyncio.run(take_screenshot(html_content, screenshot_path, vp_w, vp_h))

    # ── PPTX操作 ──────────────────────────────────────────
    print("📝 Filling PPTX template...")
    prs = Presentation(args.template)
    slide = prs.slides[0]

    # Main Message
    main_msg = data.get("main_message", "").strip()
    if main_msg:
        shape = find_shape(slide, SHAPE_MAIN_MESSAGE)
        set_placeholder_text(shape, main_msg)
        print(f"  [Main Message] {main_msg[:60]}{'...' if len(main_msg) > 60 else ''}")
    else:
        print("  ⚠ main_message が未設定です", file=sys.stderr)

    # Chart Title
    chart_title = data.get("chart_title", "").strip()
    if chart_title:
        shape = find_shape(slide, SHAPE_CHART_TITLE)
        set_placeholder_text(shape, chart_title)
        print(f"  [Chart Title]  {chart_title}")
    else:
        print("  ⚠ chart_title が未設定です", file=sys.stderr)

    # Implications (意味合い)
    implications = data.get("implications", [])
    if len(implications) != 3:
        print(f"  ⚠ implications は {len(implications)} 件です（3件推奨）", file=sys.stderr)
    fill_implications(slide, implications)

    # ── 事業モデル図の画像挿入 ─────────────────────────────
    rect_shape = find_shape(slide, SHAPE_DIAGRAM_AREA)
    if rect_shape:
        left   = rect_shape.left
        top    = rect_shape.top
        width  = rect_shape.width
        height = rect_shape.height

        # Rectangle 4 のテキストをクリアする
        if rect_shape.has_text_frame:
            for para in rect_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

        # 画像を挿入（少しパディングを入れる）
        padding = Inches(0.08)
        pic = slide.shapes.add_picture(
            screenshot_path,
            left + padding,
            top + padding,
            width - 2 * padding,
            height - 2 * padding,
        )
        print(f"  [Diagram] Image inserted into Rectangle 4 area")
    else:
        print("  ⚠ Rectangle 4 が見つかりません", file=sys.stderr)

    # ── 保存 ──────────────────────────────────────────────
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n✅ 保存しました: {args.output}")

    # スクリーンショット一時ファイル削除
    if os.path.exists(screenshot_path):
        os.unlink(screenshot_path)


if __name__ == "__main__":
    main()
