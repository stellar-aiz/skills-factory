"""
fill_business_overview.py — 事業セグメント概要スライドをPPTXネイティブオブジェクトで生成

ISSUE-004 (v0.3) Phase 2 — `business-deepdive-agent` の「事業の概要は？」論点 1 枚スライド。
customer-profile-pptx の構造を継承（左カラム: key-value テーブル / 右カラム: 業績 or KPI）。

テンプレート: business-overview-template.pptx（customer-profile-template ベース）
  - 既存テーブルを削除
  - 左側: 事業の概要テーブル（2列: ラベル | 値、ブレットポイント形式）
  - 右側: 業績チャート (mode=revenue_chart) または KPI カード (mode=kpi_cards)
  - revenue_chart 時は CAGR 矢印＋テキスト注釈

Usage:
  python fill_business_overview.py \
    --data /home/claude/business_overview_data.json \
    --template <path>/business-overview-template.pptx \
    --output /mnt/user-data/outputs/BusinessOverview_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Layout Constants ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"

PANEL_Y = Inches(1.50)

LEFT_X = Inches(0.41)
LEFT_W = Inches(5.80)

RIGHT_X = Inches(6.50)
RIGHT_W = Inches(6.40)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_BAR = RGBColor(0x4E, 0x79, 0xA7)
COLOR_LINE = RGBColor(0x33, 0x33, 0x33)
COLOR_CAGR_ARROW = RGBColor(0x33, 0x33, 0x33)
COLOR_KPI_CARD_BG = RGBColor(0xF7, 0xF7, 0xF7)
COLOR_KPI_CARD_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
COLOR_KPI_VALUE = RGBColor(0x4E, 0x79, 0xA7)

FONT_NAME_JP = "Meiryo UI"
FONT_SIZE_LABEL = Pt(14)
FONT_SIZE_VALUE = Pt(14)
FONT_SIZE_SECTION = Pt(14)

MAIN_MESSAGE_LIMIT = 65


def _validate(data):
    main_message = data.get("main_message", "")
    if not main_message:
        raise SystemExit("ERROR: main_message is required")
    if len(main_message) > MAIN_MESSAGE_LIMIT:
        raise SystemExit(
            f"ERROR: main_message exceeds {MAIN_MESSAGE_LIMIT} chars "
            f"(actual: {len(main_message)}). Shorten it.\n"
            f"  text: {main_message!r}"
        )
    if not data.get("parent_company"):
        raise SystemExit("ERROR: parent_company is required")
    if not data.get("segment_name"):
        raise SystemExit("ERROR: segment_name is required")
    overview = data.get("overview") or {}
    if not overview.get("items"):
        raise SystemExit("ERROR: overview.items is required")
    perf = data.get("performance") or {}
    mode = perf.get("mode")
    if mode not in ("revenue_chart", "kpi_cards"):
        raise SystemExit(
            f"ERROR: performance.mode must be 'revenue_chart' or 'kpi_cards' (got {mode!r})"
        )
    if mode == "revenue_chart" and not perf.get("data"):
        raise SystemExit("ERROR: performance.data is required for revenue_chart mode")
    if mode == "kpi_cards" and not perf.get("cards"):
        raise SystemExit("ERROR: performance.cards is required for kpi_cards mode")


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def remove_shape(slide, name):
    shape = find_shape(slide, name)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def add_section_title(slide, text, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.30))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = FONT_SIZE_SECTION
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.30), width, Inches(0.02)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = COLOR_TEXT
    underline.line.fill.background()
    return txBox


def build_overview_table(slide, items, left, top, width):
    """事業の概要テーブルを構築（ブレットポイント形式、枠線なし）"""
    n_rows = len(items)
    n_cols = 2
    col0_w = Inches(1.60)  # "• ラベル" column（事業向けはやや広め: 「主要製品/サービス」等）
    col1_w = width - col0_w

    min_row_h = Inches(0.35)
    table_h = min_row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = shape.table

    table.columns[0].width = col0_w
    table.columns[1].width = col1_w

    tbl_elem = shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '0', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    for tr in tbl_elem.findall(qn('a:tr')):
        tr.set('h', str(min_row_h))

    for r_idx, item in enumerate(items):
        label = item.get("label", "")
        value = item.get("value", "")
        bullet_label = f"•  {label}"
        _style_cell(table.cell(r_idx, 0), bullet_label, True, FONT_SIZE_LABEL)
        _style_cell(table.cell(r_idx, 1), value, False, FONT_SIZE_VALUE)

    print(f"  ✓ 事業概要テーブル: {n_rows}行")
    return shape


def _style_cell(cell, text, bold, font_size):
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        txBody = etree.SubElement(tc, qn("a:txBody"))

    old_bodyPr = txBody.find(qn("a:bodyPr"))
    if old_bodyPr is not None:
        txBody.remove(old_bodyPr)
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"), attrib={
        "wrap": "square",
        "lIns": "0", "rIns": "0",
        "tIns": "27432", "bIns": "27432",
        "anchor": "t",
    })
    txBody.insert(0, bodyPr)

    if txBody.find(qn("a:lstStyle")) is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.insert(1, lstStyle)

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p_elem = etree.SubElement(txBody, qn("a:p"))
    pPr = etree.SubElement(p_elem, qn("a:pPr"))
    pPr.set("algn", "l")

    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct"), attrib={"val": "100000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "0"})
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(spcAft, qn("a:spcPts"), attrib={"val": "0"})

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(int(font_size.pt * 100)),
        "b": "1" if bold else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": "333333"})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = str(text)

    old_tcPr = tc.find(qn("a:tcPr"))
    if old_tcPr is not None:
        tc.remove(old_tcPr)
    tcPr = etree.SubElement(tc, qn("a:tcPr"), attrib={
        "marL": "45720", "marR": "18288",
        "marT": "27432", "marB": "27432",
        "anchor": "t",
    })
    for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = etree.SubElement(tcPr, qn(border_name), attrib={"w": "0", "cmpd": "sng"})
        etree.SubElement(ln, qn("a:noFill"))


def build_combo_chart(slide, perf_data, left, top, width, height):
    """PowerPointネイティブ複合チャート（棒＋折れ線）"""
    from pptx.chart.data import CategoryChartData

    data = perf_data["data"]
    bar_label = perf_data.get("bar_label", "セグメント売上高")
    line_label = perf_data.get("line_label", "セグメント営業利益率")

    chart_data = CategoryChartData()
    chart_data.categories = [d["year"] for d in data]
    chart_data.add_series(bar_label, [d["revenue"] for d in data])
    chart_data.add_series(line_label, [d["op_margin"] for d in data])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = chart_frame.chart

    plotArea = chart._chartSpace.chart.plotArea
    barChart = plotArea.findall(qn('c:barChart'))[0]

    sers = barChart.findall(qn('c:ser'))
    line_ser_xml = copy.deepcopy(sers[1])
    barChart.remove(sers[1])

    lineChart = etree.SubElement(plotArea, qn('c:lineChart'))
    etree.SubElement(lineChart, qn('c:grouping'), attrib={'val': 'standard'})
    etree.SubElement(lineChart, qn('c:varyColors'), attrib={'val': '0'})
    lineChart.append(line_ser_xml)

    marker_xml = line_ser_xml.find(qn('c:marker'))
    if marker_xml is None:
        marker_xml = etree.SubElement(line_ser_xml, qn('c:marker'))
    symbol = marker_xml.find(qn('c:symbol'))
    if symbol is None:
        symbol = etree.SubElement(marker_xml, qn('c:symbol'))
    symbol.set('val', 'circle')
    sz = marker_xml.find(qn('c:size'))
    if sz is None:
        sz = etree.SubElement(marker_xml, qn('c:size'))
    sz.set('val', '9')

    catAx = plotArea.findall(qn('c:catAx'))
    valAx = plotArea.findall(qn('c:valAx'))

    sec_valAx_id = "2094734553"
    sec_catAx_id = "2094734554"

    etree.SubElement(lineChart, qn('c:axId'), attrib={'val': sec_catAx_id})
    etree.SubElement(lineChart, qn('c:axId'), attrib={'val': sec_valAx_id})

    sec_catAx = etree.SubElement(plotArea, qn('c:catAx'))
    etree.SubElement(sec_catAx, qn('c:axId'), attrib={'val': sec_catAx_id})
    scaling = etree.SubElement(sec_catAx, qn('c:scaling'))
    etree.SubElement(scaling, qn('c:orientation'), attrib={'val': 'minMax'})
    etree.SubElement(sec_catAx, qn('c:delete'), attrib={'val': '1'})
    etree.SubElement(sec_catAx, qn('c:axPos'), attrib={'val': 'b'})
    etree.SubElement(sec_catAx, qn('c:crossAx'), attrib={'val': sec_valAx_id})

    sec_valAx_elem = etree.SubElement(plotArea, qn('c:valAx'))
    etree.SubElement(sec_valAx_elem, qn('c:axId'), attrib={'val': sec_valAx_id})
    scaling2 = etree.SubElement(sec_valAx_elem, qn('c:scaling'))
    etree.SubElement(scaling2, qn('c:orientation'), attrib={'val': 'minMax'})
    etree.SubElement(sec_valAx_elem, qn('c:delete'), attrib={'val': '1'})
    etree.SubElement(sec_valAx_elem, qn('c:axPos'), attrib={'val': 'r'})
    etree.SubElement(sec_valAx_elem, qn('c:numFmt'), attrib={
        'formatCode': '0.0"%"', 'sourceLinked': '0'
    })
    etree.SubElement(sec_valAx_elem, qn('c:crossAx'), attrib={'val': sec_catAx_id})
    etree.SubElement(sec_valAx_elem, qn('c:crosses'), attrib={'val': 'max'})

    bar_ser = barChart.findall(qn('c:ser'))[0]
    spPr = bar_ser.find(qn('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(bar_ser, qn('c:spPr'))
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        sf = etree.SubElement(spPr, qn('a:solidFill'))
    for child in list(sf):
        sf.remove(child)
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': '4E79A7'})

    line_spPr = line_ser_xml.find(qn('c:spPr'))
    if line_spPr is None:
        line_spPr = etree.SubElement(line_ser_xml, qn('c:spPr'))
    ln = line_spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(line_spPr, qn('a:ln'))
    ln.set('w', '19050')
    line_sf = ln.find(qn('a:solidFill'))
    if line_sf is None:
        line_sf = etree.SubElement(ln, qn('a:solidFill'))
    for child in list(line_sf):
        line_sf.remove(child)
    etree.SubElement(line_sf, qn('a:srgbClr'), attrib={'val': '003366'})

    marker_spPr = marker_xml.find(qn('c:spPr'))
    if marker_spPr is None:
        marker_spPr = etree.SubElement(marker_xml, qn('c:spPr'))
    m_sf = etree.SubElement(marker_spPr, qn('a:solidFill'))
    etree.SubElement(m_sf, qn('a:srgbClr'), attrib={'val': '003366'})

    add_data_labels_to_ser(bar_ser, position='outEnd', num_format='0.0')
    add_data_labels_to_ser(line_ser_xml, position='t', num_format='0.0', font_color='FFFFFF')

    chart.has_legend = False

    for vax in plotArea.findall(qn('c:valAx')):
        for mg in vax.findall(qn('c:majorGridlines')):
            vax.remove(mg)
        for mg in vax.findall(qn('c:minorGridlines')):
            vax.remove(mg)
    for cax in plotArea.findall(qn('c:catAx')):
        for mg in cax.findall(qn('c:majorGridlines')):
            cax.remove(mg)

    gapWidth = barChart.find(qn('c:gapWidth'))
    if gapWidth is None:
        gapWidth = etree.SubElement(barChart, qn('c:gapWidth'))
    gapWidth.set('val', '108')

    for ax in plotArea.findall(qn('c:catAx')) + plotArea.findall(qn('c:valAx')):
        delete_elem = ax.find(qn('c:delete'))
        if delete_elem is not None and delete_elem.get('val') == '1':
            continue
        txPr = ax.find(qn('c:txPr'))
        if txPr is None:
            txPr = etree.SubElement(ax, qn('c:txPr'))
        bodyPr = txPr.find(qn('a:bodyPr'))
        if bodyPr is None:
            bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
            txPr.insert(0, bodyPr)
        bodyPr.set('rot', '-5400000')
        bodyPr.set('vert', 'horz')

        if txPr.find(qn('a:lstStyle')) is None:
            etree.SubElement(txPr, qn('a:lstStyle'))
        p = txPr.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(txPr, qn('a:p'))
        pPr = p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p, qn('a:pPr'))
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': '1100'})
        else:
            defRPr.set('sz', '1100')
        if defRPr.find(qn('a:latin')) is None:
            etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
        if defRPr.find(qn('a:ea')) is None:
            etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})

    for vax in valAx:
        del_elem = vax.find(qn('c:delete'))
        if del_elem is None:
            del_elem = etree.SubElement(vax, qn('c:delete'))
        del_elem.set('val', '1')

    chart.has_title = False

    plotArea_spPr = plotArea.find(qn('c:spPr'))
    if plotArea_spPr is None:
        plotArea_spPr = etree.SubElement(plotArea, qn('c:spPr'))
    etree.SubElement(plotArea_spPr, qn('a:noFill'))

    print(f"  ✓ 複合チャート: {len(data)}年分")
    return chart_frame


def add_data_labels_to_ser(ser_xml, position='outEnd', num_format='0.0', font_color='333333'):
    dLbls = ser_xml.find(qn('c:dLbls'))
    if dLbls is None:
        dLbls = etree.SubElement(ser_xml, qn('c:dLbls'))
    for child in list(dLbls):
        dLbls.remove(child)

    etree.SubElement(dLbls, qn('c:numFmt'), attrib={
        'formatCode': num_format, 'sourceLinked': '0'
    })
    etree.SubElement(dLbls, qn('c:showLegendKey'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showVal'), attrib={'val': '1'})
    etree.SubElement(dLbls, qn('c:showCatName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showSerName'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showPercent'), attrib={'val': '0'})
    etree.SubElement(dLbls, qn('c:showBubbleSize'), attrib={'val': '0'})

    pos_map = {'outEnd': 'outEnd', 't': 't', 'ctr': 'ctr'}
    etree.SubElement(dLbls, qn('c:dLblPos'), attrib={
        'val': pos_map.get(position, 'outEnd')
    })

    txPr = etree.SubElement(dLbls, qn('c:txPr'))
    etree.SubElement(txPr, qn('a:bodyPr'))
    etree.SubElement(txPr, qn('a:lstStyle'))
    p = etree.SubElement(txPr, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    defRPr = etree.SubElement(pPr, qn('a:defRPr'), attrib={'sz': '1200'})
    etree.SubElement(defRPr, qn('a:latin'), attrib={'typeface': FONT_NAME_JP})
    etree.SubElement(defRPr, qn('a:ea'), attrib={'typeface': FONT_NAME_JP})
    sf = etree.SubElement(defRPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': font_color})


def add_cagr_annotation(slide, perf_data, chart_left, chart_top, chart_width, chart_height):
    data = perf_data["data"]
    if len(data) < 2:
        return

    first_rev = data[0]["revenue"]
    last_rev = data[-1]["revenue"]
    max_rev = max(d["revenue"] for d in data)
    n_years = len(data) - 1
    n_cats = len(data)

    if first_rev > 0 and last_rev > 0 and n_years > 0:
        cagr = (last_rev / first_rev) ** (1.0 / n_years) - 1
        cagr_text = f"+{cagr*100:.1f}%" if cagr >= 0 else f"{cagr*100:.1f}%"
    else:
        cagr_text = "N/A"

    plot_left_margin = chart_width * 0.06
    plot_right_margin = chart_width * 0.04
    plot_top_margin = chart_height * 0.18
    plot_bottom_margin = chart_height * 0.14

    plot_left = chart_left + plot_left_margin
    plot_right = chart_left + chart_width - plot_right_margin
    plot_top = chart_top + plot_top_margin
    plot_bottom = chart_top + chart_height - plot_bottom_margin
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    cat_width = plot_w / n_cats
    first_bar_cx = plot_left + 0.5 * cat_width
    last_bar_cx = plot_left + (n_cats - 0.5) * cat_width

    axis_max = max_rev * 1.20
    first_bar_top_y = plot_bottom - (first_rev / axis_max) * plot_h
    last_bar_top_y = plot_bottom - (last_rev / axis_max) * plot_h

    gap_above = Inches(1.20)
    arrow_start_x = int(first_bar_cx)
    arrow_start_y = int(first_bar_top_y - gap_above)
    arrow_end_x = int(last_bar_cx)
    arrow_end_y = int(last_bar_top_y - gap_above)

    connector = slide.shapes.add_connector(
        1,
        arrow_start_x, arrow_start_y,
        arrow_end_x, arrow_end_y
    )
    connector.line.color.rgb = COLOR_CAGR_ARROW
    connector.line.width = Pt(1.5)

    cxnSp = connector._element
    spPr = cxnSp.find(qn('p:spPr'))
    if spPr is None:
        spPr = cxnSp.find(qn('a:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:tailEnd'), attrib={
        'type': 'triangle', 'w': 'med', 'len': 'med'
    })

    text_w = Inches(1.30)
    text_h = Inches(0.40)
    mid_x = (arrow_start_x + arrow_end_x) / 2
    mid_y = (arrow_start_y + arrow_end_y) / 2
    text_x = int(mid_x - text_w / 2)
    text_y = int(mid_y - text_h / 2)

    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, text_x, text_y, text_w, text_h
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    oval.line.color.rgb = COLOR_CAGR_ARROW
    oval.line.width = Pt(1.0)

    tf = oval.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cagr_text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP

    print(f"  ✓ CAGR注釈: {cagr_text} ({data[0]['year']}→{data[-1]['year']})")


def add_unit_label(slide, text, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.22))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_TEXT
    run.font.name = FONT_NAME_JP


def add_custom_legend(slide, perf_data, left, top, width):
    bar_label = perf_data.get("bar_label", "セグメント売上高")
    line_label = perf_data.get("line_label", "セグメント営業利益率")

    legend_w = Inches(3.40)
    legend_h = Inches(0.22)
    legend_x = left + width - legend_w

    sq_size = Inches(0.14)
    sq_y = top + (legend_h - sq_size) / 2

    bar_marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        legend_x, int(sq_y), sq_size, sq_size
    )
    bar_marker.fill.solid()
    bar_marker.fill.fore_color.rgb = COLOR_BAR
    bar_marker.line.fill.background()

    bar_text_x = legend_x + sq_size + Inches(0.06)
    bar_text_w = Inches(1.50)
    txBox1 = slide.shapes.add_textbox(int(bar_text_x), top, bar_text_w, legend_h)
    tf1 = txBox1.text_frame
    tf1.word_wrap = False
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.add_run()
    run1.text = bar_label
    run1.font.size = Pt(12)
    run1.font.color.rgb = COLOR_TEXT
    run1.font.name = FONT_NAME_JP

    line_section_x = bar_text_x + bar_text_w + Inches(0.10)
    line_w = Inches(0.30)
    line_y = top + legend_h / 2
    connector = slide.shapes.add_connector(
        1, int(line_section_x), int(line_y),
        int(line_section_x + line_w), int(line_y)
    )
    connector.line.color.rgb = RGBColor(0x00, 0x33, 0x66)
    connector.line.width = Pt(1.5)

    circle_size = Inches(0.12)
    circle_x = line_section_x + (line_w - circle_size) / 2
    circle_y = top + (legend_h - circle_size) / 2
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(circle_x), int(circle_y), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
    circle.line.fill.background()

    line_text_x = line_section_x + line_w + Inches(0.06)
    line_text_w = Inches(1.50)
    txBox2 = slide.shapes.add_textbox(int(line_text_x), top, line_text_w, legend_h)
    tf2 = txBox2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = line_label
    run2.font.size = Pt(12)
    run2.font.color.rgb = COLOR_TEXT
    run2.font.name = FONT_NAME_JP

    print(f"  ✓ カスタム凡例: {bar_label} / {line_label}")


def build_kpi_cards(slide, perf_data, left, top, width, height):
    """KPI カードグリッド（2 列 × 行可変）。各カード: KPI 名 / 値 / 補足。"""
    cards = perf_data["cards"]
    n = len(cards)
    if n < 1 or n > 6:
        raise SystemExit(f"ERROR: kpi_cards.cards must have 1-6 entries (got {n})")

    n_cols = 2
    n_rows = (n + 1) // 2
    gap = Inches(0.20)
    card_w = (width - gap) / n_cols
    card_h = (height - gap * (n_rows - 1)) / n_rows if n_rows > 1 else height

    for idx, card in enumerate(cards):
        row = idx // n_cols
        col = idx % n_cols
        cx = int(left + col * (card_w + gap))
        cy = int(top + row * (card_h + gap))

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, int(card_w), int(card_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_KPI_CARD_BG
        bg.line.color.rgb = COLOR_KPI_CARD_BORDER
        bg.line.width = Pt(0.75)
        bg.text_frame.text = ""  # avoid placeholder

        name = card.get("name", "")
        value = card.get("value", "")
        unit = card.get("unit", "")
        sub = card.get("sub", "")

        # KPI 名
        name_box = slide.shapes.add_textbox(
            cx + Inches(0.15), cy + Inches(0.10),
            int(card_w - Inches(0.30)), Inches(0.30)
        )
        tf_n = name_box.text_frame
        tf_n.word_wrap = True
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.LEFT
        run_n = p_n.add_run()
        run_n.text = name
        run_n.font.size = Pt(12)
        run_n.font.bold = True
        run_n.font.color.rgb = COLOR_TEXT
        run_n.font.name = FONT_NAME_JP

        # 値（中央大きく）
        value_h = Inches(0.70)
        value_box = slide.shapes.add_textbox(
            cx + Inches(0.15), cy + (card_h - value_h) / 2,
            int(card_w - Inches(0.30)), int(value_h)
        )
        tf_v = value_box.text_frame
        tf_v.word_wrap = False
        p_v = tf_v.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        run_v = p_v.add_run()
        run_v.text = f"{value}{unit}" if unit else str(value)
        run_v.font.size = Pt(28)
        run_v.font.bold = True
        run_v.font.color.rgb = COLOR_KPI_VALUE
        run_v.font.name = FONT_NAME_JP

        # 補足（下部）
        if sub:
            sub_box = slide.shapes.add_textbox(
                cx + Inches(0.15), cy + card_h - Inches(0.40),
                int(card_w - Inches(0.30)), Inches(0.28)
            )
            tf_s = sub_box.text_frame
            tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]
            p_s.alignment = PP_ALIGN.CENTER
            run_s = p_s.add_run()
            run_s.text = sub
            run_s.font.size = Pt(10)
            run_s.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            run_s.font.name = FONT_NAME_JP

    print(f"  ✓ KPI カード: {n}枚 ({n_cols}列×{n_rows}行)")


def add_source_label(slide, text):
    src_left = Inches(0.41)
    src_top = Inches(7.05)
    src_w = Inches(8.00)
    src_h = Inches(0.30)

    txBox = slide.shapes.add_textbox(src_left, src_top, src_w, src_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = FONT_NAME_JP
    print(f"  ✓ 出典: {text}")


def main():
    parser = argparse.ArgumentParser(description="事業セグメント概要 PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    _validate(data)

    prs = Presentation(args.template)
    slide = prs.slides[0]

    main_message = data["main_message"]
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_message)
    print(f"  ✓ Main Message: {main_message}")

    parent_company = data["parent_company"]
    segment_name = data["segment_name"]
    chart_title = data.get("chart_title") or f"{parent_company}：{segment_name}の概要"
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
    print(f"  ✓ Chart Title: {chart_title}")

    remove_shape(slide, "Table 1")

    overview = data["overview"]
    section_title_left = overview.get("section_title", "事業の概要")
    add_section_title(slide, section_title_left, LEFT_X, PANEL_Y, LEFT_W)
    build_overview_table(slide, overview["items"], LEFT_X, PANEL_Y + Inches(0.40), LEFT_W)

    perf = data["performance"]
    mode = perf["mode"]
    section_title_right = perf.get("section_title", "業績" if mode == "revenue_chart" else "主要 KPI")
    add_section_title(slide, section_title_right, RIGHT_X, PANEL_Y, RIGHT_W)

    if mode == "revenue_chart":
        unit_label = perf.get("unit_label", "")
        if unit_label:
            add_unit_label(slide, unit_label, RIGHT_X, PANEL_Y + Inches(0.35), Inches(2.50))
        add_custom_legend(slide, perf, RIGHT_X, PANEL_Y + Inches(0.35), RIGHT_W)
        chart_top = PANEL_Y + Inches(0.55)
        chart_h = Inches(4.80)
        build_combo_chart(slide, perf, RIGHT_X, chart_top, RIGHT_W, chart_h)
        add_cagr_annotation(slide, perf, RIGHT_X, chart_top, RIGHT_W, chart_h)
    else:  # kpi_cards
        cards_top = PANEL_Y + Inches(0.55)
        cards_h = Inches(4.80)
        build_kpi_cards(slide, perf, RIGHT_X, cards_top, RIGHT_W, cards_h)

    source = data.get("source", "")
    if source:
        add_source_label(slide, source)

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
