"""
fill_customer_sales_detail.py — 主要販売先詳細テーブルをHTMLで描画→スクリーンショット→PPTXに挿入

テンプレート構造（customer-sales-detail-template.pptx）:
  - Title 1            (PLACEHOLDER): メインメッセージ（上段、太字）
  - Text Placeholder 2 (PLACEHOLDER): チャートタイトル（下段）
  - Content Area       (AUTO_SHAPE):  HTML screenshot挿入先
  - Source / Source 3  (TEXT_BOX/PLACEHOLDER): 出典（左下）

Usage:
  python fill_customer_sales_detail.py \
    --data /home/claude/customer_sales_detail_data.json \
    --output /mnt/user-data/outputs/CustomerSalesDetail_output.pptx \
    --brand stellar_aiz
"""

import argparse
import asyncio
import copy
import json
import os
import sys
import tempfile
from html import escape as _esc

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── brand_resolver bootstrap ─────────────────────────────────
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair."""
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape名マッピング ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
# Source candidates (stella uses 'Source', roleup uses 'Source 3', LibreOffice 'PlaceHolder 3')
SHAPE_SOURCE_CANDIDATES = ("Source 3", "Source", "PlaceHolder 3")

DEVICE_SCALE = 2

# ── Theme-controlled module variables (stella defaults) ──
FONT_NAME_JP = "Noto Sans CJK JP, Meiryo UI, Hiragino Sans"
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
TEXT_HEX = "333333"
SOURCE_HEX = "666666"
HEADER_BG_HEX = "F0F0F0"
HEADER_BORDER_HEX = "666666"
ROW_BORDER_HEX = "E0E0E0"
ROW_ALT_BG_HEX = "FAFAFA"
ROW_OTHER_BG_HEX = "F5F5F5"
SOURCE_FONT_PT = 10  # stella default; roleup overrides via theme.font_size_source_pt
TITLE_FONT_PT = 28
SUBTITLE_FONT_PT = 14


def _apply_theme(theme):
    """Override module-level fonts/colors/sizes from theme.json."""
    global FONT_NAME_JP, COLOR_TEXT, TEXT_HEX, SOURCE_HEX
    global HEADER_BG_HEX, HEADER_BORDER_HEX, ROW_BORDER_HEX
    global SOURCE_FONT_PT, TITLE_FONT_PT, SUBTITLE_FONT_PT
    FONT_NAME_JP = theme.font_ea
    COLOR_TEXT = theme.color("text")
    TEXT_HEX = theme.hex_no_hash("text")
    SOURCE_HEX = theme.hex_no_hash("source")
    HEADER_BG_HEX = theme.hex_no_hash("label_bg")
    HEADER_BORDER_HEX = theme.hex_no_hash("label_bar")
    ROW_BORDER_HEX = theme.hex_no_hash("highlight_other")
    SOURCE_FONT_PT = theme.pt_value("font_size_source_pt")
    TITLE_FONT_PT = theme.pt_value("font_size_title_pt")
    SUBTITLE_FONT_PT = theme.pt_value("font_size_subtitle_pt")


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_source_shape(slide):
    for cand in SHAPE_SOURCE_CANDIDATES:
        s = find_shape(slide, cand)
        if s is not None:
            return s
    return None


def _make_brand_run(para, text_str, font_pt=None, bold=False, color_hex=None):
    """Create a new run with brand font/color/size; clears existing runs first."""
    for r in list(para._p.findall(qn("a:r"))):
        para._p.remove(r)

    r_elem = etree.SubElement(para._p, qn("a:r"))
    rPr_attrs = {"lang": "ja-JP"}
    if bold:
        rPr_attrs["b"] = "1"
    if font_pt is not None:
        rPr_attrs["sz"] = str(int(font_pt * 100))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib=rPr_attrs)
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color_hex or TEXT_HEX})
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_NAME_JP})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_NAME_JP})
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text_str


def set_textbox_text(shape, text, font_pt=None, color_hex=None):
    """Overwrite shape text while applying brand font/color (preserves bold from existing)."""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    bold = False
    if para.runs:
        existing_rPr = para.runs[0]._r.find(qn("a:rPr"))
        if existing_rPr is not None and existing_rPr.get("b") == "1":
            bold = True
    _make_brand_run(para, text, font_pt=font_pt, bold=bold, color_hex=color_hex)


def fmt_number(val):
    if val is None:
        return "NA"
    if isinstance(val, (int, float)):
        return f"{int(val):,}"
    return str(val)


def fmt_pct(val):
    if val is None:
        return "NA"
    if isinstance(val, (int, float)):
        return f"{val:.1f}%"
    return str(val)


def generate_html(data, viewport_width, viewport_height):
    """販売先詳細データからHTML（10列テーブル）を生成する"""
    customers = data.get("customers", [])
    unit = _esc(data.get("unit", "千円"))

    PT = viewport_width / 900.0  # 1pt ≈ viewport_width / 900 CSS px

    font_header = int(11 * PT)
    font_body = int(10.5 * PT)
    pad_cell_v = int(1.8 * PT)
    pad_cell_h = int(2.5 * PT)

    col_widths = [
        ("3%", "center"),
        ("9%", "left"),
        ("7%", "right"),
        ("5%", "right"),
        ("22%", "left"),
        ("12%", "left"),
        ("8%", "left"),
        ("11%", "right"),
        ("9%", "right"),
        ("7%", "right"),
    ]

    headers = [
        "#",
        "企業名",
        f"売上高<br><span style='font-size:{int(font_header*0.8)}px;font-weight:400;'>({unit})</span>",
        "割合",
        "事業内容",
        "本社所在地",
        "上場",
        f"直近期売上高<br><span style='font-size:{int(font_header*0.8)}px;font-weight:400;'>({unit})</span>",
        f"利益<br><span style='font-size:{int(font_header*0.8)}px;font-weight:400;'>({unit})</span>",
        "利益率",
    ]

    colgroup = "<colgroup>\n"
    for w, _ in col_widths:
        colgroup += f'  <col style="width:{w};">\n'
    colgroup += "</colgroup>"

    header_cells = ""
    for i, h in enumerate(headers):
        _, align = col_widths[i]
        header_cells += f'''
            <th style="padding:{pad_cell_v}px {pad_cell_h}px;font-size:{font_header}px;
                font-weight:700;color:#{TEXT_HEX};text-align:{align};
                border-bottom:2px solid #{HEADER_BORDER_HEX};white-space:nowrap;">{h}</th>'''

    body_rows = ""
    for i, cust in enumerate(customers):
        name = cust.get("name", "")
        revenue = cust.get("revenue", 0)
        share = cust.get("share", 0)
        rank = cust.get("rank", i + 1)
        business = cust.get("business")
        headquarters = cust.get("headquarters")
        listing = cust.get("listing")
        latest_revenue = cust.get("latest_revenue")
        profit = cust.get("profit")
        profit_margin = cust.get("profit_margin")

        is_other = (name == "その他" or rank == -1)

        rank_str = str(rank) if not is_other else ""
        rev_str = f"{int(revenue):,}" if isinstance(revenue, (int, float)) else str(revenue)
        share_str = f"{share:.1f}%" if isinstance(share, (int, float)) else str(share)

        if is_other:
            biz_str = ""
            hq_str = ""
            list_str = ""
            lr_str = ""
            prof_str = ""
            pm_str = ""
        else:
            biz_str = _esc(business) if business else ""
            hq_str = _esc(headquarters) if headquarters else ""
            list_str = _esc(listing) if listing else ""
            lr_str = fmt_number(latest_revenue)
            prof_str = fmt_number(profit)
            pm_str = fmt_pct(profit_margin)

        if is_other:
            bg = f"#{ROW_OTHER_BG_HEX}"
        elif i % 2 == 1:
            bg = f"#{ROW_ALT_BG_HEX}"
        else:
            bg = "#FFFFFF"

        border = f"1px solid #{ROW_BORDER_HEX}"
        biz_style = "white-space:normal;word-break:break-all;"

        cells = [
            (rank_str, "center", ""),
            (_esc(name), "left", "white-space:nowrap;"),
            (rev_str, "right", ""),
            (share_str, "right", ""),
            (biz_str, "left", biz_style),
            (hq_str, "left", ""),
            (list_str, "left", ""),
            (lr_str, "right", ""),
            (prof_str, "right", ""),
            (pm_str, "right", ""),
        ]

        row_cells = ""
        for j, (val, align, extra_style) in enumerate(cells):
            row_cells += f'''
                <td style="padding:{pad_cell_v}px {pad_cell_h}px;font-size:{font_body}px;
                    color:#{TEXT_HEX};text-align:{align};border-bottom:{border};{extra_style}">{val}</td>'''

        body_rows += f'''
            <tr style="background:{bg};">{row_cells}
            </tr>'''

    body_pad = int(2 * PT)
    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
html, body {{
    width:{viewport_width}px;
    height:{viewport_height}px;
    font-family:'{FONT_NAME_JP}','Noto Sans CJK JP','Meiryo UI','Hiragino Sans',sans-serif;
    background:#FFFFFF;
    padding:{body_pad}px;
    overflow:hidden;
}}
</style>
</head>
<body>
<table style="width:100%;border-collapse:collapse;
    font-family:'{FONT_NAME_JP}','Noto Sans CJK JP','Meiryo UI','Hiragino Sans',sans-serif;
    table-layout:fixed;">
    {colgroup}
    <thead>
        <tr style="background:#{HEADER_BG_HEX};">{header_cells}
        </tr>
    </thead>
    <tbody>{body_rows}
    </tbody>
</table>
</body></html>'''

    return html


async def take_screenshot(html_content, output_path, vw, vh):
    from playwright.async_api import async_playwright

    with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8") as f:
        f.write(html_content)
        html_path = f.name

    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page(
                viewport={"width": vw, "height": vh},
                device_scale_factor=DEVICE_SCALE,
            )
            await page.goto(f"file://{html_path}", wait_until="networkidle")
            await page.wait_for_timeout(500)
            await page.screenshot(path=output_path, full_page=False)
            await browser.close()
    finally:
        os.unlink(html_path)

    print(f"  ✓ Screenshot saved: {output_path}")


def insert_screenshot_into_pptx(template_path, data, screenshot_path, output_path, vw, vh):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    for s in slide.shapes:
        print(f"  Shape: '{s.name}' type={s.shape_type}")

    main_msg = data.get("main_message", "")
    msg_shape = find_shape(slide, SHAPE_MAIN_MESSAGE) or find_shape(slide, "PlaceHolder 1")
    set_textbox_text(msg_shape, main_msg, font_pt=TITLE_FONT_PT, color_hex=TEXT_HEX)

    chart_title = data.get("chart_title", "主要販売先からの完成工事売上高と割合")
    title_shape = find_shape(slide, SHAPE_CHART_TITLE) or find_shape(slide, "PlaceHolder 2")
    set_textbox_text(title_shape, chart_title, font_pt=SUBTITLE_FONT_PT, color_hex=TEXT_HEX)

    source_text = data.get("source", "")
    source_shape = find_source_shape(slide)
    if source_shape is not None:
        if source_text:
            set_textbox_text(source_shape, f"出典：{source_text}",
                             font_pt=SOURCE_FONT_PT, color_hex=SOURCE_HEX)
        else:
            set_textbox_text(source_shape, "",
                             font_pt=SOURCE_FONT_PT, color_hex=SOURCE_HEX)

    content_shape = find_shape(slide, SHAPE_CONTENT_AREA)
    if content_shape and os.path.exists(screenshot_path):
        left = content_shape.left
        top = content_shape.top
        width = content_shape.width
        height = content_shape.height
        slide.shapes.add_picture(screenshot_path, left, top, width, height)
        print(f"  ✓ Screenshot inserted at ({left},{top}) size=({width},{height})")
        # Remove the Content Area placeholder shape so brand_compliance C1
        # ('no guide rectangle') passes. The picture has already been laid down.
        _silent_remove_shape(slide, SHAPE_CONTENT_AREA)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    _finalize_pptx(output_path)
    print(f"  ✓ PPTX saved: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="主要販売先詳細テーブルスライド生成。--brand stellar_aiz / roleup を選択。",
    )
    parser.add_argument("--data", required=True, help="JSONデータファイルパス")
    parser.add_argument("--template", required=False, default=None, help="PPTXテンプレートパス (任意)")
    parser.add_argument("--output", required=True, help="出力PPTXファイルパス")
    add_brand_arg(parser)
    args = parser.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "customer-sales-detail")
    print(f"  ✓ Brand:    {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    print("=== 主要販売先詳細テーブルスライド生成 ===")

    # Content Area の実寸から viewport を逆算
    prs_probe = Presentation(template_path)
    content = find_shape(prs_probe.slides[0], SHAPE_CONTENT_AREA)
    if content is not None:
        vw = int(content.width / 914400.0 * 200)
        vh = int(content.height / 914400.0 * 200)
    else:
        vw, vh = 2504, 1080

    print(f"Step 1: HTML生成... (viewport {vw}×{vh})")
    html_content = generate_html(data, vw, vh)

    debug_html_path = os.path.join(tempfile.gettempdir(), "customer_sales_detail_debug.html")
    with open(debug_html_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    print("Step 2: スクリーンショット取得...")
    screenshot_path = os.path.join(tempfile.gettempdir(), "customer_sales_detail_screenshot.png")
    asyncio.run(take_screenshot(html_content, screenshot_path, vw, vh))

    print("Step 3: PPTX生成...")
    insert_screenshot_into_pptx(template_path, data, screenshot_path, args.output, vw, vh)

    if os.path.exists(screenshot_path):
        os.unlink(screenshot_path)

    print("=== 完了 ===")


if __name__ == "__main__":
    main()
