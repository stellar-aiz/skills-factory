"""One-shot generator: customer-sales-detail-pptx の roleup curated template を
stella template から派生して作る。

実行は本ファイルを直接 python3 で起動するだけ。生成済テンプレは
`assets/roleup/customer-sales-detail-template.pptx` に上書き保存する。

設計:
  - スライドサイズ: A4 landscape (10691813 × 7559675 EMU = 11.69 × 8.27 in)
  - Title 1, Text Placeholder 2, Content Area: 座標を A4 用に再設定
  - Source → Source 3 (placeholder 名規約 + 6pt Yu Gothic UI 茶系)
  - フォント: Yu Gothic UI 全 run
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(SKILL_DIR, "assets", "stellar_aiz", "customer-sales-detail-template.pptx")
DST = os.path.join(SKILL_DIR, "assets", "roleup", "customer-sales-detail-template.pptx")

A4_W = 10691813
A4_H = 7559675

FONT_EA = "Yu Gothic UI"
COLOR_TEXT = RGBColor(0x24, 0x1A, 0x17)
COLOR_SUBTITLE = RGBColor(0x89, 0x71, 0x41)
COLOR_SOURCE = RGBColor(0x3E, 0x3A, 0x39)


def _find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _set_run_font(run, font_name, color, size_pt=None, bold=False):
    run.font.name = font_name
    run.font.color.rgb = color
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    run.font.bold = bold
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn("a:rPr"))
    for tag in ("a:latin", "a:ea"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": font_name})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": font_name})


def main():
    prs = Presentation(SRC)
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    slide = prs.slides[0]

    for s in list(slide.shapes):
        if "think-cell" in (s.name or ""):
            sp = s._element
            sp.getparent().remove(sp)

    title = _find(slide, "Title 1")
    if title is not None:
        title.left = Inches(0.41)
        title.top = Inches(0.41)
        title.width = Inches(10.87)
        title.height = Inches(0.50)
        if title.has_text_frame:
            for para in title.text_frame.paragraphs:
                for run in para.runs:
                    _set_run_font(run, FONT_EA, COLOR_TEXT, size_pt=22, bold=True)

    sub = _find(slide, "Text Placeholder 2")
    if sub is not None:
        sub.left = Inches(0.41)
        sub.top = Inches(1.00)
        sub.width = Inches(10.87)
        sub.height = Inches(0.36)
        if sub.has_text_frame:
            for para in sub.text_frame.paragraphs:
                for run in para.runs:
                    _set_run_font(run, FONT_EA, COLOR_SUBTITLE, size_pt=12, bold=False)

    content = _find(slide, "Content Area")
    if content is not None:
        content.left = Inches(0.41)
        content.top = Inches(1.50)
        content.width = Inches(10.87)
        content.height = Inches(5.70)

    # Rename Source → Source 3 for compliance check rule
    src = _find(slide, "Source")
    if src is not None:
        src.name = "Source 3"
        src.left = Inches(0.41)
        src.top = Inches(7.30)
        src.width = Inches(10.87)
        src.height = Inches(0.40)
        if src.has_text_frame:
            tf = src.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                # Replace existing runs with a single 6pt source-styled run
                p = tf.paragraphs[0]
                existing_text = "".join(r.text for r in p.runs) or "出典: "
                for r in list(p._p.findall(qn("a:r"))):
                    p._p.remove(r)
                run = p.add_run()
                run.text = existing_text
                _set_run_font(run, FONT_EA, COLOR_SOURCE, size_pt=6, bold=False)

    if _find(slide, "Source 3") is None:
        src_tb = slide.shapes.add_textbox(
            Inches(0.41), Inches(7.30), Inches(10.87), Inches(0.40)
        )
        src_tb.name = "Source 3"
        para = src_tb.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "出典: "
        _set_run_font(run, FONT_EA, COLOR_SOURCE, size_pt=6, bold=False)

    prs.save(DST)
    print(f"✅ Generated: {DST}")
    print(f"   slide_size: {prs.slide_width} × {prs.slide_height} EMU")
    for s in slide.shapes:
        print(f"   - {s.name}")


if __name__ == "__main__":
    main()
