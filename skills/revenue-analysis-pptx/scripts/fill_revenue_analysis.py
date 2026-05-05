"""
fill_revenue_analysis.py - Revenue Analysis slide generator (v5, brand-aware)

Layout: Top = EBITDA margin line, Middle = CAGR arrow, Bottom = Grouped bar chart (Revenue + EBITDA)
Grouped bar chart uses two series: Revenue (accent_revenue_bar) and EBITDA (chart_palette[1]).

Phase 2 (ISSUE-010): brand-aware で stellar_aiz / roleup を出し分け。
"""
import argparse, json, math, os, sys

# brand_resolver bootstrap (Phase 2 — brand-aware: stellar_aiz / roleup)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import resolve_brand, add_brand_arg  # noqa: E402
from format_helpers import resolve_top_text, resolve_subtitle_text, require_source  # noqa: E402

SKILL_ID = "revenue-analysis-pptx"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_SOURCE = "Source 3"

# Defaults (stella V1 values; reassigned in main() via _apply_theme).
MARGIN_ZONE_TOP = Inches(1.80)
MARGIN_ZONE_BOTTOM = Inches(3.10)

CHART_LEFT = Inches(0.80)
CHART_WIDTH = Inches(11.73)
CHART_TOP = Inches(3.20)
CHART_HEIGHT = Inches(3.70)

SOURCE_X = Inches(0.41)
SOURCE_Y = Inches(7.05)
SOURCE_W = Inches(8.0)
SOURCE_H = Inches(0.30)

COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_BAR_REV = RGBColor(0x4E, 0x79, 0xA7)
COLOR_BAR_EBITDA = RGBColor(0xE1, 0x81, 0x2C)
COLOR_LINE = RGBColor(0x00, 0x33, 0x66)
COLOR_CAGR = RGBColor(0x33, 0x33, 0x33)
COLOR_SOURCE = RGBColor(0x66, 0x66, 0x66)

COLOR_BAR_REV_HEX = "4E79A7"
COLOR_BAR_EBITDA_HEX = "E1812C"
COLOR_DLBL_HEX = "333333"
COLOR_LINE_HEX = "003366"

FONT = "Meiryo UI"

# Font sizes in pt (overwritten by _apply_theme).
PT_AXIS = 11
PT_DLBL_REV = 11
PT_DLBL_EBITDA = 10
PT_CAT_AXIS = 12
PT_MARGIN_LBL = 12
PT_CAGR = 16
PT_LEGEND = 12
PT_UNIT = 12
PT_SOURCE = 10

PLOT_L_PCT, PLOT_R_PCT, PLOT_T_PCT, PLOT_B_PCT = 0.07, 0.02, 0.04, 0.12

# Theme module-global; populated in main() via _apply_theme(theme).
_THEME = None


def _apply_theme(theme):
    """Reassign module-level brand-aware globals from a resolved BrandTheme."""
    global _THEME
    global MARGIN_ZONE_TOP, MARGIN_ZONE_BOTTOM
    global CHART_LEFT, CHART_WIDTH, CHART_TOP, CHART_HEIGHT
    global SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H
    global COLOR_TEXT, COLOR_BAR_REV, COLOR_BAR_EBITDA, COLOR_LINE, COLOR_CAGR, COLOR_SOURCE
    global COLOR_BAR_REV_HEX, COLOR_BAR_EBITDA_HEX, COLOR_DLBL_HEX, COLOR_LINE_HEX
    global FONT
    global PT_AXIS, PT_DLBL_REV, PT_DLBL_EBITDA, PT_CAT_AXIS, PT_MARGIN_LBL
    global PT_CAGR, PT_LEGEND, PT_UNIT, PT_SOURCE

    _THEME = theme

    MARGIN_ZONE_TOP = theme.layout("margin_zone_top_in")
    MARGIN_ZONE_BOTTOM = theme.layout("margin_zone_bottom_in")
    CHART_LEFT = theme.layout("chart_left_in")
    CHART_WIDTH = theme.layout("chart_width_in")
    CHART_TOP = theme.layout("chart_top_in")
    CHART_HEIGHT = theme.layout("chart_height_in")
    SOURCE_X = theme.layout("source_x_in")
    SOURCE_Y = theme.layout("source_y_in")
    SOURCE_W = theme.layout("source_w_in")
    SOURCE_H = theme.layout("source_h_in")

    COLOR_TEXT = theme.color("text")
    COLOR_BAR_REV = theme.color("accent_revenue_bar")
    COLOR_LINE = theme.color("accent_op_margin_line")
    COLOR_CAGR = theme.color("cagr_arrow")
    COLOR_SOURCE = theme.color("source")

    COLOR_BAR_REV_HEX = theme.hex_no_hash("accent_revenue_bar").upper()
    COLOR_DLBL_HEX = theme.hex_no_hash("text").upper()
    COLOR_LINE_HEX = theme.hex_no_hash("accent_op_margin_line").upper()

    FONT = theme.font_ea

    if theme.id == "stellar_aiz":
        # V1 hardcoded values for regression-zero. Stella's chart_palette[1]
        # (#F28E2B) is close to V1's #E1812C but not identical, so pin V1 hex.
        COLOR_BAR_EBITDA = RGBColor(0xE1, 0x81, 0x2C)
        COLOR_BAR_EBITDA_HEX = "E1812C"
        PT_AXIS = 11
        PT_DLBL_REV = 11
        PT_DLBL_EBITDA = 10
        PT_CAT_AXIS = 12
        PT_MARGIN_LBL = 12
        PT_CAGR = 16
        PT_LEGEND = 12
        PT_UNIT = 12
        PT_SOURCE = 10
    else:
        # EBITDA = chart_palette[1] (#897141 beige for roleup, harmonious with
        # accent_revenue_bar #7C4C2C and accent_op_margin_line #604C3F).
        ebh = theme.chart_palette[1].lstrip("#").upper()
        COLOR_BAR_EBITDA = RGBColor(int(ebh[0:2], 16), int(ebh[2:4], 16), int(ebh[4:6], 16))
        COLOR_BAR_EBITDA_HEX = ebh
        # Roleup C4 allowed set: {22, 14, 12, 10, 6}
        body = theme.font_size_body_pt_value(skill_id=SKILL_ID)  # 10
        PT_AXIS = body
        PT_DLBL_REV = body
        PT_DLBL_EBITDA = body
        PT_CAT_AXIS = body
        PT_MARGIN_LBL = body
        PT_CAGR = 14  # subtitle/key-message size, in allowed set
        PT_LEGEND = 12
        PT_UNIT = body
        PT_SOURCE = theme.pt_value("font_size_source_pt")  # 6


def _silent_remove_shape(slide, shape_name):
    for s in list(slide.shapes):
        if s.name == shape_name:
            sp = s._element
            sp.getparent().remove(sp)


def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name: return s
    return None

def set_text(shape, text):
    if not shape: return
    tf = shape.text_frame; p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]: r.text = ""
    else:
        r = etree.SubElement(p._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t = etree.SubElement(r, qn("a:t")); t.text = text

def remove_shape(slide, name):
    s = find_shape(slide, name)
    if s: slide.shapes._spTree.remove(s._element)

def nice_max(val):
    if val <= 0: return 100
    t = val * 1.20
    mag = 10 ** math.floor(math.log10(t))
    for s in [1, 2, 2.5, 5, 10]:
        c = math.ceil(t / (s * mag)) * (s * mag)
        if c >= t: return int(c) if c == int(c) else c
    return math.ceil(t / mag) * mag

def nice_tick(axis_max):
    raw = axis_max / 5
    mag = 10 ** math.floor(math.log10(raw))
    for n in [1, 2, 2.5, 5, 10]:
        c = n * mag
        if c >= raw * 0.8: return int(c) if c == int(c) else c
    return raw

def fmt_num(v):
    return f"{int(v):,}" if v == int(v) else f"{v:,.1f}"

def chart_plot_bounds():
    cl, ct, cw, ch = int(CHART_LEFT), int(CHART_TOP), int(CHART_WIDTH), int(CHART_HEIGHT)
    return (cl + cw*PLOT_L_PCT, ct + ch*PLOT_T_PCT,
            cl + cw*(1-PLOT_R_PCT), ct + ch*(1-PLOT_B_PCT))

def bar_cx(n, i):
    """X center of the i-th category group"""
    pl, _, pr, _ = chart_plot_bounds()
    w = (pr - pl) / n
    return pl + (i + 0.5) * w

def bar_top_y(val, amax):
    _, pt, _, pb = chart_plot_bounds()
    return pb - (val / amax * (pb - pt)) if amax > 0 else pb

def add_dlabels(ser, pos='outEnd', fmt='#,##0', color=None, sz=None, bold=False):
    if color is None: color = COLOR_DLBL_HEX
    if sz is None: sz = PT_DLBL_REV * 100
    dl = ser.find(qn('c:dLbls'))
    if dl is None: dl = etree.SubElement(ser, qn('c:dLbls'))
    for c in list(dl): dl.remove(c)
    etree.SubElement(dl, qn('c:numFmt'), attrib={'formatCode': fmt, 'sourceLinked': '0'})
    for tag, val in [('showLegendKey','0'),('showVal','1'),('showCatName','0'),
                     ('showSerName','0'),('showPercent','0'),('showBubbleSize','0')]:
        etree.SubElement(dl, qn(f'c:{tag}'), attrib={'val': val})
    etree.SubElement(dl, qn('c:dLblPos'), attrib={'val': pos})
    txPr = etree.SubElement(dl, qn('c:txPr'))
    etree.SubElement(txPr, qn('a:bodyPr'))
    etree.SubElement(txPr, qn('a:lstStyle'))
    p = etree.SubElement(txPr, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    att = {'sz': str(sz)}
    if bold: att['b'] = '1'
    dr = etree.SubElement(pPr, qn('a:defRPr'), attrib=att)
    etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT})
    etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT})
    sf = etree.SubElement(dr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': color})

def set_ser_color(ser, hex_color):
    """Set fill color of a bar chart series"""
    sp = ser.find(qn('c:spPr'))
    if sp is None: sp = etree.SubElement(ser, qn('c:spPr'))
    sf = sp.find(qn('a:solidFill'))
    if sf is None: sf = etree.SubElement(sp, qn('a:solidFill'))
    for c in list(sf): sf.remove(c)
    etree.SubElement(sf, qn('a:srgbClr'), attrib={'val': hex_color})


def build_grouped_bar_chart(slide, jd):
    """Build grouped bar chart with Revenue + EBITDA in the LOWER zone."""
    from pptx.chart.data import CategoryChartData
    data = jd["data"]
    rev_label = jd.get("bar_label", "売上高")
    ebitda_label = "EBITDA"

    cd = CategoryChartData()
    cd.categories = [d["year"] for d in data]
    cd.add_series(rev_label, [d["revenue"] for d in data])
    cd.add_series(ebitda_label, [d["ebitda"] for d in data])

    cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT, cd)
    chart = cf.chart
    pa = chart._chartSpace.chart.plotArea
    bc = pa.findall(qn('c:barChart'))[0]

    # Show left Y-axis
    vax = pa.findall(qn('c:valAx'))[0]
    de = vax.find(qn('c:delete'))
    if de is None: de = etree.SubElement(vax, qn('c:delete'))
    de.set('val', '0')

    revs = [d["revenue"] for d in data]
    amax = nice_max(max(revs))
    tick = nice_tick(amax)

    sc = vax.find(qn('c:scaling'))
    if sc is None: sc = etree.SubElement(vax, qn('c:scaling'))
    for tag, val in [('c:max', str(amax)), ('c:min', '0')]:
        e = sc.find(qn(tag))
        if e is None: e = etree.SubElement(sc, qn(tag))
        e.set('val', val)
    mu = vax.find(qn('c:majorUnit'))
    if mu is None: mu = etree.SubElement(vax, qn('c:majorUnit'))
    mu.set('val', str(tick))
    nf = vax.find(qn('c:numFmt'))
    if nf is None: nf = etree.SubElement(vax, qn('c:numFmt'))
    nf.set('formatCode', '#,##0'); nf.set('sourceLinked', '0')

    # Y-axis font
    axis_sz = str(PT_AXIS * 100)
    tp = vax.find(qn('c:txPr'))
    if tp is None: tp = etree.SubElement(vax, qn('c:txPr'))
    bp = tp.find(qn('a:bodyPr'))
    if bp is None: bp = etree.SubElement(tp, qn('a:bodyPr')); tp.insert(0, bp)
    if tp.find(qn('a:lstStyle')) is None: etree.SubElement(tp, qn('a:lstStyle'))
    p = tp.find(qn('a:p'))
    if p is None: p = etree.SubElement(tp, qn('a:p'))
    pp = p.find(qn('a:pPr'))
    if pp is None: pp = etree.SubElement(p, qn('a:pPr'))
    dr = pp.find(qn('a:defRPr'))
    if dr is None: dr = etree.SubElement(pp, qn('a:defRPr'), attrib={'sz': axis_sz})
    else: dr.set('sz', axis_sz)
    if dr.find(qn('a:latin')) is None: etree.SubElement(dr, qn('a:latin'), attrib={'typeface': FONT})
    if dr.find(qn('a:ea')) is None: etree.SubElement(dr, qn('a:ea'), attrib={'typeface': FONT})
    asf = dr.find(qn('a:solidFill'))
    if asf is None: asf = etree.SubElement(dr, qn('a:solidFill'))
    for c in list(asf): asf.remove(c)
    etree.SubElement(asf, qn('a:srgbClr'), attrib={'val': COLOR_DLBL_HEX})

    # Series colors
    sers = bc.findall(qn('c:ser'))
    set_ser_color(sers[0], COLOR_BAR_REV_HEX)
    set_ser_color(sers[1], COLOR_BAR_EBITDA_HEX)

    # Data labels
    add_dlabels(sers[0], 'outEnd', '#,##0', COLOR_DLBL_HEX, PT_DLBL_REV * 100, True)
    add_dlabels(sers[1], 'outEnd', '#,##0', COLOR_DLBL_HEX, PT_DLBL_EBITDA * 100, False)

    chart.has_legend = False

    # Remove gridlines
    for v in pa.findall(qn('c:valAx')):
        for mg in v.findall(qn('c:majorGridlines')): v.remove(mg)
        for mg in v.findall(qn('c:minorGridlines')): v.remove(mg)
    for c in pa.findall(qn('c:catAx')):
        for mg in c.findall(qn('c:majorGridlines')): c.remove(mg)

    # Gap between groups (narrower for grouped)
    gw = bc.find(qn('c:gapWidth'))
    if gw is None: gw = etree.SubElement(bc, qn('c:gapWidth'))
    gw.set('val', '80')

    # Cat axis font
    cat_sz = str(PT_CAT_AXIS * 100)
    for ax in pa.findall(qn('c:catAx')):
        de2 = ax.find(qn('c:delete'))
        if de2 is not None and de2.get('val') == '1': continue
        tp2 = ax.find(qn('c:txPr'))
        if tp2 is None: tp2 = etree.SubElement(ax, qn('c:txPr'))
        bp2 = tp2.find(qn('a:bodyPr'))
        if bp2 is None: bp2 = etree.SubElement(tp2, qn('a:bodyPr')); tp2.insert(0, bp2)
        bp2.set('rot', '0'); bp2.set('vert', 'horz')
        if tp2.find(qn('a:lstStyle')) is None: etree.SubElement(tp2, qn('a:lstStyle'))
        p2 = tp2.find(qn('a:p'))
        if p2 is None: p2 = etree.SubElement(tp2, qn('a:p'))
        pp2 = p2.find(qn('a:pPr'))
        if pp2 is None: pp2 = etree.SubElement(p2, qn('a:pPr'))
        dr2 = pp2.find(qn('a:defRPr'))
        if dr2 is None: dr2 = etree.SubElement(pp2, qn('a:defRPr'), attrib={'sz': cat_sz, 'b': '1'})
        else: dr2.set('sz', cat_sz); dr2.set('b', '1')
        if dr2.find(qn('a:latin')) is None: etree.SubElement(dr2, qn('a:latin'), attrib={'typeface': FONT})
        if dr2.find(qn('a:ea')) is None: etree.SubElement(dr2, qn('a:ea'), attrib={'typeface': FONT})

    chart.has_title = False
    psp = pa.find(qn('c:spPr'))
    if psp is None: psp = etree.SubElement(pa, qn('c:spPr'))
    etree.SubElement(psp, qn('a:noFill'))

    print(f"  Grouped bar chart: {len(data)} years, axis_max={amax}, tick={tick}")
    return cf, amax


def draw_margin_line(slide, jd):
    """Draw EBITDA margin line in the UPPER zone."""
    data = jd["data"]
    n = len(data)
    margins = [round(d["ebitda"]/d["revenue"]*100, 1) if d["revenue"]>0 else 0 for d in data]
    mx, mn = max(margins), min(margins)
    rng = mx - mn

    zone_top = int(MARGIN_ZONE_TOP)
    zone_bot = int(MARGIN_ZONE_BOTTOM)
    label_space = Inches(0.30)
    usable_top = zone_top + int(label_space)
    usable_h = zone_bot - usable_top

    top_pct, bot_pct = 0.98, 0.90
    if rng < 0.5: top_pct = bot_pct = 0.95

    pts = []
    for i, m in enumerate(margins):
        cx = int(bar_cx(n, i))
        if rng >= 0.5:
            norm = (m - mn) / rng
            y = int(usable_top + usable_h * (1.0 - norm))
        else:
            y = int(usable_top + usable_h * 0.4)
        pts.append((cx, y))

    for i in range(len(pts)-1):
        x1,y1 = pts[i]; x2,y2 = pts[i+1]
        cn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        cn.line.color.rgb = COLOR_LINE; cn.line.width = Pt(2.0)

    ms = Inches(0.16)
    for x,y in pts:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x-int(ms/2), y-int(ms/2), ms, ms)
        c.fill.solid(); c.fill.fore_color.rgb = COLOR_LINE; c.line.fill.background()

    lh, lw = Inches(0.25), Inches(0.80)
    for i,(x,y) in enumerate(pts):
        tb = slide.shapes.add_textbox(x-int(lw/2), y-int(lh)-Inches(0.04), lw, lh)
        tf = tb.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"{margins[i]:.1f}%"
        r.font.size = Pt(PT_MARGIN_LBL); r.font.bold = True
        r.font.color.rgb = COLOR_LINE; r.font.name = FONT

    print(f"  Margin line: {n} pts, {mn:.1f}%-{mx:.1f}%")
    return margins


def add_cagr(slide, jd, amax):
    """Draw CAGR arrow above revenue bars."""
    data = jd["data"]
    if len(data) < 2: return
    fr, lr = data[0]["revenue"], data[-1]["revenue"]
    ny, nc = len(data)-1, len(data)
    if fr > 0 and lr > 0 and ny > 0:
        cagr = (lr/fr)**(1.0/ny) - 1
        txt = f"+{cagr*100:.1f}%" if cagr >= 0 else f"{cagr*100:.1f}%"
    else: txt = "N/A"

    x1 = int(bar_cx(nc, 0)); x2 = int(bar_cx(nc, nc-1))
    gap = Inches(0.40)
    y1 = int(bar_top_y(fr, amax) - gap)
    y2 = int(bar_top_y(lr, amax) - gap)

    cn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    cn.line.color.rgb = COLOR_CAGR; cn.line.width = Pt(1.5)
    sp = cn._element.find(qn('p:spPr'))
    if sp is None: sp = cn._element.find(qn('a:spPr'))
    ln = sp.find(qn('a:ln'))
    if ln is None: ln = etree.SubElement(sp, qn('a:ln'))
    etree.SubElement(ln, qn('a:tailEnd'), attrib={'type':'triangle','w':'med','len':'med'})

    tw, th = Inches(1.30), Inches(0.40)
    midx, midy = (x1+x2)/2, (y1+y2)/2
    ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(midx-tw/2), int(midy-th/2), tw, th)
    ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    ov.line.color.rgb = COLOR_CAGR; ov.line.width = Pt(1.0)
    tf = ov.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(PT_CAGR); r.font.bold = True
    r.font.color.rgb = COLOR_TEXT; r.font.name = FONT
    print(f"  CAGR: {txt}")


def add_legend(slide, jd, left, top, width):
    """Custom legend: ■売上高  ■EBITDA  ●━EBITDA率"""
    rev_label = jd.get("bar_label", "売上高")
    line_label = jd.get("line_label", "EBITDA率")

    lh = Inches(0.22)
    lx = left + width - Inches(5.00)  # wider for 3 items
    ss = Inches(0.14)
    sy = top + (lh - ss) / 2

    cur_x = lx

    # ■ 売上高
    bm1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(cur_x), int(sy), ss, ss)
    bm1.fill.solid(); bm1.fill.fore_color.rgb = COLOR_BAR_REV; bm1.line.fill.background()
    cur_x += ss + Inches(0.05)
    tb1 = slide.shapes.add_textbox(int(cur_x), top, Inches(0.70), lh)
    tf1 = tb1.text_frame; tf1.word_wrap = False
    p1 = tf1.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run(); r1.text = rev_label; r1.font.size = Pt(PT_LEGEND)
    r1.font.color.rgb = COLOR_TEXT; r1.font.name = FONT
    cur_x += Inches(0.85)

    # ■ EBITDA
    bm2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(cur_x), int(sy), ss, ss)
    bm2.fill.solid(); bm2.fill.fore_color.rgb = COLOR_BAR_EBITDA; bm2.line.fill.background()
    cur_x += ss + Inches(0.05)
    tb2 = slide.shapes.add_textbox(int(cur_x), top, Inches(0.80), lh)
    tf2 = tb2.text_frame; tf2.word_wrap = False
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = "EBITDA"; r2.font.size = Pt(PT_LEGEND)
    r2.font.color.rgb = COLOR_TEXT; r2.font.name = FONT
    cur_x += Inches(1.00)

    # ●━ EBITDA率
    llw = Inches(0.30); lly = top + lh / 2
    cn = slide.shapes.add_connector(1, int(cur_x), int(lly), int(cur_x + llw), int(lly))
    cn.line.color.rgb = COLOR_LINE; cn.line.width = Pt(1.5)
    cs = Inches(0.12)
    cx2 = cur_x + (llw - cs) / 2; cy2 = top + (lh - cs) / 2
    cc = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx2), int(cy2), cs, cs)
    cc.fill.solid(); cc.fill.fore_color.rgb = COLOR_LINE; cc.line.fill.background()
    cur_x += llw + Inches(0.05)
    tb3 = slide.shapes.add_textbox(int(cur_x), top, Inches(1.20), lh)
    tf3 = tb3.text_frame; tf3.word_wrap = False
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run(); r3.text = line_label; r3.font.size = Pt(PT_LEGEND)
    r3.font.color.rgb = COLOR_TEXT; r3.font.name = FONT


def add_unit(slide, text, left, top):
    tb = slide.shapes.add_textbox(left, top, Inches(3.0), Inches(0.22))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(PT_UNIT)
    r.font.color.rgb = COLOR_TEXT; r.font.name = FONT


def add_source_dynamic(slide, text):
    """Stella: dynamic textbox at SOURCE_X/Y (V1 behaviour)."""
    tb = slide.shapes.add_textbox(SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(PT_SOURCE)
    r.font.color.rgb = COLOR_SOURCE; r.font.name = FONT


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True)
    ap.add_argument(
        "--template", required=False, default=None,
        help="Optional explicit template path. If omitted, resolved from --brand "
             "(via brand_resolver.template_path).",
    )
    ap.add_argument("--output", required=True)
    add_brand_arg(ap)
    args = ap.parse_args()

    theme = resolve_brand(args.brand, SKILL_DIR)
    _apply_theme(theme)
    template_path = args.template or theme.template_path(SKILL_DIR, "revenue-analysis")
    print(f"  ✓ Brand: {theme.id} ({theme.label})")
    print(f"  ✓ Template: {template_path}")

    with open(args.data, "r", encoding="utf-8") as f:
        jd = json.load(f)

    # Roleup: source field is required (hard-fail). Stella: no-op.
    require_source(jd, theme, skill_id=SKILL_ID)

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Top / subtitle placeholder semantics differ between brands:
    #   stella: Title 1 = main_message (結論文), Text Placeholder 2 = chart_title
    #   roleup: Title 1 = chart_title (スライドタイトル), Text Placeholder 2 = main_message
    top_text = resolve_top_text(jd, theme)
    sub_text = resolve_subtitle_text(jd, theme) or "売上分析ー売上高・EBITDAの推移"
    set_text(find_shape(slide, SHAPE_MAIN_MESSAGE), top_text)
    set_text(find_shape(slide, SHAPE_CHART_TITLE), sub_text)
    print(f"  ✓ Top placeholder ({theme.top_placeholder_field()}): {top_text[:40]}")
    print(f"  ✓ Subtitle placeholder ({theme.subtitle_placeholder_field()}): {sub_text[:40]}")

    # Stella template carries Table 1 placeholder; roleup template does not.
    # remove_shape is a no-op when shape is absent.
    remove_shape(slide, "Table 1")

    # Roleup: silently remove brown guide rectangles carried by template.
    _silent_remove_shape(slide, "正方形/長方形 1")
    _silent_remove_shape(slide, "正方形/長方形 8")

    lt = MARGIN_ZONE_TOP - Inches(0.25)
    ul = jd.get("unit_label", "（単位：百万円、%）")
    if ul: add_unit(slide, ul, CHART_LEFT, lt)
    add_legend(slide, jd, CHART_LEFT, lt, CHART_WIDTH)

    margins = draw_margin_line(slide, jd)
    cf, amax = build_grouped_bar_chart(slide, jd)
    add_cagr(slide, jd, amax)

    src = jd.get("source", "")
    if src:
        if theme.id == "stellar_aiz":
            add_source_dynamic(slide, src)
        else:
            source_shape = find_shape(slide, SHAPE_SOURCE)
            if source_shape is not None:
                set_text(source_shape, src)
            else:
                add_source_dynamic(slide, src)
        print(f"  ✓ Source: {src[:40]}")

    outdir = os.path.dirname(args.output)
    if outdir: os.makedirs(outdir, exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  Output: {args.output}")

if __name__ == "__main__":
    main()
