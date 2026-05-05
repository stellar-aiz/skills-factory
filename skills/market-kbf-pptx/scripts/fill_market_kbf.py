"""
fill_market_kbf.py — 市場KBF（Key Business Factor）スライドをPPTXネイティブテーブルで生成

テンプレート構造 (market-kbf-template.pptx):
  - Title 1            (PLACEHOLDER): メインメッセージ
  - Text Placeholder 2 (PLACEHOLDER): チャートタイトル
  - Content Area       (AUTO_SHAPE):  削除してネイティブテーブルに置換
  - Source             (TEXT_BOX):    出典

方式:
  3列×4行（ヘッダー1行 + KBF×3行）のネイティブテーブルを生成。
  列1: KBF（紺背景＋白太字、強調）
  列2: 詳細（白背景、本文）
  列3: プレイヤーの抑え方の例（薄グレー背景、プレイヤー名のみ太字）

  KBFは必ず3つ固定。各KBFの player_examples は1〜5社（5社上限）。

Usage:
  python fill_market_kbf.py \
    --data /path/to/market_kbf_data.json \
    --template <SKILL_DIR>/assets/market-kbf-template.pptx \
    --output /path/to/MarketKBF_output.pptx
"""

import argparse
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


# ── Shape名マッピング（competitor-summary-template と同構造）──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE = "Source"

# ── レイアウト定数 ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

CONTENT_LEFT = Inches(0.41)
CONTENT_TOP = Inches(1.50)
CONTENT_WIDTH = Inches(12.52)
CONTENT_BOTTOM = Inches(6.90)
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP

# 列幅比率（合計1.0）
COL_RATIO_KBF = 0.20
COL_RATIO_DESC = 0.30
COL_RATIO_EXAMPLES = 0.50

# ── 色定数 ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_HEADER_BG = RGBColor(0xF0, 0xF0, 0xF0)
COLOR_KBF_BG = RGBColor(0x1F, 0x38, 0x64)        # 紺色（列1強調）
COLOR_KBF_TEXT = RGBColor(0xFF, 0xFF, 0xFF)      # 白（列1文字）
COLOR_DESC_BG = RGBColor(0xFF, 0xFF, 0xFF)       # 白（列2）
COLOR_EXAMPLES_BG = RGBColor(0xFA, 0xFA, 0xFA)   # 薄グレー（列3）

# ── フォント ──
FONT_NAME_JP = "Meiryo UI"
FONT_NAME_LATIN = "Arial"

# ── KBF 数（v0.2 で可変化: 既定 3、範囲 2〜5）──
KBF_COUNT_DEFAULT = 3
KBF_COUNT_MIN = 2
KBF_COUNT_MAX = 5


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def set_textbox_text(shape, text):
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def _add_run(p_elem, text, *, font_size, bold, text_rgb):
    """段落要素に1つのrunを追加する（フォント・色・太字を細かく制御）。"""
    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP",
        "sz": str(font_size * 100),
        "b": "1" if bold else "0",
    })
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{text_rgb[0]:02X}{text_rgb[1]:02X}{text_rgb[2]:02X}")
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", FONT_NAME_LATIN)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT_NAME_JP)
    cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", FONT_NAME_JP)
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text


def _reset_cell(cell, *, bg_color, v_align):
    if bg_color is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    if v_align == "middle":
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "bottom":
        cell.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        cell.vertical_anchor = MSO_ANCHOR.TOP
    cell.margin_left = Inches(0.10)
    cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.08)
    cell.margin_bottom = Inches(0.08)
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    return txBody


def apply_simple_cell(cell, text, *, font_size, bold,
                      bg_color=None, text_color=None,
                      align="left", v_align="top"):
    """単一テキスト（または行リスト）でセルを埋めるシンプルなセル設定。"""
    txBody = _reset_cell(cell, bg_color=bg_color, v_align=v_align)
    cell.text_frame.word_wrap = True

    if isinstance(text, list):
        lines = [str(line) for line in text]
    else:
        lines = [str(text)] if text else [""]

    align_map = {"left": "l", "center": "ctr", "right": "r"}
    algn_code = align_map.get(align, "l")

    text_rgb = text_color if text_color is not None else COLOR_TEXT

    for line in lines:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", algn_code)
        _add_run(p_elem, line, font_size=font_size, bold=bold, text_rgb=text_rgb)


def apply_player_examples_cell(cell, player_examples, *, font_size,
                               bg_color=None, text_color=None):
    """プレイヤー例セルを構築する。

    各 example について `{player}: {example}` を1段落とし、
    `{player}` のみ太字、コロン以降は通常体で出す（同一段落内2run）。
    """
    txBody = _reset_cell(cell, bg_color=bg_color, v_align="top")
    cell.text_frame.word_wrap = True

    text_rgb = text_color if text_color is not None else COLOR_TEXT

    if not player_examples:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", "l")
        _add_run(p_elem, "", font_size=font_size, bold=False, text_rgb=text_rgb)
        return

    for ex in player_examples:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", "l")
        # 段落間隔を少し広めに（行間でプレイヤーが識別しやすくなる）
        pPr.set("marB", "60000")  # EMU - 約6ptの段落後スペース

        player = ex.get("player", "")
        example_text = ex.get("example", "")
        # プレイヤー名（太字）
        _add_run(p_elem, f"{player}: ", font_size=font_size, bold=True, text_rgb=text_rgb)
        # 例文（通常体）
        _add_run(p_elem, example_text, font_size=font_size, bold=False, text_rgb=text_rgb)


MAIN_MESSAGE_MAX = 65
KBF_NAME_MAX = 15
KBF_DESCRIPTION_MAX = 120
PLAYER_EXAMPLE_MAX = 80


def validate_input(data):
    """入力JSONのバリデーション。KBF=KBF_COUNT_MIN〜KBF_COUNT_MAX（v0.2 可変化）、player_examples=1〜5、文字数制約。"""
    main_message = data.get("main_message", "")
    if len(main_message) > MAIN_MESSAGE_MAX:
        raise ValueError(
            f"main_message は {MAIN_MESSAGE_MAX} 字以内である必要があります "
            f"（受領: {len(main_message)} 字）: {main_message[:80]}..."
        )
    kbf_list = data.get("kbf_list")
    if not isinstance(kbf_list, list):
        raise ValueError("kbf_list は配列である必要があります")
    if not (KBF_COUNT_MIN <= len(kbf_list) <= KBF_COUNT_MAX):
        raise ValueError(
            f"kbf_list の要素数は {KBF_COUNT_MIN}〜{KBF_COUNT_MAX} の範囲である必要があります"
            f"（受領: {len(kbf_list)}、既定: {KBF_COUNT_DEFAULT}）"
        )
    for i, kbf in enumerate(kbf_list):
        name = kbf.get("name", "")
        if not name:
            raise ValueError(f"kbf_list[{i}].name が空です")
        if len(name) > KBF_NAME_MAX:
            raise ValueError(
                f"kbf_list[{i}].name は {KBF_NAME_MAX} 字以内（受領: {len(name)}）: {name}"
            )
        description = kbf.get("description", "")
        if not description:
            raise ValueError(f"kbf_list[{i}].description が空です")
        if len(description) > KBF_DESCRIPTION_MAX:
            raise ValueError(
                f"kbf_list[{i}].description は {KBF_DESCRIPTION_MAX} 字以内"
                f"（受領: {len(description)}）: {description[:80]}..."
            )
        examples = kbf.get("player_examples", [])
        if not isinstance(examples, list) or len(examples) < 1:
            raise ValueError(f"kbf_list[{i}].player_examples は1要素以上必要です")
        if len(examples) > 5:
            raise ValueError(
                f"kbf_list[{i}].player_examples は最大5要素まで（受領: {len(examples)}）"
            )
        for j, ex in enumerate(examples):
            example_text = ex.get("example", "")
            if len(example_text) > PLAYER_EXAMPLE_MAX:
                raise ValueError(
                    f"kbf_list[{i}].player_examples[{j}].example は {PLAYER_EXAMPLE_MAX} 字以内"
                    f"（受領: {len(example_text)}）: {example_text[:80]}..."
                )


def build_table(slide, data):
    """3列×(1+kbf_count)行のKBFテーブルを構築（v0.2 で kbf_count 可変化）。"""
    kbf_list = data["kbf_list"]
    kbf_count = len(kbf_list)
    n_rows = 1 + kbf_count  # ヘッダー1 + KBF×kbf_count
    n_cols = 3

    content_shape = find_shape(slide, SHAPE_CONTENT_AREA)
    if content_shape is not None:
        tbl_left = content_shape.left
        tbl_top = content_shape.top
        tbl_width = content_shape.width
        tbl_height = content_shape.height
        sp_tree = slide.shapes._spTree
        sp_tree.remove(content_shape._element)
    else:
        tbl_left = CONTENT_LEFT
        tbl_top = CONTENT_TOP
        tbl_width = CONTENT_WIDTH
        tbl_height = CONTENT_HEIGHT

    # 列幅
    col_kbf_w = int(tbl_width * COL_RATIO_KBF)
    col_desc_w = int(tbl_width * COL_RATIO_DESC)
    col_examples_w = tbl_width - col_kbf_w - col_desc_w

    # 行高: ヘッダー固定、データ行は等分
    header_row_h = Inches(0.40)
    available_h = tbl_height - header_row_h
    data_row_h = int(available_h / kbf_count)

    print(f"  ✓ テーブル: {n_rows}行 × {n_cols}列")
    print(f"    列幅: KBF={col_kbf_w/914400:.2f}in / 詳細={col_desc_w/914400:.2f}in / 例={col_examples_w/914400:.2f}in")
    print(f"    データ行高: {data_row_h/914400:.2f}in")

    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table_shape.name = "MarketKBFTable"
    table = table_shape.table

    # 列幅設定
    table.columns[0].width = col_kbf_w
    table.columns[1].width = col_desc_w
    table.columns[2].width = col_examples_w

    # 行高設定（XML）
    tbl_xml = table_shape._element.find('.//' + qn('a:tbl'))
    for i, tr in enumerate(tbl_xml.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(header_row_h))
        else:
            tr.set('h', str(data_row_h))

    # tblPrクリア（デフォルトスタイル除去）
    tbl_elem = table_shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # ── ヘッダー行 ──（v0.1 R5: 12→14pt）
    apply_simple_cell(
        table.cell(0, 0), "KBF",
        font_size=14, bold=True,
        bg_color=COLOR_HEADER_BG,
        align="center", v_align="middle",
    )
    apply_simple_cell(
        table.cell(0, 1), "詳細",
        font_size=14, bold=True,
        bg_color=COLOR_HEADER_BG,
        align="center", v_align="middle",
    )
    apply_simple_cell(
        table.cell(0, 2), "プレイヤーの抑え方の例",
        font_size=14, bold=True,
        bg_color=COLOR_HEADER_BG,
        align="center", v_align="middle",
    )

    # ── データ行（KBF×kbf_count、v0.2 で可変化）──（v0.1 R5: KBF名 14→16, 詳細 11→13, プレイヤー例 10→12）
    for r_idx, kbf in enumerate(kbf_list):
        row = r_idx + 1

        # 列1: KBF（紺背景＋白文字、中央寄せ・縦中央）
        apply_simple_cell(
            table.cell(row, 0), kbf.get("name", ""),
            font_size=16, bold=True,
            bg_color=COLOR_KBF_BG,
            text_color=COLOR_KBF_TEXT,
            align="center", v_align="middle",
        )

        # 列2: 詳細（白背景、左寄せ・縦中央）
        apply_simple_cell(
            table.cell(row, 1), kbf.get("description", ""),
            font_size=13, bold=False,
            bg_color=COLOR_DESC_BG,
            align="left", v_align="middle",
        )

        # 列3: プレイヤー例（薄グレー、プレイヤー名のみ太字）
        apply_player_examples_cell(
            table.cell(row, 2),
            kbf.get("player_examples", []),
            font_size=12,
            bg_color=COLOR_EXAMPLES_BG,
        )

        print(f"    KBF{row}: {kbf.get('name', '')}（例 {len(kbf.get('player_examples', []))} 件）")

    print("  ✓ テーブル生成完了")


def main():
    parser = argparse.ArgumentParser(description="市場KBF PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True, help="JSONデータファイルのパス")
    parser.add_argument("--template", required=True, help="PPTXテンプレートのパス")
    parser.add_argument("--output", required=True, help="出力PPTXファイルのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    validate_input(data)
    print(f"  データ読み込み: KBF={len(data['kbf_list'])}件 / market={data.get('chart_title', 'N/A')}")

    prs = Presentation(args.template)
    slide = prs.slides[0]

    # 1. メインメッセージ
    main_message = data.get("main_message", "")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_message)
    print(f"  ✓ Main Message: {main_message[:40]}...")

    # 2. チャートタイトル
    chart_title = data.get("chart_title", "Key Business Factor")
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
    print(f"  ✓ Chart Title: {chart_title}")

    # 3. テーブル構築
    build_table(slide, data)

    # 4. 出典
    source_text = data.get("source", "出典：各社IR、業界レポート")
    source_shape = find_shape(slide, SHAPE_SOURCE)
    if source_shape is not None:
        set_textbox_text(source_shape, source_text)
        print(f"  ✓ Source: {source_text[:40]}...")

    # 5. 出力
    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
