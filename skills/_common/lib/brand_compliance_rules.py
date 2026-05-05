"""Brand compliance rules for fill_*.py outputs.

Each rule is a callable `(prs, ctx) -> Iterable[CheckResult]` that yields one
or more CheckResult records. Rules are grouped into profiles selected by
`(skill_id, brand)` from `tools/check_brand_compliance.py`.

Profiles
--------
Currently only `pilot3 × roleup` profiles are populated. `pilot3 × stellar_aiz`
profiles are skeletons (empty list with a TODO marker). When ISSUE-010 expands
brand-aware support to more skills, add new profiles here without touching the
CLI.

Coverage of pilot3_roleup profile
----------------------------------
  C1  ガイド矩形 (`正方形/長方形 1/8`, `Content Area`) が残存していない
  C2  スライドタイトル placeholder (`Title 1`) が theme.font_size_title_pt
  C4  本文・表テキストの font.size が roleup 許容集合 {22, 14, 12, 10, 6} のみ
  C5  出典 placeholder (`Source 3`) のフォントが theme.font_size_source_pt
  C6  出典 placeholder のテキストが空でない (theme.layout_rules.source_required)
  C7  slide_size が theme.slide_size と一致 (±tolerance)
  C8  全 run の font.name (latin/ea) が theme.fonts.ea と一致 (None は許容)
  C10 chart 軸 (`c:catAx`/`c:valAx`) の defRPr/@sz が theme.font_size_body_pt
  C11 テキスト含有 shape 同士の bbox 重なり (chart shape は除外)
  C12 chart object の has_legend == False (二重凡例防止)
"""
from __future__ import annotations

import io
import json
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Iterable, Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"a": NS_A, "c": NS_C, "p": NS_P}

REPO_ROOT = Path(__file__).resolve().parents[3]

# Slide size tolerance (EMU). 5000 EMU ≈ 0.0055 inch.
SLIDE_SIZE_TOL_EMU = 5000

# Bbox overlap threshold: ratio = overlap_area / min(area_a, area_b)
BBOX_OVERLAP_THRESHOLD = 0.10

# Allowed font sizes (pt) for roleup body/text content. Tuned to match
# user spec 2026-05-04: 22 (title) / 14 (key message) / 12 (subtitle) /
# 10 (body+table) / 6 (source).
ROLEUP_ALLOWED_FONT_SIZES_PT = {22, 14, 12, 10, 6}


# ─────────────────────────────────────────
# Result types
# ─────────────────────────────────────────


@dataclass
class CheckResult:
    rule_id: str
    passed: bool
    message: str
    details: Optional[dict] = None
    severity: str = "error"  # 'error' | 'warning'


@dataclass
class CheckContext:
    pptx_path: str
    skill_id: str
    brand: str
    theme: dict  # loaded theme.json content


# ─────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────


def load_theme(brand: str) -> dict:
    path = REPO_ROOT / "skills" / "_common" / "brands" / brand / "theme.json"
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def _find_title_placeholder(prs):
    """Locate the title placeholder.

    Priority:
      1. ph type="title" (most reliable; survives LibreOffice roundtrip)
      2. shape name in ("Title 1", "PlaceHolder 1") fallback
    Returns (slide_idx, shape) or (None, None).
    """
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            ph = shape._element.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is not None and ph.get("type") == "title":
                return s_idx, shape
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in ("Title 1", "PlaceHolder 1"):
                return s_idx, shape
    return None, None


def _find_source_placeholder(prs):
    """Locate the source placeholder.

    The source placeholder has no specific ph type, so we rely on shape name.
    Roleup template uses 'Source 3'; LibreOffice roundtrip may rename it to
    'PlaceHolder 3'. Stella uses 'Source'.
    Returns (slide_idx, shape) or (None, None).
    """
    candidates = ("Source 3", "Source", "PlaceHolder 3")
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in candidates:
                return s_idx, shape
    return None, None


def _placeholder_size_pt(shape):
    """Return the effective font size (pt) for the first run-bearing paragraph,
    falling back through run rPr → paragraph defRPr → None.

    Roleup placeholders typically rely on template-level defaultStyle which
    means run-level rPr/@sz is absent. Reading the paragraph defRPr lets us
    distinguish "explicitly set incorrect size" from "inherits template size
    (correct by construction)".
    """
    for para in shape.text_frame.paragraphs:
        # 1. run rPr/@sz の最初に見つかった非 None
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt
        # 2. paragraph defRPr/@sz
        p_elem = para._p
        defRPr = p_elem.find(".//a:pPr/a:defRPr", NS)
        if defRPr is not None and defRPr.get("sz"):
            return int(defRPr.get("sz")) / 100.0
    return None


def _iter_all_runs(prs):
    """Yield (slide_idx, shape, paragraph, run) tuples for all text shapes."""
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    yield s_idx, shape, para, run


def _iter_table_runs(prs):
    """Yield (slide_idx, table_shape, row, col, run) for all table cells."""
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            yield s_idx, shape, r_idx, c_idx, run


def _bbox_intersection_area(a, b):
    al, at, aw, ah = a
    bl, bt, bw, bh = b
    ar, ab = al + aw, at + ah
    br_, bb = bl + bw, bt + bh
    ox = max(0, min(ar, br_) - max(al, bl))
    oy = max(0, min(ab, bb) - max(at, bt))
    return ox * oy


def _shape_bbox(shape):
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    return (int(shape.left), int(shape.top), int(shape.width), int(shape.height))


def _read_chart_xmls(pptx_path) -> dict:
    """Return {chart_part_name: parsed_xml_root} for all charts in the pptx."""
    out = {}
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if name.startswith("ppt/charts/chart") and name.endswith(".xml"):
                out[name] = etree.fromstring(z.read(name))
    return out


# ─────────────────────────────────────────
# Rules
# ─────────────────────────────────────────


def rule_no_guide_rectangles(prs, ctx) -> Iterable[CheckResult]:
    """C1: roleup 公式テンプレ由来のガイド矩形が出力に残っていない。"""
    forbidden = {"正方形/長方形 1", "正方形/長方形 8", "Content Area"}
    found = []
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.name in forbidden:
                found.append(f"slide{s_idx+1}/{shape.name}")
    if found:
        yield CheckResult("C1_no_guide_rect", False,
                          f"ガイド矩形が残存: {found}",
                          {"shapes": found})
    else:
        yield CheckResult("C1_no_guide_rect", True, "ガイド矩形 0 件")


def rule_title_size(prs, ctx) -> Iterable[CheckResult]:
    """C2: スライドタイトル placeholder のフォントサイズ。

    run rPr/@sz と paragraph defRPr/@sz の両方を探索。両方とも未指定なら
    template の master/layout defaultStyle を継承する設計のため PASS とする。
    """
    expected = ctx.theme["defaults"]["font_size_title_pt"]
    s_idx, shape = _find_title_placeholder(prs)
    if shape is None:
        yield CheckResult("C2_title_size", False,
                          "スライドタイトル placeholder 未検出 (title type or 'Title 1'/'PlaceHolder 1')")
        return
    actual = _placeholder_size_pt(shape)
    if actual is None:
        yield CheckResult("C2_title_size", True,
                          f"スライドタイトル size はテンプレ default 継承 (expected {expected}pt)")
        return
    if actual != expected:
        yield CheckResult("C2_title_size", False,
                          f"スライドタイトル {actual}pt (expected {expected}pt)",
                          {"slide": s_idx + 1, "shape": shape.name,
                           "actual_pt": actual, "expected_pt": expected})
    else:
        yield CheckResult("C2_title_size", True, f"スライドタイトル {expected}pt 確認")


def rule_allowed_font_sizes(prs, ctx) -> Iterable[CheckResult]:
    """C4: 全 run / 全テーブルセルのフォントサイズが roleup 許容集合のみ。"""
    allowed = ROLEUP_ALLOWED_FONT_SIZES_PT
    bad = []
    for s_idx, shape, para, run in _iter_all_runs(prs):
        if run.font.size is None:
            continue
        sz = run.font.size.pt
        if sz not in allowed:
            txt = run.text.strip()[:30]
            bad.append({"slide": s_idx + 1, "shape": shape.name,
                        "size_pt": sz, "text_sample": txt})
    for s_idx, shape, r_idx, c_idx, run in _iter_table_runs(prs):
        if run.font.size is None:
            continue
        sz = run.font.size.pt
        if sz not in allowed:
            txt = run.text.strip()[:30]
            bad.append({"slide": s_idx + 1, "shape": shape.name,
                        "row": r_idx, "col": c_idx,
                        "size_pt": sz, "text_sample": txt})
    if bad:
        yield CheckResult("C4_allowed_font_sizes", False,
                          f"許容外フォントサイズ {len(bad)} 件 (許容: {sorted(allowed, reverse=True)}pt)",
                          {"violations": bad[:20]})
    else:
        yield CheckResult("C4_allowed_font_sizes", True,
                          f"フォントサイズ全 run が許容集合内 {sorted(allowed, reverse=True)}pt")


def rule_source_size(prs, ctx) -> Iterable[CheckResult]:
    """C5: 出典 placeholder のフォントサイズ。テンプレ defaultStyle 継承時は PASS。"""
    expected = ctx.theme["defaults"]["font_size_source_pt"]
    s_idx, shape = _find_source_placeholder(prs)
    if shape is None:
        yield CheckResult("C5_source_size", False,
                          "出典 placeholder 未検出 ('Source 3'/'Source'/'PlaceHolder 3')")
        return
    actual = _placeholder_size_pt(shape)
    if actual is None:
        yield CheckResult("C5_source_size", True,
                          f"出典 size はテンプレ default 継承 (expected {expected}pt)")
        return
    if actual != expected:
        yield CheckResult("C5_source_size", False,
                          f"出典 {actual}pt (expected {expected}pt)",
                          {"slide": s_idx + 1, "shape": shape.name,
                           "actual_pt": actual, "expected_pt": expected})
    else:
        yield CheckResult("C5_source_size", True, f"出典 {expected}pt 確認")


def rule_source_required(prs, ctx) -> Iterable[CheckResult]:
    """C6: theme.layout_rules.source_required = true の場合、出典テキストが空でない。"""
    required = ctx.theme.get("layout_rules", {}).get("source_required", False)
    if not required:
        yield CheckResult("C6_source_required", True,
                          "source_required=False のため検査 skip", severity="warning")
        return
    s_idx, shape = _find_source_placeholder(prs)
    if shape is None:
        yield CheckResult("C6_source_required", False,
                          "出典 placeholder 不在 ('Source 3'/'Source'/'PlaceHolder 3')")
        return
    txt = shape.text_frame.text.strip()
    if not txt:
        yield CheckResult("C6_source_required", False,
                          f"slide{s_idx+1} の出典 placeholder '{shape.name}' のテキストが空",
                          {"slide": s_idx + 1, "shape": shape.name})
    else:
        yield CheckResult("C6_source_required", True,
                          f"出典テキスト有り: {txt[:40]!r}")


def rule_slide_size(prs, ctx) -> Iterable[CheckResult]:
    """C7: slide_size が theme.slide_size と ±tolerance 内で一致。"""
    expected_w_in = ctx.theme["slide_size"]["width_in"]
    expected_h_in = ctx.theme["slide_size"]["height_in"]
    expected_w_emu = int(expected_w_in * 914400)
    expected_h_emu = int(expected_h_in * 914400)
    actual_w = prs.slide_width
    actual_h = prs.slide_height
    if abs(actual_w - expected_w_emu) > SLIDE_SIZE_TOL_EMU or abs(actual_h - expected_h_emu) > SLIDE_SIZE_TOL_EMU:
        yield CheckResult("C7_slide_size", False,
                          f"slide_size mismatch: actual={actual_w}×{actual_h} EMU, "
                          f"expected={expected_w_emu}×{expected_h_emu} EMU (±{SLIDE_SIZE_TOL_EMU})",
                          {"actual": [actual_w, actual_h],
                           "expected": [expected_w_emu, expected_h_emu]})
    else:
        yield CheckResult("C7_slide_size", True,
                          f"slide_size OK ({expected_w_in}×{expected_h_in} in)")


def rule_font_name(prs, ctx) -> Iterable[CheckResult]:
    """C8: 全 run の font.name が theme.fonts.ea と一致 (None は許容: テンプレ default 継承)。"""
    expected = ctx.theme["fonts"]["ea"]
    bad = []
    for s_idx, shape, para, run in _iter_all_runs(prs):
        actual = run.font.name
        if actual is not None and actual != expected:
            txt = run.text.strip()[:30]
            bad.append({"slide": s_idx + 1, "shape": shape.name,
                        "actual": actual, "expected": expected,
                        "text_sample": txt})
    if bad:
        yield CheckResult("C8_font_name", False,
                          f"フォント名 mismatch {len(bad)} 件 (expected: '{expected}')",
                          {"violations": bad[:20]})
    else:
        yield CheckResult("C8_font_name", True, f"フォント名 全 run 'expected' or None")


def rule_chart_axis_font_size(prs, ctx) -> Iterable[CheckResult]:
    """C10: chart 軸 (catAx/valAx) の defRPr/@sz が theme.font_size_body_pt × 100。"""
    expected_pt = ctx.theme["defaults"]["font_size_body_pt"]
    expected_sz = str(expected_pt * 100)
    bad = []
    for chart_part, root in _read_chart_xmls(ctx.pptx_path).items():
        for ax_tag in ("c:catAx", "c:valAx"):
            for ax in root.findall(f".//{ax_tag}", NS):
                # 削除済み軸はスキップ
                de = ax.find("c:delete", NS)
                if de is not None and de.get("val") == "1":
                    continue
                drs = ax.findall(".//a:defRPr", NS)
                for dr in drs:
                    sz = dr.get("sz")
                    if sz is None:
                        continue
                    if sz != expected_sz:
                        bad.append({"chart_part": chart_part, "axis": ax_tag,
                                    "actual_sz": sz, "expected_sz": expected_sz})
    if bad:
        yield CheckResult("C10_chart_axis_font", False,
                          f"chart 軸フォントサイズ mismatch {len(bad)} 件 (expected sz='{expected_sz}' = {expected_pt}pt)",
                          {"violations": bad[:20]})
    else:
        yield CheckResult("C10_chart_axis_font", True,
                          f"chart 軸フォント全て {expected_pt}pt")


def rule_textbox_bbox_overlap(prs, ctx) -> Iterable[CheckResult]:
    """C11 (NEW A): テキスト含有 shape 同士の bbox 重なり検査。

    例外: chart shape は内部に shape を意図的にレイヤーする設計のため、
    chart 自身および chart bbox 内の shape はペア比較から除外する。
    閾値: 重なり面積 / min(面積_a, 面積_b) > BBOX_OVERLAP_THRESHOLD で fail。
    """
    overlaps = []
    for s_idx, slide in enumerate(prs.slides):
        # chart bbox 群を集める (除外ゾーン)
        chart_bboxes = []
        for shape in slide.shapes:
            if shape.has_chart:
                bb = _shape_bbox(shape)
                if bb is not None:
                    chart_bboxes.append(bb)

        # text 含有 shape を収集 (chart 自身は除外)
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_chart:
                continue
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            bb = _shape_bbox(shape)
            if bb is None:
                continue
            # chart bbox 内に「完全包含」される shape はチャート内意図的配置とみなして除外
            inside_chart = False
            for cbb in chart_bboxes:
                cl, ct, cw, ch = cbb
                sl, st, sw, sh = bb
                if sl >= cl and st >= ct and sl + sw <= cl + cw and st + sh <= ct + ch:
                    inside_chart = True
                    break
            if inside_chart:
                continue
            text_shapes.append((shape, bb, txt))

        # ペアワイズ比較
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                sa, ba, ta = text_shapes[i]
                sb, bb, tb = text_shapes[j]
                inter = _bbox_intersection_area(ba, bb)
                if inter <= 0:
                    continue
                area_a = ba[2] * ba[3]
                area_b = bb[2] * bb[3]
                min_area = min(area_a, area_b)
                if min_area == 0:
                    continue
                ratio = inter / min_area
                if ratio > BBOX_OVERLAP_THRESHOLD:
                    overlaps.append({
                        "slide": s_idx + 1,
                        "shape_a": sa.name, "text_a": ta[:25],
                        "shape_b": sb.name, "text_b": tb[:25],
                        "overlap_ratio": round(ratio, 3),
                    })
    if overlaps:
        yield CheckResult("C11_textbox_overlap", False,
                          f"テキスト shape 同士の bbox 重なり {len(overlaps)} 件 "
                          f"(閾値: 重なり面積/小面積 > {BBOX_OVERLAP_THRESHOLD})",
                          {"overlaps": overlaps[:20]})
    else:
        yield CheckResult("C11_textbox_overlap", True, "テキスト shape 重なりなし")


def rule_no_chart_builtin_legend(prs, ctx) -> Iterable[CheckResult]:
    """C12 (NEW B): chart object に builtin legend が無いこと (二重凡例防止)。

    pilot3 の roleup ではカスタム凡例を別 textbox として配置する設計のため、
    chart 自身の has_legend は False でなければならない。
    """
    bad = []
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            if shape.chart.has_legend:
                bad.append({"slide": s_idx + 1, "shape": shape.name})
    if bad:
        yield CheckResult("C12_no_chart_legend", False,
                          f"chart に builtin legend が残存 {len(bad)} 件 "
                          f"(カスタム凡例を別途配置する設計のため has_legend=False が必須)",
                          {"charts": bad})
    else:
        yield CheckResult("C12_no_chart_legend", True,
                          "chart builtin legend 全て無効")


# ─────────────────────────────────────────
# Profiles
# ─────────────────────────────────────────

# 共通 roleup ルール (pilot 3 すべてで同じ)
_COMMON_ROLEUP_RULES: list[Callable] = [
    rule_no_guide_rectangles,
    rule_title_size,
    rule_allowed_font_sizes,
    rule_source_size,
    rule_source_required,
    rule_slide_size,
    rule_font_name,
    rule_chart_axis_font_size,
    rule_textbox_bbox_overlap,
    rule_no_chart_builtin_legend,
]

# Profile registry: (skill_id, brand) -> list[rule]
PROFILES: dict[tuple, list[Callable]] = {
    ("customer-profile-pptx", "roleup"): list(_COMMON_ROLEUP_RULES),
    ("market-environment-pptx", "roleup"): list(_COMMON_ROLEUP_RULES),
    ("market-share-pptx", "roleup"): list(_COMMON_ROLEUP_RULES),
    ("positioning-map-pptx", "roleup"): [
        # positioning-map は手描き (OVAL/CONNECTOR) でチャート不在のため C10/C12 を除外
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("competitor-summary-pptx", "roleup"): [
        # competitor-summary はテーブルのみでチャート不在のため C10/C12 を除外
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("market-kbf-pptx", "roleup"): [
        # market-kbf もテーブルのみでチャート不在のため C10/C12 を除外
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("pest-analysis-pptx", "roleup"): [
        # pest-analysis は手描き shape (Rectangle + textbox) でチャート不在のため C10/C12 を除外
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("company-history-pptx", "roleup"): [
        # ch にはチャートが無いため C10/C12 を除外
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("executive-summary-pptx", "roleup"): [
        # exec-summary はチャートを持たない (5 findings 縦積みの textbox + 縦バーシェイプ)
        # ため C10/C12 を除外。残り 8 ルール (C1/C2/C4/C5/C6/C7/C8/C11) を適用。
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("section-divider-pptx", "roleup"): [
        # 中扉は装飾スライド: 巨大数字 (180pt) と SECTION ラベル textbox が
        # 完全包含で C11 を必ず fail させる + Title 1 / Source 3 placeholder を削除する設計のため
        # C2/C5/C6 適用不可、本文サイズも装飾 (180/22pt) のため C4 不適用。
        # 適用するのは C1 (ガイド矩形除去) / C7 (slide size) / C8 (font name) のみ。
        rule_no_guide_rectangles, rule_slide_size, rule_font_name,
    ],
    ("table-of-contents-pptx", "roleup"): [
        # TOC はチャート不在 (C10/C12 除外) かつ出典任意 (C6 除外)。
        # 残り 7 ルール (C1/C2/C4/C5/C7/C8/C11) を適用。
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend, rule_source_required)
    ],
    ("data-availability-pptx", "roleup"): [
        # data-availability はテーブル+textbox 構成、チャート不在のため C10/C12 を除外。
        # 残り 8 ルール (C1/C2/C4/C5/C6/C7/C8/C11) を適用。
        r for r in _COMMON_ROLEUP_RULES
        if r not in (rule_chart_axis_font_size, rule_no_chart_builtin_legend)
    ],
    ("revenue-analysis-pptx", "roleup"): list(_COMMON_ROLEUP_RULES),
    # stella 版は ISSUE-010 で stella 仕様確定後に追加 (現状 skeleton)
    ("customer-profile-pptx", "stellar_aiz"): [],   # TODO: ISSUE-010
    ("market-environment-pptx", "stellar_aiz"): [], # TODO: ISSUE-010
    ("company-history-pptx", "stellar_aiz"): [],    # TODO: ISSUE-010
    ("executive-summary-pptx", "stellar_aiz"): [],  # TODO: ISSUE-010
    ("market-share-pptx", "stellar_aiz"): [],       # TODO: ISSUE-010
    ("positioning-map-pptx", "stellar_aiz"): [],    # TODO: ISSUE-010
    ("competitor-summary-pptx", "stellar_aiz"): [], # TODO: ISSUE-010
    ("market-kbf-pptx", "stellar_aiz"): [],         # TODO: ISSUE-010
    ("pest-analysis-pptx", "stellar_aiz"): [],      # TODO: ISSUE-010
    ("section-divider-pptx", "stellar_aiz"): [],    # TODO: ISSUE-010
    ("table-of-contents-pptx", "stellar_aiz"): [],  # TODO: ISSUE-010
    ("data-availability-pptx", "stellar_aiz"): [],  # TODO: ISSUE-010
    ("revenue-analysis-pptx", "stellar_aiz"): [],   # TODO: ISSUE-010
}


def get_profile(skill_id: str, brand: str) -> list[Callable]:
    key = (skill_id, brand)
    if key not in PROFILES:
        raise KeyError(
            f"profile not found for skill={skill_id!r} brand={brand!r}. "
            f"Available: {sorted(PROFILES.keys())}"
        )
    return PROFILES[key]


def run_profile(prs, ctx: CheckContext) -> list[CheckResult]:
    rules = get_profile(ctx.skill_id, ctx.brand)
    if not rules:
        return [CheckResult(
            "PROFILE_EMPTY", False,
            f"profile ({ctx.skill_id}, {ctx.brand}) は未実装 (skeleton)。"
            f"ISSUE-010 で対応予定。",
            severity="warning",
        )]
    results = []
    for rule in rules:
        results.extend(rule(prs, ctx))
    return results
