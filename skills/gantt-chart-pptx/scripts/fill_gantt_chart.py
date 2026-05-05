"""
fill_gantt_chart.py — ガントチャートをHTML描画→スクリーンショット→PPTXに挿入
"""

import argparse

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402
import asyncio
import json
import os
import sys
import copy
import math
import tempfile
from datetime import datetime, timedelta
from html import escape

from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass


SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE  = "Text Placeholder 2"

CHART_LEFT   = 370800
CHART_TOP    = 1350000
CHART_WIDTH  = 11450400
CHART_MAX_H  = 5200000

VP_WIDTH  = 2200
DEVICE_SCALE = 2

MAX_ROWS_PER_PAGE = 14
FALLBACK_ACCENT2 = "#1A3C6E"


def parse_date(s):
    return datetime.strptime(s.replace("/", "-"), "%Y-%m-%d")


def extract_accent2(template_path):
    try:
        prs = Presentation(template_path)
        for rel in prs.slide_masters[0].part.rels.values():
            if "theme" in rel.reltype:
                theme_elem = etree.fromstring(rel.target_part.blob)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                cs = theme_elem.find(".//a:clrScheme", ns)
                if cs is not None:
                    a2 = cs.find("a:accent2", ns)
                    if a2 is not None:
                        sr = a2.find("a:srgbClr", ns)
                        if sr is not None:
                            return f"#{sr.get('val')}"
    except Exception as e:
        print(f"  WARNING: {e}", file=sys.stderr)
    return FALLBACK_ACCENT2


def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def lighten(c, f):
    r, g, b = hex_to_rgb(c)
    return f"#{int(r+(255-r)*f):02x}{int(g+(255-g)*f):02x}{int(b+(255-b)*f):02x}"


def month_ticks(s, e):
    t = []
    d = s.replace(day=1)
    while d <= e:
        if d >= s: t.append(d)
        d = d.replace(year=d.year+1, month=1) if d.month == 12 else d.replace(month=d.month+1)
    return t


def week_ticks(s, e):
    t = []
    d = s - timedelta(days=s.weekday())
    while d <= e:
        if d >= s: t.append(d)
        d += timedelta(days=7)
    return t


def auto_unit(s, e):
    return "week" if (e - s).days <= 90 else "month"


def generate_html(data, task_rows, accent2):
    ss = parse_date(data["schedule_start"]) - timedelta(days=3)
    se = parse_date(data["schedule_end"]) + timedelta(days=7)
    td = max(1, (se - ss).days)

    milestones = data.get("milestones", [])
    tu = data.get("time_unit", auto_unit(ss, se))

    bar_full  = accent2
    bar_light = lighten(accent2, 0.55)
    hdr_bg    = accent2

    row_h    = 58
    hdr_h    = 68
    ms_row_h = 72       # マイルストーン行を高くしてひし形とテキストの余裕を確保
    bar_h    = 28
    bar_top  = 15

    top_off = hdr_h + 8
    has_ms = len(milestones) > 0
    if has_ms:
        top_off += ms_row_h

    total_h = top_off + len(task_rows) * row_h + 20
    lp = 20
    cp = 80

    ticks = week_ticks(ss, se) if tu == "week" else month_ticks(ss, se)

    html = f"""<!DOCTYPE html>
<html lang="ja"><head><meta charset="UTF-8">
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{
  font-family:'Meiryo','Meiryo UI','Yu Gothic','Noto Sans CJK JP',sans-serif;
  background:white;width:{VP_WIDTH}px;height:{total_h}px;overflow:hidden;
}}
.gc{{position:relative;width:100%;height:100%}}

.th{{
  position:absolute;top:0;height:{hdr_h}px;
  display:flex;align-items:flex-end;padding-bottom:6px;
  font-size:26px;font-weight:600;color:#333;
  border-left:1px solid #bbb;padding-left:8px;
}}
.tg{{position:absolute;top:{hdr_h}px;bottom:0;width:1px;background:#e8e8e8}}

.msr{{
  position:absolute;top:{hdr_h+4}px;height:{ms_row_h}px;
  width:100%;display:flex;align-items:center;
  background:#f8f8f8;border-bottom:1px solid #ddd;
}}
.msr .rl{{
  width:{lp}%;padding:0 8px 0 16px;
  font-size:26px;font-weight:700;color:#C62828;flex-shrink:0;
}}
.msr .rc{{width:{cp}%;position:relative;height:100%;flex-shrink:0}}
.msd{{
  position:absolute;top:58%;width:16px;height:16px;
  background:#C62828;transform:translate(-50%,-50%) rotate(45deg);
  border-radius:2px;z-index:5;
}}
.msl{{
  position:absolute;top:4px;font-size:22px;color:#C62828;
  font-weight:600;white-space:nowrap;
  transform:translateX(-50%);
}}

.gr{{position:absolute;height:{row_h}px;width:100%;display:flex;align-items:center}}
.gre{{background:#fafbfc}}.gro{{background:#fff}}

.rl{{
  width:{lp}%;padding:0 10px 0 18px;
  font-size:26px;color:#1a1a1a;
  white-space:nowrap;overflow:hidden;text-overflow:ellipsis;flex-shrink:0;
}}
.rl .ow{{color:#888;font-size:24px}}

.ph .rl{{font-weight:700;font-size:26px;color:white;padding-left:18px}}

.rc{{width:{cp}%;position:relative;height:100%;flex-shrink:0}}
.tb{{position:absolute;top:{bar_top}px;height:{bar_h}px;border-radius:3px;min-width:4px}}

.ls{{position:absolute;left:{lp}%;top:0;bottom:0;width:1px;background:#ccc;z-index:3}}
.hbl{{position:absolute;left:0;right:0;top:{hdr_h}px;height:2px;background:#333;z-index:3}}
</style></head>
<body><div class="gc">
<div class="ls"></div><div class="hbl"></div>
"""

    for tick in ticks:
        pct = (tick - ss).days / td * cp + lp
        lbl = tick.strftime("%m/%d") if tu == "week" else \
              (tick.strftime("%Y/%m") if tick.month == 1 or tick == ticks[0] else tick.strftime("%m月"))
        html += f'<div class="th" style="left:{pct:.3f}%">{escape(lbl)}</div>\n'
        html += f'<div class="tg" style="left:{pct:.3f}%"></div>\n'

    if has_ms:
        html += f'<div class="msr"><div class="rl">マイルストーン</div><div class="rc">\n'
        for ms in milestones:
            md = parse_date(ms["date"])
            if ss <= md <= se:
                p = (md - ss).days / td * 100
                html += f'<div class="msl" style="left:{p:.3f}%">{escape(ms.get("name",""))}</div>\n'
                html += f'<div class="msd" style="left:{p:.3f}%"></div>\n'
        html += '</div></div>\n'

    for i, row in enumerate(task_rows):
        y = top_off + i * row_h
        ec = "gre" if i % 2 == 0 else "gro"
        if row["type"] == "phase_header":
            html += f'<div class="gr ph" style="top:{y}px;background:{hdr_bg}"><div class="rl">{escape(row["name"])}</div><div class="rc"></div></div>\n'
        else:
            ts = parse_date(row["start"]); te = parse_date(row["end"])
            prog = row.get("progress", 0); ow = row.get("owner", "")
            bl = max(0, (ts - ss).days / td * 100)
            bw = max(0.3, (te - ts).days / td * 100)
            pw = bw * prog / 100 if prog else 0
            lt = escape(row["name"])
            if ow: lt += f' <span class="ow">({escape(ow)})</span>'
            html += f'''<div class="gr {ec}" style="top:{y}px"><div class="rl">{lt}</div><div class="rc">
<div class="tb" style="left:{bl:.3f}%;width:{bw:.3f}%;background:{bar_light}"></div>
<div class="tb" style="left:{bl:.3f}%;width:{pw:.3f}%;background:{bar_full}"></div>
</div></div>\n'''

    html += "</div></body></html>"
    return html, VP_WIDTH, total_h


def build_task_rows(data):
    rows = []
    phases = data.get("phases", [])
    if phases:
        for ph in phases:
            rows.append({"type": "phase_header", "name": ph["name"]})
            for t in ph.get("tasks", []):
                r = dict(t); r["type"] = "task"; rows.append(r)
    else:
        for t in data.get("tasks", []):
            r = dict(t); r["type"] = "task"; rows.append(r)
    return rows


def split_rows_into_pages(all_rows, max_per_page):
    """フェーズ見出しが最後の行にならないようにページ分割する"""
    pages = []
    current_page = []

    for i, row in enumerate(all_rows):
        current_page.append(row)

        # ページが満杯になったら切る
        if len(current_page) >= max_per_page:
            # 最後の行がphase_headerなら、それを次のページに繰り越す
            if current_page[-1]["type"] == "phase_header":
                carried = current_page.pop()
                pages.append(current_page)
                current_page = [carried]
            else:
                pages.append(current_page)
                current_page = []

    # 残りを追加
    if current_page:
        pages.append(current_page)

    return pages


async def take_screenshot(html, out, w, h):
    from playwright.async_api import async_playwright
    hp = os.path.join(tempfile.gettempdir(), "gantt.html")
    with open(hp, "w", encoding="utf-8") as f: f.write(html)
    async with async_playwright() as p:
        br = await p.chromium.launch()
        pg = await br.new_page(viewport={"width": w, "height": h}, device_scale_factor=DEVICE_SCALE)
        await pg.goto(f"file://{hp}")
        await pg.wait_for_timeout(500)
        await pg.screenshot(path=out, full_page=False)
        await br.close()
    print(f"  Screenshot: {out} ({os.path.getsize(out)} bytes)")
    os.unlink(hp)


def find_shape(sl, nm):
    for s in sl.shapes:
        if s.name == nm: return s
    return None


def set_placeholder_text(sh, txt):
    if sh is None: return
    p = sh.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = txt
        for r in p.runs[1:]: r.text = ""
    else:
        r = etree.SubElement(p._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t = etree.SubElement(r, qn("a:t")); t.text = txt


def dup_slide(prs, tmpl):
    ns = prs.slides.add_slide(tmpl.slide_layout)
    nsp = ns.shapes._spTree
    for c in list(nsp):
        if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr"): nsp.remove(c)
    for c in tmpl.shapes._spTree:
        if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr"): nsp.append(copy.deepcopy(c))
    return ns


def pop_slide(sl, mm, ct, sp, pn, tp, hw, hh):
    set_placeholder_text(find_shape(sl, SHAPE_MAIN_MESSAGE), mm)
    dt = ct if tp <= 1 else f"{ct}（{pn}/{tp}）"
    set_placeholder_text(find_shape(sl, SHAPE_CHART_TITLE), dt)
    aspect = hh / hw
    iw = CHART_WIDTH
    ih = int(iw * aspect)
    if ih > CHART_MAX_H: ih = CHART_MAX_H; iw = int(ih / aspect)
    sl.shapes.add_picture(sp, CHART_LEFT, CHART_TOP, iw, ih)
    print(f"  [Page {pn}] {iw}x{ih} EMU")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", required=True)
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f: data = json.load(f)

    mm = data.get("main_message", "").strip()
    ct = data.get("chart_title", "").strip()
    a2 = extract_accent2(args.template)
    print(f"  Accent2: {a2}")

    rows = build_task_rows(data)
    to = [r for r in rows if r["type"] == "task"]
    if "schedule_start" not in data and to: data["schedule_start"] = min(r["start"] for r in to)
    if "schedule_end" not in data and to: data["schedule_end"] = max(r["end"] for r in to)

    print(f"  Rows: {len(rows)}")

    # フェーズ見出し+タスクがセットで同じページに収まるよう分割
    pages = split_rows_into_pages(rows, MAX_ROWS_PER_PAGE)
    tp = len(pages)
    print(f"  Pages: {tp}")
    for pi, pg in enumerate(pages):
        ph_count = sum(1 for r in pg if r["type"] == "phase_header")
        tk_count = sum(1 for r in pg if r["type"] == "task")
        print(f"    Page {pi+1}: {len(pg)} rows ({ph_count} headers + {tk_count} tasks)")

    # 各ページのスクリーンショット
    ssi = []
    for pi, chunk in enumerate(pages):
        h, vw, vh = generate_html(data, chunk, a2)
        sp = os.path.join(tempfile.gettempdir(), f"gantt_{pi}.png")
        asyncio.run(take_screenshot(h, sp, vw, vh))
        ssi.append((sp, vw, vh))

    prs = Presentation(args.template)
    tmpl = prs.slides[0]
    exs = [dup_slide(prs, tmpl) for _ in range(1, tp)]

    pop_slide(tmpl, mm, ct, ssi[0][0], 1, tp, ssi[0][1], ssi[0][2])
    for pi, es in enumerate(exs):
        s = ssi[pi+1]; pop_slide(es, mm, ct, s[0], pi+2, tp, s[1], s[2])

    prs.save(args.output)

    _finalize_pptx(args.output)

    print(f"\n  Saved: {args.output} ({tp} slide(s))")
    for sp, _, _ in ssi:
        if os.path.exists(sp): os.unlink(sp)


if __name__ == "__main__":
    main()
