"""
fill_competitor_summary.py — 競合比較サマリースライドをPPTXネイティブテーブルで生成

テンプレート構造 (competitor-summary-template.pptx):
  - Title 1            (PLACEHOLDER): メインメッセージ
  - Text Placeholder 2 (PLACEHOLDER): チャートタイトル
  - Content Area       (AUTO_SHAPE):  削除してネイティブテーブルに置換
  - Source             (TEXT_BOX):    出典

方式: 横型比較テーブル（行=比較項目、列=企業）をネイティブテーブルで生成。
      target_company を指定した場合のみ、対象会社の列をイエロー背景＋太字でハイライト。
      target_company が未指定なら competitors[] のみで全社フラット表示（強調なし）。
      3〜5競合の可変列数に対応（対象会社＋競合で計4〜6列）。

Usage:
  python fill_competitor_summary.py \
    --data /home/claude/competitor_summary_data.json \
    --template <SKILL_DIR>/assets/competitor-summary-template.pptx \
    --output /mnt/user-data/outputs/CompetitorSummary_output.pptx
"""

import argparse
import copy
import json
import os
import sys

# brand_resolver bootstrap (passive --brand acceptance until brand-aware migration)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(SKILL_DIR, "..", "_common", "lib"))
from brand_resolver import add_brand_arg  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def _finalize_pptx(path):
    """LibreOffice roundtrip to normalize OOXML so PowerPoint stops asking for repair.

    No-op if soffice is unavailable or the conversion fails; the original file
    is preserved. Added by tools/add_finalize_hook.py.
    """
    import os, shutil, subprocess, tempfile, glob
    candidates = [
        os.environ.get("SOFFICE_BIN"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/usr/bin/soffice",
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    soffice = next((c for c in candidates if c and os.path.exists(c)), None)
    if not soffice:
        return
    try:
        with tempfile.TemporaryDirectory(prefix="pptx_rt_") as tmp:
            subprocess.run(
                [soffice, f"-env:UserInstallation=file://{tmp}/prof",
                 "--headless", "--convert-to", "pptx",
                 "--outdir", tmp, str(path)],
                timeout=120, capture_output=True, check=True,
            )
            found = glob.glob(os.path.join(tmp, "*.pptx"))
            if found:
                shutil.move(found[0], str(path))
    except Exception:
        pass



# ── Shape名マッピング ──
SHAPE_MAIN_MESSAGE = "Title 1"
SHAPE_CHART_TITLE = "Text Placeholder 2"
SHAPE_CONTENT_AREA = "Content Area"
SHAPE_SOURCE = "Source"

# ── レイアウト定数 ──
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

CONTENT_LEFT = Inches(0.41)
CONTENT_TOP = Inches(1.50)
CONTENT_WIDTH = Inches(12.52)
CONTENT_BOTTOM = Inches(6.90)     # Source の上
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP  # ~5.40in

# 比較項目列（一番左）の幅比率
ITEM_COL_RATIO = 0.14

# ── 色定数 ──
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_HEADER_BG = RGBColor(0xF0, 0xF0, 0xF0)   # 比較項目列の背景（グレー）
COLOR_TARGET_BG = RGBColor(0xFF, 0xF4, 0xC2)   # 対象会社列の背景（イエロー）
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_EVEN_ROW = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)

# ── フォント ──
FONT_NAME_JP = "Meiryo UI"
FONT_NAME_LATIN = "Arial"


def find_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    print(f"  ⚠ WARNING: Shape '{name}' not found", file=sys.stderr)
    return None


def remove_shape(slide, name):
    shape = find_shape(slide, name)
    if shape is not None:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)
        print(f"  ✓ Shape '{name}' removed")


def set_textbox_text(shape, text):
    """TextBoxのテキストを上書き（既存スタイルを保持）"""
    if shape is None:
        return
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        r_elem = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ja-JP"})
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = text


def get_font_sizes(num_competitors):
    """競合社数に応じてフォントサイズを動的決定（対象+5社まで対応）

    deck_skeleton_standard.json の limits.max_competitors=5 に合わせ、
    対象＋3〜5社（4〜6列）の3段階で運用する。

    v0.2 で対象＋10社まで拡張したが、9列構成（target+8）で cell 30字制限の運用負荷が
    高いため、2026-04-29 に 5 社上限へ撤回（ISSUE-008）。
    """
    # 対象会社を含めた総列数（比較項目列を除く）
    num_companies = num_competitors + 1
    if num_companies <= 4:           # 対象＋3社
        return {"header": 14, "body": 13, "item": 14}
    elif num_companies == 5:         # 対象＋4社
        return {"header": 13, "body": 12, "item": 13}
    else:                            # 対象＋5社（上限）
        return {"header": 12, "body": 11, "item": 12}


def apply_cell_style(cell, text, *,
                     font_size=10, bold=False,
                     bg_color=None, text_color=None,
                     align="left", v_align="top", is_multiline=False):
    """セルにテキストとスタイルを適用

    Args:
        cell: pptxテーブルセル
        text: テキスト（文字列または文字列リスト）
        font_size: フォントサイズ (pt, int)
        bold: 太字
        bg_color: 背景色 (RGBColor)
        text_color: 文字色 (RGBColor)
        align: 'left' / 'center' / 'right'
        v_align: 'top' / 'middle' / 'bottom'
        is_multiline: 複数段落にするか
    """
    # 背景色
    if bg_color is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    # 垂直方向の配置
    if v_align == "middle":
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif v_align == "bottom":
        cell.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        cell.vertical_anchor = MSO_ANCHOR.TOP

    # 余白設定（読みやすさ）
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)

    # テキストフレーム取得・既存段落を全削除
    tf = cell.text_frame
    tf.word_wrap = True

    # text_frameの既存段落を全て削除して作り直す
    tc = cell._tc
    txBody = tc.find(qn("a:txBody"))
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # テキストを段落リストに正規化
    if isinstance(text, list):
        lines = [str(line) for line in text]
    else:
        lines = [str(text)] if text else [""]

    align_map = {"left": "l", "center": "ctr", "right": "r"}
    algn_code = align_map.get(align, "l")

    text_rgb = text_color if text_color is not None else COLOR_TEXT

    for line in lines:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        pPr.set("algn", algn_code)

        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
            "lang": "ja-JP",
            "sz": str(font_size * 100),
            "b": "1" if bold else "0",
        })

        # 色
        solidFill = etree.SubElement(rPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{text_rgb[0]:02X}{text_rgb[1]:02X}{text_rgb[2]:02X}")

        # 日本語フォント
        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", FONT_NAME_LATIN)
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", FONT_NAME_JP)
        cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", FONT_NAME_JP)

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line


def build_table(slide, data):
    """競合比較テーブルを動的に構築

    target_company が指定されていれば列1をイエロー＋太字で強調。
    未指定の場合は competitors[] のみを並べ、強調列なしのフラット表示。
    """
    target = data.get("target_company") or None
    has_target = bool(target and target.get("name"))
    competitors = data["competitors"]
    items = data["comparison_items"]  # ["事業内容", "本社所在地", ...]

    # 表示企業リスト: target ありなら先頭に target、なしなら competitors のみ
    companies = ([target] + list(competitors)) if has_target else list(competitors)

    # 列数計算: 比較項目列(1) + 企業列(N or 1+N)
    n_companies = len(companies)
    n_rows = len(items) + 1  # ヘッダー行 + 項目行
    n_cols = 1 + n_companies  # 比較項目 + 企業

    # Content Areaを削除
    content_shape = find_shape(slide, SHAPE_CONTENT_AREA)
    if content_shape is not None:
        tbl_left = content_shape.left
        tbl_top = content_shape.top
        tbl_width = content_shape.width
        tbl_height = content_shape.height
        sp_tree = slide.shapes._spTree
        sp_tree.remove(content_shape._element)
    else:
        tbl_left = CONTENT_LEFT
        tbl_top = CONTENT_TOP
        tbl_width = CONTENT_WIDTH
        tbl_height = CONTENT_HEIGHT

    # 列幅計算: 比較項目列14%、残りの86%を企業列で等分
    item_col_w = int(tbl_width * ITEM_COL_RATIO)
    remaining_w = tbl_width - item_col_w
    company_col_w = remaining_w // n_companies

    # 行高計算: ヘッダー行は固定、データ行は等分
    header_row_h = Inches(0.40)
    available_h = tbl_height - header_row_h
    data_row_h = int(available_h / len(items))

    # フォントサイズ決定（get_font_sizes は内部で +1 して総表示企業数に戻すので
    # 「target あり/なし」に関わらず "総表示企業数 - 1" を渡す）
    fs = get_font_sizes(n_companies - 1)

    print(f"  ✓ テーブル: {n_rows}行 × {n_cols}列")
    if has_target:
        print(f"    対象会社+競合={n_companies}社, 列幅(企業): {company_col_w/914400:.2f}in")
    else:
        print(f"    全社フラット表示={n_companies}社（強調なし）, 列幅(企業): {company_col_w/914400:.2f}in")
    print(f"    データ行高: {data_row_h/914400:.2f}in, フォント: body={fs['body']}pt")

    # テーブル追加
    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table_shape.name = "CompetitorSummaryTable"
    table = table_shape.table

    # 列幅設定
    table.columns[0].width = item_col_w
    for c in range(1, n_cols):
        table.columns[c].width = company_col_w

    # 行高さ設定（XMLレベル）
    tbl_xml = table_shape._element.find('.//' + qn('a:tbl'))
    for i, tr in enumerate(tbl_xml.findall(qn('a:tr'))):
        if i == 0:
            tr.set('h', str(header_row_h))
        else:
            tr.set('h', str(data_row_h))

    # tblPrクリア（デフォルトスタイル除去）
    tbl_elem = table_shape._element.find('.//' + qn('a:tbl'))
    old_tblPr = tbl_elem.find(qn('a:tblPr'))
    if old_tblPr is not None:
        tbl_elem.remove(old_tblPr)
    tblPr = etree.SubElement(tbl_elem, qn('a:tblPr'), attrib={
        'firstRow': '1', 'bandRow': '0'
    })
    tbl_elem.insert(0, tblPr)

    # ── ヘッダー行 ──
    # (0,0): 空白（比較項目の見出し列）
    apply_cell_style(
        table.cell(0, 0), "",
        font_size=fs["header"], bold=True,
        bg_color=COLOR_HEADER_BG,
        align="center", v_align="middle",
    )

    # (0,1+): 企業名（target ありなら列1をイエロー強調、それ以外はグレー）
    for c_idx, comp in enumerate(companies):
        is_target_col = has_target and c_idx == 0
        apply_cell_style(
            table.cell(0, 1 + c_idx),
            comp.get("name", f"競合{c_idx+1}"),
            font_size=fs["header"], bold=True,
            bg_color=COLOR_TARGET_BG if is_target_col else COLOR_HEADER_BG,
            align="center", v_align="middle",
        )

    # ── データ行 ──
    for r_idx, item in enumerate(items):
        item_key = item.get("key")  # 例: "事業内容"
        item_label = item.get("label", item_key)
        row = r_idx + 1

        # 行の背景色（偶数行はグレー）
        row_bg = COLOR_EVEN_ROW if (r_idx % 2 == 0) else COLOR_WHITE

        # (row, 0): 比較項目ラベル
        apply_cell_style(
            table.cell(row, 0), item_label,
            font_size=fs["item"], bold=True,
            bg_color=COLOR_HEADER_BG,
            align="left", v_align="middle",
        )

        # (row, 1+): 企業の値（target あり&列1のみイエロー＋bold）
        for c_idx, comp in enumerate(companies):
            is_target_col = has_target and c_idx == 0
            val = comp.get(item_key, "")
            apply_cell_style(
                table.cell(row, 1 + c_idx), val,
                font_size=fs["body"],
                bold=is_target_col,
                bg_color=COLOR_TARGET_BG if is_target_col else row_bg,
                align="left", v_align="top",
            )

        print(f"    行{row}: {item_label}")

    print(f"  ✓ テーブル生成完了")


MAIN_MESSAGE_MAX = 65
CELL_VALUE_MAX = 30
COMPETITORS_MIN = 2
COMPETITORS_MAX = 5  # deck_skeleton_standard.json limits.max_competitors と同期（2026-04-29 ISSUE-008 で 10→5 撤回）


def _validate_input(data):
    """入力JSONのバリデーション。main_message ≤65字、competitors=2〜5、各セル値 ≤30字。"""
    main_message = data.get("main_message", "")
    if len(main_message) > MAIN_MESSAGE_MAX:
        raise ValueError(
            f"main_message は {MAIN_MESSAGE_MAX} 字以内（受領: {len(main_message)}）: {main_message[:80]}..."
        )
    target = data.get("target_company") or {}
    has_target = bool(target.get("name"))
    competitors = data.get("competitors", [])
    if not (COMPETITORS_MIN <= len(competitors) <= COMPETITORS_MAX):
        raise ValueError(
            f"competitors の要素数は {COMPETITORS_MIN}〜{COMPETITORS_MAX} の範囲である必要があります"
            f"（受領: {len(competitors)}）。max_competitors は 2026-04-29 に 10→5 へ撤回（ISSUE-008）"
        )
    comparison_items = data.get("comparison_items", [])
    keys = [item.get("key") for item in comparison_items if item.get("key")]
    for k in keys:
        if k in ("name",):
            continue
        if has_target:
            v = target.get(k, "")
            if isinstance(v, str) and len(v) > CELL_VALUE_MAX:
                raise ValueError(
                    f"target_company.{k} は {CELL_VALUE_MAX} 字以内（受領: {len(v)}）: {v}"
                )
        for i, c in enumerate(competitors):
            cv = c.get(k, "")
            if isinstance(cv, str) and len(cv) > CELL_VALUE_MAX:
                raise ValueError(
                    f"competitors[{i}].{k} は {CELL_VALUE_MAX} 字以内（受領: {len(cv)}）: {cv}"
                )


def main():
    parser = argparse.ArgumentParser(description="競合比較サマリー PowerPoint ジェネレーター")
    parser.add_argument("--data", required=True, help="JSONデータファイルのパス")
    parser.add_argument("--template", required=True, help="PPTXテンプレートのパス")
    parser.add_argument("--output", required=True, help="出力PPTXファイルのパス")
    add_brand_arg(parser)  # passive: accepted but ignored until brand migration
    args = parser.parse_args()

    # JSONデータ読み込み
    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    _validate_input(data)
    n_comp = len(data.get("competitors", []))
    target_name = (data.get("target_company") or {}).get("name") or "（指定なし・強調なしモード）"
    print(f"  データ読み込み: 対象={target_name}, 競合={n_comp}社")

    # テンプレート読み込み
    prs = Presentation(args.template)
    slide = prs.slides[0]

    # 1. メインメッセージ
    main_message = data.get("main_message", "")
    set_textbox_text(find_shape(slide, SHAPE_MAIN_MESSAGE), main_message)
    print(f"  ✓ Main Message: {main_message[:40]}...")

    # 2. チャートタイトル
    chart_title = data.get("chart_title", "競合比較")
    set_textbox_text(find_shape(slide, SHAPE_CHART_TITLE), chart_title)
    print(f"  ✓ Chart Title: {chart_title}")

    # 3. テーブル構築
    build_table(slide, data)

    # 4. 出典
    source_text = data.get("source", "出典：各社HP、IR資料、東京商工リサーチ等")
    source_shape = find_shape(slide, SHAPE_SOURCE)
    if source_shape is not None:
        set_textbox_text(source_shape, source_text)
        print(f"  ✓ Source: {source_text[:40]}...")

    # 5. 出力
    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    _finalize_pptx(args.output)
    print(f"\n  ✅ 出力完了: {args.output}")


if __name__ == "__main__":
    main()
