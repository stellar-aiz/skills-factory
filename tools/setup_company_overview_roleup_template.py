#!/usr/bin/env python3
"""Generate skills/company-overview-pptx-v2/assets/roleup/company-overview-template.pptx.

Base: skills/customer-profile-pptx/assets/roleup/customer-profile-template.pptx
(A4 横 11.693×8.268, Title 1 / Text Placeholder 2 / Source 3 / 茶色ガイド ×2 / object 8 ×2)

Adds the company-overview-v2-specific shapes:
  - Overview Table (left panel: 5.23 × 5.40 in @ x=0.41, y=1.92)
    A 2×2 placeholder table with cell rPr pre-set to Yu Gothic UI 10pt so
    rebuild_table() in fill_company_overview.py can copy that rPr to new
    rows and stay within the C4 allowed-font-size set {22,14,12,10,6}.
  - Photo Caption 1 / Photo Area 1 (right top: 5.23 × 0.30 + 5.23 × 2.40 in)
  - Photo Caption 2 / Photo Area 2 (right bottom)

The brown guide rectangles (`正方形/長方形 1`, `正方形/長方形 8`) and
`object 8` decorations are kept; fill_company_overview.py removes the brown
rectangles silently and the object 8 decorations remain harmless under the
panel content.

Idempotent: re-running overwrites the output file.

Usage:
  python3 tools/setup_company_overview_roleup_template.py
"""
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(
    REPO_ROOT,
    "skills/customer-profile-pptx/assets/roleup/customer-profile-template.pptx",
)
DST = os.path.join(
    REPO_ROOT,
    "skills/company-overview-pptx-v2/assets/roleup/company-overview-template.pptx",
)

# Coordinates (inches) — cp roleup の左右 panel と source 領域に整合
TABLE_X, TABLE_Y, TABLE_W, TABLE_H = 0.41, 1.92, 5.23, 5.40

PHOTO_CAP1_X, PHOTO_CAP1_Y = 6.07, 1.92
PHOTO_CAP1_W, PHOTO_CAP1_H = 5.23, 0.30
PHOTO_AREA1_X, PHOTO_AREA1_Y = 6.07, 2.27
PHOTO_AREA1_W, PHOTO_AREA1_H = 5.23, 2.40

PHOTO_CAP2_X, PHOTO_CAP2_Y = 6.07, 4.77
PHOTO_CAP2_W, PHOTO_CAP2_H = 5.23, 0.30
PHOTO_AREA2_X, PHOTO_AREA2_Y = 6.07, 5.12
PHOTO_AREA2_W, PHOTO_AREA2_H = 5.23, 2.40

FONT_EA = "Yu Gothic UI"
FONT_LATIN = "Yu Gothic UI"
BODY_SZ = "1000"  # 10pt for table cells (C4 allowed)
COLOR_TEXT = "241A17"
COLOR_LABEL_BG = "F2E8DD"
COLOR_VALUE_BG = "FFFFFF"
COLOR_BORDER = "CDCECE"


def add_overview_table(slide):
    """2x2 placeholder table whose cell rPr is pre-set to Yu Gothic UI 10pt."""
    shape = slide.shapes.add_table(
        2, 2,
        Inches(TABLE_X), Inches(TABLE_Y),
        Inches(TABLE_W), Inches(TABLE_H),
    )
    shape.name = "Overview Table"
    table = shape.table

    # Column widths: label 35% / value 65%
    table.columns[0].width = Inches(TABLE_W * 0.35)
    table.columns[1].width = Inches(TABLE_W * 0.65)

    placeholders = [
        ("商号", "（テンプレート行 1）"),
        ("本店所在地", "（テンプレート行 2）"),
    ]

    for r_idx, (label, value) in enumerate(placeholders):
        for c_idx, txt in enumerate((label, value)):
            cell = table.cell(r_idx, c_idx)
            tc = cell._tc
            txBody = tc.find(qn("a:txBody"))
            for p in txBody.findall(qn("a:p")):
                txBody.remove(p)
            p_elem = etree.SubElement(txBody, qn("a:p"))
            pPr = etree.SubElement(p_elem, qn("a:pPr"))
            pPr.set("algn", "l")
            r_elem = etree.SubElement(p_elem, qn("a:r"))
            rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
                "lang": "ja-JP", "sz": BODY_SZ,
            })
            etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_LATIN})
            etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_EA})
            sf = etree.SubElement(rPr, qn("a:solidFill"))
            etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": COLOR_TEXT})
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = txt

            # tcPr: bg color + border + vertical center
            old_tcPr = tc.find(qn("a:tcPr"))
            if old_tcPr is not None:
                tc.remove(old_tcPr)
            tcPr = etree.SubElement(tc, qn("a:tcPr"))
            tcPr.set("anchor", "ctr")
            bg = COLOR_LABEL_BG if c_idx == 0 else COLOR_VALUE_BG
            sf2 = etree.SubElement(tcPr, qn("a:solidFill"))
            etree.SubElement(sf2, qn("a:srgbClr"), attrib={"val": bg})
            for side in ("lnL", "lnR", "lnT", "lnB"):
                ln = etree.SubElement(tcPr, qn(f"a:{side}"), attrib={
                    "w": "6350", "cap": "flat", "cmpd": "sng", "algn": "ctr",
                })
                sfb = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(sfb, qn("a:srgbClr"), attrib={"val": COLOR_BORDER})


def add_photo_area(slide, x, y, w, h, name):
    """Photo Area: AUTO_SHAPE.RECTANGLE (light beige bg, no border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF2, 0xE8, 0xDD)
    shape.line.color.rgb = RGBColor(0xCD, 0xCE, 0xCE)
    shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    shape.text_frame.text = ""


def add_photo_caption(slide, x, y, w, h, name, text="（写真キャプション）"):
    """Photo Caption: TEXT_BOX with sz=1000 (10pt) Yu Gothic UI."""
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    tb.name = name
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if p.runs:
        for r in p.runs:
            r.text = ""
    r_elem = etree.SubElement(p._p, qn("a:r"))
    rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={
        "lang": "ja-JP", "sz": BODY_SZ, "b": "1",
    })
    etree.SubElement(rPr, qn("a:latin"), attrib={"typeface": FONT_LATIN})
    etree.SubElement(rPr, qn("a:ea"), attrib={"typeface": FONT_EA})
    sf = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": COLOR_TEXT})
    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text


def main():
    if not os.path.exists(SRC):
        print(f"ERROR: source not found: {SRC}", file=sys.stderr)
        return 1

    prs = Presentation(SRC)
    slide = prs.slides[0]

    # Add company-overview-v2-specific shapes.
    add_overview_table(slide)
    add_photo_caption(slide, PHOTO_CAP1_X, PHOTO_CAP1_Y,
                      PHOTO_CAP1_W, PHOTO_CAP1_H,
                      "Photo Caption 1", "本社家屋")
    add_photo_area(slide, PHOTO_AREA1_X, PHOTO_AREA1_Y,
                   PHOTO_AREA1_W, PHOTO_AREA1_H, "Photo Area 1")
    add_photo_caption(slide, PHOTO_CAP2_X, PHOTO_CAP2_Y,
                      PHOTO_CAP2_W, PHOTO_CAP2_H,
                      "Photo Caption 2", "主要製品 / サービス")
    add_photo_area(slide, PHOTO_AREA2_X, PHOTO_AREA2_Y,
                   PHOTO_AREA2_W, PHOTO_AREA2_H, "Photo Area 2")

    os.makedirs(os.path.dirname(DST), exist_ok=True)
    prs.save(DST)
    print(f"  ✓ Saved: {DST}")

    # Verify
    prs2 = Presentation(DST)
    s2 = prs2.slides[0]
    print(f"\n  --- shapes in {os.path.basename(DST)} ---")
    for shp in s2.shapes:
        l = shp.left/914400 if shp.left else 0
        t = shp.top/914400 if shp.top else 0
        w = shp.width/914400 if shp.width else 0
        h = shp.height/914400 if shp.height else 0
        print(f"    {shp.shape_type:>5} | {shp.name!r:<30} | x={l:.2f} y={t:.2f} w={w:.2f} h={h:.2f}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
